# Copyright (C) 2025 CrispStrobe
#
# This program is free software: you can redistribute it and/or modify
# it under the terms of the GNU Affero General Public License as published
# by the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.

# /var/www/transkript_app/app.py

import gradio as gr
import os
from typing import Tuple, List, Dict, Optional
import glob

# --- FORCE FFMPEG PATH ---
# Explicitly tell Python where to find the tools
os.environ["PATH"] += os.pathsep + "/usr/bin" + os.pathsep + "/usr/local/bin"

try:
    from pydub import AudioSegment
    AudioSegment.converter = "/usr/bin/ffmpeg"
    AudioSegment.ffprobe   = "/usr/bin/ffprobe"
except ImportError:
    print ("pydub cannot be imported")
    pass

import time
import json
import logging
import sys
import base64
import requests
import tempfile
from io import BytesIO
import asyncio
import re
import hashlib
from bs4 import BeautifulSoup

logging.basicConfig(
    level=logging.INFO,  # Global level INFO to stop library spam
    format='%(asctime)s [%(levelname)s] %(name)s:%(lineno)d - %(message)s',
    handlers=[
        logging.StreamHandler(sys.stdout),
        logging.FileHandler('/var/www/transkript_app/app.log')
    ]
)
# Only set this app to debug
logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)

# Set Gradio file upload limits
os.environ['GRADIO_TEMP_DIR'] = '/tmp/gradio'
os.makedirs('/tmp/gradio', exist_ok=True)

try:
    import fastapi_poe as fp
    HAS_POE = True
except ImportError:
    logger.warning("fastapi_poe not installed. Poe API will be unavailable.")
    HAS_POE = False

# Increase max file size
gr.set_static_paths(paths=["/var/www/transkript_app/static"])

# --- IMPORTS & SETUP ---
try:
    from PIL import Image
    import openai
except ImportError:
    print("⚠️ WARNUNG: Bitte 'pip install openai pillow' ausführen.")

from crypto_utils import crypto, HAS_PQ, KeyWrapper
# Initialize the wrapper
key_wrapper = KeyWrapper()

# ==========================================
# 🗄️ DATABASE SETUP
# ==========================================

from sqlalchemy import create_engine, Column, Integer, String, Text, DateTime, ForeignKey, Boolean, UniqueConstraint
from sqlalchemy.orm import declarative_base
from sqlalchemy.orm import sessionmaker, relationship
from datetime import datetime
import bcrypt

# Database setup
DATABASE_URL = "sqlite:///akademie_suite.db"
engine = create_engine(DATABASE_URL, connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

def migrate_to_encrypted_db():
    """
    CRITICAL: Encrypt all existing plaintext data
    Run this ONCE during deployment
    """
    db = SessionLocal()
    
    try:
        logger.info("🔐 Starting database encryption migration...")
        
        # 1. Check if crypto is initialized
        if not crypto.master_key:
            logger.error("❌ Crypto not initialized! Cannot migrate.")
            return False
        
        # 2. Migrate Transcriptions
        trans = db.query(Transcription).filter(Transcription.is_encrypted == False).all()
        logger.info(f"Migrating {len(trans)} transcriptions...")
        for t in trans:
            try:
                if t.original_text:
                    t.original_text = crypto.encrypt_text(t.original_text)
                if t.translated_text:
                    t.translated_text = crypto.encrypt_text(t.translated_text)
                t.is_encrypted = True
                t.encryption_metadata = json.dumps({"algorithm": "AES-256-GCM", "version": 1})
            except Exception as e:
                logger.error(f"Failed to encrypt transcription {t.id}: {e}")
                continue  # Skip this one, continue with others
        
        # 3. Migrate Chat History
        chats = db.query(ChatHistory).filter(ChatHistory.is_encrypted == False).all()
        logger.info(f"Migrating {len(chats)} chats...")
        for c in chats:
            try:
                if c.messages:
                    c.messages = crypto.encrypt_text(c.messages)
                c.is_encrypted = True
                c.encryption_metadata = json.dumps({"algorithm": "AES-256-GCM", "version": 1})
            except Exception as e:
                logger.error(f"Failed to encrypt chat {c.id}: {e}")
                continue
        
        # 4. Migrate Vision Results
        visions = db.query(VisionResult).filter(VisionResult.is_encrypted == False).all()
        logger.info(f"Migrating {len(visions)} vision results...")
        for v in visions:
            try:
                if v.result:
                    v.result = crypto.encrypt_text(v.result)
                if v.prompt:
                    v.prompt = crypto.encrypt_text(v.prompt)
                v.is_encrypted = True
                v.encryption_metadata = json.dumps({"algorithm": "AES-256-GCM", "version": 1})
            except Exception as e:
                logger.error(f"Failed to encrypt vision {v.id}: {e}")
                continue
        
        # 5. Encrypt Generated Images (files) - Skipped by default (too large)
        images = db.query(GeneratedImage).filter(GeneratedImage.is_encrypted == False).all()
        logger.info(f"Skipping {len(images)} images (too large for migration)")
        
        db.commit()
        logger.info("✅ Migration complete!")
        return True
        
    except Exception as e:
        db.rollback()
        logger.exception(f"🔥 Migration failed: {e}")
        return False
    finally:
        db.close()

# Call on startup (only runs once due to is_encrypted flag)
def ensure_encryption():
    """Check if migration is needed"""
    db = SessionLocal()
    try:
        # Check if any unencrypted data exists
        unencrypted_count = db.query(Transcription).filter(
            Transcription.is_encrypted == False
        ).count()
        
        if unencrypted_count > 0:
            logger.warning(f"⚠️ Found {unencrypted_count} unencrypted transcriptions. Starting migration...")
            migrate_to_encrypted_db()
    except Exception as e:
        logger.error(f"Encryption check failed: {e}")
    finally:
        db.close()

# ==========================================
# 📊 DATABASE MODELS
# ==========================================

class User(Base):
    __tablename__ = "users"
    
    id = Column(Integer, primary_key=True, index=True)
    username = Column(String, unique=True, index=True, nullable=False)
    password_hash = Column(String, nullable=False)

    # FIELDS FOR KEY WRAPPING
    salt = Column(String) # Stores the random salt for PBKDF2
    encrypted_master_key = Column(Text) # Stores the encrypted User Master Key (Lockbox)

    email = Column(String, unique=True, index=True)
    is_admin = Column(Boolean, default=False)
    is_media_manager = Column(Boolean, default=False)
    created_at = Column(DateTime, default=datetime.utcnow)
    
    # Relationships
    chat_history = relationship("ChatHistory", back_populates="user", cascade="all, delete-orphan")
    transcriptions = relationship("Transcription", back_populates="user", cascade="all, delete-orphan")
    vision_results = relationship("VisionResult", back_populates="user", cascade="all, delete-orphan")
    generated_images = relationship("GeneratedImage", back_populates="user", cascade="all, delete-orphan")
    custom_prompts = relationship("CustomPrompt", back_populates="user", cascade="all, delete-orphan")
    model_preferences = relationship("UserModelPreference", back_populates="user", cascade="all, delete-orphan")
    settings = relationship("UserSettings", uselist=False, back_populates="user", cascade="all, delete-orphan")
class ChatHistory(Base):
    __tablename__ = "chat_history"

    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    provider = Column(String, nullable=False)
    model = Column(String, nullable=False)
    messages = Column(Text, nullable=False)  # JSON string
    timestamp = Column(DateTime, default=datetime.utcnow)
    title = Column(String)  # Auto-generated summary

    user = relationship("User", back_populates="chat_history")

    is_encrypted = Column(Boolean, default=False)
    encryption_metadata = Column(Text)

class Transcription(Base):
    __tablename__ = "transcriptions"

    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    provider = Column(String, nullable=False)
    model = Column(String)
    original_text = Column(Text, nullable=False)
    translated_text = Column(Text)
    language = Column(String)
    filename = Column(String)
    timestamp = Column(DateTime, default=datetime.utcnow)
    title = Column(String)  # User-defined or auto-generated

    user = relationship("User", back_populates="transcriptions")

    is_encrypted = Column(Boolean, default=False)
    encryption_metadata = Column(Text)  # JSON with algorithm info

class VisionResult(Base):
    __tablename__ = "vision_results"

    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    provider = Column(String, nullable=False)
    model = Column(String, nullable=False)
    prompt = Column(Text, nullable=False)
    result = Column(Text, nullable=False)
    image_path = Column(String)  # Store path to uploaded image
    timestamp = Column(DateTime, default=datetime.utcnow)

    user = relationship("User", back_populates="vision_results")

    is_encrypted = Column(Boolean, default=False)
    encryption_metadata = Column(Text)

class GeneratedImage(Base):
    __tablename__ = "generated_images"

    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    provider = Column(String, nullable=False)
    model = Column(String, nullable=False)
    prompt = Column(Text, nullable=False)
    image_path = Column(String, nullable=False)
    timestamp = Column(DateTime, default=datetime.utcnow)

    user = relationship("User", back_populates="generated_images")

    is_encrypted = Column(Boolean, default=False)
    encryption_metadata = Column(Text)
    encrypted_path = Column(String)  # Path to .enc file

class FileUploadMetadata(Base):
    """NEW TABLE: Track encrypted uploads"""
    __tablename__ = "file_uploads"
    
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    original_filename = Column(String, nullable=False)
    encrypted_path = Column(String, nullable=False)
    encryption_metadata = Column(Text, nullable=False)  # JSON
    upload_timestamp = Column(DateTime, default=datetime.utcnow)
    file_hash = Column(String)  # SHA256 of original
    
    user = relationship("User", backref="uploaded_files")

class CustomPrompt(Base):
    __tablename__ = "custom_prompts"

    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    name = Column(String, nullable=False)
    category = Column(String)  # "chat", "transcription", etc.
    prompt_text = Column(Text, nullable=False)
    is_shared = Column(Boolean, default=False)  # Share with other users
    timestamp = Column(DateTime, default=datetime.utcnow)

    user = relationship("User", back_populates="custom_prompts")

class UserModelPreference(Base):
    __tablename__ = "user_model_preferences"
    
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    provider = Column(String, nullable=False)
    model_id = Column(String, nullable=False)  # The actual model ID used by API
    display_name = Column(String)  # Human-readable name
    display_order = Column(Integer, default=0)  # Lower = higher priority
    is_visible = Column(Boolean, default=True)
    created_at = Column(DateTime, default=datetime.utcnow)
    
    # Relationship
    user = relationship("User", back_populates="model_preferences")
    
    # Unique constraint: one entry per user+provider+model
    __table_args__ = (
        UniqueConstraint('user_id', 'provider', 'model_id', name='_user_provider_model_uc'),
    )

class UserSettings(Base):
    """Store user preferences for app behavior"""
    __tablename__ = "user_settings"
    
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"), nullable=False, unique=True)
    
    # Attachment & Chunking Settings
    auto_chunk_enabled = Column(Boolean, default=True)
    chunk_size = Column(Integer, default=4000)
    chunk_overlap = Column(Integer, default=200)
    
    # Chat Settings
    auto_truncate_history = Column(Boolean, default=True)
    show_truncation_warning = Column(Boolean, default=True)
    
    # Display Preferences
    show_token_counts = Column(Boolean, default=False)
    compact_mode = Column(Boolean, default=False)
    
    created_at = Column(DateTime, default=datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)
    
    user = relationship("User", back_populates="settings")

class DSFARecord(Base):
    """Track Data Protection Impact Assessments"""
    __tablename__ = "dsfa_records"
    
    id = Column(Integer, primary_key=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    assessment_date = Column(DateTime, default=datetime.utcnow)
    use_case = Column(String, nullable=False)  # "transcription", "chat", etc.
    provider = Column(String, nullable=False)
    risk_level = Column(String)  # "low", "medium", "high"
    special_categories = Column(Boolean, default=False)  # Art. 9 DS-GVO
    mitigation_measures = Column(Text)
    approved_by = Column(String)
    notes = Column(Text)
    
    user = relationship("User")

def auto_dsfa_check(provider: str, use_case: str, data_type: str) -> dict:
    """
    Automatic DSFA trigger detection
    Returns: {requires_dsfa: bool, risk_level: str, reason: str}
    """
    high_risk_providers = ["OpenAI", "Anthropic", "Poe"]  # US providers
    high_risk_uses = ["transcription", "vision"]  # Likely to contain Art. 9 data
    
    requires = False
    risk = "low"
    reason = ""
    
    if provider in high_risk_providers:
        requires = True
        risk = "high"
        reason = "US-Provider ohne EU-Datenverarbeitung"
    
    if use_case in high_risk_uses and "audio" in data_type.lower():
        requires = True
        risk = "high" if risk == "high" else "medium"
        reason += " | Religiöse Inhalte wahrscheinlich (Gottesdienst)"
    
    return {
        "requires_dsfa": requires,
        "risk_level": risk,
        "reason": reason
    }

class UserConfirmation(Base):
    """Track user confirmations for warnings/dialogs"""
    __tablename__ = "user_confirmations"
    
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    confirmation_type = Column(String, nullable=False)  # "upload_warning", "us_provider_warning", etc.
    confirmation_count = Column(Integer, default=0)
    last_confirmed = Column(DateTime, default=datetime.utcnow)
    permanently_dismissed = Column(Boolean, default=False)  # After 3 times
    created_at = Column(DateTime, default=datetime.utcnow)
    
    user = relationship("User", backref="confirmations")
    
    __table_args__ = (
        UniqueConstraint('user_id', 'confirmation_type', name='_user_confirmation_uc'),
    )

def check_user_confirmation(user_id: int, confirmation_type: str) -> tuple:
    """
    Check if user needs to see confirmation dialog.
    Returns: (show_dialog: bool, count: int)
    """
    db = SessionLocal()
    try:
        confirm = db.query(UserConfirmation).filter(
            UserConfirmation.user_id == user_id,
            UserConfirmation.confirmation_type == confirmation_type
        ).first()
        
        if not confirm:
            # First time - show dialog
            return True, 0
        
        if confirm.permanently_dismissed:
            # User has seen it 3+ times, don't show anymore
            return False, confirm.confirmation_count
        
        if confirm.confirmation_count < 3:
            # Still in the 3-time window
            return True, confirm.confirmation_count
        
        # Automatically dismiss after 3 times
        confirm.permanently_dismissed = True
        db.commit()
        return False, confirm.confirmation_count
        
    finally:
        db.close()

def record_user_confirmation(user_id: int, confirmation_type: str):
    """Record that user has seen and confirmed a dialog"""
    db = SessionLocal()
    try:
        confirm = db.query(UserConfirmation).filter(
            UserConfirmation.user_id == user_id,
            UserConfirmation.confirmation_type == confirmation_type
        ).first()
        
        if not confirm:
            confirm = UserConfirmation(
                user_id=user_id,
                confirmation_type=confirmation_type,
                confirmation_count=1
            )
            db.add(confirm)
        else:
            confirm.confirmation_count += 1
            confirm.last_confirmed = datetime.utcnow()
            
            if confirm.confirmation_count >= 3:
                confirm.permanently_dismissed = True
        
        db.commit()
        logger.info(f"User {user_id} confirmed '{confirmation_type}' (count: {confirm.confirmation_count})")
        
    except Exception as e:
        db.rollback()
        logger.error(f"Error recording confirmation: {e}")
    finally:
        db.close()

def reset_user_confirmation(user_id: int, confirmation_type: str):
    """Admin function to reset confirmation tracking"""
    db = SessionLocal()
    try:
        db.query(UserConfirmation).filter(
            UserConfirmation.user_id == user_id,
            UserConfirmation.confirmation_type == confirmation_type
        ).delete()
        db.commit()
        return True
    except:
        db.rollback()
        return False
    finally:
        db.close()

# Create tables
Base.metadata.create_all(bind=engine)

# ==========================================
# 🔐 DATABASE HELPER FUNCTIONS
# ==========================================

def get_db():
    """Get database session"""
    return SessionLocal()

def hash_password(password: str) -> str:
    """Hash a password using bcrypt"""
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password: str, hashed: str) -> bool:
    """Verify a password against its hash"""
    return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))

def authenticate_user(username, password):
    """
    Authenticate user AND create keychain if needed.
    Returns a 'SafeUser' object that works after DB close.
    """
    db = get_db()
    try:
        logger.info(f"🔐 AUTH ATTEMPT: User '{username}'")
        
        user = db.query(User).filter(User.username == username).first()
        
        if not user:
            logger.error(f"❌ AUTH FAIL: User '{username}' not found in DB")
            return None, None
            
        if not verify_password(password, user.password_hash):
            logger.error(f"❌ AUTH FAIL: Password hash mismatch for '{username}'")
            return None, None
        
        logger.info(f"✅ Password verified for '{username}'. Unlocking keychain...")

        # 1. Migration Logic (User has no keys yet)
        if not user.salt or not user.encrypted_master_key:
            logger.info(f"🔄 Migrating user {username} to per-user encryption...")
            try:
                keychain = key_wrapper.create_user_keychain(password)
                
                # Wrap the GLOBAL key (legacy support)
                from cryptography.fernet import Fernet
                salt_bytes = base64.b64decode(keychain["salt"])
                wrapper_key = key_wrapper.derive_wrapper_key(password, salt_bytes)
                f = Fernet(wrapper_key)
                encrypted_umk = f.encrypt(crypto.global_key)
                
                user.salt = keychain["salt"]
                user.encrypted_master_key = encrypted_umk.decode('utf-8')
                db.commit()
                
                umk = crypto.global_key
            except Exception as e:
                logger.error(f"Migration failed: {e}")
                return None, None
        else:
            # 2. Normal Unlock
            try:
                umk = key_wrapper.unlock_user_keychain(
                    password, 
                    user.salt, 
                    user.encrypted_master_key
                )
            except Exception as e:
                logger.error(f"❌ KEYCHAIN ERROR: {e}")
                umk = None
        
        if not umk:
            logger.error(f"❌ AUTH FAIL: Could not decrypt Master Key for '{username}'")
            return None, None

        logger.info(f"🔓 Login successful for '{username}'")

        # 3. Create safe object
        class SafeUser: pass
        
        safe_user = SafeUser()
        safe_user.id = user.id
        safe_user.username = user.username
        safe_user.is_admin = user.is_admin
        safe_user.is_media_manager = getattr(user, 'is_media_manager', False)
        
        return safe_user, umk

    except Exception as e:
        logger.exception(f"🔥 Auth Critical Error: {e}")
        return None, None
    finally:
        db.close()

def create_default_users():
    """Create default users if they don't exist"""
    db = SessionLocal()
    try:
        # Check if users exist
        if db.query(User).count() == 0:
            # Create admin user
            admin = User(
                username="admin123",
                password_hash=hash_password("ÄndereDasSofort!"),
                email="adminmail@yourdomain.de",
                is_admin=True
            )
            db.add(admin)

            # Create regular user
            user = User(
                username="user123",
                password_hash=hash_password("ÄndereDasAuchGleich!"),
                email="usermail@yourdomain.de",
                is_admin=False
            )
            db.add(user)

            db.commit()
            logger.info("✅ Default users created")
        else:
            logger.info("Users already exist, skipping creation")
    except Exception as e:
        db.rollback()
        logger.exception(f"Error creating default users: {str(e)}")
    finally:
        db.close()

def save_transcription(user_id: int, provider: str, model: str, original: str,
                       translated: str = None, language: str = None, filename: str = None, 
                       title: str = None, user_state=None): # <-- Add user_state arg
    db = SessionLocal()
    try:
        # GET KEY FROM SESSION
        umk = user_state.get('umk') if user_state else None
        if not umk: 
            # Fallback for system operations if needed, but risky
            logger.warning("No user session key found, falling back to global")
            umk = crypto.global_key

        encrypted_original = crypto.encrypt_text(original, key=umk) if original else None
        encrypted_translated = crypto.encrypt_text(translated, key=umk) if translated else None
            
        trans = Transcription(
            user_id=user_id,
            provider=provider,
            model=model or "N/A",
            original_text=encrypted_original,
            translated_text=encrypted_translated,
            language=language,
            filename=filename,
            title=title or f"Transkript {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            is_encrypted=True,
            encryption_metadata=json.dumps({"algorithm": "AES-256-GCM", "version": 2})
        )
        db.add(trans)
        db.commit()
        return trans.id
    except Exception as e:
        db.rollback()
        logger.exception(f"Error saving transcription: {str(e)}")
        raise
    finally:
        db.close()

def get_decrypted_transcription(trans_id: int, user_id: int, user_state=None):
    """Load and decrypt using user session key"""
    db = SessionLocal()
    try:
        trans = db.query(Transcription).filter(
            Transcription.id == trans_id,
            Transcription.user_id == user_id
        ).first()
        
        if not trans: return None
        
        # GET KEY
        umk = user_state.get('umk') if user_state else crypto.global_key
        
        if trans.is_encrypted:
            try:
                # Try decrypting with Session Key
                decrypted = crypto.decrypt_text(trans.original_text, key=umk)
                
                # Check validity (AES-GCM auth tag will fail if key is wrong)
                if decrypted == "[Decryption Failed]" and umk != crypto.global_key:
                    # Fallback: Data might be old (encrypted with global key)
                    decrypted = crypto.decrypt_text(trans.original_text, key=crypto.global_key)
                
                trans.original_text = decrypted
                
                # 4. Handle Translation Text (same logic)
                if trans.translated_text:
                    decrypted_trans = crypto.decrypt_text(trans.translated_text, key=umk)
                    if decrypted_trans == "[Decryption Failed]" and umk != crypto.global_key:
                        decrypted_trans = crypto.decrypt_text(trans.translated_text, key=crypto.global_key)
                    trans.translated_text = decrypted_trans

                    
            except Exception as e:
                logger.error(f"Decryption failed: {e}")
                trans.original_text = "[Fehler: Daten konnten nicht entschlüsselt werden]"
        
        return trans
    finally:
        db.close()

def get_decrypted_chat(chat_id: int, user_id: int, user_state=None):
    db = SessionLocal()
    try:
        chat = db.query(ChatHistory).filter(ChatHistory.id == chat_id, ChatHistory.user_id == user_id).first()
        if not chat: return None
        
        # GET KEY
        umk = user_state.get('umk') if user_state else None
        
        if chat.is_encrypted:
            # Try decrypting with Session Key
            try:
                # 2. Try decrypting with User Key first
                decrypted_json = crypto.decrypt_text(chat.messages, key=umk)
                
                # 3. If that failed (returns specific error string), try Global Key
                # This handles legacy data migration automatically
                if decrypted_json == "[Decryption Failed]" and umk != crypto.global_key:
                    # Fallback: Try global key (for data created before migration)
                    decrypted_json = crypto.decrypt_text(chat.messages, key=crypto.global_key)
                
                chat.messages = decrypted_json
            except:
                chat.messages = "[]"
        
        return chat
    finally:
        db.close()

def change_password_secure(user_id, old_password, new_password):
    """Changes password and re-wraps the Master Key without touching data"""
    db = SessionLocal()
    try:
        user = db.query(User).filter(User.id == user_id).first()
        if not user: return False, "User missing"
        
        # Verify old password hash first
        if not verify_password(old_password, user.password_hash):
            return False, "Altes Passwort falsch"
            
        # RE-WRAP MAGIC
        if user.salt and user.encrypted_master_key:
            try:
                new_keychain = key_wrapper.rewrap_keychain(
                    old_password, 
                    new_password, 
                    user.salt, 
                    user.encrypted_master_key
                )
                
                user.salt = new_keychain["salt"]
                user.encrypted_master_key = new_keychain["encrypted_master_key"]
            except ValueError:
                return False, "Kryptographischer Fehler beim Umschlüsseln."
        else:
            # Legacy user migration happens here if they change password!
            # If user has no keychain yet, create one using the Global Key 
            # (Assuming existing data is currently encrypted with Global Key)
            pass 

        # Standard password update
        user.password_hash = hash_password(new_password)
        db.commit()
        return True, "✅ Passwort geändert & Key neu verschlüsselt (Daten bleiben sicher)"
    except Exception as e:
        db.rollback()
        return False, f"Fehler: {e}"
    finally:
        db.close()

def save_chat_history(user_id: int, provider: str, model: str, messages: list, title: str = None, user_state=None):   
    """Save chat conversation to database"""
    db = SessionLocal()
    try:
        # Grab the key from the session state
        umk = user_state.get('umk') if user_state else crypto.global_key
        
        messages_json = json.dumps(messages)
        encrypted_messages = crypto.encrypt_text(messages_json, key=umk)

        chat = ChatHistory(
            user_id=user_id,
            provider=provider,
            model=model,
            messages=encrypted_messages,
            title=title or f"Chat {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            is_encrypted=True,
            encryption_metadata=json.dumps({"algorithm": "AES-256-GCM-UserKey", "version": 2})
        )
        db.add(chat)
        db.commit()
        db.refresh(chat)  # Refresh to get the ID
        chat_id = chat.id  # Get ID BEFORE closing session
        return chat_id
    except Exception as e:
        db.rollback()
        logger.exception(f"Error saving chat history: {str(e)}")
        raise
    finally:
        db.close()

def generate_ai_label(provider: str, model: str) -> str:
    """Generate standardized label for AI-generated content"""
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
    compliance = PROVIDERS.get(provider, {}).get("badge", "Unbekannt")
    
    return f"""
---
**🤖 KI-Generiert**
- **Provider:** {provider} ({compliance})
- **Modell:** {model}
- **Datum:** {timestamp}

*Dieser Inhalt wurde mit KI-Unterstützung erstellt und durch Menschen geprüft.*
---
"""

def save_transcription(user_id: int, provider: str, model: str, original: str,
                       translated: str = None, language: str = None, filename: str = None, 
                       title: str = None, user_state=None): # <-- Add user_state arg
    db = SessionLocal()
    try:
        # GET KEY FROM SESSION
        umk = user_state.get('umk') if user_state else None
        if not umk: 
            # Fallback for system operations if needed, but risky
            logger.warning("No user session key found, falling back to global")
            umk = crypto.global_key

        encrypted_original = crypto.encrypt_text(original, key=umk) if original else None
        encrypted_translated = crypto.encrypt_text(translated, key=umk) if translated else None
            
        trans = Transcription(
            user_id=user_id,
            provider=provider,
            model=model or "N/A",
            original_text=encrypted_original,
            translated_text=encrypted_translated,
            language=language,
            filename=filename,
            title=title or f"Transkript {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            is_encrypted=True,
            encryption_metadata=json.dumps({"algorithm": "AES-256-GCM", "version": 2})
        )
        db.add(trans)
        db.commit()
        return trans.id
    except Exception as e:
        db.rollback()
        logger.exception(f"Error saving transcription: {str(e)}")
        raise
    finally:
        db.close()

def get_decrypted_vision(vision_id: int, user_id: int, user_state=None):
    """Load and decrypt vision result"""
    db = SessionLocal()
    try:
        vis = db.query(VisionResult).filter(
            VisionResult.id == vision_id,
            VisionResult.user_id == user_id
        ).first()
        
        if not vis:
            return None
        
        # Get key from session
        umk = user_state.get('umk') if user_state else crypto.global_key
        
        if vis.is_encrypted:
            try:
                # Decrypt with user key, fallback to global
                vis.result = crypto.decrypt_text(vis.result, key=umk)
                if vis.result == "[Decryption Failed]" and umk != crypto.global_key:
                    vis.result = crypto.decrypt_text(vis.result, key=crypto.global_key)
                
                if vis.prompt:
                    vis.prompt = crypto.decrypt_text(vis.prompt, key=umk)
                    if vis.prompt == "[Decryption Failed]":
                        vis.prompt = crypto.decrypt_text(vis.prompt, key=crypto.global_key)
            except Exception as e:
                logger.error(f"Vision decryption failed: {e}")
                vis.result = "[Fehler: Entschlüsselung fehlgeschlagen]"
        
        return vis
    finally:
        db.close()

def save_vision_result(user_id: int, provider: str, model: str, prompt: str, result: str, image_path: str = None, user_state=None):
    """Save vision analysis to database"""
    db = SessionLocal()
    try:
        # Get Key
        umk = user_state.get('umk') if user_state else crypto.global_key
        
        # Encrypt result and prompt manually before creating object
        enc_result = crypto.encrypt_text(result, key=umk)
        enc_prompt = crypto.encrypt_text(prompt, key=umk)
        
        vision = VisionResult(
            user_id=user_id,
            provider=provider,
            model=model,
            prompt=enc_prompt,
            result=enc_result,
            is_encrypted=True,
            image_path=image_path,
            encryption_metadata=json.dumps({"algorithm": "AES-256-GCM-UserKey", "version": 2})
        )
        db.add(vision)
        db.commit()
        db.refresh(vision)
        vision_id = vision.id
        return vision_id
    except Exception as e:
        db.rollback()
        logger.exception(f"Error saving vision result: {str(e)}")
        raise
    finally:
        db.close()

def save_generated_image(user_id: int, provider: str, model: str, prompt: str, image_path: str, user_state=None):
    """Save generated image to database"""
    db = SessionLocal()
    try:
        umk = user_state.get('umk') if user_state else crypto.global_key
        enc_prompt = crypto.encrypt_text(prompt, key=umk)
        img = GeneratedImage(
            user_id=user_id,
            provider=provider,
            model=model,
            prompt=enc_prompt,
            is_encrypted=True,
            image_path=image_path
        )
        db.add(img)
        db.commit()
        db.refresh(img)
        img_id = img.id
        return img_id
    except Exception as e:
        db.rollback()
        logger.exception(f"Error saving generated image: {str(e)}")
        raise
    finally:
        db.close()

def get_user_transcriptions(user_id: int, limit: int = 50):
    """Get user's transcription history"""
    db = get_db()
    results = db.query(Transcription).filter(
        Transcription.user_id == user_id
    ).order_by(Transcription.timestamp.desc()).limit(limit).all()
    db.close()
    return results

def get_user_custom_prompts(user_id: int, category: str = None):
    """Get user's custom prompts"""
    db = get_db()
    query = db.query(CustomPrompt).filter(CustomPrompt.user_id == user_id)
    if category:
        query = query.filter(CustomPrompt.category == category)
    results = query.order_by(CustomPrompt.timestamp.desc()).all()
    db.close()
    return results

def save_custom_prompt(user_id: int, name: str, prompt_text: str, category: str = "general", is_shared: bool = False):
    """Save a custom prompt template"""
    db = SessionLocal()
    try:
        prompt = CustomPrompt(
            user_id=user_id,
            name=name,
            category=category,
            prompt_text=prompt_text,
            is_shared=is_shared
        )
        db.add(prompt)
        db.commit()
        db.refresh(prompt)
        prompt_id = prompt.id
        return prompt_id
    except Exception as e:
        db.rollback()
        logger.exception(f"Error saving custom prompt: {str(e)}")
        raise
    finally:
        db.close()

# Add after existing database functions

def delete_transcription(trans_id: int, user_id: int):
    """Delete a transcription safely"""
    db = get_db()
    try:
        trans = db.query(Transcription).filter(Transcription.id == trans_id, Transcription.user_id == user_id).first()
        if trans:
            db.delete(trans)
            db.commit()
            return True
        return False
    except Exception as e:
        db.rollback()
        logger.exception(f"DB Error: {e}")
        return False
    finally:
        db.close()

def delete_chat_history(chat_id: int, user_id: int):
    """Delete a chat history safely"""
    db = get_db()
    try:
        chat = db.query(ChatHistory).filter(ChatHistory.id == chat_id, ChatHistory.user_id == user_id).first()
        if chat:
            db.delete(chat)
            db.commit()
            return True
        return False
    except Exception as e:
        db.rollback()
        logger.exception(f"DB Error: {e}")
        return False
    finally:
        db.close()

def delete_vision_result(vision_id: int, user_id: int):
    """Delete a vision result safely"""
    db = get_db()
    try:
        vision = db.query(VisionResult).filter(VisionResult.id == vision_id, VisionResult.user_id == user_id).first()
        if vision:
            db.delete(vision)
            db.commit()
            return True
        return False
    except Exception as e:
        db.rollback()
        logger.exception(f"DB Error: {e}")
        return False
    finally:
        db.close()


def get_user_settings(user_id: int):
    """Get or create user settings"""
    db = get_db()
    try:
        settings = db.query(UserSettings).filter(UserSettings.user_id == user_id).first()
        
        if not settings:
            # Create default settings
            settings = UserSettings(user_id=user_id)
            db.add(settings)
            db.commit()
            db.refresh(settings)
        
        return settings
    finally:
        db.close()

def update_user_settings(user_id: int, **kwargs):
    """Update user settings"""
    db = SessionLocal()
    try:
        settings = db.query(UserSettings).filter(UserSettings.user_id == user_id).first()
        
        if not settings:
            settings = UserSettings(user_id=user_id)
            db.add(settings)
        
        # Update provided fields
        for key, value in kwargs.items():
            if hasattr(settings, key):
                setattr(settings, key, value)
        
        settings.updated_at = datetime.utcnow()
        db.commit()
        
        return True, "✅ Einstellungen gespeichert"
    except Exception as e:
        db.rollback()
        logger.exception(f"Error updating settings: {e}")
        return False, f"🔥 Fehler: {str(e)}"
    finally:
        db.close()

def delete_generated_image(img_id: int, user_id: int):
    """Delete a generated image safely"""
    db = get_db()
    try:
        img = db.query(GeneratedImage).filter(GeneratedImage.id == img_id, GeneratedImage.user_id == user_id).first()
        if img:
            if img.image_path and os.path.exists(img.image_path):
                try: os.remove(img.image_path)
                except: pass
            db.delete(img)
            db.commit()
            return True
        return False
    except Exception as e:
        db.rollback()
        logger.exception(f"DB Error: {e}")
        return False
    finally:
        db.close()

def delete_custom_prompt(prompt_id: int, user_id: int):
    """Delete a custom prompt"""
    db = get_db()
    prompt = db.query(CustomPrompt).filter(
        CustomPrompt.id == prompt_id,
        CustomPrompt.user_id == user_id
    ).first()
    if prompt:
        db.delete(prompt)
        db.commit()
        db.close()
        return True
    db.close()
    return False

def get_user_chat_history(user_id: int, limit: int = 50):
    """Get user's chat history"""
    db = get_db()
    results = db.query(ChatHistory).filter(
        ChatHistory.user_id == user_id
    ).order_by(ChatHistory.timestamp.desc()).limit(limit).all()
    db.close()
    return results

def get_user_vision_results(user_id: int, limit: int = 50):
    """Get user's vision results"""
    db = get_db()
    results = db.query(VisionResult).filter(
        VisionResult.user_id == user_id
    ).order_by(VisionResult.timestamp.desc()).limit(limit).all()
    db.close()
    return results

def get_user_generated_images(user_id: int, limit: int = 50):
    """Get user's generated images"""
    db = get_db()
    results = db.query(GeneratedImage).filter(
        GeneratedImage.user_id == user_id
    ).order_by(GeneratedImage.timestamp.desc()).limit(limit).all()
    db.close()
    return results

def ensure_database_schema():
    """Ensure database schema is up to date"""
    try:
        # This creates tables that don't exist
        Base.metadata.create_all(bind=engine)
        
        # Check for UserSettings table specifically
        db = SessionLocal()
        try:
            # Try to query UserSettings - will fail if table doesn't exist
            db.query(UserSettings).first()
            logger.info("✅ Database schema is up to date")
        except Exception as e:
            logger.warning(f"UserSettings table issue: {e}")
            # Force recreation
            UserSettings.__table__.create(bind=engine, checkfirst=True)
            logger.info("✅ Created UserSettings table")
        finally:
            db.close()
            
    except Exception as e:
        logger.exception(f"Database schema check failed: {e}")

def migrate_add_media_manager_column():
    """Add is_media_manager column to existing users table"""
    try:
        from sqlalchemy import inspect, text
        
        inspector = inspect(engine)
        columns = [col['name'] for col in inspector.get_columns('users')]
        
        if 'is_media_manager' not in columns:
            logger.info("Adding is_media_manager column to users table...")
            with engine.connect() as conn:
                conn.execute(text("ALTER TABLE users ADD COLUMN is_media_manager BOOLEAN DEFAULT 0"))
                conn.commit()
            logger.info("✅ Migration complete")
        else:
            logger.info("is_media_manager column already exists")
    except Exception as e:
        logger.error(f"Migration failed: {e}")

# Auto-Migration
ensure_database_schema()
migrate_add_media_manager_column()

# Initialize default users
create_default_users()

# ==========================================
# 🎬 YOUTUBE CHANNEL WHITELIST
# ==========================================

YOUTUBE_CHANNEL_WHITELIST = [
    # Format: Channel-Identifier (wird automatisch normalisiert)
    "AkademieKanal",           # /user/AkademieKanal
    "@theologisches-forum",     # /@theologisches-forum
    "grenzfragen",              # /grenzfragen
    
    # Ihr könnt auch Channel-IDs verwenden (am sichersten):
    # "UCxyz123abc...",         # /channel/UCxyz123abc...
    
    # Weitere genehmigte Kanäle hier hinzufügen
]

def extract_youtube_channel_identifier(url):
    """
    Extract channel identifier from various YouTube URL formats.
    
    Supports:
    - youtube.com/user/USERNAME
    - youtube.com/@HANDLE
    - youtube.com/c/CUSTOMNAME
    - youtube.com/channel/CHANNEL_ID
    - youtube.com/CUSTOMNAME (direct)
    - youtu.be/VIDEO (extracts from video metadata if needed)
    """
    import re
    from urllib.parse import urlparse, parse_qs
    
    try:
        parsed = urlparse(url)
        path = parsed.path
        
        # Pattern 1: /user/USERNAME
        match = re.search(r'/user/([^/?&#]+)', path)
        if match:
            return match.group(1)
        
        # Pattern 2: /@HANDLE
        match = re.search(r'/@([^/?&#]+)', path)
        if match:
            return f"@{match.group(1)}"
        
        # Pattern 3: /c/CUSTOMNAME
        match = re.search(r'/c/([^/?&#]+)', path)
        if match:
            return match.group(1)
        
        # Pattern 4: /channel/CHANNEL_ID
        match = re.search(r'/channel/([^/?&#]+)', path)
        if match:
            return match.group(1)
        
        # Pattern 5: Direct custom URL (e.g., /grenzfragen)
        match = re.search(r'/([^/?&#]+)$', path)
        if match and match.group(1) not in ['watch', 'embed', 'v', 'shorts']:
            return match.group(1)
        
        # Pattern 6: Video URL - try to get channel from API
        if '/watch' in path or 'youtu.be' in parsed.netloc:
            try:
                import yt_dlp
                
                cookie_file = "/var/www/transkript_app/cookies.txt"
                
                ydl_opts = {
                    'quiet': True,
                    'no_warnings': True,
                    'extract_flat': True,
                    # WICHTIG: User-Agent setzen, damit wir wie ein Browser aussehen
                    'user_agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                }
                
                if os.path.exists(cookie_file):
                    ydl_opts['cookiefile'] = cookie_file

                with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                    info = ydl.extract_info(url, download=False)
                    
                    channel_id = info.get('channel_id')
                    if channel_id: return channel_id
                    
                    uploader_id = info.get('uploader_id')
                    if uploader_id: return uploader_id
                    
                    channel_url = info.get('channel_url', '')
                    if channel_url:
                        return extract_youtube_channel_identifier(channel_url)
                        
            except Exception as e:
                logger.warning(f"Could not extract channel from video URL: {e}")
        
        return None
        
    except Exception as e:
        logger.error(f"Error extracting channel identifier: {e}")
        return None

def is_youtube_channel_whitelisted(url):
    """Check if YouTube URL is from a whitelisted channel"""
    if not is_youtube_url(url):
        return True  # Non-YouTube URLs are allowed by default (can be changed)
    
    channel_id = extract_youtube_channel_identifier(url)
    
    if not channel_id:
        logger.warning(f"Could not extract channel from URL: {url}")
        return False
    
    # Normalize for comparison
    channel_id_lower = channel_id.lower()
    
    for whitelisted in YOUTUBE_CHANNEL_WHITELIST:
        whitelisted_lower = whitelisted.lower()
        
        # Direct match or partial match (for handles with/without @)
        if (channel_id_lower == whitelisted_lower or 
            channel_id_lower == f"@{whitelisted_lower}" or
            f"@{channel_id_lower}" == whitelisted_lower):
            logger.info(f"✅ Channel '{channel_id}' is whitelisted")
            return True
    
    logger.warning(f"❌ Channel '{channel_id}' not in whitelist")
    return False

def download_from_url(url, download_video=False, save_to_storage=True, user_state=None):
    """
    Download audio/video from URL with role-based permissions and whitelist.
    
    Permissions:
    - Admin: All URLs
    - Media Manager: All URLs
    - Regular User: Only whitelisted YouTube channels + non-YouTube URLs

    ⚠️ RECHTLICHER HINWEIS:
    Diese Funktion darf nur für eigene oder lizenzierte Inhalte verwendet werden.
    Der Nutzer trägt die volle Verantwortung für die Einhaltung der Urheberrechte.
    """

    import subprocess
    import requests
    from urllib.parse import urlparse
    
    # --- AUTHENTICATION CHECK ---
    if not user_state or not user_state.get("id"):
        return False, None, "❌ Nicht authentifiziert"
    
    username = user_state.get("username", "unknown")
    is_admin = user_state.get("is_admin", False)
    is_media_manager = user_state.get("is_media_manager", False)
    
    # --- PERMISSION CHECK ---
    if is_youtube_url(url):
        # Admins and Media Managers can download everything
        if not (is_admin or is_media_manager):
            # Regular users: Check whitelist
            if not is_youtube_channel_whitelisted(url):
                channel_id = extract_youtube_channel_identifier(url)
                logger.warning(f"🚫 BLOCKED: User '{username}' tried to download from non-whitelisted channel: {channel_id}")
                return False, None, (
                    f"❌ Kanal nicht freigegeben.\n\n"
                    f"Dieser YouTube-Kanal ist nicht in der Whitelist.\n"
                    f"Erkannter Kanal: {channel_id or 'unbekannt'}\n\n"
                    f"Freigegebene Kanäle:\n" + 
                    "\n".join(f"  • {ch}" for ch in YOUTUBE_CHANNEL_WHITELIST[:5]) +
                    f"\n\n💡 Kontaktiere einen Administrator oder Medienverwalter."
                )
            
            # Whitelisted channel - log and proceed
            logger.info(f"✅ ALLOWED: User '{username}' downloading from whitelisted channel")
        else:
            # Admin/Media Manager - log but allow
            role = "Admin" if is_admin else "Medienverwalter"
            logger.info(f"✅ ALLOWED: {role} '{username}' downloading: {url}")
    
    # --- AUDIT LOG ---
    logger.warning(f"📥 DOWNLOAD: User='{username}' | Role={'Admin' if is_admin else 'MediaMgr' if is_media_manager else 'User'} | URL={url} | Video={download_video}")
    
    # Determine output directory
    if save_to_storage and user_state.get("username"):
        output_dir = get_user_storage_path(user_state["username"])
        if not output_dir:
            output_dir = tempfile.gettempdir()
            logger.warning("Storage path unavailable, using temp")
    else:
        output_dir = tempfile.gettempdir()
    
    os.makedirs(output_dir, exist_ok=True)
    
    # --- METHOD 1: YT-DLP (Best for YouTube/streaming) ---
    if is_youtube_url(url):
        logger.info(f"Attempting yt-dlp download for YouTube: {url}")
        result = download_with_yt_dlp(url, output_dir, download_video)
        
        if result and os.path.exists(result):
            size_mb = os.path.getsize(result) / (1024*1024)
            return True, result, f"✅ yt-dlp hat heruntergeladen: ({size_mb:.1f} MB)"
        else:
            # STOP HERE! Do not fall back to FFmpeg/Direct for YouTube URLs
            # because they will just download the 'Sign In' HTML page.
            return False, None, "❌ YouTube-Download fehlgeschlagen (Cookie/Login erforderlich)"
    
    # 2. Non-YouTube Handling (Flexible)
    # Try yt-dlp first (it handles Twitch, Vimeo, direct files too)
    try:
        result = download_with_yt_dlp(url, output_dir, download_video)
        if result and os.path.exists(result):
            return True, result, "✅ yt-dlp erfolg"
    except: pass

    # Fallback 1: FFmpeg (Good for direct streams)
    if not download_video:
        try:
            logger.info(f"Attempting ffmpeg download: {url}")
            result = download_with_ffmpeg(url, output_dir)
            if result and os.path.exists(result):
                return True, result, "✅ ffmpeg erfolg"
        except Exception as e:
            logger.warning(f"ffmpeg failed: {e}")

    # Fallback 2: Direct HTTP (Only for actual files)
    try:
        logger.info(f"Attempting direct HTTP download: {url}")
        result = download_with_requests(url, output_dir)
        if result and os.path.exists(result):
            # Sanity check: Ensure it's not a tiny HTML error page
            if os.path.getsize(result) < 5000:
                with open(result, 'r', errors='ignore') as f:
                    head = f.read(500)
                    if "<!DOCTYPE html" in head or "<html" in head:
                         os.remove(result)
                         return False, None, "❌ Fehler: URL lieferte HTML statt Audio"
            
            return True, result, "✅ HTTP download erfolg"
    except Exception as e:
        logger.warning(f"HTTP download failed: {e}")

    return False, None, "❌ Alle Methoden fehlgeschlagen"

# ==========================================
# 📦 STORAGE BOX HELPER
# ==========================================
STORAGE_MOUNT_POINT = "/mnt/akademie_storage"

def get_user_storage_path(username):
    """
    Returns the user-specific storage path.
    Creates directory if it doesn't exist.
    """
    if not username:
        logger.warning("get_user_storage_path called without username")
        return None
        
    user_path = os.path.join(STORAGE_MOUNT_POINT, username)
    
    # Create directory with proper permissions
    try:
        os.makedirs(user_path, exist_ok=True)
        logger.info(f"✅ User storage path ready: {user_path}")
    except Exception as e:
        logger.error(f"❌ Failed to create storage path for {username}: {e}")
        return None
    
    return user_path

def ensure_user_storage_dirs(username):
    """
    Creates user-specific storage directory.
    Called at login time.
    """
    if not username:
        logger.warning("ensure_user_storage_dirs called without username")
        return False
    
    user_path = get_user_storage_path(username)
    
    if user_path and os.path.exists(user_path):
        logger.info(f"✅ Storage directory verified for user '{username}': {user_path}")
        return True
    
    logger.error(f"❌ Storage directory creation failed for user '{username}'")
    return False


def copy_storage_file_to_temp(file_path):
    """
    Copies a file from Storage Box to local temp for processing.
    Returns path to local temp file.
    """
    if not file_path: return None
    
    # Security check: Ensure path is within mount point
    abs_path = os.path.abspath(file_path)
    if not abs_path.startswith(os.path.abspath(STORAGE_MOUNT_POINT)):
        raise ValueError("Zugriff verweigert: Datei liegt außerhalb der Storage Box")
        
    if not os.path.exists(abs_path):
        raise ValueError("Datei existiert nicht")
        
    # Copy to temp
    filename = os.path.basename(abs_path)
    temp_dir = tempfile.gettempdir()
    local_dest = os.path.join(temp_dir, f"sb_{int(time.time())}_{filename}")
    
    import shutil
    shutil.copy2(abs_path, local_dest)
    return local_dest

def get_storage_root(user_state=None):
    """
    Returns the user-specific storage root path.
    NO LONGER returns global mount point.
    """
    if not user_state or not user_state.get("username"):
        logger.warning("get_storage_root called without valid user_state")
        return None
    
    username = user_state.get("username")
    user_path = get_user_storage_path(username)
    
    logger.debug(f"get_storage_root for '{username}': {user_path}")
    return user_path

# ==========================================
# AUDIO HELPERS
# ==========================================

import shutil

JOB_STATE_DIR = "/var/www/transkript_app/jobs"
os.makedirs(JOB_STATE_DIR, exist_ok=True)

def create_job_manifest(job_id, audio_path, provider, model, chunks, lang, prompt, temp, user_id):
    """Save job state to disk with User ID"""
    manifest = {
        "job_id": job_id,
        "user_id": user_id,  # Track ownership
        "audio_path": audio_path,
        "provider": provider,
        "model": model,
        "lang": lang,
        "prompt": prompt,
        "temp": temp,
        "created_at": time.time(),
        "chunks": [{"path": p, "status": "pending"} for p in chunks],
        "transcript_parts": [""] * len(chunks)
    }
    with open(os.path.join(JOB_STATE_DIR, f"{job_id}.json"), "w") as f:
        json.dump(manifest, f)
    return manifest

def update_job_chunk_status(job_id, chunk_index, status, text=None):
    """Update status of a specific chunk"""
    path = os.path.join(JOB_STATE_DIR, f"{job_id}.json")
    if not os.path.exists(path): return
    
    with open(path, "r") as f:
        manifest = json.load(f)
    
    manifest["chunks"][chunk_index]["status"] = status
    if text is not None:
        manifest["transcript_parts"][chunk_index] = text
        
    with open(path, "w") as f:
        json.dump(manifest, f)

def get_failed_jobs(user_id):
    """List incomplete jobs filtered by User ID"""
    jobs = []
    if not os.path.exists(JOB_STATE_DIR): return []
    
    for f in os.listdir(JOB_STATE_DIR):
        if f.endswith(".json"):
            try:
                with open(os.path.join(JOB_STATE_DIR, f), "r") as jf:
                    data = json.load(jf)
                    
                    # Check User ID match
                    if data.get("user_id") != user_id:
                        continue
                        
                    # Check if any chunk is not 'done'
                    pending = sum(1 for c in data["chunks"] if c["status"] != "done")
                    if pending > 0:
                        jobs.append([
                            data["job_id"], 
                            datetime.fromtimestamp(data["created_at"]).strftime("%Y-%m-%d %H:%M"),
                            f"{pending}/{len(data['chunks'])} offen",
                            data["provider"]
                        ])
            except: pass
    return jobs

def get_file_hash(filepath):
    """Generate SHA256 hash for file to track progress"""
    sha256_hash = hashlib.sha256()
    with open(filepath, "rb") as f:
        # Read and update hash string value in blocks of 4K
        for byte_block in iter(lambda: f.read(4096), b""):
            sha256_hash.update(byte_block)
    return sha256_hash.hexdigest()

def cleanup_chunks(chunk_dir):
    """Safely remove temporary chunk directory"""
    try:
        if chunk_dir and os.path.exists(chunk_dir):
            shutil.rmtree(chunk_dir)
            logger.info(f"Cleaned up chunk directory: {chunk_dir}")
    except Exception as e:
        logger.warning(f"Failed to cleanup chunks: {e}")

def split_audio_into_chunks(audio_path, chunk_minutes=10):
    """
    Splits audio into segments to bypass API size/duration limits.
    Returns: (List of file paths, temp_directory_path)
    """
    try:
        # --- FIX: Ensure imports exist ---
        import math 
        from pydub import AudioSegment
        
        # Create unique temp dir
        job_id = int(time.time())
        chunk_dir = os.path.join(tempfile.gettempdir(), f"transcription_chunks_{job_id}")
        os.makedirs(chunk_dir, exist_ok=True)
        
        logger.info(f"Loading audio for splitting: {audio_path}")
        audio = AudioSegment.from_file(audio_path)
        
        duration_ms = len(audio)
        chunk_length_ms = chunk_minutes * 60 * 1000
        total_chunks = math.ceil(duration_ms / chunk_length_ms)
        
        logger.info(f"Audio duration: {duration_ms/1000/60:.2f} min. Splitting into {total_chunks} chunks.")
        
        chunk_paths = []
        
        for i in range(total_chunks):
            start_ms = i * chunk_length_ms
            end_ms = min((i + 1) * chunk_length_ms, duration_ms)
            
            chunk = audio[start_ms:end_ms]
            
            # Export with low compression to keep size down but quality up
            chunk_filename = os.path.join(chunk_dir, f"chunk_{i:03d}.mp3")
            chunk.export(chunk_filename, format="mp3", bitrate="128k")
            chunk_paths.append(chunk_filename)
            
        return chunk_paths, chunk_dir
        
    except Exception as e:
        logger.exception(f"Error splitting audio: {e}")
        return None, None
    
def split_audio(filepath, chunk_length_ms=600000): # 10 minutes default (safe for Mistral 15m limit)
    """Split audio into chunks and return list of temp file paths"""
    try:
        audio = AudioSegment.from_file(filepath)
        chunks = []
        duration_ms = len(audio)
        
        # Create temp dir for chunks
        chunk_dir = os.path.join(tempfile.gettempdir(), "audio_chunks")
        os.makedirs(chunk_dir, exist_ok=True)
        
        base_name = os.path.splitext(os.path.basename(filepath))[0]
        
        for i, start_ms in enumerate(range(0, duration_ms, chunk_length_ms)):
            end_ms = min(start_ms + chunk_length_ms, duration_ms)
            chunk = audio[start_ms:end_ms]
            
            # Export chunk
            chunk_path = os.path.join(chunk_dir, f"{base_name}_part_{i}.mp3")
            chunk.export(chunk_path, format="mp3")
            chunks.append(chunk_path)
            
        return chunks
    except Exception as e:
        logger.error(f"Error splitting audio: {e}")
        return None

# ==========================================
# URL DOWNLOAD HELPERS
# ==========================================

def is_youtube_url(url):
    """Check if URL is from YouTube"""
    youtube_patterns = [
        r'(youtube\.com/watch\?v=)',
        r'(youtu\.be/)',
        r'(youtube\.com/embed/)',
        r'(youtube\.com/v/)'
    ]
    return any(re.search(pattern, url, re.IGNORECASE) for pattern in youtube_patterns)

def download_with_yt_dlp(url, output_dir, keep_video=False):
    """Download with yt-dlp using cookies, Node.js and optimized client settings"""
    try:
        import yt_dlp
        
        cookie_file = "/var/www/transkript_app/cookies.txt"
        
        ydl_opts = {
            'outtmpl': os.path.join(output_dir, '%(title)s.%(ext)s'),
            'quiet': False,
            'no_warnings': False,
            'nocheckcertificate': True,
            'ignoreerrors': True,
            'no_color': True,
            
            # WICHTIG: 'js_runtimes' Zeile LÖSCHEN! 
            # yt-dlp findet das installierte Node.js automatisch.
            
            # iOS Client erzwingen (beste Umgehung für Bot-Sperre derzeit)
            'extractor_args': {'youtube': {'player_client': ['ios']}},
        }

        # Cookies laden
        if os.path.exists(cookie_file):
            logger.info(f"🍪 Loading cookies from: {cookie_file}")
            ydl_opts['cookiefile'] = cookie_file
        else:
            logger.warning("⚠️ No cookies.txt found!")

        # Format Einstellungen
        if keep_video:
            ydl_opts.update({
                'format': 'bestvideo+bestaudio/best',
                'merge_output_format': 'mp4',
            })
        else:
            ydl_opts.update({
                'format': 'bestaudio/best',
                'postprocessors': [{
                    'key': 'FFmpegExtractAudio',
                    'preferredcodec': 'mp3',
                    'preferredquality': '192',
                }],
            })
        
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            
            if not info:
                logger.error("❌ yt-dlp failed: Download blocked or invalid URL.")
                return None

            title = info.get('title', 'video_download')
            safe_title = re.sub(r'[<>:"/\\|?*\u0000-\u001F\u007F-\u009F]', '_', title)
            ext = 'mp4' if keep_video else 'mp3'
            
            # Datei suchen
            expected_file = os.path.join(output_dir, f"{safe_title}.{ext}")
            if os.path.exists(expected_file):
                return expected_file
            
            # Fallback Suche
            import glob
            search_pattern = os.path.join(output_dir, f"{safe_title}*.{ext}")
            found = glob.glob(search_pattern)
            if found:
                return found[0]

        return None
        
    except Exception as e:
        logger.error(f"yt-dlp error: {e}")
        return None

def download_with_ffmpeg(url, output_dir):
    """Download audio using ffmpeg directly"""
    try:
        import subprocess
        
        # Generate output filename
        output_file = os.path.join(output_dir, f"download_{int(time.time())}.mp3")
        
        cmd = [
            'ffmpeg',
            '-i', url,
            '-vn',  # No video
            '-acodec', 'libmp3lame',
            '-ab', '192k',
            '-y',  # Overwrite
            output_file
        ]
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            timeout=300,  # 5 minutes
            text=True
        )
        
        if result.returncode == 0 and os.path.exists(output_file):
            return output_file
        
        logger.warning(f"ffmpeg stderr: {result.stderr[:500]}")
        return None
        
    except FileNotFoundError:
        logger.error("ffmpeg not found in PATH")
        return None
    except subprocess.TimeoutExpired:
        logger.error("ffmpeg timeout")
        return None
    except Exception as e:
        logger.error(f"ffmpeg error: {e}")
        return None

def download_with_requests(url, output_dir):
    """Simple HTTP download for direct file URLs"""
    try:
        import requests
        from urllib.parse import urlparse, unquote
        
        # Get filename from URL
        parsed = urlparse(url)
        filename = os.path.basename(unquote(parsed.path))
        
        if not filename or '.' not in filename:
            filename = f"download_{int(time.time())}.mp3"
        
        output_file = os.path.join(output_dir, filename)
        
        response = requests.get(url, stream=True, timeout=60)
        response.raise_for_status()
        
        with open(output_file, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        
        if os.path.exists(output_file) and os.path.getsize(output_file) > 0:
            return output_file
        
        return None
        
    except Exception as e:
        logger.error(f"HTTP download error: {e}")
        return None

# ==========================================
# 👥 USER MANAGEMENT FUNCTIONS
# ==========================================

def create_user(username, password, email, is_admin=False):
    """Create a new user with their own Encryption Keychain"""
    db = SessionLocal()
    try:
        if db.query(User).filter(User.username == username).first():
            return False, "❌ Benutzername existiert bereits"

        # 1. Generate Keychain (Returns: salt, encrypted_master_key, decrypted_master_key)
        keychain = key_wrapper.create_user_keychain(password)

        new_user = User(
            username=username,
            password_hash=hash_password(password),
            email=email,
            is_admin=is_admin,
            # Store the Lockbox
            salt=keychain["salt"],
            encrypted_master_key=keychain["encrypted_master_key"]
        )
        db.add(new_user)
        db.commit()
        return True, f"✅ Benutzer '{username}' erstellt (Secure Keychain Active)"
    except Exception as e:
        db.rollback()
        logger.exception(f"Error creating user: {str(e)}")
        return False, f"🔥 Fehler: {str(e)}"
    finally:
        db.close()

def delete_user(user_id, current_user_id):
    """Delete a user (cannot delete self)"""
    db = SessionLocal()
    try:
        if user_id == current_user_id:
            return False, "❌ Du kannst dich nicht selbst löschen"
        
        user = db.query(User).filter(User.id == user_id).first()
        if not user:
            return False, "❌ Benutzer nicht gefunden"
        
        username = user.username
        db.delete(user)
        db.commit()
        
        logger.info(f"User deleted: {username} (ID: {user_id})")
        return True, f"✅ Benutzer '{username}' gelöscht"
    except Exception as e:
        db.rollback()
        logger.exception(f"Error deleting user: {str(e)}")
        return False, f"🔥 Fehler: {str(e)}"
    finally:
        db.close()

def rename_user(user_id, new_username):
    """Rename a user"""
    db = SessionLocal()
    try:
        # Check if new username already exists
        existing = db.query(User).filter(User.username == new_username).first()
        if existing and existing.id != user_id:
            return False, "❌ Benutzername existiert bereits"
        
        user = db.query(User).filter(User.id == user_id).first()
        if not user:
            return False, "❌ Benutzer nicht gefunden"
        
        old_username = user.username
        user.username = new_username
        db.commit()
        
        logger.info(f"User renamed: {old_username} → {new_username}")
        return True, f"✅ Benutzer umbenannt: '{old_username}' → '{new_username}'"
        
    except Exception as e:
        db.rollback()
        logger.exception(f"Error renaming user: {str(e)}")
        return False, f"🔥 Fehler: {str(e)}"
    finally:
        db.close()

def reset_user_password(user_id, new_password, admin_aware=False):
    """Admin password reset - WARNING: User data will become inaccessible!"""
    if not admin_aware:
        return False, "⚠️ WARNUNG: Passwort-Reset macht verschlüsselte Daten unzugänglich. Nutzer muss Daten neu erstellen."
    
    db = SessionLocal()
    try:
        user = db.query(User).filter(User.id == user_id).first()
        if not user:
            return False, "❌ Benutzer nicht gefunden"
        
        # Generate NEW keychain (old data becomes inaccessible)
        new_keychain = key_wrapper.create_user_keychain(new_password)
        
        user.password_hash = hash_password(new_password)
        user.salt = new_keychain["salt"]
        user.encrypted_master_key = new_keychain["encrypted_master_key"]
        
        db.commit()
        
        return True, "✅ Passwort zurückgesetzt. WICHTIG: Alte Daten sind nicht mehr zugänglich!"
    except Exception as e:
        db.rollback()
        logger.exception(f"Error resetting password: {str(e)}")
        return False, f"🔥 Fehler: {str(e)}"
    finally:
        db.close()

def toggle_admin_status(user_id, current_user_id):
    """Toggle admin status for a user (cannot change own status)"""
    db = SessionLocal()
    try:
        # Prevent changing own admin status
        if user_id == current_user_id:
            return False, "❌ Du kannst deinen eigenen Admin-Status nicht ändern"
        
        user = db.query(User).filter(User.id == user_id).first()
        if not user:
            return False, "❌ Benutzer nicht gefunden"
        
        # Toggle admin status
        user.is_admin = not user.is_admin
        db.commit()
        
        status = "Admin" if user.is_admin else "Normaler Benutzer"
        logger.info(f"Admin status changed for {user.username}: {status}")
        return True, f"✅ '{user.username}' ist jetzt: {status}"
        
    except Exception as e:
        db.rollback()
        logger.exception(f"Error toggling admin status: {str(e)}")
        return False, f"🔥 Fehler: {str(e)}"
    finally:
        db.close()

def toggle_media_manager_status(user_id, current_user_id):
    """Toggle media manager status for a user (only admins can do this)"""
    db = SessionLocal()
    try:
        # Prevent changing own media manager status
        if user_id == current_user_id:
            return False, "❌ Du kannst deinen eigenen Medienverwalter-Status nicht ändern"
        
        user = db.query(User).filter(User.id == user_id).first()
        if not user:
            return False, "❌ Benutzer nicht gefunden"
        
        # Toggle status
        user.is_media_manager = not user.is_media_manager
        db.commit()
        
        status = "Medienverwalter" if user.is_media_manager else "Normaler Benutzer"
        logger.info(f"Media manager status changed for {user.username}: {status}")
        return True, f"✅ '{user.username}' ist jetzt: {status}"
        
    except Exception as e:
        db.rollback()
        logger.exception(f"Error toggling media manager status: {str(e)}")
        return False, f"🔥 Fehler: {str(e)}"
    finally:
        db.close()

def get_all_users():
    """Get all users for admin panel (updated for media manager column)"""
    db = SessionLocal()
    try:
        users = db.query(User).order_by(User.created_at.desc()).all()
        data = []
        for u in users:
            role = "👑 Admin" if u.is_admin else ("🎬 Medienverwalter" if u.is_media_manager else "👤 User")
            data.append([
                u.id,
                u.username,
                u.email or "—",
                role,  # Updated to show media manager
                u.created_at.strftime("%Y-%m-%d %H:%M")
            ])
        return data
    except Exception as e:
        logger.exception(f"Error getting users: {str(e)}")
        return []
    finally:
        db.close()

def get_all_users():
    """Get all users for admin panel"""
    db = SessionLocal()
    try:
        users = db.query(User).order_by(User.created_at.desc()).all()
        data = []
        for u in users:
            data.append([
                u.id,
                u.username,
                u.email or "—",
                "✅ Admin" if u.is_admin else "👤 User",
                u.created_at.strftime("%Y-%m-%d %H:%M")
            ])
        return data
    except Exception as e:
        logger.exception(f"Error getting users: {str(e)}")
        return []
    finally:
        db.close()

def update_user_email(user_id, new_email):
    """Update user's email address"""
    db = SessionLocal()
    try:
        # Check if email already exists
        if new_email:
            existing = db.query(User).filter(User.email == new_email).first()
            if existing and existing.id != user_id:
                return False, "❌ E-Mail wird bereits verwendet"
        
        user = db.query(User).filter(User.id == user_id).first()
        if not user:
            return False, "❌ Benutzer nicht gefunden"
        
        user.email = new_email
        db.commit()
        
        logger.info(f"Email updated for user: {user.username}")
        return True, f"✅ E-Mail für '{user.username}' aktualisiert"
        
    except Exception as e:
        db.rollback()
        logger.exception(f"Error updating email: {str(e)}")
        return False, f"🔥 Fehler: {str(e)}"
    finally:
        db.close()


# ==========================================
# 🎯 MODEL PREFERENCES MANAGEMENT
# ==========================================

def fetch_available_models(provider, api_key=None):
    """
    Ruft verfügbare Modelle ab. 
    Für Poe: Nutzt die offizielle API und sortiert nach Modalitäten (Text, Bild, etc.).
    """
    try:
        if provider == "Poe":
            api_key = api_key or API_KEYS.get("POE")
            if not api_key:
                return None, "API Key erforderlich"

            try:
                # Direkter Request an die Poe API, da wir die Metadaten brauchen
                headers = {"Authorization": f"Bearer {api_key}"}
                response = requests.get("https://api.poe.com/v1/models", headers=headers, timeout=15)
                
                if response.status_code == 200:
                    data = response.json()
                    model_list = []
                    
                    # Wir kategorisieren die Modelle für die Rückgabe
                    # (Standardmäßig geben wir alle zurück, aber fügen Infos hinzu)
                    for model in data.get("data", []):
                        arch = model.get("architecture", {})
                        inputs = arch.get("input_modalities", [])
                        outputs = arch.get("output_modalities", [])
                        
                        # Bestimme den Typ für die Anzeige
                        type_label = "Chat"
                        if "image" in outputs:
                            type_label = "Bild-Gen"
                        elif "video" in outputs:
                            type_label = "Video-Gen"
                        elif "audio" in outputs:
                            type_label = "Audio-Gen"
                        elif "image" in inputs and "text" in outputs:
                            type_label = "Vision/Chat"
                        
                        model_list.append({
                            "id": model["id"],
                            "name": f"{model.get('metadata', {}).get('display_name', model['id'])} ({type_label})",
                            "type": type_label, # Hilfreich für Filterung später
                            "context_length": model.get("context_length", 0)
                        })
                    
                    # Sortieren nach Name
                    model_list.sort(key=lambda x: x["name"])
                    return model_list, None
                else:
                    return None, f"Poe API Fehler: {response.status_code} - {response.text}"
            except Exception as e:
                return None, f"Poe Verbindungsfehler: {str(e)}"

        elif provider == "Scaleway":
            # Scaleway uses OpenAI-compatible API
            api_key = api_key or API_KEYS.get("SCALEWAY")
            if not api_key:
                return None, "API Key erforderlich"
            
            try:
                client = openai.OpenAI(base_url="https://api.scaleway.ai/v1", api_key=api_key)
                models = client.models.list()
                model_list = []
                for model in models.data:
                    model_list.append({
                        "id": model.id,
                        "name": model.id,
                        "owned_by": getattr(model, 'owned_by', 'scaleway')
                    })
                return model_list, None
            except Exception as e:
                return None, f"Fehler: {str(e)}"
        
        elif provider == "Nebius":
            # Nebius uses OpenAI-compatible API
            api_key = api_key or API_KEYS.get("NEBIUS")
            if not api_key:
                return None, "API Key erforderlich"
            
            try:
                client = openai.OpenAI(base_url="https://api.tokenfactory.nebius.com/v1", api_key=api_key)
                models = client.models.list()
                model_list = []
                for model in models.data:
                    model_list.append({
                        "id": model.id,
                        "name": model.id,
                        "owned_by": getattr(model, 'owned_by', 'nebius')
                    })
                return model_list, None
            except Exception as e:
                return None, f"Fehler: {str(e)}"
        
        elif provider == "Mistral":
            # Mistral API
            api_key = api_key or API_KEYS.get("MISTRAL")
            if not api_key:
                return None, "API Key erforderlich"
            
            try:
                client = openai.OpenAI(base_url="https://api.mistral.ai/v1", api_key=api_key)
                models = client.models.list()
                model_list = []
                for model in models.data:
                    model_list.append({
                        "id": model.id,
                        "name": model.id,
                        "owned_by": getattr(model, 'owned_by', 'mistral')
                    })
                return model_list, None
            except Exception as e:
                return None, f"Fehler: {str(e)}"
        
        elif provider == "OpenRouter":
            # OpenRouter models list
            try:
                response = requests.get("https://openrouter.ai/api/v1/models", timeout=10)
                if response.status_code == 200:
                    data = response.json()
                    model_list = []
                    for model in data.get("data", []):
                        model_list.append({
                            "id": model["id"],
                            "name": model.get("name", model["id"]),
                            "context_length": model.get("context_length", 0)
                        })
                    return model_list, None
                return None, f"Fehler: Status {response.status_code}"
            except Exception as e:
                return None, f"Fehler: {str(e)}"
        
        elif provider == "Groq":
            # Groq API - use REST endpoint instead of client library
            api_key = api_key or API_KEYS.get("GROQ")
            if not api_key:
                return None, "API Key erforderlich"
            
            try:
                headers = {
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                }
                response = requests.get("https://api.groq.com/openai/v1/models", headers=headers, timeout=10)
                
                if response.status_code == 200:
                    data = response.json()
                    model_list = []
                    for model in data.get("data", []):
                        model_list.append({
                            "id": model["id"],
                            "name": model["id"],
                            "owned_by": model.get("owned_by", "groq")
                        })
                    return model_list, None
                return None, f"Fehler: Status {response.status_code}"
            except Exception as e:
                return None, f"Fehler: {str(e)}"
        
        else:
            # For providers without /models endpoint, return configured models
            provider_config = PROVIDERS.get(provider, {})
            chat_models = provider_config.get("chat_models", [])
            model_list = [{"id": m, "name": m} for m in chat_models]
            return model_list, None
    
    except Exception as e:
        logger.exception(f"Error fetching models for {provider}: {str(e)}")
        return None, f"Fehler: {str(e)}"

def get_user_model_preferences(user_id, provider):
    """Get user's model preferences for a provider"""
    db = SessionLocal()
    try:
        prefs = db.query(UserModelPreference).filter(
            UserModelPreference.user_id == user_id,
            UserModelPreference.provider == provider
        ).order_by(UserModelPreference.display_order).all()
        
        return prefs
    except Exception as e:
        logger.exception(f"Error getting user model preferences: {str(e)}")
        return []
    finally:
        db.close()

def save_user_model_preferences(user_id, provider, model_configs):
    """
    Save user's model preferences
    model_configs: list of dicts with keys: model_id, display_name, is_visible, display_order
    """
    db = SessionLocal()
    try:
        # Delete existing preferences for this provider
        db.query(UserModelPreference).filter(
            UserModelPreference.user_id == user_id,
            UserModelPreference.provider == provider
        ).delete()
        
        # Add new preferences
        for config in model_configs:
            pref = UserModelPreference(
                user_id=user_id,
                provider=provider,
                model_id=config["model_id"],
                display_name=config.get("display_name", config["model_id"]),
                display_order=config.get("display_order", 0),
                is_visible=config.get("is_visible", True)
            )
            db.add(pref)
        
        db.commit()
        logger.info(f"Saved {len(model_configs)} model preferences for user {user_id}, provider {provider}")
        return True, "✅ Einstellungen gespeichert"
        
    except Exception as e:
        db.rollback()
        logger.exception(f"Error saving model preferences: {str(e)}")
        return False, f"🔥 Fehler: {str(e)}"
    finally:
        db.close()

def get_user_visible_models(user_id, provider):
    """Get list of visible models for user, in order"""
    prefs = get_user_model_preferences(user_id, provider)
    
    if not prefs:
        # No preferences set, return all models from provider config
        provider_config = PROVIDERS.get(provider, {})
        return provider_config.get("chat_models", [])
    
    # Return visible models in order
    visible_models = [p.model_id for p in prefs if p.is_visible]
    return visible_models if visible_models else [prefs[0].model_id] if prefs else []

def get_default_model_for_user(user_id, provider):
    """Get user's default model (first in order)"""
    visible = get_user_visible_models(user_id, provider)
    return visible[0] if visible else None
        
# ==========================================
# ⚙️ ZENTRALE KONFIGURATION
# ==========================================

DEFAULT_CHAT_PROVIDER = "Mistral"
DEFAULT_CHAT_MODEL = "mistral-large-latest"

# API Keys (aus Environment oder direkt hier eintragen)
API_KEYS = {
    "GLADIA": os.environ.get("GLADIA_API_KEY", "your_key"),
    "SCALEWAY": os.environ.get("SCALEWAY_API_KEY", "your_key"),
    "NEBIUS": os.environ.get("NEBIUS_API_KEY", "your_key"),
    "MISTRAL": os.environ.get("MISTRAL_API_KEY", "your_key"),
    "OPENROUTER": os.environ.get("OPENROUTER_API_KEY", "your_key"),
    "GROQ": os.environ.get("GROQ_API_KEY", "your_key"),
    "POE": os.environ.get("POE_API_KEY", "your_poe_key_here"),
    "DEEPGRAM": os.environ.get("DEEPGRAM_API_KEY", "your_key"), 
    "ASSEMBLYAI": os.environ.get("ASSEMBLYAI_API_KEY", "your_key"),
    "OPENAI": os.environ.get("OPENAI_API_KEY", "your_key"),
    "COHERE": os.environ.get("COHERE_API_KEY", "your_key"),
    "TOGETHER": os.environ.get("TOGETHER_API_KEY", "your_key"),
    "OVH": os.environ.get("OVH_API_KEY", "your_key"),
    "CEREBRAS": os.environ.get("CEREBRAS_API_KEY", "your_key"),
    "GOOGLEAI": os.environ.get("GOOGLEAI_API_KEY", "your_key"),
    "ANTHROPIC": os.environ.get("ANTHROPIC_API_KEY", "your_key"),
}

# Provider-Datenbank (Modelle, Endpoints, Compliance)
PROVIDERS = {
    "Scaleway": {
        "base_url": "https://api.scaleway.ai/v1",
        "key_name": "SCALEWAY",
        "badge": "🇫🇷 <b>DSGVO-Konform</b>",
        "chat_models": [
            "gpt-oss-120b", 
            "mistral-small-3.2-24b-instruct-2506", 
            "gemma-3-27b-it", 
            "qwen3-235b-a22b-instruct-2507", 
            "llama-3.3-70b-instruct", 
            "deepseek-r1-distill-llama-70b"
        ],
        "vision_models": ["pixtral-12b-2409", "mistral-small-3.1-24b-instruct-2503"],
        "audio_models": ["whisper-large-v3"],
        "image_models": ["pixtral-12b-2409"],
        "context_limits": {
            "gpt-oss-120b": 32768,
            "mistral-small-3.2-24b-instruct-2506": 32768,
            "gemma-3-27b-it": 96000,
            "qwen3-235b-a22b-instruct-2507": 131072,
            "llama-3.3-70b-instruct": 131072,
            "deepseek-r1-distill-llama-70b": 8192,
            "pixtral-12b-2409": 32768,
            "mistral-small-3.1-24b-instruct-2503": 96000,
            "whisper-large-v3": 16384,
        }
    },
    
    "Gladia": {
        "base_url": "https://api.gladia.io/v2",
        "key_name": "GLADIA",
        "badge": "🇫🇷 <b>DSGVO-Konform</b>",
        "audio_models": ["gladia-v2"],
        "context_limits": {
            "gladia-v2": 1000000 
        }
    },
    
    "Nebius": {
        "base_url": "https://api.tokenfactory.nebius.com/v1",
        "key_name": "NEBIUS",
        "badge": "🇪🇺 <b>DSGVO-Konform</b>",
        "chat_models": [
            "deepseek-ai/DeepSeek-R1-0528",
            "nvidia/Llama-3_1-Nemotron-Ultra-253B-v1",
            "openai/gpt-oss-120b",
            "moonshotai/Kimi-K2-Instruct",
            "moonshotai/Kimi-K2-Thinking",
            "zai-org/GLM-4.5",
            "meta-llama/Llama-3.3-70B-Instruct"
        ],
        "image_models": ["black-forest-labs/flux-schnell", "black-forest-labs/flux-dev"],
        "vision_models": ["google/gemma-3-27b-it", "Qwen/Qwen2.5-VL-72B-Instruct", "nvidia/Nemotron-Nano-V2-12b"],
        "context_limits": {
            "deepseek-ai/DeepSeek-R1-0528": 163840,
            "nvidia/Llama-3_1-Nemotron-Ultra-253B-v1": 131072,
            "openai/gpt-oss-120b": 32768,
            "moonshotai/Kimi-K2-Instruct": 128000,
            "moonshotai/Kimi-K2-Thinking": 128000,
            "zai-org/GLM-4.5": 128000,
            "meta-llama/Llama-3.3-70B-Instruct": 131072,
            "black-forest-labs/flux-schnell": 4096,
            "black-forest-labs/flux-dev": 4096,
        }
    },
    
    "Mistral": {
        "base_url": "https://api.mistral.ai/v1",
        "key_name": "MISTRAL",
        "badge": "🇫🇷 <b>DSGVO-Konform</b>",
        "chat_models": [
            "mistral-large-latest",
            "mistral-medium-2508",
            "magistral-medium-2509",
            "open-mistral-nemo-2407"
        ],
        "vision_models": ["pixtral-large-2411", "pixtral-12b-2409", "mistral-ocr-latest"],
        "audio_models": ["voxtral-mini-latest"],
        "context_limits": {
            "mistral-large-latest": 128000,
            "mistral-medium-2508": 128000,
            "magistral-medium-2509": 128000,
            "open-mistral-nemo-2407": 128000,
            "pixtral-large-2411": 128000,
            "pixtral-12b-2409": 32768,
            "mistral-ocr-latest": 32768,
            "voxtral-mini-latest": 16384,
        }
    },
    
    "Deepgram": {
        "base_url": "https://api.eu.deepgram.com/v1",
        "key_name": "DEEPGRAM",
        "badge": "🇪🇺 <b>EU-Server, US-Firma</b>",
        "audio_models": ["nova-3-general", "nova-2-general", "nova-2"],
        "context_limits": {
            "nova-3-general": 16384,
            "nova-2-general": 16384,
            "nova-2": 16384,
        }
    },
    
    "AssemblyAI": {
        "base_url": "https://api.eu.assemblyai.com/v2",
        "key_name": "ASSEMBLYAI",
        "badge": "🇪🇺 <b>EU-Server, US-Firma</b>",
        "audio_models": ["universal", "slam-1"],
        "context_limits": {
            "universal": 16384,
            "slam-1": 16384,
        }
    },
    
    "OpenRouter": {
        "base_url": "https://openrouter.ai/api/v1",
        "key_name": "OPENROUTER",
        "badge": "🇺🇸 <b>US-Server</b>!",
        "chat_models": [
            # 1M+ Context
            "google/gemini-2.0-pro-exp-02-05:free",
            "google/gemini-2.0-flash-thinking-exp:free",
            "google/gemini-2.0-flash-exp:free",
            "google/gemini-2.5-pro-exp-03-25:free",
            "google/gemini-flash-1.5-8b-exp",
            # 100K+ Context
            "deepseek/deepseek-r1-zero:free",
            "deepseek/deepseek-r1:free",
            "deepseek/deepseek-v3-base:free",
            "deepseek/deepseek-chat-v3-0324:free",
            "deepseek/deepseek-chat:free",
            "google/gemma-3-4b-it:free",
            "google/gemma-3-12b-it:free",
            "qwen/qwen2.5-vl-72b-instruct:free",
            "nvidia/llama-3.1-nemotron-70b-instruct:free",
            "meta-llama/llama-3.2-1b-instruct:free",
            "meta-llama/llama-3.2-11b-vision-instruct:free",
            "meta-llama/llama-3.1-8b-instruct:free",
            "mistralai/mistral-nemo:free",
            # 64K-100K Context
            "mistralai/mistral-small-3.1-24b-instruct:free",
            "google/gemma-3-27b-it:free",
            "qwen/qwen2.5-vl-3b-instruct:free",
            "qwen/qwen-2.5-vl-7b-instruct:free",
            # 32K-64K Context
            "google/learnlm-1.5-pro-experimental:free",
            "qwen/qwq-32b:free",
            "google/gemini-2.0-flash-thinking-exp-1219:free",
            "bytedance-research/ui-tars-72b:free",
            "google/gemma-3-1b-it:free",
            "mistralai/mistral-small-24b-instruct-2501:free",
            "qwen/qwen-2.5-coder-32b-instruct:free",
            "qwen/qwen-2.5-72b-instruct:free",
            # 8K-32K Context
            "meta-llama/llama-3.2-3b-instruct:free",
            "qwen/qwq-32b-preview:free",
            "deepseek/deepseek-r1-distill-qwen-32b:free",
            "qwen/qwen2.5-vl-32b-instruct:free",
            "deepseek/deepseek-r1-distill-llama-70b:free",
            "qwen/qwen-2-7b-instruct:free",
            "google/gemma-2-9b-it:free",
            "mistralai/mistral-7b-instruct:free",
            "microsoft/phi-3-mini-128k-instruct:free",
            "meta-llama/llama-3-8b-instruct:free",
            "meta-llama/llama-3.3-70b-instruct:free",
            # 4K Context
            "huggingfaceh4/zephyr-7b-beta:free",
        ],
        "vision_models": [
            "google/gemini-2.0-pro-exp-02-05:free",
            "google/gemini-2.0-flash-thinking-exp:free",
            "google/gemini-2.0-flash-exp:free",
            "google/gemini-2.5-pro-exp-03-25:free",
            "google/gemini-flash-1.5-8b-exp",
            "qwen/qwen2.5-vl-72b-instruct:free",
            "meta-llama/llama-3.2-11b-vision-instruct:free",
            "mistralai/mistral-small-3.1-24b-instruct:free",
            "google/gemma-3-27b-it:free",
            "qwen/qwen2.5-vl-3b-instruct:free",
            "qwen/qwen-2.5-vl-7b-instruct:free",
            "bytedance-research/ui-tars-72b:free",
            "qwen/qwen2.5-vl-32b-instruct:free",
        ],
        "audio_models": [
            "google/gemini-2.0-flash-lite-001",
            "mistralai/voxtral-small-24b-2507",
            "google/gemini-2.5-flash-lite"
        ],
        "image_models": [
            "google/gemini-2.5-flash-image",
            "openai/gpt-5-image-mini",
            "google/gemini-3-pro-image-preview",
            "black-forest-labs/flux.2-pro",
            "black-forest-labs/flux.2-flex"
        ],
        "context_limits": {
            # 1M+ Context
            "google/gemini-2.0-pro-exp-02-05:free": 2000000,
            "google/gemini-2.0-flash-thinking-exp:free": 1048576,
            "google/gemini-2.0-flash-exp:free": 1048576,
            "google/gemini-2.5-pro-exp-03-25:free": 1000000,
            "google/gemini-flash-1.5-8b-exp": 1000000,
            # 100K+ Context
            "deepseek/deepseek-r1-zero:free": 163840,
            "deepseek/deepseek-r1:free": 163840,
            "deepseek/deepseek-v3-base:free": 131072,
            "deepseek/deepseek-chat-v3-0324:free": 131072,
            "deepseek/deepseek-chat:free": 131072,
            "google/gemma-3-4b-it:free": 131072,
            "google/gemma-3-12b-it:free": 131072,
            "qwen/qwen2.5-vl-72b-instruct:free": 131072,
            "nvidia/llama-3.1-nemotron-70b-instruct:free": 131072,
            "meta-llama/llama-3.2-1b-instruct:free": 131072,
            "meta-llama/llama-3.2-11b-vision-instruct:free": 131072,
            "meta-llama/llama-3.1-8b-instruct:free": 131072,
            "mistralai/mistral-nemo:free": 128000,
            # 64K-100K Context
            "mistralai/mistral-small-3.1-24b-instruct:free": 96000,
            "google/gemma-3-27b-it:free": 96000,
            "qwen/qwen2.5-vl-3b-instruct:free": 64000,
            "qwen/qwen-2.5-vl-7b-instruct:free": 64000,
            # 32K-64K Context
            "google/learnlm-1.5-pro-experimental:free": 40960,
            "qwen/qwq-32b:free": 40000,
            "google/gemini-2.0-flash-thinking-exp-1219:free": 40000,
            "bytedance-research/ui-tars-72b:free": 32768,
            "google/gemma-3-1b-it:free": 32768,
            "mistralai/mistral-small-24b-instruct-2501:free": 32768,
            "qwen/qwen-2.5-coder-32b-instruct:free": 32768,
            "qwen/qwen-2.5-72b-instruct:free": 32768,
            # 8K-32K Context
            "meta-llama/llama-3.2-3b-instruct:free": 20000,
            "qwen/qwq-32b-preview:free": 16384,
            "deepseek/deepseek-r1-distill-qwen-32b:free": 16000,
            "qwen/qwen2.5-vl-32b-instruct:free": 8192,
            "deepseek/deepseek-r1-distill-llama-70b:free": 8192,
            "qwen/qwen-2-7b-instruct:free": 8192,
            "google/gemma-2-9b-it:free": 8192,
            "mistralai/mistral-7b-instruct:free": 8192,
            "microsoft/phi-3-mini-128k-instruct:free": 8192,
            "meta-llama/llama-3-8b-instruct:free": 8192,
            "meta-llama/llama-3.3-70b-instruct:free": 8000,
            # 4K Context
            "huggingfaceh4/zephyr-7b-beta:free": 4096,
            # Audio/Image
            "google/gemini-2.0-flash-lite-001": 1000000,
            "mistralai/voxtral-small-24b-2507": 32768,
            "google/gemini-2.5-flash-lite": 1000000,
            "google/gemini-2.5-flash-image": 1000000,
            "openai/gpt-5-image-mini": 128000,
            "google/gemini-3-pro-image-preview": 1000000,
            "black-forest-labs/flux.2-pro": 4096,
            "black-forest-labs/flux.2-flex": 4096,
        }
    },
    
    "Groq": {
        "base_url": "https://api.groq.com/openai/v1",
        "key_name": "GROQ",
        "badge": "🇺🇸 <b>US-Server</b>!",
        "chat_models": [
            "deepseek-r1-distill-llama-70b",
            "deepseek-r1-distill-qwen-32b",
            "gemma2-9b-it",
            "llama-3.1-8b-instant",
            "llama-3.2-1b-preview",
            "llama-3.2-3b-preview",
            "llama-3.2-11b-vision-preview",
            "llama-3.2-90b-vision-preview",
            "llama-3.3-70b-specdec",
            "llama-3.3-70b-versatile",
            "llama-guard-3-8b",
            "llama3-8b-8192",
            "llama3-70b-8192",
            "mistral-saba-24b",
            "qwen-2.5-32b",
            "qwen-2.5-coder-32b",
            "qwen-qwq-32b",
        ],
        "audio_models": [
            "distil-whisper-large-v3-en",
            "whisper-large-v3",
            "whisper-large-v3-turbo"
        ],
        "vision_models": [
            "llama-3.2-11b-vision-preview",
            "llama-3.2-90b-vision-preview"
        ],
        "context_limits": {
            "deepseek-r1-distill-llama-70b": 8192,
            "deepseek-r1-distill-qwen-32b": 8192,
            "gemma2-9b-it": 8192,
            "llama-3.1-8b-instant": 131072,
            "llama-3.2-1b-preview": 131072,
            "llama-3.2-3b-preview": 131072,
            "llama-3.2-11b-vision-preview": 131072,
            "llama-3.2-90b-vision-preview": 131072,
            "llama-3.3-70b-specdec": 131072,
            "llama-3.3-70b-versatile": 131072,
            "llama-guard-3-8b": 8192,
            "llama3-8b-8192": 8192,
            "llama3-70b-8192": 8192,
            "mistral-saba-24b": 32768,
            "qwen-2.5-32b": 32768,
            "qwen-2.5-coder-32b": 32768,
            "qwen-qwq-32b": 32768,
            "distil-whisper-large-v3-en": 16384,
            "whisper-large-v3": 16384,
            "whisper-large-v3-turbo": 16384,
        }
    },
    
    "Poe": {
        "base_url": "https://api.poe.com/v1",
        "key_name": "POE",
        "badge": "🌐 <b>US-Server</b>!",
        "chat_models": [
            "gpt-5.1-instant",
            "claude-sonnet-4.5",
            "gemini-3-pro",
            "gpt-5.1",
            "gpt-4o",
            "claude-3.5-sonnet",
            "deepseek-r1",
            "grok-4"
        ],
        "vision_models": [
            "claude-sonnet-4.5",
            "gpt-5.1",
            "gemini-3-pro",
            "gpt-4o",
            "claude-3.5-sonnet"
        ],
        "image_models": [
            "gpt-image-1",
            "flux-pro-1.1-ultra",
            "ideogram-v3",
            "dall-e-3",
            "playground-v3"
        ],
        "audio_models": [
            "elevenlabs-v3",
            "sonic-3.0"
        ],
        "video_models": [
            "kling-2.5-turbo-pro",
            "runway-gen-4-turbo",
            "veo-3.1"
        ],
        "supports_system": True,
        "supports_streaming": True,
        "context_limits": {
            "gpt-5.1-instant": 128000,
            "claude-sonnet-4.5": 200000,
            "gemini-3-pro": 2000000,
            "gpt-5.1": 128000,
            "gpt-4o": 128000,
            "claude-3.5-sonnet": 200000,
            "deepseek-r1": 163840,
            "grok-4": 131072,
            "gpt-image-1": 4096,
            "flux-pro-1.1-ultra": 4096,
            "ideogram-v3": 4096,
            "dall-e-3": 4096,
            "playground-v3": 4096,
            "elevenlabs-v3": 4096,
            "sonic-3.0": 4096,
            "kling-2.5-turbo-pro": 4096,
            "runway-gen-4-turbo": 4096,
            "veo-3.1": 4096,
        }
    },
    
    "OpenAI": {
        "base_url": "https://api.openai.com/v1",
        "key_name": "OPENAI",
        "badge": "🇺🇸 <b>US-Server</b>!",
        "chat_models": [
            "gpt-3.5-turbo",
            "gpt-3.5-turbo-0125",
            "gpt-4",
            "gpt-4-turbo",
            "gpt-4o",
            "gpt-4o-mini",
            "o1-preview",
            "o1-mini"
        ],
        "vision_models": [
            "gpt-4-turbo",
            "gpt-4o",
            "gpt-4o-mini",
            "o1-preview",
            "o1-mini"
        ],
        "context_limits": {
            "gpt-3.5-turbo": 16385,
            "gpt-3.5-turbo-0125": 16385,
            "gpt-3.5-turbo-instruct": 4096,
            "gpt-4": 8192,
            "gpt-4-turbo": 128000,
            "gpt-4o": 128000,
            "gpt-4o-mini": 128000,
            "o1-preview": 128000,
            "o1-mini": 128000,
        }
    },
    
    "Cohere": {
        "base_url": "https://api.cohere.ai/v1",
        "key_name": "COHERE",
        "badge": "🇺🇸 <b>US-Server</b>!",
        "chat_models": [
            "command-r-plus-08-2024",
            "command-r-plus",
            "command-r-08-2024",
            "command-r",
            "command",
            "c4ai-aya-expanse-8b",
            "c4ai-aya-expanse-32b",
        ],
        "context_limits": {
            "command-r-plus-08-2024": 131072,
            "command-r-plus-04-2024": 131072,
            "command-r-plus": 131072,
            "command-r-08-2024": 131072,
            "command-r-03-2024": 131072,
            "command-r": 131072,
            "command": 4096,
            "command-nightly": 131072,
            "command-light": 4096,
            "command-light-nightly": 4096,
            "c4ai-aya-expanse-8b": 8192,
            "c4ai-aya-expanse-32b": 131072,
        }
    },
    
    "Together": {
        "base_url": "https://api.together.xyz/v1",
        "key_name": "TOGETHER",
        "badge": "🇺🇸 <b>US-Server</b>!",
        "chat_models": [
            "meta-llama/Meta-Llama-3.1-8B-Instruct-Turbo",
            "deepseek-ai/DeepSeek-R1-Distill-Llama-70B-free",
            "meta-llama/Llama-3.3-70B-Instruct-Turbo-Free",
        ],
        "vision_models": ["meta-llama/Llama-Vision-Free"],
        "context_limits": {
            "meta-llama/Meta-Llama-3.1-8B-Instruct-Turbo": 131072,
            "deepseek-ai/DeepSeek-R1-Distill-Llama-70B-free": 8192,
            "meta-llama/Llama-Vision-Free": 8192,
            "meta-llama/Llama-3.3-70B-Instruct-Turbo-Free": 8192,
        }
    },
    
    "OVH": {
        "base_url": "https://llama-3-1-70b-instruct.endpoints.kepler.ai.cloud.ovh.net/api/openai_compat/v1",
        "key_name": "OVH",
        "badge": "🇫🇷 <b>DSGVO-Konform</b>",
        "chat_models": [
            "ovh/codestral-mamba-7b-v0.1",
            "ovh/deepseek-r1-distill-llama-70b",
            "ovh/llama-3.1-70b-instruct",
            "ovh/llama-3.1-8b-instruct",
            "ovh/llama-3.3-70b-instruct",
            "ovh/mistral-7b-instruct-v0.3",
            "ovh/mistral-nemo-2407",
            "ovh/mixtral-8x7b-instruct",
            "ovh/qwen2.5-coder-32b-instruct",
        ],
        "vision_models": [
            "ovh/llava-next-mistral-7b",
            "ovh/qwen2.5-vl-72b-instruct"
        ],
        "context_limits": {
            "ovh/codestral-mamba-7b-v0.1": 131072,
            "ovh/deepseek-r1-distill-llama-70b": 8192,
            "ovh/llama-3.1-70b-instruct": 131072,
            "ovh/llama-3.1-8b-instruct": 131072,
            "ovh/llama-3.3-70b-instruct": 131072,
            "ovh/llava-next-mistral-7b": 8192,
            "ovh/mistral-7b-instruct-v0.3": 32768,
            "ovh/mistral-nemo-2407": 131072,
            "ovh/mixtral-8x7b-instruct": 32768,
            "ovh/qwen2.5-coder-32b-instruct": 32768,
            "ovh/qwen2.5-vl-72b-instruct": 131072,
        }
    },
    
    "Cerebras": {
        "base_url": "https://api.cerebras.ai/v1",
        "key_name": "CEREBRAS",
        "badge": "🇺🇸 <b>US-Server</b>!",
        "chat_models": [
            "llama3.1-8b",
            "llama-3.3-70b"
        ],
        "context_limits": {
            "llama3.1-8b": 8192,
            "llama-3.3-70b": 8192,
        }
    },
    
    "GoogleAI": {
        "base_url": "https://generativelanguage.googleapis.com/v1beta",
        "key_name": "GOOGLEAI",
        "badge": "🇺🇸 <b>US-Server</b>?",
        "chat_models": [
            "gemini-1.0-pro",
            "gemini-1.5-flash",
            "gemini-1.5-pro",
            "gemini-2.0-pro",
            "gemini-2.5-pro"
        ],
        "vision_models": [
            "gemini-1.5-pro",
            "gemini-1.0-pro",
            "gemini-1.5-flash",
            "gemini-2.0-pro",
            "gemini-2.5-pro"
        ],
        "context_limits": {
            "gemini-1.0-pro": 32768,
            "gemini-1.5-flash": 1000000,
            "gemini-1.5-pro": 1000000,
            "gemini-2.0-pro": 2000000,
            "gemini-2.5-pro": 2000000,
        }
    },
    
    "Anthropic": {
        "base_url": "https://api.anthropic.com/v1",
        "key_name": "ANTHROPIC",
        "badge": "🇺🇸 <b>US-Server</b>!",
        "chat_models": [
            "claude-3-7-sonnet-20250219",
            "claude-3-5-sonnet-20241022",
            "claude-3-5-haiku-20240307",
            "claude-3-opus-20240229",
        ],
        "vision_models": [
            "claude-3-7-sonnet-20250219",
            "claude-3-5-sonnet-20241022",
            "claude-3-5-haiku-20240307",
            "claude-3-opus-20240229",
        ],
        "context_limits": {
            "claude-3-7-sonnet-20250219": 128000,
            "claude-3-5-sonnet-20241022": 200000,
            "claude-3-5-haiku-20240307": 200000,
            "claude-3-5-sonnet-20240620": 200000,
            "claude-3-opus-20240229": 200000,
            "claude-3-haiku-20240307": 200000,
            "claude-3-sonnet-20240229": 200000,
        }
    },
    
    "HuggingFace": {
        "base_url": "https://api-inference.huggingface.co/models",
        "key_name": "HUGGINGFACE",
        "badge": "🌐 <b>US-Server</b>?",
        "chat_models": [
            "microsoft/phi-3-mini-4k-instruct",
            "microsoft/Phi-3-mini-128k-instruct",
            "HuggingFaceH4/zephyr-7b-beta",
            "deepseek-ai/DeepSeek-Coder-V2-Instruct",
            "mistralai/Mistral-7B-Instruct-v0.3",
            "NousResearch/Nous-Hermes-2-Mixtral-8x7B-DPO",
            "microsoft/Phi-3.5-mini-instruct",
            "google/gemma-2-2b-it",
            "Qwen/Qwen2.5-7B-Instruct",
            "tiiuae/falcon-7b-instruct",
            "Qwen/QwQ-32B-preview",
        ],
        "vision_models": [
            "Qwen/Qwen2.5-VL-7B-Instruct",
            "Qwen/qwen2.5-vl-3b-instruct",
            "Qwen/qwen2.5-vl-32b-instruct",
            "Qwen/qwen2.5-vl-72b-instruct",
        ],
        "context_limits": {
            "microsoft/phi-3-mini-4k-instruct": 4096,
            "microsoft/Phi-3-mini-128k-instruct": 131072,
            "HuggingFaceH4/zephyr-7b-beta": 8192,
            "deepseek-ai/DeepSeek-Coder-V2-Instruct": 8192,
            "mistralai/Mistral-7B-Instruct-v0.3": 32768,
            "NousResearch/Nous-Hermes-2-Mixtral-8x7B-DPO": 32768,
            "microsoft/Phi-3.5-mini-instruct": 4096,
            "google/gemma-2-2b-it": 2048,
            "openai-community/gpt2": 1024,
            "microsoft/phi-2": 2048,
            "TinyLlama/TinyLlama-1.1B-Chat-v1.0": 2048,
            "Qwen/Qwen2.5-7B-Instruct": 131072,
            "tiiuae/falcon-7b-instruct": 8192,
            "Qwen/QwQ-32B-preview": 32768,
            "Qwen/Qwen2.5-VL-7B-Instruct": 64000,
            "Qwen/qwen2.5-vl-3b-instruct": 64000,
            "Qwen/qwen2.5-vl-32b-instruct": 8192,
            "Qwen/qwen2.5-vl-72b-instruct": 131072,
        }
    },
}

# EU-ONLY MODE: Set to True to disable US providers
EU_ONLY_MODE = True

RESTRICTED_PROVIDERS = {
    "OpenAI": "🇺🇸 US-Server",
    "Anthropic": "🇺🇸 US-Server", 
    "OpenRouter": "🇺🇸 US-Server",
    "Groq": "🇺🇸 US-Server",
    "Poe": "🇺🇸 US-Server",
    "Cohere": "🇺🇸 US-Server",
    "Together": "🇺🇸 US-Server",
    "Cerebras": "🇺🇸 US-Server",
}

def is_provider_allowed(provider_name: str) -> tuple:
    """
    Returns: (allowed: bool, reason: str)
    """
    if not EU_ONLY_MODE:
        return True, ""
    
    if provider_name in RESTRICTED_PROVIDERS:
        return False, f"⛔ {provider_name} ist im EU-Modus nicht erlaubt (Grund: {RESTRICTED_PROVIDERS[provider_name]})"
    
    return True, ""

def get_compliance_html(provider):
    """Holt Badge-Text direkt aus der Provider-Config (Single Row)."""
    p_data = PROVIDERS.get(provider, {})
    return p_data.get("badge", "❓ Unbekannt")

def get_compliance_html_raw(provider):
    """Gibt nur den Text/Icon zurück. Styling erfolgt im UI-Update."""
    # Ensure emojis are first characters
    badges = {
        "Scaleway": "&#127467;&#127479; <b>DSGVO</b><br>(FR)",             # 🇫🇷
        "Nebius": "&#127466;&#127482; <b>DSGVO</b><br>(NL)",               # 🇪🇺
        "Mistral": "&#127467;&#127479; <b>DSGVO</b><br>(FR)",              # 🇫🇷
        "Gladia": "&#127467;&#127479; <b>DSGVO</b><br>(FR)",               # 🇫🇷
        "OpenRouter": "&#127482;&#127480; <b>US-Host</b><br>(nicht EU)",   # 🇺🇸
        "Groq": "&#127482;&#127480; <b>US-Host</b><br>(nicht EU)",         # 🇺🇸
        "Deepgram": "&#127466;&#127482; <b>US-Firma</b><br>(EU-Server)",   # 🇪🇺
        "AssemblyAI": "&#127466;&#127482; <b>US-Firma</b><br>(EU-Server)", # 🇪🇺
    }
    # Using <br> allows nicer wrapping in the new small badge box
    return badges.get(provider, "❓ Unbekannt")

# Gladia Spezial-Config
GLADIA_CONFIG = {
    "url": "https://api.gladia.io/v2",
    "vocab": [
        "Christian Ströbele", "Konstanze Jüngling", "Fabian Jaskolla", "Jesus Christus", "Amen", "Halleluja",
        "Evangelium", "Predigt", "Liturgie", "Gottesdienst", "Pfarrei",
        "Diözese", "Kirchenvorstand", "Fürbitten", "Akademie",
        "Tagungshaus", "Compliance", "Synode", "Ökumene"
    ]
}

def login_user(username, password):
    user, umk = authenticate_user(username, password)
    
    if user and umk:
        # Save UMK in Session State (RAM Only)
        # Never write this 'new_state' to disk/logs!
        new_state = {
            "id": user.id,
            "username": user.username,
            "is_admin": user.is_admin,
            "is_media_manager": user.is_media_manager,
            "umk": umk # THE KEY
        }
        
        return (True, f"✅ Willkommen, {user.username}!", gr.update(visible=True), gr.update(visible=False), new_state)
    
    return (False, "❌ Login fehlgeschlagen", gr.update(visible=False), gr.update(visible=True), {})

def logout_user():
    """Logout function returns empty state"""
    # Return empty dictionary to reset session_state
    empty_state = {"id": None, "username": None, "is_admin": False}
    return f"👋 Auf Wiedersehen!", gr.update(visible=False), gr.update(visible=True), empty_state

# ==========================================
# TOKEN MANAGEMENT & CHUNKING
# ==========================================

def estimate_tokens(text: str) -> int:
    """
    Conservative token estimation.
    English is ~4 chars/token, but code/German is closer to 2.5-3.
    We use 3 chars/token as a safe buffer to prevent API errors.
    """
    if not text: return 0
    return len(str(text)) // 3

def get_model_context_limit(provider: str, model: str) -> int:
    """Get context window size for a specific model"""
    provider_data = PROVIDERS.get(provider, {})
    
    # Check if provider has model-specific limits
    if "context_limits" in provider_data:
        return provider_data["context_limits"].get(model, 4096)
    
    # Default limits by provider
    defaults = {
        "Scaleway": 32000,
        "Nebius": 128000,
        "Mistral": 128000,
        "OpenRouter": 128000,
        "Groq": 8192,
        "Poe": 128000,
        "Deepgram": 16384,  # Audio context
        "AssemblyAI": 16384
    }
    
    return defaults.get(provider, 4096)

def check_content_fits_context(content: str, provider: str, model: str, reserve_tokens: int = 1000) -> Tuple[bool, int, int]:
    """
    Check if content fits within model's context window.
    
    Returns:
        (fits: bool, estimated_tokens: int, max_tokens: int)
    """
    estimated = estimate_tokens(content)
    limit = get_model_context_limit(provider, model)
    usable_limit = limit - reserve_tokens
    
    return (estimated <= usable_limit, estimated, limit)

def prune_messages(messages: list, model_limit: int, max_output_tokens: int = 1000) -> list:
    """
    Smartly trims conversation history with VERBOSE CLI LOGGING.
    """
    # 1. Calculate Safe Budget
    # Reserve tokens for output + specific safety buffer (e.g. 500)
    safety_buffer = 500
    safe_input_limit = model_limit - max_output_tokens - safety_buffer
    
    # Sanity check
    if safe_input_limit < 1000: 
        safe_input_limit = 2000 # Force a minimum floor if config is weird
        print(f"[CTX] ⚠️ Warning: calculated input limit was too low. Forced to {safe_input_limit}")

    print(f"\n[CTX] 📊 CONTEXT CALCULATION:")
    print(f"[CTX]    Model Limit: {model_limit}")
    print(f"[CTX]    Output Reserve: -{max_output_tokens}")
    print(f"[CTX]    Safety Buffer: -{safety_buffer}")
    print(f"[CTX]    ==========================")
    print(f"[CTX]    AVAILABLE INPUT BUDGET: {safe_input_limit} tokens")

    current_tokens = 0
    kept_indices = []
    
    # 2. Mandatory: System Prompt (Index 0)
    if messages and messages[0]["role"] == "system":
        t = estimate_tokens(messages[0]["content"])
        current_tokens += t
        kept_indices.append(0)
        print(f"[CTX]    + {t:5d} tok | (Required) System Prompt")
    
    # 3. Mandatory: Latest User Message (Last Index)
    if len(messages) > 1:
        last_idx = len(messages) - 1
        if last_idx not in kept_indices:
            t = estimate_tokens(messages[last_idx]["content"])
            current_tokens += t
            kept_indices.append(last_idx)
            print(f"[CTX]    + {t:5d} tok | (Required) Latest User Message")

    # 4. Fill remaining budget with history (Newest -> Oldest)
    remaining_indices = [i for i in range(len(messages)) if i not in kept_indices]
    remaining_indices.reverse() # Start from newest history items

    dropped_count = 0
    
    for i in remaining_indices:
        role = messages[i]["role"]
        # Snippet for log (first 30 chars)
        content_preview = str(messages[i]["content"])[:30].replace('\n', ' ') + "..."
        msg_tokens = estimate_tokens(messages[i]["content"])
        
        if current_tokens + msg_tokens <= safe_input_limit:
            current_tokens += msg_tokens
            kept_indices.append(i)
            print(f"[CTX]    + {msg_tokens:5d} tok | (Kept) {role}: {content_preview}")
        else:
            # Once we hit the limit, we effectively drop everything older
            print(f"[CTX]    - {msg_tokens:5d} tok | (DROP) {role}: {content_preview} [BUDGET FULL]")
            dropped_count += 1
            # We don't break immediately if you want to try fitting smaller older messages, 
            # BUT usually for chat continuity, you stop once you hit a block.
            # We break here to maintain conversational continuity.
            break 
            
    # 5. Reconstruct and Sort
    final_messages = [messages[i] for i in sorted(kept_indices)]
    
    print(f"[CTX]    ==========================")
    print(f"[CTX]    Total Used: {current_tokens} / {safe_input_limit}")
    print(f"[CTX]    Messages: {len(final_messages)} kept, {len(messages) - len(final_messages)} dropped.\n")
    
    return final_messages

def split_content_into_chunks(text: str, max_tokens: int = 4000, overlap: int = 200) -> List[str]:
    """
    Split text into overlapping chunks that fit within token limits.
    
    Args:
        text: Input text to split
        max_tokens: Maximum tokens per chunk (in characters: tokens * 4)
        overlap: Character overlap between chunks
        
    Returns:
        List of text chunks
    """
    max_chars = max_tokens * 4
    overlap_chars = overlap * 4
    
    paragraphs = text.split('\n\n')
    chunks = []
    current_chunk = []
    current_size = 0
    
    for para in paragraphs:
        para_size = len(para)
        
        if current_size + para_size <= max_chars:
            current_chunk.append(para)
            current_size += para_size + 2
        else:
            if current_chunk:
                chunk_text = '\n\n'.join(current_chunk)
                chunks.append(chunk_text)
                
                if overlap_chars > 0 and chunks:
                    overlap_text = chunks[-1][-overlap_chars:]
                    current_chunk = [overlap_text, para]
                    current_size = len(overlap_text) + para_size + 2
                else:
                    current_chunk = [para]
                    current_size = para_size
            else:
                sentences = para.split('. ')
                current_chunk = [sentences[0]]
                current_size = len(sentences[0])
    
    if current_chunk:
        chunks.append('\n\n'.join(current_chunk))
    
    return chunks

# ==========================================
# 🛠️ HELPER FUNCTIONS
# ==========================================

def is_provider_implemented(provider_name):
    """Check if a provider is fully implemented"""
    implemented_providers = {
        "Scaleway": True,
        "Nebius": True,
        "Mistral": True,
        "OpenRouter": True,
        "Groq": True,
        "Poe": True,
        "Deepgram": True,
        "AssemblyAI": True,
        # New providers - mark as not implemented yet
        "OpenAI": False,
        "Cohere": False,
        "Together": False,
        "OVH": False,
        "Cerebras": False,
        "GoogleAI": False,
        "Anthropic": False,
        "HuggingFace": False,  
    }
    return implemented_providers.get(provider_name, False)

def get_provider_status(provider_name):
    """Get status message for provider"""
    if is_provider_implemented(provider_name):
        return ""
    return "⚠️ In Entwicklung - Noch nicht verfügbar"

def get_client(provider_name, api_key_override=None):
    """Factory: Erstellt einen OpenAI-Client für JEDEN Provider"""
    
    # Check if implemented
    if not is_provider_implemented(provider_name):
        raise NotImplementedError(
            f"Provider '{provider_name}' ist noch nicht implementiert. "
            f"Bitte wählen Sie einen anderen Provider."
        )
    
    conf = PROVIDERS.get(provider_name)
    if not conf:
        raise ValueError(f"Unbekannter Provider: {provider_name}")

    key = api_key_override if api_key_override else API_KEYS.get(conf["key_name"])
    if not key or key == "your_key":
        raise ValueError(
            f"Kein API Key für {provider_name} gefunden. "
            f"Bitte in Umgebungsvariablen oder Einstellungen konfigurieren."
        )

    return openai.OpenAI(base_url=conf["base_url"], api_key=key)

def get_provider_choices_with_status():
    """Get provider choices with implementation status"""
    choices = []
    for provider_name in PROVIDERS.keys():
        if is_provider_implemented(provider_name):
            choices.append(provider_name)
        else:
            choices.append(f"{provider_name} (⚠️ Bald verfügbar)")
    return choices

def call_poe_sync(messages, model, api_key):
    """
    Synchronously call Poe API by running async code
    """
    import asyncio
    
    async def call_poe_async():
        poe_messages = []
        for msg in messages:
            role = msg["role"]
            content = msg["content"]
            
            # Skip system messages - Poe doesn't support them
            if role == "system":
                continue
            
            # Convert 'assistant' to 'bot' for Poe
            if role == "assistant":
                role = "bot"
            
            # Handle multimodal content
            if isinstance(content, list):
                text_parts = [item.get("text", "") for item in content if item.get("type") == "text"]
                content = " ".join(text_parts)
            
            # Convert to Poe ProtocolMessage
            poe_messages.append(fp.ProtocolMessage(role=role, content=content))
        
        # Call Poe API and collect response
        full_response = ""
        try:
            async for partial in fp.get_bot_response(
                messages=poe_messages,
                bot_name=model,
                api_key=api_key
            ):
                full_response += partial.text
        except Exception as e:
            raise Exception(f"Poe API error: {str(e)}")
        
        return full_response
    
    # Run the async function synchronously
    try:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        result = loop.run_until_complete(call_poe_async())
        loop.close()
        return result
    except Exception as e:
        raise Exception(f"Poe Fehler: {str(e)}")

def encode_image(image_path):
    if not image_path: return None
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def format_duration(seconds):
    if seconds is None: return "00:00"
    m, s = divmod(seconds, 60)
    return f"{int(m):02d}:{int(s):02d}"

def run_deepgram_transcription(audio_path, model, lang, diar, key):
    """
    Handle Deepgram transcription using the EU endpoint.
    Uses a single request for simplicity, relying on Deepgram's direct support.
    """
    import requests
    import json
    
    API_URL = "https://api.eu.deepgram.com/v1/listen"
    
    # KORRIGIERTE LOGIK: Verwende UI-Key oder falle auf systemweit konfigurierten Key zurück
    api_key = key if key else API_KEYS.get("DEEPGRAM")

    if not api_key or api_key == "your_key":
        yield "❌ Kein Deepgram Key gefunden oder konfiguriert.", ""
        return

    logs = "🚀 Starte Deepgram Upload (EU Endpoint)..."
    yield logs, ""
    
    headers = {
        "Authorization": f"Token {api_key}",
        "Content-Type": "application/octet-stream"
    }

    # Deepgram parameters: use Nova-2/3 by default
    query_params = {
        "model": model or "nova-2-general", 
        "smart_format": "true",
        "language": lang if lang and lang != "auto" else "de", # Default to German
        "punctuate": "true",
        "diarize": "true" if diar else "false",
        "paragraphs": "true" if diar else "false", # Paragraphs/Utterances is usually required for diarization output
        "utterances": "true" if diar else "false",
    }
    
    try:
        with open(audio_path, "rb") as audio_file:
            response = requests.post(
                API_URL, 
                headers=headers, 
                params=query_params, 
                data=audio_file,
                timeout=300 # 5 minutes for upload and transcription
            )
        
        response.raise_for_status()
        result = response.json()
        
        logs += "\n✅ Upload & Transkription fertig."
        yield logs, ""
        
        # --- Parsing Results ---
        final_text = ""
        if diar and result.get('results', {}).get('utterances'):
            utterances = result['results']['utterances']
            for utt in utterances:
                speaker = utt.get('speaker', 0)
                transcript = utt.get('transcript', '').strip()
                start_sec = utt.get('start', 0) # Seconds
                
                time_code = format_duration(start_sec)
                final_text += f"[{time_code}] Sprecher {speaker}: {transcript}\n"
            
        else:
            # Standard Transcript
            transcript = result.get('results', {}).get('channels', [{}])[0].get('alternatives', [{}])[0].get('transcript', '')
            final_text = transcript
            
        if not final_text.strip():
            final_text = "❌ Transkript ist leer! Bitte Sprache prüfen."

        yield logs + "\n🎉 Fertig!", final_text.strip()
        
    except requests.exceptions.RequestException as e:
        error_msg = f"Deepgram API Fehler: {e}"
        if 'response' in locals() and response.status_code == 400:
            try:
                error_msg += f"\nDetails: {response.json().get('err_msg', response.text)}"
            except: pass
        
        logger.exception(f"Deepgram Error: {e}")
        yield logs + f"\n🔥 Abbruch: {error_msg}", ""
    except Exception as e:
        logger.exception(f"Deepgram Error: {e}")
        yield logs + f"\n🔥 Abbruch: {str(e)}", ""

def run_assemblyai_transcription(audio_path, model, lang, diar, key):
    """
    Handle AssemblyAI transcription using the EU endpoint.
    Requires: Upload -> Job Submission -> Polling (Generator Pattern).
    """
    import requests
    import time
    import json
    
    EU_BASE_URL = "https://api.eu.assemblyai.com/v2"

    # KORRIGIERTE LOGIK: Verwende UI-Key oder falle auf systemweit konfigurierten Key zurück
    api_key = key if key else API_KEYS.get("ASSEMBLYAI")
    
    if not api_key or api_key == "your_key":
        yield "❌ Kein AssemblyAI Key gefunden oder konfiguriert.", ""
        return

    logs = "🚀 Starte AssemblyAI Upload (EU Endpoint)..."
    yield logs, ""
    
    headers = {"Authorization": api_key}
    
    try:
        # A. Upload
        upload_url = f"{EU_BASE_URL}/upload"
        upload_headers = headers.copy()
        upload_headers["Content-Type"] = "application/octet-stream"

        with open(audio_path, "rb") as audio_file:
            response = requests.post(upload_url, headers=upload_headers, data=audio_file, timeout=300)
        response.raise_for_status()
        upload_result = response.json()
        final_audio_url = upload_result["upload_url"]
        
        logs += "\n✅ Upload erfolgreich. Starte Job..."
        yield logs, ""
        
        # B. Job Submission
        submit_url = f"{EU_BASE_URL}/transcript"
        
        # Use recommended defaults for general, accurate transcription
        request_data = {
            "audio_url": final_audio_url,
            "language_code": lang if lang and lang != "auto" else "de",
            "punctuate": True,
            "format_text": True,
            "speaker_labels": diar, # Diarization
            "speech_models": [model or "universal"], # Use selected model
        }

        response = requests.post(submit_url, headers=headers, json=request_data)
        response.raise_for_status()
        submit_result = response.json()
        transcript_id = submit_result["id"]
        
        logs += f"\n✅ Job ID: {transcript_id}. Starte Polling..."
        yield logs, ""
        
        # C. Polling
        get_url = f"{EU_BASE_URL}/transcript/{transcript_id}"
        poll_count = 0
        
        while True:
            time.sleep(5)
            poll_count += 1
            
            poll_response = requests.get(get_url, headers=headers)
            poll_response.raise_for_status()
            poll_result = poll_response.json()
            
            status = poll_result.get("status")
            
            if status == "completed":
                # Extract Transcript
                transcript = poll_result.get('text', '')
                final_text = ""
                
                if diar and poll_result.get('utterances'):
                    # Process diarized utterances
                    for utt in poll_result['utterances']:
                        speaker = utt.get('speaker') # AssemblyAI returns "A", "B", etc.
                        text = utt.get('text', '').strip()
                        start_sec = utt.get('start', 0) / 1000 # Convert ms to seconds
                        
                        time_code = format_duration(start_sec)
                        final_text += f"[{time_code}] Sprecher {speaker}: {text}\n"
                else:
                    final_text = transcript
                    
                if not final_text.strip():
                    final_text = "❌ Transkript ist leer! Bitte Sprache/Modell prüfen."

                logs += f"\n✅ Status: COMPLETED nach {poll_count * 5} Sekunden."
                yield logs + "\n🎉 Fertig!", final_text.strip()
                return

            elif status == "error":
                raise Exception(f"AssemblyAI Error: {poll_result.get('error')}")

            if poll_count % 3 == 0:
                logs += f"\n⏳ Status: {status}... Polling {poll_count}..."
                yield logs, ""
                
    except requests.exceptions.RequestException as e:
        error_msg = f"AssemblyAI API Fehler: {e}"
        try:
            if 'response' in locals():
                error_msg += f"\nDetails: {response.text}"
        except: pass
        logger.exception(f"AssemblyAI Error: {e}")
        yield logs + f"\n🔥 Abbruch: {error_msg}", ""
    except Exception as e:
        logger.exception(f"AssemblyAI Error: {e}")
        yield logs + f"\n🔥 Abbruch: {str(e)}", ""

# ==========================================
# 1. CHAT LOGIK
# ==========================================

def run_chat(message, history, provider, model, temp, system_prompt, key, r_effort, r_tokens, user_state):
    """
    Enhanced chat runner with strict context logging and content extraction fix
    """
    # VALIDATION
    allowed, reason = is_provider_allowed(provider)
    if not allowed:
        yield f"⛔ **Zugriff verweigert**\n\n{reason}\n\nBitte wählen Sie einen EU-Provider (Scaleway, Mistral, Nebius)."
        return
    
    import re
    import json

    print(f"[DEBUG] 🚀 run_chat started. History length: {len(history)}")

    # Clean provider name
    clean_provider = provider.replace(" ⚠️", "").strip()
    
    if not is_provider_implemented(clean_provider):
        yield f"⚠️ **{clean_provider} ist noch nicht verfügbar**"
        return

    # --- HELPER: Extract text from Gradio's multimodal list format ---
    def extract_clean_content(content):
        if isinstance(content, list):
            # Handle format: [{'type': 'text', 'text': '...'}]
            text_parts = [str(item.get("text", "")) for item in content if item.get("type") == "text"]
            return " ".join(text_parts)
        return str(content) if content else ""
    # -----------------------------------------------------------------
    
    try:
        # 1. Build Full Message List
        raw_messages = []
        
        # System Prompt
        if system_prompt and system_prompt.strip():
            raw_messages.append({"role": "system", "content": str(system_prompt)})
        
        # Chat History
        for msg in history:
            # FIX: Properly extract text if content is a list/dict
            content = extract_clean_content(msg.get("content"))
            raw_messages.append({"role": msg["role"], "content": content})
            
        # Current Message
        # FIX: Also extract content for the current message
        current_content = extract_clean_content(message)
        raw_messages.append({"role": "user", "content": current_content})

        # 2. Get Context Limit Logic
        context_limit = get_model_context_limit(provider, model)
        
        # Determine max output
        max_output = 4096 
        if r_tokens > 0: max_output = int(r_tokens)
        
        # 3. PRUNE (This calls our new verbose function)
        final_messages = prune_messages(raw_messages, context_limit, max_output)
        
        # 4. Warnings logic
        original_len = len(raw_messages)
        final_len = len(final_messages)
        truncated_count = original_len - final_len
        
        warning_prefix = ""
        if truncated_count > 0:
            warning_prefix = f"⚠️ *Context-Limit ({context_limit} T): {truncated_count} älteste Nachrichten entfernt.*\n\n"

        # 5. Initialize Client
        client = get_client(provider, key)
        
        # 6. Setup Parameters
        params = {
            "model": model,
            "messages": final_messages,
            "stream": True
        }
        
        # Provider nuances...
        extra_body = {}
        if provider == "OpenRouter":
            reasoning_config = {}
            if r_tokens > 0:
                reasoning_config["max_tokens"] = int(r_tokens)
            elif r_effort != "default":
                reasoning_config["effort"] = r_effort
            
            if reasoning_config:
                extra_body["reasoning"] = reasoning_config
                extra_body["include_reasoning"] = True

        elif provider == "Scaleway":
            if r_effort and r_effort != "default":
                params["reasoning_effort"] = r_effort
            if r_tokens > 0:
                params["max_completion_tokens"] = int(r_tokens)
            params["temperature"] = float(temp)

        elif provider == "Mistral":
            if r_tokens > 0:
                params["max_tokens"] = int(r_tokens)
            params["temperature"] = float(temp)

        else:
            params["temperature"] = float(temp)
            if r_tokens > 0:
                params["max_tokens"] = int(r_tokens)

        if extra_body:
            params["extra_body"] = extra_body

        # 7. Execute Stream
        stream = client.chat.completions.create(**params)
        
        full_response = ""
        reasoning_buffer = ""
        is_thinking = False
        
        for chunk in stream:
            if chunk.choices and len(chunk.choices) > 0:
                delta = chunk.choices[0].delta
                
                # --- A. Reasoning Handling ---
                new_reasoning = ""
                if hasattr(delta, 'reasoning') and delta.reasoning:
                    new_reasoning = delta.reasoning
                elif hasattr(delta, 'reasoning_content') and delta.reasoning_content:
                    new_reasoning = delta.reasoning_content
                elif hasattr(delta, 'reasoning_details') and delta.reasoning_details:
                    if isinstance(delta.reasoning_details, str):
                        new_reasoning = delta.reasoning_details
                    elif isinstance(delta.reasoning_details, list):
                        for detail in delta.reasoning_details:
                            if detail.get('type') == 'reasoning.text':
                                new_reasoning += detail.get('text', '')

                if new_reasoning:
                    reasoning_buffer += new_reasoning
                    display_thought = f"<details open><summary>💭 Gedankengang ({len(reasoning_buffer)} zeichen)</summary>\n\n{reasoning_buffer}\n</details>\n\n"
                    yield warning_prefix + display_thought + full_response
                    warning_prefix = "" 
                    continue

                # --- B. Content Handling ---
                if delta.content:
                    val = delta.content
                    if "<think>" in val:
                        is_thinking = True
                        val = val.replace("<think>", "")
                    elif "</think>" in val:
                        is_thinking = False
                        val = val.replace("</think>", "")
                    
                    if is_thinking:
                        reasoning_buffer += val
                        display_thought = f"<details open><summary>💭 Gedankengang ({len(reasoning_buffer)} zeichen)</summary>\n\n{reasoning_buffer}\n</details>\n\n"
                        yield warning_prefix + display_thought + full_response
                        warning_prefix = ""
                    else:
                        full_response += val
                        if reasoning_buffer:
                            display_thought = f"<details><summary>💭 Gedankengang ({len(reasoning_buffer)} zeichen)</summary>\n\n{reasoning_buffer}\n</details>\n\n"
                            yield warning_prefix + display_thought + full_response
                        else:
                            yield warning_prefix + full_response
                        warning_prefix = ""
                
    except Exception as e:
        err_str = str(e)
        if any(phrase in err_str.lower() for phrase in ["context length", "maximum context", "too many tokens"]):
            yield f"🔥 **Kritischer Fehler: Context-Limit.**\n\nTrotz automatischer Kürzung war die Anfrage zu groß.\nTechnischer Fehler: {str(e)}"
            return
        
        logger.exception(f"Chat error with {provider}: {str(e)}")
        yield f"🔥 Fehler ({provider}): {str(e)}"
        
# ==========================================
# 2. VISION LOGIK
# ==========================================

def run_vision(image, prompt, provider, model, key, user_state):
    # --- SECURITY CHECK ---
    if not user_state or not user_state.get("id"):
        return "⛔ Nicht autorisiert. Bitte anmelden."
    # ----------------------
    if not image: return "❌ Bitte Bild hochladen."
    
    try:
        client = get_client(provider, key)
        b64_img = encode_image(image)
        
        messages = [{
            "role": "user",
            "content": [
                {"type": "text", "text": prompt},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64_img}"}}
            ]
        }]
        
        # Standard OpenAI Vision call works for Poe now too
        response = client.chat.completions.create(
            model=model, messages=messages, max_tokens=1000
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"🔥 Vision Fehler: {str(e)}"

# ==========================================
# 3. TRANSKRIPTION
# ==========================================

def smart_format(utterances, show_tc, smart_merge, group_para, show_spk):
    """V6 Formatierungs-Logik für Gladia"""
    if not utterances: return ""
    out = ""; buf = ""; buf_start = 0; buf_spk = -1; last_end = 0

    for turn in utterances:
        txt = turn.get("text", "").strip()
        start = turn.get("start", 0); end = turn.get("end", 0); spk = turn.get("speaker", 0)

        if buf == "": buf_start = start; buf_spk = spk

        change = (spk != buf_spk)
        pause = start - last_end
        new_para = group_para and (pause > 1.5)

        if change or (new_para and buf):
            prefix = ""
            if show_tc: prefix += f"[{format_duration(buf_start)}] "
            if show_spk: prefix += f"Sprecher {buf_spk}: "
            out += f"{prefix}{buf}\n"
            if new_para: out += "\n"
            buf = txt; buf_start = start; buf_spk = spk
        else:
            if smart_merge: buf += " " + txt if buf else txt
            else:
                prefix = ""
                if show_tc: prefix += f"[{format_duration(start)}] "
                if show_spk: prefix += f"Sprecher {spk}: "
                out += f"{prefix}{txt}\n"; buf = ""
        last_end = end

    if buf:
        prefix = ""
        if show_tc: prefix += f"[{format_duration(buf_start)}] "
        if show_spk: prefix += f"Sprecher {buf_spk}: "
        out += f"{prefix}{buf}\n"
    return out

def run_mistral_transcription(audio_path, model, lang, api_key):
    """
    Handle Mistral transcription with chunking and state tracking.
    """
    import os
    
    client = openai.OpenAI(
        base_url="https://api.mistral.ai/v1", 
        api_key=api_key
    )
    
    # 1. Hash file to check/create progress
    file_hash = get_file_hash(audio_path)
    logger.info(f"Processing file hash: {file_hash}")
    
    # 2. Split file (Mistral limit ~15 mins)
    # We use 10 min chunks to be safe and allow for overhead
    logger.info("Splitting audio for Mistral (limit ~15m)...")
    chunks = split_audio(audio_path, chunk_length_ms=10 * 60 * 1000)
    
    if not chunks:
        yield "❌ Fehler beim Aufteilen der Audiodatei."
        return

    full_transcript = ""
    total_chunks = len(chunks)
    
    yield f"📂 Datei in {total_chunks} Teile zerlegt. Starte Upload & Transkription..."

    # 3. Process Chunks
    for i, chunk_path in enumerate(chunks):
        try:
            current_step = i + 1
            yield f"⏳ Verarbeite Teil {current_step}/{total_chunks}..."
            
            # Open file in binary mode
            with open(chunk_path, "rb") as f:
                # Mistral uses the standard OpenAI-compatible audio endpoint structure
                # Note: 'timestamp_granularities' is unique to Mistral but passed via extra_body or standard params depending on SDK version.
                # Since we use the OpenAI client for compatibility, we use standard params.
                
                params = {
                    "model": model or "voxtral-mini-latest",
                    "file": f,
                    "response_format": "json"
                }
                
                # Mistral specific: language param
                if lang and lang != "auto":
                    params["language"] = lang

                transcript = client.audio.transcriptions.create(**params)
                
                text_part = transcript.text
                full_transcript += text_part + " "
                
                logger.info(f"Chunk {current_step} complete.")
                yield f"✅ Teil {current_step}/{total_chunks} fertig."

        except Exception as e:
            logger.exception(f"Error processing chunk {i}: {e}")
            yield f"🔥 Fehler bei Teil {current_step}: {str(e)}"
            # We don't abort, try next chunk? Or abort? Usually better to abort to avoid broken texts.
            return 
        finally:
            # Cleanup chunk
            if os.path.exists(chunk_path):
                os.remove(chunk_path)

    yield full_transcript.strip()

def run_chunked_api_transcription(client, model, chunk_paths, lang, prompt, temp, job_id=None):
    """
    Iterates through chunks with State Tracking.
    """
    import time
    
    # Load manifest if resuming
    start_index = 0
    full_transcript_parts = [""] * len(chunk_paths)
    
    if job_id:
        try:
            with open(os.path.join(JOB_STATE_DIR, f"{job_id}.json"), "r") as f:
                manifest = json.load(f)
                # Find first non-done chunk
                for i, c in enumerate(manifest["chunks"]):
                    if c["status"] == "done":
                        full_transcript_parts[i] = manifest["transcript_parts"][i]
                    else:
                        if start_index == 0 and i > 0: start_index = i
        except: pass

    total = len(chunk_paths)
    MAX_RETRIES = 3
    BASE_DELAY = 10 
    
    for i in range(start_index, total):
        chunk_path = chunk_paths[i]
        step = i + 1
        chunk_success = False
        retry_count = 0
        
        # If resuming, check if file still exists (tmp files might be gone if server restarted)
        if not os.path.exists(chunk_path):
            yield f"❌ Fehler: Chunk-Datei {step} nicht gefunden. Job kann nicht fortgesetzt werden."
            return

        while not chunk_success and retry_count < MAX_RETRIES:
            try:
                yield f"⏳ Verarbeite Teil {step}/{total}..."
                
                with open(chunk_path, "rb") as f:
                    args = {
                        "model": model,
                        "file": f,
                        "response_format": "json"
                    }
                    if lang and lang != "auto": args["language"] = lang
                    if temp is not None: args["temperature"] = float(temp)
                    
                    # Context handling
                    if i == 0 and prompt:
                        args["prompt"] = prompt
                    elif i > 0:
                        # Construct context from previous successful parts
                        prev_text = "".join(full_transcript_parts[:i])
                        if prev_text: args["prompt"] = prev_text[-220:]

                    resp = client.audio.transcriptions.create(**args)
                    text_part = resp.text if hasattr(resp, 'text') else str(resp)
                    
                    # Update State
                    full_transcript_parts[i] = text_part + " "
                    if job_id:
                        update_job_chunk_status(job_id, i, "done", text_part + " ")
                    
                    yield f"✅ Teil {step}/{total} erledigt."
                    chunk_success = True
            
            except Exception as e:
                error_str = str(e)
                if "429" in error_str or "rate_limit" in error_str.lower():
                    retry_count += 1
                    wait = BASE_DELAY * (2 ** retry_count)
                    yield f"⚠️ Rate Limit (429). Warte {wait}s..."
                    time.sleep(wait)
                else:
                    logger.error(f"Critical error in chunk {step}: {e}")
                    full_transcript_parts[i] = f" [Fehler in Teil {step}] "
                    if job_id: update_job_chunk_status(job_id, i, "failed")
                    break 
        
        if not chunk_success:
            full_transcript_parts[i] = f" [Timeout Teil {step}] "
            if job_id: update_job_chunk_status(job_id, i, "failed")

    yield "".join(full_transcript_parts).strip()
        
def run_transcription(audio, provider, model, lang, whisper_temp, whisper_prompt, diar, trans, target, key, chunk_opt=True, chunk_len=10, user_id=None):
    """
    Unified transcription router.
    - Gladia: Handles Files (Upload -> URL) OR Direct URLs (YouTube/etc).
    - Others: Requires local file.
    """
    logger.info("=" * 50)
    logger.info(f"TRANSCRIPTION START: {provider} | Model: {model} | Input: {audio}")

    # --- 1. Validation ---
    is_url = isinstance(audio, str) and audio.startswith("http")
    
    if not audio:
        yield "❌ Keine Datei oder URL.", "", ""
        return

    # Check file existence only if it's NOT a URL
    if not is_url and not os.path.exists(audio):
        yield f"❌ Datei nicht gefunden: {audio}", "", ""
        return

    if not is_url:
        try:
            file_size = os.path.getsize(audio)
            if file_size == 0:
                yield "❌ Datei ist leer (0 Bytes).", "", ""
                return
        except Exception as e:
            logger.error(f"File check error: {e}")

    # --- 2. BRANCH A: GLADIA (Native URL Support) ---
    if provider == "Gladia":
        logger.info("Using Gladia ...")
        
        api_key = key if key else API_KEYS.get("GLADIA", "")
        if not api_key:
            yield "❌ Kein Gladia Key gefunden.", "", ""
            return

        headers = {"x-gladia-key": api_key, "accept": "application/json"}
        
        # A. Handle Input (File Upload OR Direct URL)
        audio_url_for_api = ""
        
        if is_url:
            # Direct URL (YouTube, etc.)
            audio_url_for_api = audio
            logs = f"🚀 Sende URL direkt an Gladia API ({audio})..."
            yield logs, "", ""
        else:
            # File Upload
            logs = "🚀 Starte Gladia Upload..."
            yield logs, "", ""
            try:
                fname = os.path.basename(audio)
                with open(audio, 'rb') as f:
                    r = requests.post(
                        f"{GLADIA_CONFIG['url']}/upload",
                        headers=headers,
                        files={'audio': (fname, f, 'audio/wav')}, # Mime type sniffing is auto handled usually
                        timeout=600 
                    )
                if r.status_code != 200:
                    raise Exception(f"Upload failed: {r.text}")
                
                audio_url_for_api = r.json().get("audio_url")
                logs += "\n✅ Upload fertig. Starte Job..."
                yield logs, "", ""
            except Exception as e:
                yield f"🔥 Upload Fehler: {str(e)}", "", ""
                return

        try:
            # B. Job Config
            vocab_list = [{"value": w} for w in GLADIA_CONFIG.get('vocab', [])]
            
            payload = {
                "audio_url": audio_url_for_api,
                "language_config": {
                    "code_switching": (lang == "auto"),
                    "languages": [] if lang == "auto" else [lang]
                },
                "diarization": diar,
                "diarization_config": {"min_speakers": 1, "max_speakers": 10} if diar else None,
                "custom_vocabulary": True,
                "custom_vocabulary_config": {"vocabulary": vocab_list},
                "translation": trans,
                "translation_config": {
                    "target_languages": [target],
                    "model": "base",
                    "match_original_utterances": True
                } if trans else None
            }

            # C. Start Job
            r = requests.post(f"{GLADIA_CONFIG['url']}/pre-recorded", headers=headers, json=payload)
            if r.status_code != 201:
                raise Exception(f"Job start failed: {r.text}")
            
            result_url = r.json().get("result_url")
            job_id = r.json().get("id", "unknown")
            logs += f"\n🆔 Job ID: {job_id}"
            
            # D. Polling
            poll_count = 0
            start_t = time.time()
            
            while True:
                time.sleep(5)
                poll_count += 1
                elapsed = time.time() - start_t
                
                if poll_count % 2 == 0: 
                    yield f"{logs}\n⏳ Verarbeite... ({format_duration(elapsed)})", "", ""

                if elapsed > 7200: # 2 Hours max (YouTube videos can be long)
                    raise Exception("Timeout nach 120 Minuten")

                r = requests.get(result_url, headers=headers)
                if r.status_code != 200: continue
                
                data = r.json()
                status = data.get("status")
                
                if status == "done":
                    break
                elif status == "error":
                    raise Exception(f"Gladia Error: {json.dumps(data)}")

            # E. Process Result
            res = data.get("result", {})
            
            transcription = res.get("transcription", {})
            utterances = transcription.get("utterances", [])
            final_text = smart_format(utterances, True, True, True, diar)
            if not final_text: 
                final_text = transcription.get("full_transcript", "")

            trans_text = ""
            if trans:
                translation = res.get("translation", {})
                t_res = translation.get("results", [])
                if t_res:
                    t_utt = t_res[0].get("utterances", [])
                    trans_text = smart_format(t_utt, True, True, True, diar)
            
            yield f"{logs}\n🎉 Fertig!", final_text, trans_text

        except Exception as e:
            logger.exception(f"Gladia Error: {e}")
            yield f"🔥 Fehler: {str(e)}", "", ""
            
    # --- BRANCH B, C, D (Requires Local File) ---
    else:
        # Prevent URL usage for providers that don't support it natively
        if is_url:
             yield f"❌ Der Provider '{provider}' unterstützt keine direkten URLs. Bitte Datei herunterladen oder Gladia wählen.", "", ""
             return        
    
        # --- 2. BRANCH B: DEEPGRAM (Native Single-Shot EU) ---
        if provider == "Deepgram":
            logger.info("Using Deepgram Native Sync Flow (EU)")
            full_text = ""
            # model, lang, diar aus den UI-Optionen verwenden
            for log, text in run_deepgram_transcription(audio, model, lang, diar, key):
                if text: full_text = text
                yield log, full_text, "(Keine Übersetzung verfügbar)"
            return # Ende der Funktion

        # --- 2. BRANCH C: ASSEMBLYAI (Native Async EU) ---
        elif provider == "AssemblyAI":
            logger.info("Using AssemblyAI Native Async Flow (EU)")
            full_text = ""
            # model, lang, diar aus den UI-Optionen verwenden
            for log, text in run_assemblyai_transcription(audio, model, lang, diar, key):
                if text: full_text = text
                yield log, full_text, "(Keine Übersetzung verfügbar)"
            return # Ende der Funktion

        # --- 3. BRANCH D: GENERIC CHUNKING (Mistral, Scaleway, Groq) ---
        else:
            logger.info(f"Using Generic Provider: {provider}")
            
            try:
                client = get_client(provider, key)
            except Exception as e:
                yield f"🔥 Client Fehler: {str(e)}", "", ""
                return

            if not model:
                conf = PROVIDERS.get(provider, {})
                model = conf.get("audio_models", ["whisper-large-v3"])[0]

            logs = f"🚀 Starte {provider} ({model})..."
            yield logs, "", ""

            chunks = []
            chunk_dir = None # Default if not splitting

            try:
                # --- OPTIONAL CHUNKING LOGIC ---
                if chunk_opt:
                    yield f"{logs}\n✂️ Teile Audio (alle {chunk_len} Min)...", "", ""
                    chunks, chunk_dir = split_audio_into_chunks(audio, chunk_minutes=int(chunk_len))
                    
                    if not chunks:
                        yield "❌ Fehler beim Aufteilen der Datei.", "", ""
                        return
                    logs += f"\n📂 {len(chunks)} Teile erstellt."
                else:
                    # No chunking: Pass original file as single item list
                    yield f"{logs}\n⚠️ Chunking deaktiviert. Sende Originaldatei...", "", ""
                    if os.path.getsize(audio) > 25 * 1024 * 1024:
                        logger.warning("File > 25MB and chunking disabled. API might fail.")
                        logs += "\n⚠️ WARNUNG: Datei > 25MB. Upload könnte fehlschlagen."
                    chunks = [audio]

                # --- CREATE JOB MANIFEST (For Resume Capability) ---
                job_id = int(time.time())
                create_job_manifest(job_id, audio, provider, model, chunks, lang, whisper_prompt, whisper_temp, user_id)
                logs += f"\n🆔 Job-ID: {job_id} (Für Resume gespeichert)"
                yield logs, "", ""

                # D. Run Sequential Processing
                full_text = ""
                
                # Pass job_id to allow state saving during processing
                transcriber = run_chunked_api_transcription(
                    client, model, chunks, lang, whisper_prompt, whisper_temp, job_id=job_id
                )

                for update in transcriber:
                    if len(update) < 300 and (update.startswith("⏳") or update.startswith("✅") or update.startswith("⚠️")):
                        logs += f"\n{update}"
                        yield logs, full_text, ""
                    else:
                        full_text = update

                yield logs + "\n🎉 Fertig!", full_text, "(Keine Übersetzung verfügbar)"

            except Exception as e:
                logger.exception(f"Provider Error: {e}")
                yield logs + f"\n🔥 Abbruch: {str(e)}", "", ""
            
            finally:
                # E. Cleanup (Only if we actually created chunks)
                if chunk_dir:
                    cleanup_chunks(chunk_dir)

                
def run_and_save_transcription(audio, provider, model, lang, w_temp, w_prompt, diar, trans, target, key, chunk_opt, chunk_len, dg_lang, dg_diar, aa_lang, aa_diar, url_input, dl_video, dl_destination, force_local_dl, user_state):
    # for URLs: If User is Admin/Manager AND checked the box -> Download. Otherwise -> Force Gladia direct mode.
    
    # --- SECURITY CHECK ---
    if not user_state or not user_state.get("id"):
        yield "⛔ Nicht autorisiert. Bitte anmelden.", "", ""
        return
    
    user_id = user_state["id"]
    is_privileged = user_state.get("is_admin", False) or user_state.get("is_media_manager", False)
    
    # ----------------------
    
    # --- HANDLE URL INPUT ---
    actual_audio_path = audio
    cleanup_downloaded = False
    
    if url_input and url_input.strip():
        url = url_input.strip()
        
        # YouTube warning
        if is_youtube_url(url):
            yield "⚠️ YouTube-URL erkannt. Rechte prüfen!\n", "", ""
            time.sleep(1)

        # LOGIC:
        # 1. Privileged users CAN force local download (yt-dlp).
        # 2. Normal users MUST use Gladia Direct.
        # 3. If Provider is NOT Gladia, local download is mandatory (only Privileged can do it).

        should_download = False
        
        if is_privileged and force_local_dl:
            # Admin explicitly requested download
            should_download = True
            yield "📥 [Admin] Starte lokalen Download via yt-dlp...\n", "", ""
        
        elif provider != "Gladia":
            # Other providers require a file
            if is_privileged:
                should_download = True
                yield f"📥 [Info] Provider '{provider}' benötigt lokale Datei. Starte Download...\n", "", ""
            else:
                yield "⛔ Zugriff verweigert: Nur Gladia unterstützt direkte URLs für Standard-Nutzer. Bitte Provider 'Gladia' wählen.", "", ""
                return
        else:
            # Provider is Gladia and NOT forced download -> Use Direct URL
            yield "🔗 Sende URL direkt an Gladia (Kein Download)...\n", "", ""
            actual_audio_path = url # Pass URL string directly
        
        # EXECUTE DOWNLOAD IF NEEDED
        if should_download:
            # Determine save location
            save_to_storage = (dl_destination == "Storage Box")
            
            success, downloaded_path, dl_msg = download_from_url(
                url=url,
                download_video=dl_video,
                save_to_storage=save_to_storage,
                user_state=user_state
            )
            
            if not success:
                yield f"❌ Download fehlgeschlagen: {dl_msg}", "", ""
                return
            
            yield f"{dl_msg}\n", "", ""
            actual_audio_path = downloaded_path
            
            # Mark for cleanup if saved to temp
            if not save_to_storage:
                cleanup_downloaded = True
    
    # --- PROCEED WITH TRANSCRIPTION ---
    
    # 1. Prepare parameters (Legacy mapping)
    final_lang = lang
    final_diar = diar
    final_trans = trans
    final_target = target

    if provider == "Deepgram":
        final_lang = dg_lang
        final_diar = dg_diar
        final_trans = False
        final_target = None
    elif provider == "AssemblyAI":
        final_lang = aa_lang
        final_diar = aa_diar
        final_trans = False
        final_target = None
    
    try:
        # 2. Run transcription
        result = None
        for result in run_transcription(
            actual_audio_path, provider, model, 
            final_lang, w_temp, w_prompt, 
            final_diar, final_trans, final_target, 
            key, chunk_opt, chunk_len,
            user_id=user_id
        ):
            yield result

        # 3. Save to database
        if user_id and result and len(result) > 1 and result[1]:
            # Don't save empty/error results
            if "❌" in result[1] or not result[1].strip():
                return

            logger.info("Auto-saving transcription...")
            # For URLs, filename is the URL, for files, it's basename
            filename = url_input if (url_input and not should_download) else (os.path.basename(actual_audio_path) if actual_audio_path else "Audio")
            
            try:
                trans_id = save_transcription(
                    user_id=user_id,
                    provider=provider,
                    model=model or "N/A",
                    original=result[1],
                    translated=result[2] if len(result) > 2 else None,
                    language=final_lang,
                    filename=filename,
                    user_state=user_state # Pass key
                )
                updated_log = result[0] + f"\n\n💾 Automatisch gespeichert (ID: {trans_id})"
                yield (updated_log, result[1], result[2] if len(result) > 2 else "")

            except Exception as save_error:
                logger.error(f"DB Save error: {save_error}")
                yield (result[0] + f"\n⚠️ DB Fehler: {save_error}", result[1], result[2])

    except Exception as e:
        logger.exception(f"Orchestrator error: {e}")
        yield f"🔥 Kritischer Fehler: {str(e)}", "", ""
    
    finally:
        # Cleanup downloaded file if it was temporary
        if cleanup_downloaded and actual_audio_path and os.path.exists(actual_audio_path):
            try:
                os.remove(actual_audio_path)
                logger.info(f"Cleaned up temp file: {actual_audio_path}")
            except: pass

# ==========================================
# 4. BILDGENERIERUNG 
# ==========================================

def run_image_gen(prompt, provider, model, width, height, steps, key, user_state):
    # --- SECURITY CHECK ---
    if not user_state or not user_state.get("id"):
        return None, "⛔ Nicht autorisiert. Bitte anmelden."
    # ----------------------
    
    import requests 
    import base64
    import tempfile
    import re
    import time
    
    try:
        client = get_client(provider, key)
        
        # --- SPECIAL CASE: POE (Chat-to-Image) ---
        if provider == "Poe":
            response = client.chat.completions.create(
                model=model,
                messages=[{"role": "user", "content": prompt}],
                stream=False
            )
            response_text = response.choices[0].message.content
            
            # Extract Image URL
            match = re.search(r'!\[.*?\]\((https?://.*?)\)', response_text)
            if not match:
                match = re.search(r'\[.*?\]\((https?://.*?)\)', response_text)
            if not match:
                match = re.search(r'(https?://[^\s]+)', response_text)
            
            if not match:
                return None, f"❌ Kein Bild gefunden. Antwort: {response_text[:200]}"
            
            image_url = match.group(1).rstrip(".,;)")
            
            # Download with Retry
            for attempt in range(1, 4):
                try:
                    r = requests.get(image_url, timeout=15)
                    if r.status_code == 200:
                        tfile = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                        tfile.write(r.content)
                        tfile.close()
                        return tfile.name, "✅ Erfolg (Poe)"
                    elif r.status_code in [403, 404]:
                        time.sleep(2)
                        continue
                    else:
                        return None, f"❌ Download Fehler: {r.status_code}"
                except Exception as e:
                    logger.warning(f"Download attempt {attempt} failed: {e}")
                    time.sleep(2)
            
            return None, "❌ Bild konnte nach 3 Versuchen nicht geladen werden."
        
        # --- SPECIAL CASE: OPENROUTER (Chat Completions with Modalities) ---
        if provider == "OpenRouter":
            response = client.chat.completions.create(
                model=model,
                messages=[{"role": "user", "content": prompt}],
                modalities=["image", "text"],
                stream=False
            )
            
            message = response.choices[0].message
            
            if not hasattr(message, 'images') or not message.images:
                content = getattr(message, 'content', '')
                return None, f"❌ Keine Bilder generiert. Antwort: {content[:200]}"
            
            # FIX: Handle Dictionary vs Object access
            first_image = message.images[0]
            image_data_url = None
            
            if isinstance(first_image, dict):
                # It's a dictionary (logs showed this)
                image_data_url = first_image.get('image_url', {}).get('url')
            elif hasattr(first_image, 'image_url'):
                # It's an object
                image_data_url = first_image.image_url.url
            
            if not image_data_url or not image_data_url.startswith('data:image/'):
                return None, f"❌ Ungültiges Bildformat: {str(image_data_url)[:50]}"
            
            base64_data = image_data_url.split('base64,', 1)[1]
            img_data = base64.b64decode(base64_data)
            
            tfile = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            tfile.write(img_data)
            tfile.close()
            return tfile.name, "✅ Erfolg (OpenRouter)"
        
        # --- SPECIAL CASE: SCALEWAY ---
        if provider == "Scaleway":
            return None, "❌ Scaleway: Bildgenerierung derzeit nicht unterstützt. Bitte Nebius verwenden."
        
        # --- STANDARD PROVIDER: NEBIUS ---
        params = {
            "model": model,
            "prompt": prompt,
            "response_format": "b64_json",
        }
        
        if provider == "Nebius":
            params["extra_body"] = {
                "width": width, 
                "height": height, 
                "num_inference_steps": steps,
                "response_extension": "jpg"
            }
            
        response = client.images.generate(**params)
        
        img_data = base64.b64decode(response.data[0].b64_json)
        tfile = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        tfile.write(img_data)
        tfile.close()
        return tfile.name, "✅ Erfolg"
        
    except Exception as e:
        logger.exception(f"Image Gen Error: {str(e)}")
        return None, f"🔥 Fehler: {str(e)}"
    
# ==========================================
# 📋 PREDEFINED PROMPT TEMPLATES
# ==========================================

TRANSCRIPT_PROMPTS = {
    "Veranstaltungsrückblick": """Schreibe auf der Grundlage dieses automatisch erstellten Transkripts einen professionellen Veranstaltungsrückblick für unsere Website.

Berücksichtige dabei:
- Hauptthemen und Kernaussagen
- Wichtige Diskussionspunkte
- Atmosphäre und Teilnehmerfeedback (falls erwähnt)
- Schlussfolgerungen und Ausblicke

Transkript:
{transcript}

Zusätzliche Hinweise:
{notes}""",

    "Zusammenfassung": """Erstelle eine prägnante Zusammenfassung dieses Transkripts.

Gliedere nach:
1. Hauptthema und Zielsetzung
2. Kernaussagen (3-5 Punkte)
3. Wichtigste Erkenntnisse
4. Offene Fragen oder Ausblick

Transkript:
{transcript}

Zusätzliche Hinweise:
{notes}""",

    "Pressemitteilung": """Verfasse eine Pressemitteilung basierend auf diesem Veranstaltungstranskript.

Die Pressemitteilung sollte:
- Einen aufmerksamkeitsstarken Titel haben
- Die 5 W-Fragen beantworten (Wer, Was, Wann, Wo, Warum)
- Zitate von Rednern einbinden
- Maximal 300 Wörter umfassen

Transkript:
{transcript}

Zusätzliche Hinweise:
{notes}""",

    "Social Media Posts": """Erstelle 3 verschiedene Social Media Posts basierend auf diesem Transkript:

1. LinkedIn Post (professionell, 150 Wörter)
2. Twitter/X Thread (3-4 Tweets)
3. Instagram Caption (ansprechend, mit Emojis)

Transkript:
{transcript}

Zusätzliche Hinweise:
{notes}""",

    "Protokoll": """Erstelle ein formelles Protokoll dieser Veranstaltung basierend auf dem Transkript.

Struktur:
1. Datum, Ort, Teilnehmer
2. Tagesordnung/Themen
3. Diskussionspunkte und Beschlüsse
4. Offene Aufgaben und Verantwortliche
5. Nächste Schritte

Transkript:
{transcript}

Zusätzliche Hinweise:
{notes}""",

    "FAQ Generierung": """Analysiere dieses Transkript und erstelle daraus eine FAQ-Liste mit 8-10 häufig gestellten Fragen und detaillierten Antworten.

Die FAQs sollten:
- Klar und verständlich formuliert sein
- Die wichtigsten Informationen abdecken
- Für Website-Besucher nützlich sein

Transkript:
{transcript}

Zusätzliche Hinweise:
{notes}""",

    "Zitate-Sammlung": """Extrahiere die wichtigsten und aussagekräftigsten Zitate aus diesem Transkript.

Für jedes Zitat gib an:
- Sprecher (falls bekannt)
- Kontext
- Warum dieses Zitat bedeutsam ist

Transkript:
{transcript}

Zusätzliche Hinweise:
{notes}""",

    "Blogbeitrag": """Schreibe einen ansprechenden Blogbeitrag basierend auf diesem Veranstaltungstranskript.

Der Blogbeitrag sollte:
- Einen einladenden Einstieg haben
- Die Hauptinhalte verständlich aufbereiten
- Persönliche Eindrücke einbinden
- Mit einem Call-to-Action enden
- Ca. 500-800 Wörter umfassen

Transkript:
{transcript}

Zusätzliche Hinweise:
{notes}""",

    "Eigener Prompt": """{transcript}

{notes}"""
}

# ==========================================
# 📱 PWA CONFIGURATION (Meta Tags & JS Only)
# ==========================================
PWA_HEAD = """
<meta name="viewport" content="width=device-width, initial-scale=1, maximum-scale=1, user-scalable=no">
<meta name="theme-color" content="#ffffff">
<meta name="mobile-web-app-capable" content="yes">
<link rel="manifest" href="/manifest.json" crossorigin="use-credentials">
<script src="/static/pwa.js" defer></script>
"""

# ==========================================
# 🎨 CUSTOM CSS (Passed to css= argument)
# ==========================================
# ==========================================
# 🎨 CUSTOM CSS (Aggressive Mobile Fix)
# ==========================================
CUSTOM_CSS = """
/* 1. NUCLEAR RESET: Kill all spacing & focus rings */
:root, body, .gradio-container {
    --section-gap: 0px !important;
    --block-padding: 0px !important;
    --container-radius: 0px !important;
    --block-radius: 0px !important;
    --input-focus-ring-color: transparent !important;
    --input-focus-border-color: transparent !important;
    --shadow-spread: 0px !important;
    --block-border-width: 0px !important;
    --block-shadow: none !important;
    background-color: #f9fafb !important;
}

/* Force container to fill screen */
.gradio-container {
    max-width: 100% !important;
    width: 100% !important;
    margin: 0 !important;
    padding: 0 !important;
    min-height: 100vh !important;
}

/* Remove internal padding from the main wrapper */
.gradio-container > .main {
    padding: 0 !important;
    gap: 0 !important;
}

/* 2. HEADER */
.compact-header {
    background: white;
    border-bottom: 1px solid #e5e7eb;
    padding: 8px 12px !important;
    min-height: 50px;
}

/* 3. MOBILE OPTIMIZATIONS (Max Width 768px) */
@media (max-width: 768px) {
    
    /* --- TABS: Icon Only --- */
    /* We target buttons inside the container with class 'icon-nav' */
    .icon-nav button {
        font-size: 0 !important;       /* Hide text */
        padding: 15px 0 !important;    /* Taller touch target */
        display: flex;
        justify-content: center;
        align-items: center;
    }
    
    /* Restore Emoji */
    .icon-nav button::first-letter {
        font-size: 1.5rem !important;
        visibility: visible !important;
    }
    
    /* Selected State */
    .icon-nav button.selected {
        border-bottom: 3px solid #2563eb !important;
        background: #f3f4f6 !important;
    }

    /* --- BUTTONS: Icon Only --- */
    .mobile-icon-only {
        font-size: 0 !important;
        padding: 0 !important;
        display: flex;
        justify-content: center;
        align-items: center;
        min-height: 45px !important;
        width: 100% !important;
    }
    
    .mobile-icon-only::first-letter {
        font-size: 1.4rem !important;
        visibility: visible !important;
    }

    /* --- CHAT HEIGHT --- */
    #chat_window {
        height: 70vh !important;
        max-height: 70vh !important;
        overflow-y: auto !important;
    }

    /* Hide footer */
    footer { display: none !important; }
}
"""

# ==========================================
# UI HELPERS
# ==========================================

def update_transcript_chat_ui(prov, force_all=False, user_state=None):
    """
    Same logic as update_c_ui but for the transcript "Send to Chat" section
    """
    # Remove suffix if present
    clean_prov = prov.replace(" ⚠️", "").strip()

    # Check implementation
    if not is_provider_implemented(clean_prov):
        return (
            gr.update(choices=[], value=None),
            f"⚠️ {clean_prov} ist noch nicht verfügbar"
        )
    
    p_data = PROVIDERS.get(prov, {})
    badge = p_data.get("badge", "")
    api_key = API_KEYS.get(p_data.get("key_name"))
    
    # 1. Fetch
    models, error = fetch_available_models(prov, api_key)
    choices = []
    
    if models:
        # Filter for Chat types if Poe
        if prov == "Poe":
            choices = [(m["name"], m["id"]) for m in models if "Chat" in m.get("type", "Chat")]
        else:
            choices = [(m["name"], m["id"]) for m in models]
    
    # Fallback
    if not choices:
        choices = [(m, m) for m in p_data.get("chat_models", [])]

    # 2. Filter via Helper
    final_choices = get_filtered_model_choices(prov, choices, force_all, user_state)
    
    # 3. Smart default selection
    default_val = None
    if final_choices:
        # Try to use DEFAULT_CHAT_MODEL if available
        model_ids = [choice[1] for choice in final_choices]
        if DEFAULT_CHAT_MODEL in model_ids:
            default_val = DEFAULT_CHAT_MODEL
        else:
            default_val = final_choices[0][1]
    
    #return gr.update(choices=final_choices, value=default_val), badge
    return gr.update(choices=final_choices, value=default_val), f"<div class='custom-badge'>{badge}</div>"


# --- CHAT UI UPDATE ---
def update_c_ui(prov, force_all=False, user_state=None):
    # Remove suffix if present
    clean_prov = prov.replace(" ⚠️", "").strip()

    # Check implementation
    if not is_provider_implemented(clean_prov):
        return (
            gr.update(choices=[], value=None),
            f"⚠️ {clean_prov} ist noch nicht verfügbar. Bitte wählen Sie einen anderen Provider."
        )
    
    p_data = PROVIDERS.get(prov, {})
    badge = p_data.get("badge", "")
    api_key = API_KEYS.get(p_data.get("key_name"))
    
    # 1. Fetch
    models, error = fetch_available_models(prov, api_key)
    choices = []
    
    if models:
        # Filter for Chat types if Poe
        if prov == "Poe":
            choices = [(m["name"], m["id"]) for m in models if "Chat" in m.get("type", "Chat")]
        else:
            choices = [(m["name"], m["id"]) for m in models]
    
    # Fallback
    if not choices:
        choices = [(m, m) for m in p_data.get("chat_models", [])]

    # 2. Filter via Helper
    final_choices = get_filtered_model_choices(prov, choices, force_all, user_state)
    
    default_val = final_choices[0][1] if final_choices else None
    #return gr.update(choices=final_choices, value=default_val), badge
    return gr.update(choices=final_choices, value=default_val), f"<div class='custom-badge'>{badge}</div>"

# --- VISION UI UPDATE ---

# --- VISION UI UPDATE ---
def update_v_ui(prov, force_all=False, user_state=None):
    styled_badge = f"<div class='custom-badge'>{get_compliance_html(prov)}</div>"

    p_data = PROVIDERS.get(prov, {})
    api_key = API_KEYS.get(p_data.get("key_name"))
    
    # 2. Fetch Models
    models, error = fetch_available_models(prov, api_key)
    choices = []
    
    if models and prov == "Poe":
        choices = [(m["name"], m["id"]) for m in models if "Vision" in m.get("type", "")]
    elif not models or prov != "Poe":
        choices = [(m, m) for m in p_data.get("vision_models", [])]

    # 3. Filter via Helper
    final_choices = get_filtered_model_choices(prov, choices, force_all, user_state)
    default_val = final_choices[0][1] if final_choices else None
    
    return styled_badge, gr.update(choices=final_choices, value=default_val)

def update_v_ui_old(prov, force_all=False, user_state=None):
    # 1. Handle Badge Styling (Dark background, White text)
    raw_html = get_compliance_html(prov)
    styled_badge = f"""
    <div style="
        background-color: #374151; 
        color: #ffffff !important; 
        padding: 0 12px; 
        border-radius: 8px; 
        border: 1px solid #4b5563; 
        display: flex; 
        align-items: center; 
        justify-content: center;
        height: 42px; 
        font-size: 0.9em;
        white-space: nowrap;
        overflow: hidden;
    ">
        {raw_html}
    </div>
    """

    p_data = PROVIDERS.get(prov, {})
    api_key = API_KEYS.get(p_data.get("key_name"))
    
    # 2. Fetch Models
    models, error = fetch_available_models(prov, api_key)
    choices = []
    
    if models and prov == "Poe":
        choices = [(m["name"], m["id"]) for m in models if "Vision" in m.get("type", "")]
    elif not models or prov != "Poe":
        choices = [(m, m) for m in p_data.get("vision_models", [])]

    # 3. Filter via Helper
    final_choices = get_filtered_model_choices(prov, choices, force_all, user_state)
    default_val = final_choices[0][1] if final_choices else None
    
    # Return Badge and Model Update
    return styled_badge, gr.update(choices=final_choices, value=default_val)

# --- IMAGE UI UPDATE ---
def update_g_ui(prov, force_all=False, user_state=None):
    styled_badge = f"<div class='custom-badge'>{get_compliance_html(prov)}</div>"

    p_data = PROVIDERS.get(prov, {})
    api_key = API_KEYS.get(p_data.get("key_name"))
    
    # 2. Fetch Models
    models, error = fetch_available_models(prov, api_key)
    choices = []
    
    if models and prov == "Poe":
        choices = [(f"{m['name']}", m['id']) for m in models if m.get('type') in ["Bild-Gen", "Video-Gen"]]
    
    if not choices:
        choices = [(m, m) for m in p_data.get("image_models", [])]
        
    # 3. Filter via Helper
    final_choices = get_filtered_model_choices(prov, choices, force_all, user_state)
    default_val = final_choices[0][1] if final_choices else None
    
    return styled_badge, gr.update(choices=final_choices, value=default_val)

def update_g_ui_old(prov, force_all=False, user_state=None):
    # 1. Handle Badge Styling
    raw_html = get_compliance_html(prov)
    styled_badge = f"""
    <div style="
        background-color: #374151; 
        color: #ffffff !important; 
        padding: 0 12px; 
        border-radius: 8px; 
        border: 1px solid #4b5563; 
        display: flex; 
        align-items: center; 
        justify-content: center;
        height: 42px; 
        font-size: 0.9em;
        white-space: nowrap;
        overflow: hidden;
    ">
        {raw_html}
    </div>
    """

    p_data = PROVIDERS.get(prov, {})
    api_key = API_KEYS.get(p_data.get("key_name"))
    
    # 2. Fetch Models
    models, error = fetch_available_models(prov, api_key)
    choices = []
    
    if models and prov == "Poe":
        choices = [(f"{m['name']}", m['id']) for m in models if m.get('type') in ["Bild-Gen", "Video-Gen"]]
    
    if not choices:
        choices = [(m, m) for m in p_data.get("image_models", [])]
        
    # 3. Filter via Helper
    final_choices = get_filtered_model_choices(prov, choices, force_all, user_state)
    default_val = final_choices[0][1] if final_choices else None
    
    return styled_badge, gr.update(choices=final_choices, value=default_val)

# --- TRANSCRIPTION UI UPDATE ---
# Signature: (badge, t_model_update, gladia_vis, whisper_vis, deepgram_vis, assemblyai_vis)
def update_t_ui(prov, force_all=False):
    # Retrieve badge text
    badge_html = get_compliance_html(prov)
    
    # 2. Handle Models
    is_whisper = prov in ["Mistral", "Scaleway", "Groq"]
    show_model_dropdown = prov in ["Deepgram", "AssemblyAI"] or is_whisper
    
    choices = []
    default_val = None

    if show_model_dropdown:
        p_data = PROVIDERS.get(prov, {})
        raw_list = p_data.get("audio_models", [])
        choices = [(m, m) for m in raw_list]
        if choices:
            default_val = choices[0][1]

    # Return updates - Badge wrapped in the custom div
    return (
        f"<div class='custom-badge'>{badge_html}</div>",
        gr.update(visible=show_model_dropdown, choices=choices, value=default_val),
        gr.update(visible=is_whisper) 
    )
def update_t_ui_old(prov, force_all=False):
    # 1. Handle Badge Formatting
    raw_html = get_compliance_html(prov)
    
    # CSS: Multi-line support, dark background, centered
    styled_badge = f"""
    <div style="
        background-color: #374151; 
        color: #ffffff !important; 
        padding: 8px 12px; 
        border-radius: 8px; 
        border: 1px solid #4b5563; 
        margin-top: 10px; 
        display: flex; 
        align-items: center; 
        justify-content: center;
        text-align: center;
        font-size: 0.9em;
    ">
        {raw_html}
    </div>
    """

    # 2. Handle Models
    is_whisper = prov in ["Mistral", "Scaleway", "Groq"]
    show_model_dropdown = prov in ["Deepgram", "AssemblyAI"] or is_whisper
    
    choices = []
    default_val = None

    if show_model_dropdown:
        p_data = PROVIDERS.get(prov, {})
        raw_list = p_data.get("audio_models", [])
        choices = [(m, m) for m in raw_list]
        if choices:
            default_val = choices[0][1]

    # Return updates
    return (
        styled_badge,
        gr.update(visible=show_model_dropdown, choices=choices, value=default_val),
        gr.update(visible=is_whisper) 
    )
    
# --- Chat functions ---
def user_msg(msg, hist):
    # 1. CLI DEBUG LOG
    print(f"\n[DEBUG] 👤 user_msg called at {time.time()}")
    
    if not msg or not msg.strip():
        print("[DEBUG] ❌ Empty message ignored")
        return "", hist
    
    # 2. Simple Deduplication (Prevents double visual bubble)
    if hist and len(hist) > 0:
        if hist[-1]['role'] == 'user' and hist[-1]['content'] == msg:
            print(f"[DEBUG] 🚫 DUPLICATE user message detected. Ignoring.")
            return "", hist

    return "", hist + [{"role": "user", "content": msg}]

def bot_msg(hist, prov, mod, temp, sys, key, r_effort, r_tokens, user_state):
    """
    Execute chat with Debounce Lock to prevent double-generation
    """
    current_time = time.time()
    
    # --- CLI DEBUG LOG ---
    print(f"[DEBUG] 🤖 bot_msg called at {current_time}")

    # 1. DEBOUNCE LOCK (The Fix)
    # Check if we ran this function for this user less than 1.0 second ago
    last_run = user_state.get('last_chat_time', 0)
    if current_time - last_run < 1.0:
        print(f"[DEBUG] 🛑 DEBOUNCE: Ignoring double trigger ({current_time - last_run:.4f}s)")
        yield hist
        return

    # Update timestamp
    user_state['last_chat_time'] = current_time

    # 2. Standard Validation
    if not hist: 
        print("[DEBUG] ❌ History is empty")
        yield hist; return

    if hist[-1]["role"] != "user": 
        print(f"[DEBUG] ❌ Last message is '{hist[-1]['role']}', not 'user'. Skipping.")
        yield hist; return

    # 3. Prepare Generation
    last_user_msg = hist[-1]["content"]
    hist.append({"role": "assistant", "content": ""})

    raw_context = hist[:-2]
    clean_context = []
    for m in raw_context:
        role = m["role"]
        if role in ["bot", "model"]: role = "assistant"
        clean_context.append({"role": role, "content": m["content"]})

    try:
        print(f"[DEBUG] 🚀 Calling LLM: {prov} / {mod}...")
        
        # Pass user_state to run_chat
        generator = run_chat(last_user_msg, clean_context, prov, mod, temp, sys, key, r_effort, r_tokens, user_state)
        
        for chunk in generator:
            hist[-1]["content"] = chunk
            yield hist
            
    except Exception as e:
        print(f"[DEBUG] 🔥 ERROR: {str(e)}")
        hist[-1]["content"] = f"🔥 Wrapper Fehler: {str(e)}"
        yield hist

def save_chat(hist, prov, mod, user_state):
    """Save current chat to database"""
    try:
        if not user_state or not user_state.get("id"):
            logger.warning("Save chat failed: User not logged in")
            return "❌ Bitte anmelden"

        user_id = user_state["id"]
        logger.info(f"Attempting to save chat for user {user_id}")

        if not hist or len(hist) == 0:
            logger.warning("Save chat failed: Empty chat history")
            return "❌ Kein Chat zum Speichern"

        # Generate title
        first_content = hist[0].get("content", "") if isinstance(hist[0], dict) else str(hist[0])
        title = first_content[:50] + "..." if len(first_content) > 50 else first_content

        logger.info(f"Saving chat with title: {title}")
        chat_id = save_chat_history(user_id, prov, mod, hist, title)
        logger.info(f"Chat saved successfully with ID: {chat_id}")

        return f"✅ Chat gespeichert (ID: {chat_id})"

    except Exception as e:
        logger.exception(f"Error saving chat: {str(e)}")
        return f"🔥 Fehler beim Speichern: {str(e)}"
    
def load_chat_list(user_state):
    """Load user's chat history (legacy format)"""
    if not user_state or not user_state.get("id"):
        return [["Bitte anmelden", "", "", ""]]

    chats = get_user_chat_history(user_state["id"])
    data = []
    for chat in chats:
        data.append([
            chat.id,
            chat.timestamp.strftime("%Y-%m-%d %H:%M"),
            chat.title or "Ohne Titel",
            chat.model
        ])
    return data if data else [["Keine Chats vorhanden", "", "", ""]]

def load_chat_list_with_state(user_state=None):
    """Load user's chat history returning both Display Data and State Data"""
    if not user_state or not user_state.get("id"):
        return [["Bitte anmelden", "", "", ""]], []

    chats = get_user_chat_history(user_state["id"])
    
    # Clean data for State (ID, Date, Title, Model)
    state_data = [[c.id, c.timestamp.strftime("%Y-%m-%d %H:%M"), c.title or "Ohne Titel", c.model] for c in chats]
    
    if not state_data:
        return [["Keine Chats vorhanden", "", "", ""]], []
        
    return state_data, state_data

def select_chat_row(evt: gr.SelectData, state_data):
    """Smart Selection for Chat List"""
    try:
        if not state_data:
            return None
            
        row_idx = evt.index[0]
        if row_idx < len(state_data):
            real_row = state_data[row_idx]
            # ID is in column 0
            return int(real_row[0])
    except Exception as e:
        logger.error(f"Selection error: {e}")
    return None

# ==========================================
# 📄 UNIVERSAL CONTENT EXTRACTOR (Robust & Granular)
# ==========================================
import mimetypes
import chardet
import shutil
import subprocess
import traceback
import tempfile
from io import StringIO

# --- Granular Imports ---
HAS_OCR = False
HAS_PDF_IMG = False
HAS_DOCX = False
HAS_FITZ = False
HAS_PDF_READER = False
HAS_PANDAS = False

try:
    import pytesseract
    HAS_OCR = True
except ImportError: logger.warning("⚠️ pytesseract missing")

try:
    from pdf2image import convert_from_path
    HAS_PDF_IMG = True
except ImportError: logger.warning("⚠️ pdf2image missing")

try:
    from docx import Document
    HAS_DOCX = True
except ImportError: logger.warning("⚠️ python-docx missing")

try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except ImportError: logger.warning("⚠️ PyMuPDF (fitz) missing")

try:
    from pypdf import PdfReader
    HAS_PDF_READER = True
except ImportError: logger.warning("⚠️ pypdf missing")

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError: logger.warning("⚠️ pandas missing")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError: 
    HAS_PPTX = False
    logger.warning("⚠️ python-pptx missing")

import zipfile
from xml.etree import ElementTree as ET

class UniversalExtractor:
    @staticmethod
    def extract(filepath):
        if not filepath or not os.path.exists(filepath):
            return "❌ Datei nicht gefunden."
            
        mime_type, _ = mimetypes.guess_type(filepath)
        ext = os.path.splitext(filepath)[1].lower()
        filename = os.path.basename(filepath)
        
        logger.info(f"🔍 Extracting: {filename} (Type: {ext})")
        
        try:
            # 1. Images (OCR)
            if (mime_type and mime_type.startswith('image/')) or ext in ['.png', '.jpg', '.jpeg', '.webp', '.bmp', '.tiff']:
                return UniversalExtractor._extract_image(filepath)
            
            # 2. PDF
            elif ext == '.pdf':
                return UniversalExtractor._extract_pdf(filepath)
            
            # 3. Modern Word
            elif ext == '.docx':
                if HAS_DOCX:
                    try: 
                        return UniversalExtractor._extract_docx(filepath)
                    except Exception as e: 
                        logger.warning(f"Docx extraction failed: {e}")
                return UniversalExtractor._extract_with_cli_tool(filepath)

            # 4. Excel and Powerpoint
            elif ext in ['.xls', '.xlsx', '.csv']:
                return UniversalExtractor._extract_excel(filepath)
            
            elif ext == '.pptx':
                return UniversalExtractor._extract_pptx(filepath)

            # 5. HTML (NEW: Multiple fallbacks)
            elif ext in ['.html', '.htm']:
                return UniversalExtractor._extract_html(filepath)

            # 6. Ebooks & Legacy Docs
            elif ext in ['.epub', '.mobi', '.azw3', '.fb2', '.doc', '.odt', '.rtf']:
                return UniversalExtractor._extract_with_cli_tool(filepath)
                
            # 7. Text/Code Fallback
            else:
                return UniversalExtractor._extract_plain_text(filepath)
                
        except Exception as e:
            logger.error(f"Critical Extractor Error: {e}")
            logger.error(traceback.format_exc())
            return f"[Systemfehler: {str(e)}]"
        
    @staticmethod
    def _extract_pptx(path):
        """
        Robust PPTX extraction preserving slide order, notes, and comments.
        """
        if not HAS_PPTX:
            return "[❌ 'python-pptx' fehlt. Bitte `pip install python-pptx` installieren.]"

        try:
            prs = Presentation(path)
            output = [f"[PPTX Extraction: {os.path.basename(path)}]\n"]
            
            # XML Namespaces
            ns = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            }

            # 1. Extract Comments (Zip Logic)
            all_comments = {}
            try:
                with zipfile.ZipFile(path, 'r') as zf:
                    # Iterate purely based on slide count, assuming standard structure
                    for i in range(1, len(prs.slides) + 1):
                        rel_file = f'ppt/slides/_rels/slide{i}.xml.rels'
                        if rel_file in zf.namelist():
                            with zf.open(rel_file) as rf:
                                rel_tree = ET.parse(rf)
                                for rel in rel_tree.findall(".//r:Relationship", ns):
                                    tgt = rel.get('Target')
                                    if tgt and "comments" in tgt:
                                        # Resolve path (../comments/comment1.xml -> ppt/comments/comment1.xml)
                                        fname = os.path.basename(tgt)
                                        comment_path = f'ppt/comments/{fname}'
                                        
                                        if comment_path in zf.namelist():
                                            slide_comments = all_comments.setdefault(i, [])
                                            with zf.open(comment_path) as cf:
                                                c_tree = ET.parse(cf)
                                                # Look for comment text
                                                for comment in c_tree.findall('.//p:cm', ns):
                                                    text_node = comment.find('.//p:text', ns)
                                                    if text_node is not None and text_node.text:
                                                        slide_comments.append(text_node.text)
            except Exception as e:
                logger.warning(f"PPTX Comment extraction warning: {e}")

            # 2. Extract Slides (Text + Notes)
            for i, slide in enumerate(prs.slides):
                slide_num = i + 1
                slide_buffer = []
                
                # -- Main Text (Sorted Top-to-Bottom) --
                # python-pptx shapes usually have .top, but we check safety
                shapes = [s for s in slide.shapes if s.has_text_frame]
                shapes.sort(key=lambda s: (s.top if hasattr(s, 'top') else 0, s.left if hasattr(s, 'left') else 0))
                
                for shape in shapes:
                    shape_text = []
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            shape_text.append(paragraph.text.strip())
                    if shape_text:
                        slide_buffer.append("\n".join(shape_text))

                # -- Speaker Notes --
                if slide.has_notes_slide:
                    notes_tf = slide.notes_slide.notes_text_frame
                    if notes_tf and notes_tf.text.strip():
                        slide_buffer.append(f"\n[Speaker Notes]:\n{notes_tf.text.strip()}")

                # -- Formatting Output --
                output.append(f"## Slide {slide_num}")
                if slide_buffer:
                    output.append("\n\n".join(slide_buffer))
                else:
                    output.append("(No text content)")

                # -- Add Comments --
                if slide_num in all_comments:
                    output.append("\n[Comments]:")
                    for c in all_comments[slide_num]:
                        output.append(f"- {c}")
                
                output.append("\n---\n")

            final_text = "\n".join(output)
            
            if len(final_text) < 50:
                return "[⚠️ PPTX scheint leer zu sein oder enthält nur Bilder]"
                
            return final_text

        except Exception as e:
            logger.error(f"PPTX Error: {e}")
            return f"[❌ Fehler beim Lesen der PPTX: {str(e)}]"

    @staticmethod
    def _extract_html(path):
        """
        Extract text from HTML with multiple fallback methods.
        Priority: BeautifulSoup > CLI tools > Plain text
        """
        logger.info(f"Attempting HTML extraction: {path}")
        
        # Method 1: BeautifulSoup (Most reliable)
        try:
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style", "meta", "link"]):
                script.decompose()
            
            # Get text
            text = soup.get_text(separator='\n', strip=True)
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            if len(text.strip()) > 50:
                logger.info(f"✅ BeautifulSoup extraction successful: {len(text)} chars")
                return f"[HTML extracted via BeautifulSoup]:\n\n{text}"
            else:
                logger.warning("BeautifulSoup returned minimal text, trying fallback...")
                
        except Exception as e:
            logger.warning(f"BeautifulSoup extraction failed: {e}")
        
        # Method 2: CLI tools (ebook-convert, pandoc)
        try:
            cli_result = UniversalExtractor._extract_with_cli_tool(path)
            if cli_result and not cli_result.startswith("[") and len(cli_result.strip()) > 50:
                logger.info(f"✅ CLI extraction successful: {len(cli_result)} chars")
                return cli_result
            logger.warning("CLI extraction returned minimal/error text")
        except Exception as e:
            logger.warning(f"CLI extraction failed: {e}")
        
        # Method 3: Plain text fallback (last resort)
        try:
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                raw_html = f.read()
            
            if len(raw_html.strip()) > 100:
                logger.warning("⚠️ Using raw HTML as fallback (not cleaned)")
                return f"[Raw HTML - Cleaning failed]:\n\n{raw_html[:10000]}"
                
        except Exception as e:
            logger.error(f"Even plain text reading failed: {e}")
        
        return "[❌ HTML konnte nicht verarbeitet werden. Datei möglicherweise beschädigt.]"

    
    @staticmethod
    def _extract_with_cli_tool(path):
        """Uses ebook-converter, calibre, or pandoc with error handling"""
        # Check common paths explicitly
        candidates = [
            "ebook-converter", 
            "/var/www/transkript_app/venv/bin/ebook-converter",
            "ebook-convert", 
            "pandoc"
        ]
        
        tool = None
        for c in candidates:
            if shutil.which(c) or os.path.exists(c):
                tool = c
                break
        
        if not tool:
            logger.warning(f"CLI Tools missing. PATH: {os.environ.get('PATH')}")
            return "[⚠️ Kein Konverter (ebook-converter/pandoc) gefunden. Installieren Sie: apt-get install calibre pandoc]"

        try:
            with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
                out_path = tmp.name

            cmd = [tool, path, out_path]
            if "pandoc" in tool:
                cmd = [tool, path, "-t", "plain", "-o", out_path]

            logger.info(f"Running CLI: {cmd}")
            
            # Run with timeout and capture stderr
            result = subprocess.run(
                cmd, 
                capture_output=True, 
                text=True, 
                timeout=120
            )
            
            # Check for errors but be lenient (some tools return non-zero but still produce output)
            if result.returncode != 0:
                logger.warning(f"CLI returned code {result.returncode}. Stderr: {result.stderr[:500]}")
                
                # Specific error: libxml2 version mismatch
                if "libxml2" in result.stderr or "html5-parser" in result.stderr:
                    logger.error("❌ ebook-convert has library conflict. Using fallback...")
                    return f"[❌ ebook-convert Bibliothekskonflikt - Verwende alternative Methode]"

            # Try to read output even if return code was non-zero
            if os.path.exists(out_path):
                with open(out_path, 'r', encoding='utf-8', errors='replace') as f:
                    text = f.read()
                os.remove(out_path)
                
                if text.strip():
                    logger.info(f"✅ CLI extraction successful: {len(text)} chars")
                    return text

            # If we got here, extraction failed
            error_msg = result.stderr[:500] if result.stderr else "Unknown error"
            logger.error(f"CLI extraction produced no output. Error: {error_msg}")
            return f"[❌ CLI Konvertierung fehlgeschlagen: {error_msg}]"

        except subprocess.TimeoutExpired:
            logger.error("CLI tool timed out after 120 seconds")
            return "[❌ Konvertierung Timeout (>120s). Datei zu groß?]"
        except Exception as e:
            logger.error(f"CLI tool exception: {e}")
            logger.error(traceback.format_exc())
            return f"[❌ CLI Fehler: {str(e)}]"

    @staticmethod
    def _extract_image(path):
        """Extract text from image with detailed error reporting"""
        if not HAS_OCR:
            return "[❌ OCR Library 'pytesseract' fehlt. Bitte installieren: pip install pytesseract]"
        
        try:
            # Check if tesseract binary exists
            if not shutil.which("tesseract"):
                return "[❌ Systemfehler: 'tesseract' Binary fehlt. Bitte installieren: apt-get install tesseract-ocr tesseract-ocr-deu]"
            
            logger.info(f"Running OCR on: {path}")
            
            # Try to open image first
            try:
                img = Image.open(path)
                logger.info(f"Image opened successfully: {img.size}, mode: {img.mode}")
            except Exception as e:
                return f"[❌ Bild konnte nicht geöffnet werden: {str(e)}]"
            
            # Run OCR
            text = pytesseract.image_to_string(img, lang='deu+eng', config='--psm 1')
            
            logger.info(f"OCR complete. Text length: {len(text)}")
            
            if text.strip():
                return f"[OCR Extraktion]:\n\n{text}"
            else:
                return "[⚠️ OCR: Kein Text im Bild erkannt. Bild könnte leer sein oder nur Grafiken enthalten.]"
                
        except pytesseract.TesseractNotFoundError:
            return "[❌ Tesseract nicht gefunden. Bitte installieren: apt-get install tesseract-ocr tesseract-ocr-deu]"
        except Exception as e:
            logger.error(f"OCR error: {e}")
            logger.error(traceback.format_exc())
            return f"[❌ OCR Fehler: {str(e)}]"

    @staticmethod
    def _extract_pdf(path):
        text = ""
        # A. PyMuPDF
        if HAS_FITZ:
            try:
                with fitz.open(path) as doc:
                    for page in doc: text += page.get_text() + "\n"
                if len(text.strip()) > 50: return text
                logger.info("PyMuPDF returned empty text (Scanned PDF?)")
            except Exception as e: logger.warning(f"PyMuPDF failed: {e}")

        # B. PyPDF
        if HAS_PDF_READER:
            try:
                reader = PdfReader(path)
                pypdf_text = ""
                for page in reader.pages: pypdf_text += (page.extract_text() or "") + "\n"
                if len(pypdf_text.strip()) > 50: return f"[PyPDF]:\n{pypdf_text}"
            except Exception as e: logger.warning(f"PyPDF failed: {e}")

        # C. CLI Tool (ebook-converter)
        try:
            cli_text = UniversalExtractor._extract_with_cli_tool(path)
            if len(cli_text.strip()) > 50 and "Fehler" not in cli_text: return cli_text
        except: pass

        # D. OCR (Scan)
        if HAS_PDF_IMG and HAS_OCR:
            try:
                logger.info("Attempting OCR on PDF...")
                return UniversalExtractor._extract_scanned_pdf(path)
            except Exception as e: return f"[OCR Fehler: {e}]"
        
        return "[PDF unlesbar: Keine Text-Ebene und OCR Tools fehlen]"

    @staticmethod
    def _extract_scanned_pdf(path):
        text = "[OCR Scan Modus]\n"
        images = convert_from_path(path, first_page=1, last_page=5)
        for i, img in enumerate(images):
            text += pytesseract.image_to_string(img, lang='deu+eng') + "\n"
        return text

    @staticmethod
    def _extract_docx(path):
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])

    @staticmethod
    def _extract_excel(path):
        if not HAS_PANDAS: return "[Pandas fehlt]"
        try:
            if path.endswith('.csv'): df = pd.read_csv(path)
            else: df = pd.read_excel(path)
            return df.to_markdown(index=False)
        except Exception as e: return f"[Excel Fehler: {e}]"

    @staticmethod
    def _extract_plain_text(path):
        try:
            with open(path, 'rb') as f:
                raw = f.read(50000)
                enc = chardet.detect(raw)['encoding'] or 'utf-8'
            with open(path, 'r', encoding=enc, errors='replace') as f:
                return f.read()
        except Exception as e: return f"[Read Error: {e}]"

def undo_last_attachment(hist):
    """Removes the last message if it was an attachment"""
    if not hist or len(hist) == 0:
        return hist, "❌ Keine Nachrichten"
    
    last_msg = hist[-1]
    # Check if the last message is from user and looks like an attachment header
    if last_msg["role"] == "user":
        content = last_msg["content"]
        if content.startswith("[Datei:") or content.startswith("[Transkript") or content.startswith("[Vision"):
            hist.pop()
            return hist, "✅ Letzten Anhang entfernt"
            
    return hist, "⚠️ Letzte Nachricht war kein Anhang"


SENSITIVE_FILE_WARNING = """
⚠️ **DATENSCHUTZ-WARNUNG**

**Bitte KEINE der folgenden Dateien hochladen:**
- ❌ Personalakten oder Bewerbungen
- ❌ Seelsorge-Protokolle oder Beichtgespräche  
- ❌ Medizinische Unterlagen
- ❌ Dokumente mit Gesundheitsdaten
- ❌ Interne strategische Dokumente (Verschlusssachen)
- ❌ Passwörter oder Zugangsdaten

**Art. 9 DS-GVO:** Besondere Kategorien (Religion, Gesundheit) sind besonders zu schützen!

✅ **Erlaubt:** Öffentliche Texte, anonymisierte Daten, allgemeine Fachliteratur
"""

def check_sensitive_filename(filename: str) -> tuple:
    """
    Returns: (is_suspicious, warning_message)
    """
    suspicious_keywords = [
        "personal", "bewerbung", "lebenslauf", "cv", "resume",
        "seelsorge", "beichte", "patient", "diagnose", "arzt",
        "geheim", "vertraulich", "verschluss", "confidential",
        "passwort", "password", "credentials", "token"
    ]
    
    lower_name = filename.lower()
    for keyword in suspicious_keywords:
        if keyword in lower_name:
            return True, f"⚠️ Dateiname enthält sensiblen Begriff: '{keyword}'"
    
    return False, ""

def attach_content_to_chat(hist, attach_type, attach_id, custom_text, uploaded_files, 
                          sb_files, user_state):
    """
    Attach content to chat with automatic chunking support.
    Uses saved user settings for chunking behavior.
    """
    # SECURITY CHECK 1: Authentication
    if not user_state or not user_state.get("id"):
        return hist, "❌ Bitte anmelden"
    
    user_id = user_state["id"]

    # SECURITY CHECK: File upload warnings (with tracking)
    if attach_type == "Datei uploaden" and uploaded_files:
        # Check if user needs warning
        show_warning, count = check_user_confirmation(user_id, "upload_warning")
        
        if show_warning:
            files_list = uploaded_files if isinstance(uploaded_files, list) else [uploaded_files]
            for file_obj in files_list:
                if hasattr(file_obj, 'name'):
                    is_suspicious, warning = check_sensitive_filename(os.path.basename(file_obj.name))
                    if is_suspicious:
                        # Record confirmation
                        record_user_confirmation(user_id, "upload_warning")
                        
                        # Show warning with counter
                        remaining = 3 - count - 1
                        if remaining > 0:
                            warning_suffix = f"\n\n_(Diese Warnung wird noch {remaining}x angezeigt)_"
                        else:
                            warning_suffix = "\n\n_(Dies ist die letzte Warnung. Zukünftig wird sie nicht mehr angezeigt.)_"
                        
                        return hist, f"{SENSITIVE_FILE_WARNING}\n\n{warning}{warning_suffix}"
    
    # Load user settings
    settings = get_user_settings(user_id)
    auto_chunk = settings.auto_chunk_enabled
    chunk_size = settings.chunk_size
    chunk_overlap = settings.chunk_overlap
    
    full_content_to_add = ""
    status_messages = []  # Consistent variable name

    def process_text_content(text, label, source_type="text"):
        """
        Process any text content with optional chunking.
        
        Args:
            text: The text content to process
            label: Display label (e.g., filename, "Transkript #123")
            source_type: Type of source ("file", "transcript", "vision", "custom")
        
        Returns:
            (full_block, status_msg)
        """
        if not text or not text.strip():
            return "", f"❌ {label}: Leer"
        
        estimated_tokens = estimate_tokens(text)
        
        # Check if chunking is needed
        if auto_chunk and estimated_tokens > chunk_size:
            chunks = split_content_into_chunks(text, chunk_size, chunk_overlap)
            
            # Create formatted output with chunks
            full_block = f"\n\n=== 📎 {label} - {len(chunks)} Teile ===\n"
            for i, chunk in enumerate(chunks, 1):
                full_block += f"\n--- Teil {i}/{len(chunks)} ---\n{chunk}\n"
            
            status = f"📄 {label}: {len(chunks)} Teile ({estimated_tokens:,} tokens)"
            return full_block, status
        
        else:
            # No chunking needed - but still truncate if extremely long
            truncated = False
            if len(text) > 150000:
                text = text[:150000] + "\n... [Gekürzt wegen Länge]"
                truncated = True
            
            full_block = f"\n\n=== 📎 {label} ===\n{text}\n"
            
            if truncated:
                status = f"📄 {label}: OK ({estimated_tokens:,} tokens, gekürzt)"
            else:
                status = f"📄 {label}: OK ({estimated_tokens:,} tokens)"
            
            return full_block, status

    def process_file_path(path, source_label):
        """Extract and process file content"""
        try:
            fname = os.path.basename(path)
            
            # Extract content
            extracted = UniversalExtractor.extract(path)
            
            # Better error detection: Check for actual error indicators
            is_error = (
                not extracted or 
                extracted.startswith("❌") or 
                extracted.startswith("[❌") or
                extracted.startswith("[⚠️") or
                (extracted.startswith("[") and "fehlt" in extracted.lower()) or
                (extracted.startswith("[") and "error" in extracted.lower()) or
                (extracted.startswith("[") and "nicht gefunden" in extracted.lower())
            )
            
            if is_error:
                # Extraction failed or returned error
                logger.warning(f"Extraction failed for {fname}: {extracted[:200]}")
                return "", f"❌ {fname}: Konnte nicht verarbeitet werden"
            
            # Check if we got actual content (not just headers)
            # Remove common headers to check actual content
            content_to_check = extracted
            for header in ["[HTML extracted via BeautifulSoup]:", "[OCR Extraktion]:", "[OCR]:", "[PyPDF]:"]:
                content_to_check = content_to_check.replace(header, "")
            
            if len(content_to_check.strip()) < 10:
                logger.warning(f"Extracted content too short for {fname}: {len(content_to_check)} chars")
                return "", f"❌ {fname}: Datei enthält zu wenig Text"
            
            # Process extracted text
            return process_text_content(extracted, f"Datei: {fname} ({source_label})", "file")
            
        except Exception as e:
            logger.error(f"Error processing file {path}: {e}")
            return "", f"❌ {os.path.basename(path)}: Fehler bei Verarbeitung"

    try:
        # === 1. UPLOADED FILES ===
        if attach_type == "Datei uploaden" and uploaded_files:
            files_list = uploaded_files if isinstance(uploaded_files, list) else [uploaded_files]
            
            for file_obj in files_list:
                if hasattr(file_obj, 'name') and file_obj.name:
                    content_block, status_msg = process_file_path(file_obj.name, "Upload")
                    full_content_to_add += content_block
                    status_messages.append(status_msg)
                else:
                    status_messages.append("❌ Ungültige Datei")

        # === 2. STORAGE BOX FILES ===
        elif attach_type == "Storage Box Datei" and sb_files:
            files_list = sb_files if isinstance(sb_files, list) else [sb_files]
            
            for f_path in files_list:
                try:
                    # Normalize path
                    if not f_path.startswith("/"): 
                        f_path = os.path.join(STORAGE_MOUNT_POINT, f_path)
                    
                    if not os.path.exists(f_path):
                        status_messages.append(f"❌ {os.path.basename(f_path)}: Nicht gefunden")
                        continue
                    
                    # Copy to temp for processing
                    local_temp = copy_storage_file_to_temp(f_path)
                    content_block, status_msg = process_file_path(local_temp, "Cloud")
                    full_content_to_add += content_block
                    status_messages.append(status_msg)
                    
                    # Cleanup
                    try: 
                        os.remove(local_temp)
                    except: 
                        pass
                        
                except Exception as e:
                    logger.error(f"Storage box file error: {e}")
                    status_messages.append(f"❌ {os.path.basename(f_path)}: Fehler")

        # === 3. TRANSCRIPT ===
        elif attach_type == "Transkript":
            if not attach_id:
                return hist, "❌ Transkript-ID fehlt"
            
            try:
                db = get_db()
                trans = db.query(Transcription).filter(
                    Transcription.id == int(attach_id), 
                    Transcription.user_id == user_id
                ).first()
                db.close()
                
                if trans and trans.original_text:
                    content_block, status_msg = process_text_content(
                        trans.original_text,
                        f"Transkript #{trans.id}",
                        "transcript"
                    )
                    full_content_to_add = content_block
                    status_messages.append(status_msg)
                else:
                    return hist, f"❌ Transkript #{attach_id} nicht gefunden"
                    
            except Exception as e:
                logger.error(f"Transcript attachment error: {e}")
                return hist, f"❌ Fehler beim Laden des Transkripts: {str(e)}"

        # === 4. VISION RESULT ===
        elif attach_type == "Vision-Ergebnis":
            if not attach_id:
                return hist, "❌ Vision-ID fehlt"
            
            try:
                # USE DECRYPTION FUNCTION
                vis = get_decrypted_vision(int(attach_id), user_id, user_state)
                
                if vis and vis.result:
                    content_block, status_msg = process_text_content(
                        vis.result,  # ✓ Now decrypted
                        f"Vision #{vis.id}",
                        "vision"
                    )
                    full_content_to_add = content_block
                    status_messages.append(status_msg)
                else:
                    return hist, f"❌ Vision-Ergebnis #{attach_id} nicht gefunden"
                    
                    
            except Exception as e:
                logger.error(f"Vision attachment error: {e}")
                return hist, f"❌ Fehler beim Laden des Vision-Ergebnisses: {str(e)}"

        # === 5. CUSTOM TEXT ===
        elif attach_type == "Eigener Text":
            if not custom_text or not custom_text.strip():
                return hist, "❌ Kein Text eingegeben"
            
            try:
                content_block, status_msg = process_text_content(
                    custom_text,
                    "Eigener Text",
                    "custom"
                )
                full_content_to_add = content_block
                status_messages.append(status_msg)
                
            except Exception as e:
                logger.error(f"Custom text attachment error: {e}")
                return hist, f"❌ Fehler beim Verarbeiten des Textes: {str(e)}"

        else:
            return hist, "❌ Ungültiger Anhang-Typ oder keine Daten"

    except Exception as e:
        logger.exception(f"Attachment Error: {e}")
        return hist, f"🔥 Kritischer Fehler: {str(e)}"

    # === UPDATE CHAT HISTORY ===
    if not hist:
        hist = []
    
    if full_content_to_add and full_content_to_add.strip():
        # Decide on display format based on length
        content_length = len(full_content_to_add)
        
        if content_length > 2000:
            # Use collapsible details for long content
            final_msg_content = (
                f"<details><summary>📎 Angehängter Inhalt "
                f"({content_length:,} Zeichen)</summary>\n\n"
                f"{full_content_to_add}\n\n</details>"
            )
        else:
            final_msg_content = full_content_to_add
        
        # Add to history
        hist.append({"role": "user", "content": final_msg_content})
        
        # Create status message
        if status_messages:
            status = "✅ Angehängt:\n" + "\n".join(f"  • {msg}" for msg in status_messages)
        else:
            status = "✅ Inhalt angehängt"
        
        return hist, status
    
    return hist, "❌ Kein Inhalt zum Anhängen gefunden"

def get_user_prompt_choices(user_state):
    """Get list of user's custom prompt names for dropdown"""
    if not user_state or not user_state.get("id"): 
        return gr.update(choices=[], value=None)
    
    prompts = get_user_custom_prompts(user_state["id"])
    names = [p.name for p in prompts]
    
    # FIX: Return gr.update to set choices, not value
    return gr.update(choices=names, value=None)

def insert_custom_prompt(prompt_name, current_msg, user_state):
    """Insert selected custom prompt text"""
    if not user_state or not user_state.get("id") or not prompt_name: 
        return current_msg
    
    db = get_db()
    p = db.query(CustomPrompt).filter(CustomPrompt.user_id == user_state["id"], CustomPrompt.name == prompt_name).first()
    db.close()
    
    if p:
        if current_msg:
            return current_msg + "\n\n" + p.prompt_text
        return p.prompt_text
    return current_msg

def load_single_chat(chat_id, user_state=None):
    """Load a specific chat into the UI"""
    if not user_state or not user_state.get("id"):
        return None, "❌ Bitte anmelden"
    
    if not chat_id:
        return None, "❌ Ungültige ID"

    # USE THE DECRYPTION FUNCTION!
    chat = get_decrypted_chat(int(chat_id), user_state["id"], user_state=user_state) 
    
    if chat:
        try:
            # chat.messages is already decrypted by get_decrypted_chat()
            messages = json.loads(chat.messages)
            if isinstance(messages, list) and len(messages) > 0:
                return messages, f"✅ Chat '{chat.title}' geladen"
            else:
                return None, "⚠️ Chat-Format veraltet"
        except json.JSONDecodeError as e:
            logger.error(f"JSON decode error for chat {chat_id}: {e}")
            return None, f"🔥 Ladefehler: Ungültiges JSON-Format"
        except Exception as e:
            logger.error(f"Error loading chat {chat_id}: {e}")
            return None, f"🔥 Ladefehler: {str(e)}"
            
    return None, "❌ Chat nicht gefunden"

def delete_chat(chat_id, user_state=None):
    """Delete a chat and update both list and state"""
    def get_fresh_data():
        return load_chat_list_with_state(user_state)

    if not user_state or not user_state.get("id") or not chat_id:
        d, s = get_fresh_data()
        return "❌ Fehler/Auth", d, s

    if delete_chat_history(int(chat_id), user_state["id"]):
        d, s = get_fresh_data()
        return "✅ Chat gelöscht", d, s
        
    d, s = get_fresh_data()
    return "❌ Fehler beim Löschen", d, s

def clear_chat():
    """Clear current chat"""
    return [], ""

def get_filtered_model_choices(provider, available_models_list, force_all, user_state):
    """
    Helper to filter models based on user preferences.
    """
    # If no login, return all models
    if not user_state or not user_state.get("id"):
        return available_models_list

    # 1. Get User Preferences (Ordered ID list)
    pref_ids = get_user_visible_models(user_state["id"], provider)
    
    # If no preferences set, return everything
    if not pref_ids:
        return available_models_list

    val_to_name = {v: k for k, v in available_models_list}
    
    # 2. Separate into Favorites and Others
    fav_choices = []
    seen_ids = set()
    
    for pid in pref_ids:
        if pid in val_to_name:
            fav_choices.append((val_to_name[pid], pid))
            seen_ids.add(pid)
            
    rest_choices = []
    for name, val in available_models_list:
        if val not in seen_ids:
            rest_choices.append((name, val))
            
    # 3. Decision Logic
    if not force_all and fav_choices:
        return fav_choices
    
    return fav_choices + rest_choices

# --- Functions for Transkription ---

# Send to Chat functionality
def send_transcript_to_chat(transcript, template, notes, custom_prompt, provider, model, api_key, user_state):
    # --- SECURITY CHECK ---
    if not user_state or not user_state.get("id"):
        return "⛔ Nicht autorisiert. Bitte anmelden."
    # ----------------------

    """Process transcript with selected prompt and return result"""
    if not transcript or transcript.strip() == "":
        return "❌ Kein Transkript vorhanden."

    # Build the full prompt
    if template == "Eigener Prompt":
        if not custom_prompt or custom_prompt.strip() == "":
            return "❌ Bitte eigenen Prompt eingeben."
        full_prompt = TRANSCRIPT_PROMPTS[template].format(
            transcript=transcript,
            notes=custom_prompt
        )
    else:
        prompt_template_text = TRANSCRIPT_PROMPTS.get(template, TRANSCRIPT_PROMPTS["Zusammenfassung"])
        full_prompt = prompt_template_text.format(
            transcript=transcript,
            notes=notes if notes else "(keine weiteren Hinweise)"
        )

    # Call chat API
    try:
        client = get_client(provider, api_key)
        status = f"🤖 Verarbeite mit {provider}/{model}...\n"

        messages = [
            {"role": "system", "content": "Du bist ein professioneller Redakteur und Content-Spezialist für kirchliche und akademische Veranstaltungen."},
            {"role": "user", "content": full_prompt}
        ]

        response = client.chat.completions.create(
            model=model,
            messages=messages,
            temperature=0.7,
            max_tokens=3000
        )

        result = response.choices[0].message.content
        status += f"✅ Fertig!\n\n---\n\n{result}"
        return status

    except Exception as e:
        return f"🔥 Fehler: {str(e)}"

def generate_and_handle_ui(prompt, provider, model, width, height, steps, key, user_state):

    """Generates image and updates ALL UI components"""
    img_path, status = run_image_gen(prompt, provider, model, width, height, steps, key, user_state)

    if img_path:
        # Copy to a more permanent location for download
        import shutil
        download_dir = "/tmp/gradio_downloads"
        os.makedirs(download_dir, exist_ok=True)
        download_path = os.path.join(download_dir, f"image_{int(time.time())}.jpg")
        shutil.copy2(img_path, download_path)
        
        return (
            img_path,                   # g_out (Preview)
            status,                     # g_stat
            img_path,                   # g_img_path (State)
            download_path,              # g_download_file (Direct file link)
            gr.update(visible=True),    # g_save_btn
            ""                          # g_save_status (Reset)
        )
    else:
        return (
            None, 
            status, 
            None,
            None,                       # g_download_file
            gr.update(visible=False), 
            ""
        )

# --- DB SAVE WRAPPER ---
def process_gallery_save(img_path, provider, prompt, model, user_state):
    """
    Encrypts the image file bytes before saving to persistent storage.
    Fixed: Properly formats the Key for Fernet.
    """
    try:
        import base64 # Ensure import
        
        if not user_state or not user_state.get("id"):
            return "❌ Bitte anmelden", gr.update(visible=True)
        
        user_id = user_state["id"]
        # Get User Master Key
        umk = user_state.get('umk') if user_state else crypto.global_key
        
        if not img_path or not os.path.exists(img_path):
            return "❌ Datei nicht gefunden (Session abgelaufen?)", gr.update(visible=True)

        # 1. Prepare Paths
        permanent_dir = "/var/www/transkript_app/generated_images"
        os.makedirs(permanent_dir, exist_ok=True)
        
        # We append .enc to signify it's encrypted
        filename = f"img_{int(time.time())}_{os.path.basename(img_path)}.enc"
        permanent_path = os.path.join(permanent_dir, filename)

        # 2. Encrypt & Save File
        with open(img_path, "rb") as f_in:
            file_data = f_in.read()
            
        # FIX: Fernet requires URL-safe Base64 encoded key
        # If umk is raw bytes (32 bytes), encode it.
        key_for_fernet = umk
        if isinstance(key_for_fernet, bytes) and len(key_for_fernet) == 32:
            key_for_fernet = base64.urlsafe_b64encode(key_for_fernet)
            
        from cryptography.fernet import Fernet
        f = Fernet(key_for_fernet)
        encrypted_data = f.encrypt(file_data)
        
        with open(permanent_path, "wb") as f_out:
            f_out.write(encrypted_data)

        # 3. Save Metadata to DB
        img_id = save_generated_image(
            user_id=int(user_id), 
            provider=str(provider), 
            model=str(model), 
            prompt=str(prompt), 
            image_path=str(permanent_path),
            user_state=user_state 
        )
        
        return f"✅ Verschlüsselt gespeichert (ID: {img_id})", gr.update(visible=False)

    except Exception as e:
        logger.exception(f"Gallery Save Error: {e}")
        return f"🔥 Fehler: {str(e)}", gr.update(visible=True)

def manual_save_transcription(original, translated, provider, model, lang, user_state, filename="manual_save.mp3"):
    """Manually save transcription to database"""
    try:
        if not user_state or not user_state.get("id"):
            return "❌ Bitte anmelden"

        if not original or original.strip() == "":
            return "❌ Kein Transkript zum Speichern"

        trans_id = save_transcription(
            user_id=user_state["id"],
            provider=provider,
            model=model or "N/A",
            original=original,
            translated=translated if translated and translated != "(Whisper: Keine Übersetzung verfügbar)" else None,
            language=lang,
            filename=filename,
            user_state=user_state
        )

        return f"✅ Transkript gespeichert (ID: {trans_id})"

    except Exception as e:
        logger.exception(f"Error manually saving transcription: {str(e)}")
        return f"🔥 Fehler: {str(e)}"
    
    
# ==========================================
# 🖥️ GUI BUILDER
# ==========================================

mobile_css_fix = """
@media (max-width: 768px) {
    /* 1. Target the tabs container */
    .tab-nav > button {
        display: flex !important;
        justify-content: center !important;
        align-items: center !important;
        font-size: 0 !important; /* Hide text by shrinking font to 0 */
        padding: 10px 0 !important;
    }

    /* 2. Make the icon visible again */
    /* Gradio icons are usually SVGs or images inside the button */
    .tab-nav > button > * {
        font-size: 1.5rem !important; /* Restore size for children (icons) */
    }

    /* 3. If icons are text-based (emojis at start of string), use your old trick */
    /* This targets the first letter (the emoji) and makes it visible */
    .tab-nav > button::first-letter {
        font-size: 1.5rem !important;
        visibility: visible !important;
    }
}
"""

# NEW
with gr.Blocks(
    title="Akademie KI Suite", 
    theme=gr.themes.Soft(), 
    head=PWA_HEAD,      # Only Meta tags/JS here
    css=CUSTOM_CSS      # CSS goes here to override Theme defaults
) as demo:
    # 0. INJECT PWA SCRIPTS (Required for Gradio 6+)
    gr.HTML("""
    <script>
    // 💾 PERSISTENCE V2
    window.saveCredsV2 = function(u, p) { 
        if (u && p) { 
            localStorage.setItem("ak_user", u); 
            localStorage.setItem("ak_pass", p); 
        } 
    };

    window.clearCredsV2 = function() { 
        localStorage.removeItem("ak_user"); 
        localStorage.removeItem("ak_pass"); 
    };

    window.getCredsV2 = function() {
        const u = localStorage.getItem("ak_user"); 
        const p = localStorage.getItem("ak_pass");
        if (u && p) {
            return JSON.stringify([u, p]);
        }
        return "";  
    };
    </script>
    """, visible=False)
        
    # 1. Define the Session "Backpack" (Stores data per browser tab)
    session_state = gr.State({"id": None, "username": None, "is_admin": False})

    # Set higher file size limits
    gr.set_static_paths(paths=["/var/www/transkript_app/static"])

    # === 1. COMPACT HEADER ===
    with gr.Row(elem_classes="compact-header", equal_height=True):
        
        # Col 1: Title
        with gr.Column(scale=4, min_width=100):
            gr.Markdown("### ⛪ KI Toolkit")
        
        # Col 2: User Status (Right Aligned via ID)
        with gr.Column(scale=2, min_width=100): 
            # Added ID here to target with CSS for right-alignment
            with gr.Row(elem_id="user-status-row"):
                login_status = gr.Markdown("👤", show_label=False)
            
        # Col 3: Logout Button
        with gr.Column(scale=0, min_width=120):
            logout_btn = gr.Button("🚪 Abmelden", size="sm", elem_classes="mobile-icon-only btn-secondary")

    # Login Screen (Unchanged)
    login_screen = gr.Column(visible=True)
    with login_screen:
        gr.Markdown("## 🔐 Anmeldung erforderlich")
        with gr.Row():
            with gr.Column(scale=1): pass
            with gr.Column(scale=1):
                login_username = gr.Textbox(label="Benutzername", placeholder="user123")
                login_password = gr.Textbox(label="Passwort", type="password", placeholder="***")
                login_btn = gr.Button("🔓 Anmelden", variant="primary")
                login_message = gr.Markdown("")
            with gr.Column(scale=1): pass

    # Main App
    main_app = gr.Column(visible=False)
    with main_app:

        with gr.Tabs(elem_classes="icon-nav"):

            # --- TAB 1: CHAT ---
            with gr.TabItem("💬 Chat", id="chat_tab") as chat_tab:
                
                # 1. ACCORDION (Reduced height via CSS)
                with gr.Accordion("⚙️ Modellauswahl & Status", open=False):
                    
                    # 1.1 CONTROLS ROW (Inside Accordion)
                    # Scales: Prov(3), Mod(4), Alle(1), Badge(2)
                    with gr.Row(variant="panel", equal_height=True, elem_classes="compact-row"):
                        
                        def get_allowed_chat_providers():
                            """Return only allowed AND implemented providers"""
                            # 1. Get all providers that support chat
                            all_providers = [p for p in PROVIDERS.keys() if "chat_models" in PROVIDERS[p]]
                            
                            # 2. Filter for implemented providers only
                            implemented_providers = [p for p in all_providers if is_provider_implemented(p)]
                            
                            # 3. Apply EU restriction if enabled
                            if EU_ONLY_MODE:
                                return [p for p in implemented_providers if p not in RESTRICTED_PROVIDERS]
                                
                            return implemented_providers

                        c_prov = gr.Dropdown(
                            choices=get_allowed_chat_providers(),  # ← UPDATED LOGIC APPLIED HERE
                            value=DEFAULT_CHAT_PROVIDER,
                            label="Anbieter",
                            show_label=True, container=True,
                            scale=3, min_width=80
                        )
                        
                        # Model
                        c_model = gr.Dropdown(
                            choices=PROVIDERS[DEFAULT_CHAT_PROVIDER]["chat_models"],
                            value=DEFAULT_CHAT_MODEL,
                            label="Modell", show_label=True, container=True,
                            scale=4, min_width=100
                        )
                        
                        # "Alle" Button (Adjusted Scale)
                        c_load_all = gr.Button("🌍 Alle", scale=1, min_width=50, elem_classes="mobile-icon-only btn-secondary")
                        
                        # Badge (Adjusted Scale)
                        with gr.Column(scale=2, min_width=100, elem_classes="badge-col"):
                            c_badge = gr.HTML(
                                value=f"<div class='custom-badge'>{get_compliance_html(DEFAULT_CHAT_PROVIDER)}</div>"
                            )

                        # Update logic
                        c_prov.change(
                            lambda p, s: update_c_ui(p, force_all=False, user_state=s),
                            inputs=[c_prov, session_state], 
                            outputs=[c_model, c_badge]
                        )
                        c_load_all.click(
                            lambda p, s: update_c_ui(p, force_all=True, user_state=s), 
                            inputs=[c_prov, session_state], 
                            outputs=[c_model, c_badge]
                        )

                # 2. CHAT WINDOW
                c_bot = gr.Chatbot(
                    height=500, 
                    #type="messages", 
                    #show_copy_button=True,
                    elem_id="chat_window"
                )
                
                # 3. INPUT
                c_msg = gr.Textbox(placeholder="Nachricht...", show_label=False, lines=3, max_lines=5)

                # 4. ACTION BUTTONS
                with gr.Row(equal_height=True, elem_classes="compact-row"):
                    c_btn = gr.Button("📤 Senden", variant="primary", scale=2, elem_id="btn-send", elem_classes="mobile-icon-only")
                    c_stop_btn = gr.Button("🛑 Stop", variant="stop", scale=1, min_width=40, elem_classes="mobile-icon-only btn-secondary")
                    c_save_btn = gr.Button("💾 Speichern", scale=1, min_width=40, elem_classes="mobile-icon-only btn-secondary")
                    c_clear_btn = gr.Button("🗑️ Neu", scale=1, min_width=40, elem_classes="mobile-icon-only btn-secondary")

                c_save_status = gr.Markdown("")

                # Settings Accordions
                # NO 'icon-nav' class here, so text will remain visible on mobile
                with gr.Accordion("⚙️ Einstellungen & Verlauf", open=False):
                    with gr.Tabs():
                        # Settings Tab
                        with gr.TabItem("⚙️ Config"):
                            c_key = gr.Textbox(label="API Key (Optional)", type="password")
                            c_sys = gr.Textbox(label="System Rolle", value="Du bist ein hilfreicher Assistent.", lines=2)
                            
                            gr.Markdown("**🧠 Reasoning**")
                            with gr.Row():
                                c_reasoning_effort = gr.Dropdown(
                                    choices=["default", "low", "medium", "high"],
                                    value="default",
                                    label="Effort",
                                    scale=1
                                )
                                c_reasoning_tokens = gr.Slider(
                                    0, 32000, value=0, step=1024,
                                    label="Token Budget",
                                    scale=2
                                )
                            c_temp = gr.Slider(0, 2, value=0.7, label="Temperatur", step=0.1)

                        # History Tab
                        with gr.TabItem("📚 Verlauf"):
                            c_history_state = gr.State([]) 
                            with gr.Row():
                                refresh_chats_btn = gr.Button("🔄", size="sm", scale=0)
                                delete_chat_btn = gr.Button("🗑️", variant="stop", size="sm", scale=0)
                            
                            old_chats = gr.Dataframe(
                                headers=["ID", "Datum", "Titel", "Modell"], 
                                value=[[None, "", "", ""]], 
                                label="Gespeicherte Chats", 
                                interactive=False, 
                                #height=200, 
                                wrap=True
                            )
                            
                            with gr.Row():
                                load_chat_id = gr.Number(label="ID", precision=0, scale=1)
                                load_chat_btn = gr.Button("📥 Laden", variant="primary", scale=1)
                            
                            chat_load_status = gr.Markdown("")

                            # History Events
                            refresh_chats_btn.click(load_chat_list_with_state, inputs=[session_state], outputs=[old_chats, c_history_state])
                            
                            old_chats.select(select_chat_row, inputs=[c_history_state], outputs=[load_chat_id])
                            
                            load_chat_btn.click(load_single_chat, inputs=[load_chat_id, session_state], outputs=[c_bot, chat_load_status])
                            
                            delete_chat_btn.click(delete_chat, inputs=[load_chat_id, session_state], outputs=[chat_load_status, old_chats, c_history_state])

                        # Attachments Tab
                        with gr.TabItem("📎 Anhang"):
                            gr.Markdown("**📝 Vorlagen**")
                            with gr.Row():
                                c_prompt_select = gr.Dropdown(choices=[], label="Vorlage", scale=2)
                                c_prompt_refresh = gr.Button("🔄", scale=0, size="sm")
                            c_insert_prompt_btn = gr.Button("⬇️ Einfügen", size="sm")
                            
                            gr.Markdown("---")
                            
                            attach_type = gr.Radio(
                                ["Transkript", "Vision-Ergebnis", "Eigener Text", "Datei uploaden", "Storage Box Datei"],
                                value="Transkript",
                                label="Typ"
                            )
                            
                            attach_id = gr.Number(label="ID", precision=0, visible=True)
                            attach_custom = gr.Textbox(label="Text", lines=3, visible=False)
                            attach_file = gr.File(label="Datei", visible=False, file_count="multiple", type="filepath")
                            
                            with gr.Group(visible=False) as sb_group:
                                attach_sb_browser = gr.FileExplorer(root_dir=STORAGE_MOUNT_POINT, glob="**/*", height=200)
                                sb_refresh_btn = gr.Button("🔄", size="sm")

                            with gr.Row():
                                attach_btn = gr.Button("➕ Anhängen", variant="secondary")
                                undo_attach_btn = gr.Button("↩️ Undo", variant="stop")
                            
                            attach_status = gr.Markdown("")

                            # Toggle Logic
                            def toggle_attach_inputs(atype):
                                return (
                                    gr.update(visible=atype in ["Transkript", "Vision-Ergebnis"]),
                                    gr.update(visible=atype == "Eigener Text"),
                                    gr.update(visible=atype == "Datei uploaden"),
                                    gr.update(visible=atype == "Storage Box Datei")
                                )
                            
                            attach_type.change(toggle_attach_inputs, attach_type, [attach_id, attach_custom, attach_file, sb_group])
                            
                            def refresh_chat_sb(): return gr.update(value=None)
                            sb_refresh_btn.click(refresh_chat_sb, outputs=attach_sb_browser)
                            
                            attach_btn.click(attach_content_to_chat, inputs=[c_bot, attach_type, attach_id, attach_custom, attach_file, attach_sb_browser, session_state], outputs=[c_bot, attach_status])
                            undo_attach_btn.click(undo_last_attachment, inputs=[c_bot], outputs=[c_bot, attach_status])

                            # Prompt Logic
                            c_prompt_refresh.click(get_user_prompt_choices, inputs=[session_state], outputs=c_prompt_select)
                            c_insert_prompt_btn.click(insert_custom_prompt, inputs=[c_prompt_select, c_msg, session_state], outputs=[c_msg])

                c_save_status = gr.Markdown("")

                # 2. History (Closed)
                with gr.Accordion("📚 Alte Chats laden", open=False):
                    c_history_state = gr.State([]) 
                    
                    with gr.Row():
                        refresh_chats_btn = gr.Button("🔄 Liste aktualisieren", size="sm")
                        delete_chat_btn = gr.Button("🗑️ Löschen", variant="stop", size="sm")
                    
                    old_chats = gr.Dataframe(
                        headers=["ID", "Datum", "Titel", "Modell"], 
                        value=[[None, "", "", ""]], 
                        label="Gespeicherte Chats", 
                        interactive=False, 
                        #height=200, 
                        wrap=True
                    )
                    
                    with gr.Row():
                        load_chat_id = gr.Number(label="Chat-ID", precision=0)
                    
                    load_chat_btn = gr.Button("📥 Chat laden", variant="primary")
                    chat_load_status = gr.Markdown("")

                    # --- Events for History ---
                    
                    # 1. Load list on click/tab select
                    chat_tab.select(
                        load_chat_list_with_state, 
                        inputs=[session_state], 
                        outputs=[old_chats, c_history_state]
                    )
                    refresh_chats_btn.click(
                        load_chat_list_with_state, 
                        inputs=[session_state],
                        outputs=[old_chats, c_history_state]
                    )
                    
                    # 2. Select row to get ID (Using State to avoid Index Errors)
                    old_chats.select(
                        select_chat_row, 
                        inputs=[c_history_state], 
                        outputs=[load_chat_id]
                    )
                    
                    # 3. Load Chat Logic
                    load_chat_btn.click(
                        load_single_chat, 
                        inputs=[load_chat_id], 
                        outputs=[c_bot, chat_load_status]
                    )
                    
                    # 4. Delete Logic (Now updates State + DataFrame to prevent 'backend' errors on next click)
                    delete_chat_btn.click(
                        delete_chat, 
                        inputs=[load_chat_id], 
                        outputs=[chat_load_status, old_chats, c_history_state]
                    )

                # 3. Attachments & Prompts
                with gr.Accordion("📎 Inhalt & Prompts", open=False):
                    
                    gr.Markdown("**📝 Vorlagen**")
                    with gr.Row():
                        c_prompt_select = gr.Dropdown(choices=[], label="Vorlage wählen", scale=2)
                        c_prompt_refresh = gr.Button("🔄", scale=0, size="sm")
                    c_insert_prompt_btn = gr.Button("⬇️ In Textfeld einfügen", size="sm")
                    
                    gr.Markdown("---")
                    
                    gr.Markdown("**📎 Anhang**")
                    gr.Markdown("_Hinweis: Mehrere Dateien möglich. Inhalte werden extrahiert (OCR/Text)._", visible=True)
                    
                    attach_type = gr.Radio(
                        ["Transkript", "Vision-Ergebnis", "Eigener Text", "Datei uploaden", "Storage Box Datei"],
                        value="Transkript",
                        label="Typ"
                    )
                    
                    attach_id = gr.Number(label="ID", precision=0, visible=True)
                    attach_custom = gr.Textbox(label="Text", lines=3, visible=False)
                    
                    attach_file = gr.File(
                        label="Dateien wählen", 
                        visible=False, 
                        file_count="multiple", 
                        type="filepath"
                    )
                    
                    with gr.Group(visible=False) as sb_group:
                        gr.Markdown("Dateien auf Server:")
                        attach_sb_browser = gr.FileExplorer(
                            root_dir=STORAGE_MOUNT_POINT,
                            glob="**/*",
                            height=200,
                            file_count="multiple"
                        )
                        sb_refresh_btn = gr.Button("🔄", size="sm")

                    with gr.Row():
                        attach_btn = gr.Button("➕ Anhängen", variant="secondary")
                        undo_attach_btn = gr.Button("↩️ Rückgängig", variant="stop")
                    
                    attach_status = gr.Markdown("")

                    # Toggle Logic
                    def toggle_attach_inputs(atype):
                        return (
                            gr.update(visible=atype in ["Transkript", "Vision-Ergebnis"]),
                            gr.update(visible=atype == "Eigener Text"),
                            gr.update(visible=atype == "Datei uploaden"),
                            gr.update(visible=atype == "Storage Box Datei")
                        )
                    
                    attach_type.change(
                        toggle_attach_inputs, 
                        attach_type, 
                        [attach_id, attach_custom, attach_file, sb_group]
                    )
                    
                    def refresh_chat_sb(): return gr.update(value=None)
                    sb_refresh_btn.click(refresh_chat_sb, outputs=attach_sb_browser)
                    
                    # Attach Event
                    attach_btn.click(
                        attach_content_to_chat, 
                        inputs=[c_bot, attach_type, attach_id, attach_custom, attach_file, attach_sb_browser, session_state], 
                        outputs=[c_bot, attach_status]
                    )
                    
                    # Undo Event
                    undo_attach_btn.click(
                        undo_last_attachment,
                        inputs=[c_bot],
                        outputs=[c_bot, attach_status]
                    )

                # --- EVENT WIRING ---

                # Chat Execution (With Stop)
                submit_event = c_msg.submit(user_msg, [c_msg, c_bot], [c_msg, c_bot], queue=False).then(
                    bot_msg, 
                    # Add session_state to the inputs list
                    [c_bot, c_prov, c_model, c_temp, c_sys, c_key, c_reasoning_effort, c_reasoning_tokens, session_state], 
                    c_bot
                )

                click_event = c_btn.click(user_msg, [c_msg, c_bot], [c_msg, c_bot], queue=False).then(
                    bot_msg, 
                    # Add session_state to the inputs list
                    [c_bot, c_prov, c_model, c_temp, c_sys, c_key, c_reasoning_effort, c_reasoning_tokens, session_state], 
                    c_bot
                )
                
                # Stop Button
                c_stop_btn.click(fn=None, cancels=[submit_event, click_event])

                # Save & Clear
                c_save_btn.click(save_chat, [c_bot, c_prov, c_model, session_state], c_save_status)
                c_clear_btn.click(lambda: ([], ""), outputs=[c_bot, c_save_status])

                # History Logic
                chat_tab.select(load_chat_list_with_state, inputs=[session_state], outputs=[old_chats, c_history_state])
                refresh_chats_btn.click(load_chat_list_with_state, inputs=[session_state], outputs=[old_chats, c_history_state])
                old_chats.select(select_chat_row, inputs=[c_history_state], outputs=[load_chat_id])
                load_chat_btn.click(load_single_chat, inputs=[load_chat_id, session_state], outputs=[c_bot, chat_load_status])
                delete_chat_btn.click(delete_chat, inputs=[load_chat_id, session_state], outputs=[chat_load_status, old_chats, c_history_state])

                # Attachment Logic
                attach_btn.click(
                    attach_content_to_chat, 
                    # Added session_state
                    inputs=[c_bot, attach_type, attach_id, attach_custom, attach_file, attach_sb_browser, session_state], 
                    outputs=[c_bot, attach_status]
                )

                # Prompt Logic
                c_prompt_refresh.click(
                    get_user_prompt_choices, 
                    inputs=[session_state], 
                    outputs=c_prompt_select
                )
                chat_tab.select(
                    get_user_prompt_choices, 
                    inputs=[session_state], 
                    outputs=c_prompt_select
                )

                c_insert_prompt_btn.click(
                    insert_custom_prompt, 
                    inputs=[c_prompt_select, c_msg, session_state], 
                    outputs=[c_msg]
                )
                
            # --- TAB 2: TRANSKRIPTION ---
            with gr.TabItem("🎙️ Transkription"):
                with gr.Row():
                    with gr.Column():
                        # --- INPUT SELECTION: Upload vs Storage Box ---
                        with gr.Tabs() as input_source_tabs:
                            with gr.TabItem("📤 Upload"):
                                t_audio = gr.Audio(type="filepath", label="Datei hochladen")
                            
                            with gr.TabItem("📦 Storage Box"):
                                gr.Markdown("Wähle Datei aus Cloud-Speicher:")
                                t_storage_browser = gr.FileExplorer(
                                    root_dir=STORAGE_MOUNT_POINT,
                                    glob="**/*", 
                                    #height=300,
                                    label="Dateien durchsuchen"
                                )
                                with gr.Row():
                                    t_refresh_sb_btn = gr.Button("🔄 Aktualisieren", size="sm", scale=0)
                                    t_load_sb_btn = gr.Button("✅ Diese verwenden", variant="secondary", scale=1)
                                t_sb_status = gr.Markdown("")

                            with gr.TabItem("🔗 URL Download"):
                                gr.Markdown("""
                                ### ✅ Freigegebene Kanäle:
                                """)
                                
                                # Dynamic whitelist display
                                whitelist_md = "\n".join([f"- `{ch}`" for ch in YOUTUBE_CHANNEL_WHITELIST])
                                gr.Markdown(whitelist_md)
                                
                                gr.Markdown("""
                                ---
                                ### ⚠️ Nur eigene oder lizenzierte Inhalte herunterladen
                                """)
                                
                                t_url_input = gr.Textbox(
                                    label="URL",
                                    placeholder="https://www.youtube.com/watch?v=... oder direkte Audio-URL",
                                    lines=1
                                )

                                # --- NEW: Admin Options ---
                                with gr.Group(visible=False) as admin_dl_group:
                                    gr.Markdown("🔒 **Admin / Medienverwalter Optionen**")
                                    with gr.Row():
                                        t_force_dl = gr.Checkbox(
                                            label="📥 Lokal herunterladen (yt-dlp)", 
                                            value=False,
                                            info="Erzwingt Download auf den Server. Nötig für Mistral/Whisper. Bei Gladia optional."
                                        )
                                        t_dl_video = gr.Checkbox(
                                            label="🎥 Video behalten",
                                            value=False,
                                            interactive=True
                                        )
                                    
                                    t_dl_destination = gr.Radio(
                                        choices=["Storage Box", "VPS Temp"],
                                        value="Storage Box",
                                        label="Speicherort"
                                    )
                                # --------------------------
                                
                                # Dynamic Logic: Show Admin Group if User is privileged
                                def toggle_admin_dl_options(user_state):
                                    if user_state and (user_state.get("is_admin") or user_state.get("is_media_manager")):
                                        return gr.update(visible=True)
                                    return gr.update(visible=False)
                                
                                # Trigger check when tab is selected or login changes
                                input_source_tabs.select(
                                    toggle_admin_dl_options, 
                                    inputs=[session_state], 
                                    outputs=admin_dl_group
                                )

                        # Logic: Storage Box Selection
                        def use_storage_file(selected_files):
                            if not selected_files:
                                return None, "❌ Keine Datei ausgewählt"
                            f_path = selected_files[0] if isinstance(selected_files, list) else selected_files
                            if not f_path.startswith("/"):
                                f_path = os.path.join(STORAGE_MOUNT_POINT, f_path)
                            try:
                                local_temp = copy_storage_file_to_temp(f_path)
                                return local_temp, f"✅ Geladen: {os.path.basename(f_path)}"
                            except Exception as e:
                                return None, f"🔥 Fehler: {str(e)}"

                        def refresh_explorer():
                            return gr.update(value=None) 

                        t_load_sb_btn.click(use_storage_file, inputs=t_storage_browser, outputs=[t_audio, t_sb_status])
                        t_refresh_sb_btn.click(refresh_explorer, outputs=t_storage_browser)

                        # --- MAIN CONTROLS ---
                        gr.Markdown("### 🎛️ Auswahl")
                        with gr.Group():
                            with gr.Row(elem_classes="compact-row"):
                                t_prov = gr.Dropdown(
                                    choices=["Gladia", "Deepgram", "AssemblyAI", "Mistral", "Scaleway", "Groq"], 
                                    value="Deepgram", 
                                    label="Engine",
                                    scale=2
                                )
                                with gr.Column(scale=0, min_width=210):
                                    t_diar = gr.Checkbox(
                                        value=True, 
                                        label="🎭 Sprecher",
                                        scale=1,
                                    )
                            
                                    t_badge = gr.HTML(
                                        value=f"<div class='custom-badge'>{get_compliance_html('Deepgram')}</div>"
                                    )

                        # --- SETTINGS ACCORDION ---
                        with gr.Accordion("⚙️ Einstellungen", open=False):
                            
                            # 1. Language & Model
                            with gr.Row():
                                t_lang = gr.Dropdown(
                                    [("Auto-Erkennung", "auto"), ("Deutsch", "de"), ("Englisch", "en")], 
                                    value="de", 
                                    label="Sprache"
                                )
                                
                                with gr.Row():
                                    t_model = gr.Dropdown(choices=[], value=None, label="Modell", scale=3, visible=False)
                                    t_refresh_models = gr.Button("🔄", size="sm", scale=0, variant="secondary")

                            # 2. Translation Settings
                            with gr.Row():
                                t_trans = gr.Checkbox(False, label="🌍 Übersetzen")
                                t_target = gr.Dropdown(
                                    [("Deutsch", "de"), ("Englisch", "en")], 
                                    value="en", 
                                    label="Zielsprache", 
                                    visible=False
                                )

                            # 3. Whisper Options (Hidden Group)
                            with gr.Group(visible=False) as w_options_group:
                                gr.Markdown("---")
                                gr.Markdown("**Erweiterte Optionen (Whisper)**")
                                with gr.Row():
                                    w_chunk_opt = gr.Checkbox(value=True, label="✂️ Chunking")
                                    w_chunk_len = gr.Number(value=10, label="Minuten", precision=0, minimum=1)
                                    w_temp = gr.Slider(0, 1, value=0, step=0.1, label="Temperatur")
                                w_prompt = gr.Textbox(label="Kontext-Prompt", placeholder="Optionaler Kontext...")

                        # Hidden state for API key
                        t_key = gr.State(value="") 
                        
                        # Start Button
                        t_btn = gr.Button("▶️ Transkription starten", variant="primary", size="lg")
                        t_log = gr.Textbox(label="Status Log", lines=3)

                        # --- DYNAMIC UI LOGIC ---
                        
                        # 1. Provider Change -> Update Badge, Model, and Whisper Options
                        t_prov.change(
                            fn=update_t_ui, 
                            inputs=t_prov, 
                            outputs=[t_badge, t_model, w_options_group]
                        )
                        
                        # 2. Translation Toggle
                        def toggle_translation(chk):
                            return gr.update(visible=chk), gr.update(visible=chk)
                        
                        # 3. Model Refresh Logic
                        t_refresh_models.click(
                            lambda p: update_t_ui(p, force_all=True)[1],
                            inputs=t_prov, 
                            outputs=t_model
                        )

                    with gr.Column():
                        # OUTPUTS
                        t_orig = gr.Textbox(label="📄 Transkript", lines=15)
                        t_trsl = gr.Textbox(label="🌍 Übersetzung", lines=15, visible=False)

                        # Wire visibility toggle for translation output
                        t_trans.change(toggle_translation, inputs=t_trans, outputs=[t_target, t_trsl])

                        with gr.Row():
                            t_save_btn = gr.Button("💾 Transkript speichern", variant="secondary")
                            t_save_status = gr.Markdown("")

                # --- SEND TO CHAT SECTION ---
                with gr.Accordion("💬 An Chat senden", open=False) as send_to_chat_section:
                    gr.Markdown("### Weiterverarbeitung")

                    with gr.Row():
                        with gr.Column(scale=1):
                            prompt_template = gr.Dropdown(
                                choices=list(TRANSCRIPT_PROMPTS.keys()),
                                value="Veranstaltungsrückblick",
                                label="📋 Prompt-Vorlage"
                            )
                        with gr.Column(scale=2):
                            with gr.Row():
                                chat_provider = gr.Dropdown(
                                    choices=[p for p in PROVIDERS.keys() if "chat_models" in PROVIDERS[p] and is_provider_implemented(p)],
                                    value=DEFAULT_CHAT_PROVIDER,
                                    label="🤖 Provider",
                                    scale=2
                                )
                                chat_model_for_transcript = gr.Dropdown(
                                    choices=PROVIDERS[DEFAULT_CHAT_PROVIDER]["chat_models"],
                                    value=DEFAULT_CHAT_MODEL,
                                    label="Modell",
                                    scale=3
                                )
                                chat_load_all_btn = gr.Button("🌍 Alle", scale=0, size="sm")
                    
                    # Badge display
                    chat_badge = gr.HTML(value=PROVIDERS[DEFAULT_CHAT_PROVIDER]["badge"])

                    additional_notes = gr.Textbox(
                        label="📝 Zusätzliche Hinweise",
                        placeholder="Erwähne Kooperationspartner, betone ...",
                        lines=2
                    )

                    custom_prompt_input = gr.Textbox(
                        label="✏️ Eigener Prompt",
                        lines=3,
                        visible=False
                    )

                    send_to_chat_btn = gr.Button("💬 An Chat senden", variant="primary")
                    send_status = gr.Markdown("")

                # --- LOGIC WIRING (Backend) ---
                
                # Chat Provider/Model Updates with filtering
                chat_provider.change(
                    fn=lambda p, s: update_transcript_chat_ui(p, force_all=False, user_state=s),
                    inputs=[chat_provider, session_state],
                    outputs=[chat_model_for_transcript, chat_badge]
                )

                chat_load_all_btn.click(
                    fn=lambda p, s: update_transcript_chat_ui(p, force_all=True, user_state=s),
                    inputs=[chat_provider, session_state],
                    outputs=[chat_model_for_transcript, chat_badge]
                )

                # Custom Prompt Visibility
                prompt_template.change(
                    fn=lambda t: gr.update(visible=(t == "Eigener Prompt")), 
                    inputs=prompt_template, 
                    outputs=custom_prompt_input
                )

                # Main Transcription Execution
                t_btn.click(
                    fn=run_and_save_transcription, 
                    inputs=[
                        t_audio, t_prov, t_model, 
                        t_lang, w_temp, w_prompt, 
                        t_diar, t_trans, t_target, t_key,
                        w_chunk_opt, w_chunk_len,
                        t_lang, t_diar, t_lang, t_diar,
                        t_url_input, t_dl_video, t_dl_destination, t_force_dl,
                        session_state 
                    ], 
                    outputs=[t_log, t_orig, t_trsl]
                )

                # Manual Save
                t_save_btn.click(
                    fn=manual_save_transcription,
                    inputs=[t_orig, t_trsl, t_prov, t_model, t_lang, session_state],
                    outputs=t_save_status
                )

                # Send to Chat Execution
                send_to_chat_btn.click(
                    fn=send_transcript_to_chat,
                    inputs=[
                        t_orig,
                        prompt_template,
                        additional_notes,
                        custom_prompt_input,
                        chat_provider,
                        chat_model_for_transcript,
                        t_key,
                        session_state
                    ],
                    outputs=send_status
                )

            # --- TAB 3: VISION ---
            with gr.TabItem("👁️ Vision"):
                with gr.Row():
                    with gr.Column():
                        
                        # --- INPUT SELECTION: Upload vs Storage Box ---
                        with gr.Tabs():
                            with gr.TabItem("📤 Upload"):
                                v_img = gr.Image(type="filepath", label="Bild hochladen", height=300)
                            
                            with gr.TabItem("📦 Storage Box"):
                                gr.Markdown("Wähle ein Bild aus dem Cloud-Speicher:")
                                v_storage_browser = gr.FileExplorer(
                                    root_dir=STORAGE_MOUNT_POINT,
                                    glob="**/*",
                                    #height=300,
                                    label="Bilder durchsuchen"
                                )
                                with gr.Row():
                                    v_refresh_sb_btn = gr.Button("🔄 Aktualisieren", size="sm", scale=0)
                                    v_load_sb_btn = gr.Button("✅ Dieses Bild verwenden", variant="secondary", scale=1)
                                v_sb_status = gr.Markdown("")

                        # Logic: Storage Box Selection (Vision)
                        def use_storage_image(selected_files):
                            if not selected_files:
                                return None, "❌ Kein Bild ausgewählt"
                            f_path = selected_files[0] if isinstance(selected_files, list) else selected_files
                            if not f_path.startswith("/"):
                                f_path = os.path.join(STORAGE_MOUNT_POINT, f_path)
                            try:
                                local_temp = copy_storage_file_to_temp(f_path)
                                return local_temp, f"✅ Geladen: {os.path.basename(f_path)}"
                            except Exception as e:
                                return None, f"🔥 Fehler: {str(e)}"

                        def refresh_v_explorer():
                            return gr.update(value=None)

                        v_load_sb_btn.click(use_storage_image, inputs=v_storage_browser, outputs=[v_img, v_sb_status])
                        v_refresh_sb_btn.click(refresh_v_explorer, outputs=v_storage_browser)


                        # --- SELECTION ROW ---
                        gr.Markdown("### 🎛️ Auswahl")
                        with gr.Group():
                            with gr.Row(equal_height=True):
                                v_prov = gr.Dropdown(
                                    ["Scaleway", "Mistral", "Nebius", "OpenRouter", "Poe"], 
                                    value="Scaleway", 
                                    label="Provider", 
                                    scale=1
                                )
                                v_model = gr.Dropdown(
                                    PROVIDERS["Scaleway"]["vision_models"], 
                                    value="pixtral-12b-2409", 
                                    label="Modell", 
                                    allow_custom_value=True,
                                    scale=2
                                )
                                v_load_all = gr.Button("🔄", scale=0, size="sm", variant="secondary")
                            
                            # Badge Row
                            v_badge = gr.HTML(value=get_compliance_html("Scaleway"))

                        # Inputs
                        v_prompt = gr.Textbox(label="Frage", value="Beschreibe dieses Bild detailliert.", lines=2)
                        
                        # Hidden Key (Pass empty string or handle in backend)
                        v_key = gr.State(value="") 
                        
                        v_btn = gr.Button("👁️ Analysieren", variant="primary", size="lg")
                    
                    with gr.Column():
                        v_out = gr.Markdown(label="Ergebnis")
                        
                # UI Updates
                v_prov.change(
                    lambda p: update_v_ui(p, force_all=False), 
                    inputs=v_prov, 
                    outputs=[v_badge, v_model]
                )
                v_load_all.click(
                    lambda p: update_v_ui(p, force_all=True), 
                    inputs=v_prov, 
                    outputs=[v_badge, v_model]
                )
                
                # Execution
                v_btn.click(run_vision, [v_img, v_prompt, v_prov, v_model, v_key, session_state], v_out)
            

            # --- TAB 4: BILDERZEUGUNG ---
            with gr.TabItem("🎨 Bilderzeugung"):
                with gr.Row():
                    with gr.Column():
                        # Prompt Input
                        g_prompt = gr.Textbox(
                            label="Prompt", 
                            placeholder="Eine futuristische Kathedrale aus Glas und Licht...", 
                            lines=3
                        )
                        
                        # --- SELECTION ROW ---
                        gr.Markdown("### 🎛️ Auswahl")
                        with gr.Group():
                            with gr.Row(equal_height=True):
                                g_provider = gr.Dropdown(
                                    ["Nebius", "Scaleway", "OpenRouter", "Poe"], 
                                    value="Nebius", 
                                    label="Provider", 
                                    scale=1
                                )
                                g_model = gr.Dropdown(
                                    PROVIDERS["Nebius"]["image_models"], 
                                    value="black-forest-labs/flux-schnell", 
                                    label="Modell",
                                    scale=2
                                )
                                g_load_all = gr.Button("🔄", scale=0, size="sm", variant="secondary")

                            # Badge Row
                            g_badge = gr.HTML(value=get_compliance_html("Nebius"))

                        # --- SETTINGS ACCORDION ---
                        with gr.Accordion("⚙️ Einstellungen", open=False):
                            with gr.Row():
                                g_w = gr.Slider(256, 1440, value=1024, step=64, label="Breite")
                                g_h = gr.Slider(256, 1440, value=768, step=64, label="Höhe")
                            g_steps = gr.Slider(1, 50, value=4, step=1, label="Schritte")

                        # Hidden Key
                        g_key = gr.State(value="")

                        g_btn = gr.Button("🎨 Generieren", variant="primary", size="lg")
                        g_stat = gr.Textbox(label="Status", interactive=False, visible=True)
                        
                    with gr.Column():
                        g_out = gr.Image(label="Ergebnis", type="filepath", 
                                         #show_download_button=False, 
                                         height=400)
                        
                        with gr.Row():
                             g_download_file = gr.File(label="Download", scale=1)
                             g_save_btn = gr.Button("💾 In Storage Box speichern", visible=False, scale=1, variant="secondary")
                        
                        g_save_status = gr.Markdown("")

                # State for path
                g_img_path = gr.State(value=None)
                
                # UI Updates
                g_provider.change(
                    lambda p: update_g_ui(p, force_all=False), 
                    inputs=g_provider, 
                    outputs=[g_badge, g_model]
                )
                g_load_all.click(
                    lambda p: update_g_ui(p, force_all=True), 
                    inputs=g_provider, 
                    outputs=[g_badge, g_model]
                )
                
                # Execution Logic
                g_btn.click(
                    generate_and_handle_ui, 
                    inputs=[g_prompt, g_provider, g_model, g_w, g_h, g_steps, g_key, session_state],  # ← ADD session_state
                    outputs=[g_out, g_stat, g_img_path, g_download_file, g_save_btn, g_save_status]
                )

                # Save to Storage Box Logic
                g_save_btn.click(
                    process_gallery_save,
                    inputs=[g_img_path, g_provider, g_prompt, g_model, session_state], 
                    outputs=[g_save_status, g_save_btn]
                )

            # --- TAB 5: VERLAUF & VERWALTUNG ---
            with gr.TabItem("📚 Verlauf & Verwaltung", id="tab_management"):
                
                gr.Markdown("### ⚙️ Verwaltung")
                
                # Optional: Helper to make tables look full
                def pad_data(data, width, min_rows=6):
                    while len(data) < min_rows:
                        row = [""] * width 
                        data.append(row)
                    return data

                # Helper: Robust Decryption for Lists
                def decrypt_for_display(encrypted_text, user_state):
                    if not encrypted_text: return ""
                    # 1. Get Keys
                    umk = user_state.get('umk') if user_state else crypto.global_key
                    
                    try:
                        # 2. Try User Key first
                        return crypto.decrypt_text(encrypted_text, key=umk)
                    except:
                        # 3. Fallback to Global Key (Legacy Support)
                        # This causes the "MAC Check Failed" log, but allows reading old data
                        try:
                            res = crypto.decrypt_text(encrypted_text, key=crypto.global_key)
                            if res == "[Decryption Failed]": return "🔒 [Verschlüsselt]"
                            return res
                        except:
                            return "🔒 [Datenfehler]"

                with gr.Tabs() as history_tabs:

                    # =========================================================
                    # 1. TRANSCRIPTIONS
                    # =========================================================
                    with gr.TabItem("🎙️ Transkriptions-Verlauf") as trans_tab:
                        trans_state = gr.State([])

                        # 1. TABLE
                        with gr.Group():
                            gr.Markdown("#### 📋 Gespeicherte Transkripte")
                            trans_history = gr.Dataframe(
                                headers=["ID", "Datum", "Titel", "Provider", "Sprache"],
                                value=[],  # Start Empty
                                #value=[[None, "", "", "", ""]] * 6,
                                interactive=False,
                                wrap=True,
                                datatype=["number", "str", "str", "str", "str"],
                                column_widths=["10%", "20%", "40%", "15%", "15%"]
                            )

                        # 2. CONTROLS
                        with gr.Row(variant="panel", equal_height=True):
                            with gr.Column(scale=1):
                                trans_id_input = gr.Number(label="Ausgewählte ID", precision=0, minimum=0)
                            with gr.Column(scale=0, min_width=50):
                                refresh_trans_btn = gr.Button("🔄", size="lg")
                            with gr.Column(scale=2):
                                with gr.Row():
                                    load_trans_btn = gr.Button("📄 Laden", variant="secondary", size="lg")
                                    trans_to_chat_btn = gr.Button("📨 An Chat senden", variant="primary", size="lg")
                            with gr.Column(scale=1):
                                delete_trans_btn = gr.Button("🗑️ Löschen", variant="stop", size="lg")

                        # 3. PREVIEW
                        loaded_trans_display = gr.Textbox(label="Inhalt", lines=8, max_lines=15)
                        trans_action_status = gr.Markdown("")

                        # --- LOGIC ---
                        def load_trans_data(user_state=None):
                            """Load transcriptions filtered by current user."""
                            # 1. Strict Auth Check
                            if not user_state or not isinstance(user_state, dict) or not user_state.get("id"):
                                return [], []
                            
                            user_id = user_state["id"]
                            
                            try:
                                db = SessionLocal()
                                # FILTER BY USER_ID
                                t_list = db.query(Transcription).filter(
                                    Transcription.user_id == user_id
                                ).order_by(Transcription.timestamp.desc()).limit(50).all()
                                db.close()
                                
                                clean_data = []
                                for t in t_list:
                                    clean_data.append([
                                        t.id, 
                                        t.timestamp.strftime("%Y-%m-%d %H:%M"), 
                                        t.title or "—", 
                                        t.provider, 
                                        t.language or "?"
                                    ])
                                    
                                return clean_data, clean_data
                                
                            except Exception as e:
                                logger.exception(f"Error loading transcriptions: {e}")
                                return [], []
                            
                        def load_single_trans(tid, user_state):
                            if not tid or not user_state or not user_state.get("id"): 
                                return gr.update(), "❌"
                            
                            t = get_decrypted_transcription(int(tid), user_state["id"], user_state)
                            if t:
                                return (t.original_text, f"✅ Geladen: {t.title}")
                            else:
                                return ("", "❌ Nicht gefunden")
                            
                        def select_trans_row(evt: gr.SelectData, state_data, user_state):
                            try:
                                if not user_state or not user_state.get("id"):
                                    return 0, "", "❌ Bitte anmelden"

                                row_idx = evt.index[0]
                                if state_data and row_idx < len(state_data):
                                    real_row = state_data[row_idx]
                                    t_id = int(real_row[0]) 
                                    
                                    content, status = load_single_trans(t_id, user_state)
                                    return t_id, content, status
                            except Exception as e: 
                                logger.error(f"Select Error: {e}")
                            
                            return 0, "", ""

                        # FIX: del_trans now returns exactly 4 values to match outputs
                        def del_trans(tid, user_state):
                            empty_df, empty_state = pad_data([], 5), []
                            
                            if not user_state or not user_state.get("id"):
                                return "", "❌ Auth Fehler", empty_df, empty_state
                            
                            if not tid:
                                d, s = load_trans_data(user_state)
                                return "", "⚠️ Keine ID", d, s

                            if delete_transcription(int(tid), user_state["id"]):
                                d, s = load_trans_data(user_state)
                                return "", "✅ Gelöscht", d, s
                            
                            d, s = load_trans_data(user_state)
                            return "", "❌ Fehler", d, s

                        # Wiring
                        refresh_trans_btn.click(load_trans_data, inputs=[session_state], outputs=[trans_history, trans_state])
                        trans_tab.select(load_trans_data, inputs=[session_state], outputs=[trans_history, trans_state])

                        trans_history.select(
                            select_trans_row, 
                            inputs=[trans_state, session_state], 
                            outputs=[trans_id_input, loaded_trans_display, trans_action_status]
                        )
                        trans_id_input.change(load_single_trans, inputs=[trans_id_input, session_state], outputs=[loaded_trans_display, trans_action_status])
                        
                        delete_trans_btn.click(
                            del_trans, 
                            inputs=[trans_id_input, session_state], 
                            outputs=[loaded_trans_display, trans_action_status, trans_history, trans_state]
                        )
                        
                        # FIX: Direct reference to c_msg (defined in Chat Tab)
                        trans_to_chat_btn.click(lambda x: x, inputs=loaded_trans_display, outputs=c_msg)

                    # =========================================================
                    # 2. GENERATED IMAGES
                    # =========================================================
                    with gr.TabItem("🎨 Generierte Bilder") as images_tab:
                        img_state = gr.State([])

                        with gr.Group():
                            gr.Markdown("#### 🖼️ Bild-Historie")
                            images_history = gr.Dataframe(
                                headers=["ID", "Datum", "Prompt", "Modell"],
                                value=[], # Start Empty
                                #value=[[None, "", "", ""]] * 6, 
                                interactive=False,
                                wrap=True,
                                datatype=["number", "str", "str", "str"],
                                column_widths=["10%", "20%", "50%", "20%"]
                            )

                        with gr.Row(variant="panel", equal_height=True):
                            with gr.Column(scale=1):
                                img_id_input = gr.Number(label="Ausgewählte ID", precision=0, minimum=0)
                            with gr.Column(scale=0, min_width=50):
                                refresh_images_btn = gr.Button("🔄", size="lg")
                            with gr.Column(scale=2):
                                with gr.Row():
                                    load_img_btn = gr.Button("🖼️ Laden", variant="secondary", size="lg")
                                    img_to_chat_btn = gr.Button("📨 Prompt an Chat", variant="primary", size="lg")
                            with gr.Column(scale=1):
                                delete_img_btn = gr.Button("🗑️ Löschen", variant="stop", size="lg")

                        with gr.Row():
                            loaded_img_display = gr.Image(label="Vorschau", height=300, type="filepath", interactive=False)
                            with gr.Column():
                                loaded_img_prompt = gr.Textbox(label="Prompt", lines=10)
                                img_action_status = gr.Markdown("")

                        # --- LOGIC ---
                        def load_img_data(user_state=None):
                            """
                            Load images and force decryption of the prompt for the table view.
                            """
                            # 1. Strict Auth Check
                            if not user_state or not isinstance(user_state, dict) or not user_state.get("id"):
                                return [], []
                            
                            try:
                                # Get Session Key
                                umk = user_state.get('umk') if user_state else crypto.global_key
                                user_id = user_state["id"]
                                
                                db = SessionLocal()
                                i_list = db.query(GeneratedImage).filter(
                                    GeneratedImage.user_id == user_id
                                ).order_by(GeneratedImage.timestamp.desc()).limit(50).all()
                                db.close()
                                
                                clean = []
                                
                                for i in i_list:
                                    # 2. Decrypt Prompt for Display
                                    raw_prompt = str(i.prompt)
                                    dec_prompt = raw_prompt


                                    # FIX: Always try to decrypt, don't check for specific header
                                    try:
                                        # Try User Key
                                        candidate = crypto.decrypt_text(raw_prompt, key=umk)
                                        if candidate and candidate != "[Decryption Failed]":
                                            dec_prompt = candidate
                                        else:
                                            # Fallback for legacy data (Global Key)
                                            if umk != crypto.global_key:
                                                legacy = crypto.decrypt_text(raw_prompt, key=crypto.global_key)
                                                if legacy and legacy != "[Decryption Failed]":
                                                    dec_prompt = legacy
                                    except Exception:
                                        # If it crashes, it's likely already plaintext or totally broken
                                        pass
                                    
                                    # Truncate for table view
                                    display_prompt = (dec_prompt[:75] + '...') if len(dec_prompt) > 75 else dec_prompt
                                    
                                    clean.append([
                                        i.id, 
                                        i.timestamp.strftime("%Y-%m-%d"), 
                                        display_prompt, 
                                        i.model
                                    ])
                                    
                                return clean, clean
                                
                            except Exception as e:
                                logger.error(f"Img Load Error: {e}") 
                                return [], []

                        def load_single_img(tid, user_state=None):
                            """
                            Loads DB record, decrypts the PROMPT and the IMAGE file.
                            Fixed: Key encoding for image decryption & blind prompt decryption.
                            """
                            if not tid or not user_state or not user_state.get("id"): 
                                return None, "", "❌"
                            
                            try:
                                import base64
                                
                                db = SessionLocal()
                                img = db.query(GeneratedImage).filter(
                                    GeneratedImage.id == int(tid),
                                    GeneratedImage.user_id == user_state["id"]
                                ).first()
                                db.close()
                                
                                if not img or not img.image_path or not os.path.exists(img.image_path):
                                    return None, "", "❌ Nicht gefunden"

                                # Get Key
                                umk = user_state.get('umk') if user_state else crypto.global_key

                                # 1. Decrypt Prompt (Blind attempt)
                                raw_prompt = str(img.prompt)
                                final_prompt = raw_prompt
                                logger.info(f"Raw prompt {raw_prompt}.")

                                try:
                                    candidate = crypto.decrypt_text(raw_prompt, key=umk)
                                    logger.info(f"Decrypted prompt {candidate}.")
                                    if candidate and candidate != "[Decryption Failed]":
                                        final_prompt = candidate
                                    elif umk != crypto.global_key:
                                        # Fallback
                                        legacy = crypto.decrypt_text(raw_prompt, key=crypto.global_key)
                                        logger.info(f"Decrypted (legacy) prompt: {legacy}.")
                                        if legacy and legacy != "[Decryption Failed]":
                                            final_prompt = legacy
                                except Exception:
                                    pass

                                # 2. Decrypt Image File (if needed)
                                display_path = img.image_path
                                
                                if img.image_path.endswith(".enc"):
                                    try:
                                        with open(img.image_path, "rb") as f:
                                            enc_data = f.read()
                                        
                                        # FIX: Ensure Key is B64 for Fernet
                                        key_for_fernet = umk
                                        if isinstance(key_for_fernet, bytes) and len(key_for_fernet) == 32:
                                            key_for_fernet = base64.urlsafe_b64encode(key_for_fernet)

                                        from cryptography.fernet import Fernet
                                        f_eng = Fernet(key_for_fernet)
                                        dec_data = f_eng.decrypt(enc_data)
                                        
                                        # Write to temp file
                                        original_name = os.path.basename(img.image_path).replace(".enc", "")
                                        if "." not in original_name: original_name += ".jpg"
                                        
                                        temp_decrypted_path = os.path.join(tempfile.gettempdir(), f"dec_{original_name}")
                                        
                                        with open(temp_decrypted_path, "wb") as f_out:
                                            f_out.write(dec_data)
                                            
                                        display_path = temp_decrypted_path
                                        
                                    except Exception as e:
                                        logger.error(f"Image decryption failed: {e}")
                                        return None, final_prompt, f"❌ Bild-Fehler: {str(e)}"

                                return display_path, final_prompt, f"✅ Geladen"

                            except Exception as e:
                                logger.error(f"Load Single Image Error: {e}")
                                return None, "", f"🔥 Fehler: {str(e)}"

                        def select_img_row(evt: gr.SelectData, state_data, user_state):
                            try:
                                if not user_state or not user_state.get("id"):
                                    return 0, None, "", "❌ Bitte anmelden"

                                row_idx = evt.index[0]
                                if state_data and row_idx < len(state_data):
                                    real_row = state_data[row_idx]
                                    tid = int(real_row[0])
                                    
                                    path, prmt, stat = load_single_img(tid, user_state)
                                    return tid, path, prmt, stat
                            except: pass
                            return 0, None, "", ""

                        def del_img(tid, user_state):
                            empty_df, empty_state = pad_data([], 4), []
                            if not user_state or not user_state.get("id"):
                                return None, "", "❌ Auth", empty_df, empty_state

                            delete_generated_image(int(tid or 0), user_state["id"])
                            d, s = load_img_data(user_state)
                            return None, "", "✅ Gelöscht", d, s

                        refresh_images_btn.click(load_img_data, inputs=[session_state], outputs=[images_history, img_state])
                        images_tab.select(load_img_data, inputs=[session_state], outputs=[images_history, img_state])

                        img_id_input.change(load_single_img, inputs=[img_id_input, session_state], outputs=[loaded_img_display, loaded_img_prompt, img_action_status])
                        
                        delete_img_btn.click(
                            del_img, 
                            inputs=[img_id_input, session_state], 
                            outputs=[loaded_img_display, loaded_img_prompt, img_action_status, images_history, img_state]
                        )
                        
                        images_history.select(
                            select_img_row, 
                            inputs=[img_state, session_state], 
                            outputs=[img_id_input, loaded_img_display, loaded_img_prompt, img_action_status]
                        )

                        # FIX: Direct reference to c_msg
                        img_to_chat_btn.click(lambda x: x, inputs=loaded_img_prompt, outputs=c_msg)

                    # =========================================================
                    # 3. CUSTOM PROMPTS
                    # =========================================================
                    with gr.TabItem("✏️ Eigene Prompt-Vorlagen") as prompts_tab:
                        prompt_state = gr.State([])
                        with gr.Row():
                            with gr.Column(scale=1):
                                gr.Markdown("#### ✨ Neue Vorlage")
                                new_prompt_name = gr.Textbox(label="Name")
                                new_prompt_category = gr.Dropdown(["Transkription", "Chat", "Vision", "Allgemein"], value="Chat", label="Kategorie")
                                new_prompt_text = gr.Textbox(label="Text", lines=12)
                                save_prompt_btn = gr.Button("💾 Speichern", variant="primary")
                                save_prompt_status = gr.Markdown("")

                            with gr.Column(scale=1):
                                gr.Markdown("#### 📂 Gespeicherte Vorlagen")
                                saved_prompts = gr.Dataframe(
                                    headers=["ID", "Name", "Kategorie"],
                                    value=[[None, "", ""]] * 6,
                                    interactive=False,
                                    wrap=True,
                                    datatype=["number", "str", "str"],
                                    column_widths=["15%", "50%", "35%"]
                                )

                                with gr.Row(variant="panel", equal_height=True):
                                    with gr.Column(scale=1):
                                        prompt_id_load = gr.Number(label="ID", precision=0, minimum=0)
                                    with gr.Column(scale=0):
                                        refresh_prompts_btn = gr.Button("🔄")
                                    with gr.Column(scale=1):
                                        prompt_to_chat_btn = gr.Button("📨 An Chat", variant="secondary")
                                    with gr.Column(scale=0):
                                        delete_prompt_btn = gr.Button("🗑️", variant="stop")

                                loaded_prompt_display = gr.Textbox(label="Vorschau", lines=5)

                        # --- LOGIC ---
                        def load_prompts_data(user_state=None):
                            if not user_state or not user_state.get("id"): 
                                return pad_data([], 3), []
                            p_list = get_user_custom_prompts(user_state["id"])
                            clean = [[p.id, p.name, p.category] for p in p_list]
                            return pad_data(list(clean), 3), clean

                        def load_single_prompt(tid, user_state=None):
                            if not tid or not user_state or not user_state.get("id"): return ""
                            db = SessionLocal()
                            p = db.query(CustomPrompt).filter(CustomPrompt.id == int(tid), CustomPrompt.user_id == user_state["id"]).first()
                            db.close()
                            return p.prompt_text if p else ""

                        def select_prompt_row(evt: gr.SelectData, state_data, user_state):
                            try:
                                if not user_state or not user_state.get("id"):
                                    return 0, "❌"

                                if state_data and evt.index[0] < len(state_data):
                                    tid = int(state_data[evt.index[0]][0])
                                    return tid, load_single_prompt(tid, user_state)
                            except: pass
                            return 0, ""

                        def save_p(n, c, t, user_state=None):
                            if not user_state or not user_state.get("id"):
                                return "❌ Auth Fehler", [], []
                            
                            save_custom_prompt(user_state["id"], n, t, c.lower())
                            d, s = load_prompts_data(user_state)
                            return "✅ Gespeichert", d, s

                        def del_p(tid, user_state=None):
                            empty_df, empty_state = pad_data([], 3), []
                            if not user_state or not user_state.get("id"):
                                return "❌ Auth Fehler", empty_df, empty_state

                            delete_custom_prompt(int(tid or 0), user_state["id"])
                            d, s = load_prompts_data(user_state)
                            return "", d, s

                        save_prompt_btn.click(save_p, inputs=[new_prompt_name, new_prompt_category, new_prompt_text, session_state], outputs=[save_prompt_status, saved_prompts, prompt_state])
                        refresh_prompts_btn.click(load_prompts_data, inputs=[session_state], outputs=[saved_prompts, prompt_state])
                        prompts_tab.select(load_prompts_data, inputs=[session_state], outputs=[saved_prompts, prompt_state])

                        prompt_id_load.change(load_single_prompt, inputs=[prompt_id_load, session_state], outputs=loaded_prompt_display)
                        
                        delete_prompt_btn.click(
                            del_p, 
                            inputs=[prompt_id_load, session_state], 
                            outputs=[loaded_prompt_display, saved_prompts, prompt_state]
                        )
                        saved_prompts.select(
                            select_prompt_row,
                            inputs=[prompt_state, session_state],
                            outputs=[prompt_id_load, loaded_prompt_display]
                        )
                                                
                        # FIX: Direct reference to c_msg
                        prompt_to_chat_btn.click(lambda x: x, inputs=loaded_prompt_display, outputs=c_msg)

                    
                    # =========================================================
                    # 4. RESUME FAILED JOBS
                    # =========================================================
                    with gr.TabItem("🔄 Abgebrochene Uploads") as jobs_tab:
                        gr.Markdown("### 🚧 Unvollständige Transkriptionen fortsetzen")
                        
                        failed_jobs_table = gr.Dataframe(
                            headers=["Job ID", "Datum", "Status", "Provider"],
                            value=[["", "", "", ""]],
                            interactive=False,
                            label="Offene Jobs"
                        )
                        
                        with gr.Row():
                            refresh_jobs_btn = gr.Button("🔄 Liste aktualisieren")
                            resume_job_id_input = gr.Number(label="Job ID zum Fortsetzen", precision=0)
                            resume_btn = gr.Button("▶️ Job Fortsetzen", variant="primary")
                        
                        resume_log = gr.Textbox(label="Resume Log", lines=5)
                        resume_result = gr.Textbox(label="Ergebnis", lines=10)

                        # Logic
                        def list_failed_jobs(user_state):
                            if not user_state or not user_state.get("id"):
                                return [["Bitte anmelden", "", "", ""]]
                            data = get_failed_jobs(user_state["id"]) # <--- PASS ID
                            return data if data else [["Keine offenen Jobs", "", "", ""]]

                        def resume_job_process(jid):
                            """Resume logic"""
                            try:
                                jid = int(jid)
                                path = os.path.join(JOB_STATE_DIR, f"{jid}.json")
                                if not os.path.exists(path): return "❌ Job nicht gefunden", ""
                                
                                with open(path, "r") as f:
                                    job = json.load(f)
                                
                                client = get_client(job["provider"])
                                chunk_paths = [c["path"] for c in job["chunks"]]
                                
                                logs = f"🚀 Setze Job {jid} fort..."
                                full_text = ""
                                
                                transcriber = run_chunked_api_transcription(
                                    client, job["model"], chunk_paths, job["lang"], 
                                    job["prompt"], job["temp"], job_id=jid
                                )
                                
                                for update in transcriber:
                                    if len(update) < 300:
                                        logs += f"\n{update}"
                                        yield logs, full_text
                                    else:
                                        full_text = update
                                        
                                yield logs + "\n🎉 Job abgeschlossen!", full_text
                                
                            except Exception as e:
                                yield f"🔥 Fehler: {e}", ""

                        # Wiring
                        refresh_jobs_btn.click(list_failed_jobs, inputs=[session_state], outputs=failed_jobs_table)
                        
                        # --- FIX: AUTO-LOAD ON TAB SELECT ---
                        jobs_tab.select(list_failed_jobs, inputs=[session_state], outputs=failed_jobs_table)

                        resume_btn.click(resume_job_process, inputs=resume_job_id_input, outputs=[resume_log, resume_result])
                        
                        # Auto-fill ID on click
                        def select_job(evt: gr.SelectData, data):
                            try: return int(data.iloc[evt.index[0], 0]) 
                            except: return 0
                        
                        failed_jobs_table.select(select_job, failed_jobs_table, resume_job_id_input)

                    # --- TAB 7: MODEL PREFERENCES ---
                    with gr.TabItem("🎯 Modell-Einstellungen"):
                        gr.Markdown("## 🎯 Bevorzugte Modelle verwalten")
                        gr.Markdown("Wähle welche Modelle in den Dropdown-Menüs erscheinen sollen und lege die Reihenfolge fest.")
                        
                        
                        def get_implemented_provider_choices():
                            """Return list of providers that are marked as implemented"""
                            return [p for p in PROVIDERS.keys() if is_provider_implemented(p)]
                        
                        with gr.Row():
                            pref_provider = gr.Dropdown(
                                choices=get_implemented_provider_choices(),
                                value="Scaleway",
                                label="Provider auswählen"
                            )
                        
                        with gr.Row():
                            fetch_models_btn = gr.Button("🔄 Verfügbare Modelle laden", variant="secondary")
                            fetch_status = gr.Markdown("")
                        
                        with gr.Row():
                            with gr.Column(scale=1):
                                gr.Markdown("### 📋 Verfügbare Modelle")
                                available_models_list = gr.Dataframe(
                                    headers=["Modell-ID", "Anzeigename"],
                                    value=[["", ""]],
                                    interactive=False,
                                    wrap=True,
                                    #height=400
                                )
                            
                            with gr.Column(scale=1):
                                gr.Markdown("### ✅ Deine Auswahl")
                                gr.Markdown("*Erste Modell = Standard. Drag & Drop zum Sortieren.*")
                                
                                selected_models_state = gr.State([])
                                
                                selected_models_display = gr.Dataframe(
                                    headers=["Reihenfolge", "Modell-ID", "Anzeigename", "Sichtbar"],
                                    value=[["", "", "", ""]],
                                    interactive=False,
                                    wrap=True,
                                    #height=400,
                                    datatype=["number", "str", "str", "bool"]
                                )
                                
                                with gr.Row():
                                    save_prefs_btn = gr.Button("💾 Einstellungen speichern", variant="primary")
                                    reset_prefs_btn = gr.Button("🔄 Zurücksetzen", variant="secondary")
                                
                                save_prefs_status = gr.Markdown("")
                        
                        with gr.Accordion("⚙️ Modell-Verwaltung", open=True):
                            with gr.Row():
                                with gr.Column():
                                    model_to_add = gr.Textbox(
                                        label="Modell-ID hinzufügen",
                                        placeholder="z.B. llama-3.3-70b-instruct"
                                    )
                                    model_display_name = gr.Textbox(
                                        label="Anzeigename (optional)",
                                        placeholder="z.B. Llama 3.3 70B"
                                    )
                                    add_model_btn = gr.Button("➕ Hinzufügen", variant="secondary")
                                
                                with gr.Column():
                                    model_to_remove_idx = gr.Number(
                                        label="Position zum Entfernen (Reihenfolge-Nummer)",
                                        precision=0,
                                        value=1
                                    )
                                    remove_model_btn = gr.Button("➖ Entfernen", variant="secondary")
                                
                                with gr.Column():
                                    move_from_idx = gr.Number(
                                        label="Von Position",
                                        precision=0,
                                        value=1
                                    )
                                    move_to_idx = gr.Number(
                                        label="Zu Position",
                                        precision=0,
                                        value=2
                                    )
                                    move_model_btn = gr.Button("↕️ Verschieben", variant="secondary")
                            
                            model_mgmt_status = gr.Markdown("")
                        
                        gr.Markdown("""
                        ### 💡 Tipps
                        
                        - **Erstes Modell** = Standard-Modell für diesen Provider
                        - **Unsichtbare Modelle** werden nicht in Dropdown-Menüs angezeigt
                        - Klicke auf "🔄 Verfügbare Modelle laden" um die neuesten Modelle vom Provider zu laden
                        - Änderungen gelten sofort nach dem Speichern
                        """)
                        
                        # =========================================================
                        # EVENT HANDLERS FOR MODEL PREFERENCES
                        # =========================================================
                        
                        def fetch_models_for_provider(provider, user_state=None):
                            """Fetch and display available models"""
                            if not user_state or not user_state.get("id"):
                                return [["", ""]], "❌ Bitte anmelden", gr.update()
                            
                            provider_key = API_KEYS.get(provider.lower(), "")
                            models, error = fetch_available_models(provider, provider_key)
                            
                            if error: return [["", ""]], f"❌ {error}", gr.update()
                            if not models: return [["", ""]], "⚠️ Keine Modelle gefunden", gr.update()
                            
                            model_data = [[m["id"], m.get("name", m["id"])] for m in models]
                            return model_data, f"✅ {len(models)} Modelle geladen", gr.update()

                        def load_user_preferences(provider, user_state=None):
                            """Load user's saved preferences for provider"""
                            if not user_state or not user_state.get("id"):
                                return [["", "", "", ""]], []
                            
                            prefs = get_user_model_preferences(user_state["id"], provider)
                            
                            if not prefs:
                                # Defaults
                                default_models = PROVIDERS.get(provider, {}).get("chat_models", [])
                                display_data = []
                                state_data = []
                                for i, model_id in enumerate(default_models, 1):
                                    display_data.append([i, model_id, model_id, True])
                                    state_data.append({
                                        "model_id": model_id,
                                        "display_name": model_id,
                                        "is_visible": True,
                                        "display_order": i
                                    })
                                return display_data, state_data
                            
                            display_data = []
                            state_data = []
                            for i, pref in enumerate(prefs, 1):
                                display_data.append([i, pref.model_id, pref.display_name or pref.model_id, "✅" if pref.is_visible else "❌"])
                                state_data.append({
                                    "model_id": pref.model_id,
                                    "display_name": pref.display_name or pref.model_id,
                                    "is_visible": pref.is_visible,
                                    "display_order": i
                                })
                            return display_data, state_data
                        
                        def add_model_to_selection(provider, model_id, display_name, current_state):
                            """Add a model to user's selection"""
                            if not model_id:
                                return gr.update(), current_state, "❌ Modell-ID erforderlich"
                            
                            # Check if already exists
                            if any(m["model_id"] == model_id for m in current_state):
                                return gr.update(), current_state, "⚠️ Modell bereits in Auswahl"
                            
                            # Add to state
                            new_model = {
                                "model_id": model_id,
                                "display_name": display_name or model_id,
                                "is_visible": True,
                                "display_order": len(current_state) + 1
                            }
                            current_state.append(new_model)
                            
                            # Update display
                            display_data = []
                            for i, m in enumerate(current_state, 1):
                                display_data.append([
                                    i,
                                    m["model_id"],
                                    m["display_name"],
                                    "✅" if m["is_visible"] else "❌"
                                ])
                            
                            return display_data, current_state, f"✅ '{model_id}' hinzugefügt"
                        
                        def remove_model_from_selection(idx, current_state):
                            """Remove a model from selection"""
                            if not current_state or idx < 1 or idx > len(current_state):
                                return gr.update(), current_state, "❌ Ungültige Position"
                            
                            removed = current_state.pop(idx - 1)
                            
                            # Update display
                            display_data = []
                            for i, m in enumerate(current_state, 1):
                                display_data.append([
                                    i,
                                    m["model_id"],
                                    m["display_name"],
                                    "✅" if m["is_visible"] else "❌"
                                ])
                            
                            return display_data, current_state, f"✅ '{removed['model_id']}' entfernt"
                        
                        def move_model_in_selection(from_idx, to_idx, current_state):
                            """Move a model in the selection order"""
                            if not current_state:
                                return gr.update(), current_state, "❌ Keine Modelle vorhanden"
                            
                            if from_idx < 1 or from_idx > len(current_state):
                                return gr.update(), current_state, "❌ Ungültige Ausgangsposition"
                            
                            if to_idx < 1 or to_idx > len(current_state):
                                return gr.update(), current_state, "❌ Ungültige Zielposition"
                            
                            # Move item
                            item = current_state.pop(from_idx - 1)
                            current_state.insert(to_idx - 1, item)
                            
                            # Update display
                            display_data = []
                            for i, m in enumerate(current_state, 1):
                                display_data.append([
                                    i,
                                    m["model_id"],
                                    m["display_name"],
                                    "✅" if m["is_visible"] else "❌"
                                ])
                            
                            return display_data, current_state, f"✅ Modell von Position {from_idx} zu {to_idx} verschoben"
                        
                        def save_preferences(provider, current_state, user_state=None):
                            if not user_state or not user_state.get("id"):
                                return "❌ Bitte anmelden"
                            
                            if not current_state:
                                return "⚠️ Keine Modelle ausgewählt"
                            
                            # Update display_order
                            for i, model in enumerate(current_state):
                                model["display_order"] = i
                            
                            success, message = save_user_model_preferences(user_state["id"], provider, current_state)
                            return message
                        
                        def reset_to_defaults(provider):
                            """Reset to default provider models"""
                            default_models = PROVIDERS.get(provider, {}).get("chat_models", [])
                            
                            display_data = []
                            state_data = []
                            for i, model_id in enumerate(default_models, 1):
                                display_data.append([i, model_id, model_id, "✅"])
                                state_data.append({
                                    "model_id": model_id,
                                    "display_name": model_id,
                                    "is_visible": True,
                                    "display_order": i
                                })
                            
                            return display_data, state_data, "✅ Auf Standard zurückgesetzt"
                        
                        # Wire up event handlers
                        fetch_models_btn.click(
                            fetch_models_for_provider,
                            inputs=[pref_provider, session_state],
                            outputs=[available_models_list, fetch_status, selected_models_display]
                        )

                        pref_provider.change(
                            load_user_preferences,
                            inputs=[pref_provider, session_state],
                            outputs=[selected_models_display, selected_models_state]
                        )
                        
                        add_model_btn.click(
                            add_model_to_selection,
                            inputs=[pref_provider, model_to_add, model_display_name, selected_models_state],
                            outputs=[selected_models_display, selected_models_state, model_mgmt_status]
                        )
                        
                        remove_model_btn.click(
                            remove_model_from_selection,
                            inputs=[model_to_remove_idx, selected_models_state],
                            outputs=[selected_models_display, selected_models_state, model_mgmt_status]
                        )
                        
                        move_model_btn.click(
                            move_model_in_selection,
                            inputs=[move_from_idx, move_to_idx, selected_models_state],
                            outputs=[selected_models_display, selected_models_state, model_mgmt_status]
                        )
                        
                        save_prefs_btn.click(
                            save_preferences,
                            inputs=[pref_provider, selected_models_state, session_state],
                            outputs=[save_prefs_status]
                        )
                        
                        reset_prefs_btn.click(
                            reset_to_defaults,
                            inputs=[pref_provider],
                            outputs=[selected_models_display, selected_models_state, save_prefs_status]
                        )
                    
                    
                    with gr.TabItem("⚙️ Einstellungen") as settings_tab:  # <-- Add 'as settings_tab'
                        gr.Markdown("### 🎛️ Persönliche Einstellungen")
                        
                        with gr.Group():
                            gr.Markdown("#### 📎 Dateianhänge & Chunking")
                            
                            setting_auto_chunk = gr.Checkbox(
                                value=True,
                                label="Automatisches Chunking für große Dateien",
                                info="Teilt große Dateien automatisch in kleinere Stücke"
                            )
                            
                            setting_chunk_size = gr.Slider(
                                minimum=1000,
                                maximum=32000,
                                value=4000,
                                step=1000,
                                label="Chunk-Größe (in Tokens)",
                                info="Größe jedes Teils bei automatischer Teilung"
                            )
                            
                            setting_overlap = gr.Slider(
                                minimum=0,
                                maximum=500,
                                value=200,
                                step=50,
                                label="Überlappung zwischen Chunks",
                                info="Wie viele Tokens zwischen Teilen überlappen"
                            )
                            
                            gr.Markdown("---")
                            gr.Markdown("#### 💬 Chat-Verhalten")
                            
                            setting_auto_truncate = gr.Checkbox(
                                value=True,
                                label="Automatisches Kürzen bei Context-Limit",
                                info="Entfernt alte Nachrichten automatisch wenn Context voll"
                            )
                            
                            setting_show_warning = gr.Checkbox(
                                value=True,
                                label="Warnung bei gekürztem Verlauf anzeigen",
                                info="Zeigt Info-Meldung wenn Nachrichten entfernt wurden"
                            )
                            
                            setting_show_tokens = gr.Checkbox(
                                value=False,
                                label="Token-Zähler anzeigen",
                                info="Zeigt geschätzte Token-Anzahl in Status-Nachrichten"
                            )
                        
                        with gr.Group():
                            gr.Markdown("#### 🤖 Modell-Context-Limits prüfen")
                            
                            with gr.Row():
                                check_provider = gr.Dropdown(
                                    choices=list(PROVIDERS.keys()),
                                    value="Scaleway",
                                    label="Provider"
                                )
                                check_model = gr.Dropdown(
                                    choices=[],
                                    label="Modell"
                                )
                            
                            check_btn = gr.Button("📊 Context-Limit prüfen")
                            limit_display = gr.Markdown("")
                        
                        with gr.Row():
                            save_settings_btn = gr.Button("💾 Einstellungen speichern", variant="primary", size="lg")
                            reset_settings_btn = gr.Button("🔄 Auf Standard zurücksetzen", variant="secondary", size="lg")
                        
                        settings_status = gr.Markdown("")

                    # Wire up the event handlers

                    # 1. Load user settings when tab is opened
                    def load_user_settings_ui(user_state):
                        """Load settings from database"""
                        if not user_state or not user_state.get("id"):
                            # Return defaults
                            return [True, 4000, 200, True, True, False, ""]
                        
                        try:
                            settings = get_user_settings(user_state["id"])
                            return [
                                settings.auto_chunk_enabled,
                                settings.chunk_size,
                                settings.chunk_overlap,
                                settings.auto_truncate_history,
                                settings.show_truncation_warning,
                                settings.show_token_counts,
                                ""  # Clear status
                            ]
                        except Exception as e:
                            logger.error(f"Error loading settings: {e}")
                            return [True, 4000, 200, True, True, False, f"⚠️ Fehler beim Laden: {str(e)}"]

                    settings_tab.select(
                        load_user_settings_ui,
                        inputs=[session_state],
                        outputs=[
                            setting_auto_chunk,
                            setting_chunk_size,
                            setting_overlap,
                            setting_auto_truncate,
                            setting_show_warning,
                            setting_show_tokens,
                            settings_status
                        ]
                    )

                    # 2. Update model dropdown when provider changes
                    def update_check_model_dropdown(provider):
                        """Update model choices based on provider"""
                        models = PROVIDERS.get(provider, {}).get("chat_models", [])
                        if models:
                            return gr.update(choices=models, value=models[0])
                        return gr.update(choices=[], value=None)

                    check_provider.change(
                        update_check_model_dropdown,
                        inputs=[check_provider],
                        outputs=[check_model]
                    )

                    # 3. Show context limit info
                    def show_limit_info(provider, model):
                        """Display context limit information for selected model"""
                        if not model:
                            return "❌ Bitte Modell auswählen"
                        
                        try:
                            limit = get_model_context_limit(provider, model)
                            char_estimate = limit * 4
                            
                            info = f"""
                    ### 📊 Context-Limit für `{model}`

                    - **Token-Limit:** {limit:,} Tokens
                    - **Zeichen (ca.):** {char_estimate:,} Zeichen
                    - **Reserviert für Antwort:** ~1,000 Tokens
                    - **Nutzbar für Eingabe:** ~{(limit - 1000):,} Tokens

                    💡 **Tipp:** Bei großen Dokumenten automatisches Chunking aktivieren!
                    """
                            return info
                        except Exception as e:
                            return f"❌ Fehler: {str(e)}"

                    check_btn.click(
                        show_limit_info,
                        inputs=[check_provider, check_model],
                        outputs=[limit_display]
                    )

                    # 4. Save settings to database
                    def save_user_settings_ui(auto_chunk, chunk_size, overlap, auto_trunc, 
                                            show_warn, show_tokens, user_state):
                        """Save settings to database"""
                        if not user_state or not user_state.get("id"):
                            return "❌ Bitte anmelden"
                        
                        try:
                            success, msg = update_user_settings(
                                user_state["id"],
                                auto_chunk_enabled=bool(auto_chunk),
                                chunk_size=int(chunk_size),
                                chunk_overlap=int(overlap),
                                auto_truncate_history=bool(auto_trunc),
                                show_truncation_warning=bool(show_warn),
                                show_token_counts=bool(show_tokens)
                            )
                            return msg
                        except Exception as e:
                            logger.error(f"Error saving settings: {e}")
                            return f"🔥 Fehler beim Speichern: {str(e)}"

                    save_settings_btn.click(
                        save_user_settings_ui,
                        inputs=[
                            setting_auto_chunk,
                            setting_chunk_size,
                            setting_overlap,
                            setting_auto_truncate,
                            setting_show_warning,
                            setting_show_tokens,
                            session_state
                        ],
                        outputs=[settings_status]
                    )

                    # 5. Reset to defaults
                    def reset_to_defaults(user_state):
                        """Reset settings to default values"""
                        if not user_state or not user_state.get("id"):
                            return [True, 4000, 200, True, True, False, "❌ Bitte anmelden"]
                        
                        try:
                            # Delete user settings to trigger default creation
                            db = SessionLocal()
                            db.query(UserSettings).filter(UserSettings.user_id == user_state["id"]).delete()
                            db.commit()
                            db.close()
                            
                            return [True, 4000, 200, True, True, False, "✅ Auf Standard zurückgesetzt"]
                        except Exception as e:
                            logger.error(f"Error resetting settings: {e}")
                            return [True, 4000, 200, True, True, False, f"🔥 Fehler: {str(e)}"]

                    reset_settings_btn.click(
                        reset_to_defaults,
                        inputs=[session_state],
                        outputs=[
                            setting_auto_chunk,
                            setting_chunk_size,
                            setting_overlap,
                            setting_auto_truncate,
                            setting_show_warning,
                            setting_show_tokens,
                            settings_status
                        ]
                    )
                        
            # --- TAB 6: USER MANAGEMENT (ADMIN ONLY) ---
            with gr.TabItem("👥 Benutzerverwaltung", visible=False) as admin_tab:
                gr.Markdown("## 👥 Benutzerverwaltung")
                gr.Markdown("*Nur für Administratoren*")
                
                with gr.Tabs():
                    # =========================================================
                    # CREATE NEW USER
                    # =========================================================
                    with gr.TabItem("➕ Neuer Benutzer"):
                        gr.Markdown("### Neuen Benutzer erstellen")
                        
                        with gr.Row():
                            with gr.Column(scale=1):
                                new_user_username = gr.Textbox(
                                    label="Benutzername",
                                    placeholder="max.mustermann"
                                )
                                new_user_password = gr.Textbox(
                                    label="Passwort",
                                    placeholder="Mindestens 8 Zeichen",
                                    type="password"
                                )
                                new_user_password_confirm = gr.Textbox(
                                    label="Passwort bestätigen",
                                    placeholder="Passwort wiederholen",
                                    type="password"
                                )
                                new_user_email = gr.Textbox(
                                    label="E-Mail (Optional)",
                                    placeholder="max@example.com"
                                )
                                new_user_is_admin = gr.Checkbox(
                                    label="Als Administrator erstellen",
                                    value=False
                                )
                                
                                create_user_btn = gr.Button("➕ Benutzer erstellen", variant="primary", size="lg")
                                create_user_status = gr.Markdown("")
                            
                            with gr.Column(scale=1):
                                gr.Markdown("""
                                ### 📋 Hinweise
                                
                                **Benutzername:**
                                - Muss eindeutig sein
                                - Keine Leerzeichen
                                - Empfohlen: kleinbuchstaben
                                
                                **Passwort:**
                                - Mindestens 8 Zeichen empfohlen
                                - Wird sicher verschlüsselt (bcrypt)
                                
                                **Administrator:**
                                - Admins können:
                                - Alle Benutzer verwalten
                                - Andere Admins erstellen
                                - Alle Daten sehen

                                """)
                    
                    # =========================================================
                    # MANAGE EXISTING USERS
                    # =========================================================
                    with gr.TabItem("⚙️ Benutzer verwalten"):
                        gr.Markdown("### Bestehende Benutzer verwalten")
                        
                        with gr.Row():
                            refresh_users_btn = gr.Button("🔄 Liste aktualisieren", size="sm")
                        
                        users_table = gr.Dataframe(
                            headers=["ID", "Benutzername", "E-Mail", "Rolle", "Erstellt"],
                            value=[["", "", "", "", ""]],
                            interactive=False,
                            wrap=True,
                            #height=400,
                            datatype=["number", "str", "str", "str", "str"],
                            column_widths=["10%", "25%", "25%", "20%", "20%"]
                        )
                        
                        with gr.Row():
                            selected_user_id = gr.Number(
                                label="Benutzer-ID",
                                precision=0,
                                value=0
                            )
                        
                        with gr.Tabs():
                            # RENAME USER
                            with gr.TabItem("✏️ Umbenennen"):
                                with gr.Row():
                                    rename_new_username = gr.Textbox(
                                        label="Neuer Benutzername",
                                        placeholder="neuer.name"
                                    )
                                    rename_user_btn = gr.Button("✏️ Umbenennen", variant="secondary")
                                rename_status = gr.Markdown("")
                            
                            # RESET PASSWORD
                            with gr.TabItem("🔑 Passwort zurücksetzen"):
                                with gr.Row():
                                    reset_new_password = gr.Textbox(
                                        label="Neues Passwort",
                                        placeholder="Neues Passwort",
                                        type="password"
                                    )
                                    reset_confirm_password = gr.Textbox(
                                        label="Passwort bestätigen",
                                        placeholder="Passwort wiederholen",
                                        type="password"
                                    )
                                reset_password_btn = gr.Button("🔑 Passwort zurücksetzen", variant="secondary")
                                reset_password_status = gr.Markdown("")
                            
                            # UPDATE EMAIL
                            with gr.TabItem("📧 E-Mail ändern"):
                                with gr.Row():
                                    update_email_input = gr.Textbox(
                                        label="Neue E-Mail",
                                        placeholder="neue@email.de"
                                    )
                                    update_email_btn = gr.Button("📧 E-Mail aktualisieren", variant="secondary")
                                update_email_status = gr.Markdown("")
                            
                            # TOGGLE ADMIN
                            with gr.TabItem("⬆️⬇️ Admin-Status"):
                                gr.Markdown("""
                                ### Admin-Status umschalten
                                
                                **Achtung:** 
                                - Du kannst deinen eigenen Status nicht ändern
                                - Admins haben volle Kontrolle über die App
                                """)
                                toggle_admin_btn = gr.Button("⬆️⬇️ Admin-Status umschalten", variant="secondary")
                                toggle_admin_status = gr.Markdown("")

                            with gr.TabItem("🎬 Medienverwalter"):
                                gr.Markdown("""
                                ### 🎬 Medienverwalter-Status umschalten:
                                - User kann alle YouTube-URLs herunterladen (nicht nur Whitelist)
                                """)
                                toggle_media_btn = gr.Button("🎬 Medienverwalter-Status umschalten", variant="secondary")
                                toggle_media_status = gr.Markdown("")

                            # Event Handler
                            toggle_media_btn.click(
                                toggle_media_manager_status,
                                inputs=[selected_user_id, session_state],
                                outputs=[toggle_media_status, users_table]
                            )
                            
                            # DELETE USER
                            with gr.TabItem("🗑️ Benutzer löschen"):
                                gr.Markdown("""
                                ### ⚠️ WARNUNG: Benutzer löschen
                                
                                Diese Aktion:
                                - Löscht den Benutzer **permanent**
                                - Löscht alle zugehörigen Daten (Chats, Transkripte, etc.)
                                - **Kann nicht rückgängig gemacht werden**
                                - Du kannst dich nicht selbst löschen
                                """)
                                
                                with gr.Row():
                                    delete_confirm = gr.Textbox(
                                        label="Bestätigung",
                                        placeholder="Tippe 'LÖSCHEN' zur Bestätigung"
                                    )
                                    delete_user_btn = gr.Button("🗑️ BENUTZER LÖSCHEN", variant="stop")
                                delete_user_status = gr.Markdown("")
                        
                        # Select user from table
                        def select_user_from_table(evt: gr.SelectData):
                            """Select user by clicking on table row"""
                            try:
                                # Get the row data
                                row_idx = evt.index[0]
                                # Load fresh data
                                users_data = get_all_users()
                                if row_idx < len(users_data):
                                    user_id = users_data[row_idx][0]
                                    return int(user_id)
                            except Exception as e:
                                logger.exception(f"Error selecting user: {str(e)}")
                            return gr.update()
                        
                        users_table.select(select_user_from_table, outputs=selected_user_id)

                    with gr.TabItem("📝 Kanal-Whitelist"):
                        gr.Markdown("### 📝 YouTube-Kanal Whitelist verwalten")
                        
                        current_whitelist = gr.Dataframe(
                            headers=["Kanal-Identifier"],
                            value=[[ch] for ch in YOUTUBE_CHANNEL_WHITELIST],
                            interactive=False,
                            #height=300
                        )
                        
                        with gr.Row():
                            new_channel = gr.Textbox(
                                label="Neuen Kanal hinzufügen",
                                placeholder="z.B. @meinkanal oder UCxyz123..."
                            )
                            add_channel_btn = gr.Button("➕ Hinzufügen", variant="primary")
                        
                        with gr.Row():
                            remove_channel = gr.Textbox(
                                label="Kanal entfernen",
                                placeholder="Exakte Schreibweise eingeben"
                            )
                            remove_channel_btn = gr.Button("➖ Entfernen", variant="stop")
                        
                        whitelist_status = gr.Markdown("")
                        
                        gr.Markdown("""
                        **💡 Hinweis:** Änderungen werden in `app.py` gespeichert und 
                        benötigen einen Neustart des Services.
                        
                        Unterstützte Formate:
                        - `AkademieKanal` (User-URL)
                        - `@theologisches-forum` (Handle)
                        - `UCxyz123abc...` (Channel-ID)
                        """)

                    def add_to_whitelist(channel_name):
                        """Add channel to whitelist (requires file write)"""
                        if not channel_name or channel_name.strip() == "":
                            return "❌ Kanal-Name erforderlich", gr.update()
                        
                        channel_name = channel_name.strip()
                        
                        if channel_name in YOUTUBE_CHANNEL_WHITELIST:
                            return "⚠️ Kanal bereits in Whitelist", gr.update()
                        
                        try:
                            # Update in-memory list
                            YOUTUBE_CHANNEL_WHITELIST.append(channel_name)
                            
                            # Update config file
                            config_path = "/var/www/transkript_app/app.py"
                            with open(config_path, 'r', encoding='utf-8') as f:
                                content = f.read()
                            
                            # Find and update the whitelist section
                            import re
                            pattern = r'(YOUTUBE_CHANNEL_WHITELIST = \[)(.*?)(\])'
                            
                            def replace_whitelist(match):
                                entries = ',\n    '.join([f'"{ch}"' for ch in YOUTUBE_CHANNEL_WHITELIST])
                                return f'{match.group(1)}\n    {entries}\n{match.group(3)}'
                            
                            new_content = re.sub(pattern, replace_whitelist, content, flags=re.DOTALL)
                            
                            with open(config_path, 'w', encoding='utf-8') as f:
                                f.write(new_content)
                            
                            # Update display
                            updated_data = [[ch] for ch in YOUTUBE_CHANNEL_WHITELIST]
                            
                            return "✅ Hinzugefügt. Bitte Service neustarten: `systemctl restart akademie_suite`", gr.update(value=updated_data)
                            
                        except Exception as e:
                            logger.error(f"Error updating whitelist: {e}")
                            return f"❌ Fehler: {str(e)}", gr.update()

                    add_channel_btn.click(add_to_whitelist, inputs=[new_channel], outputs=[whitelist_status, current_whitelist])
                
                # =========================================================
                # EVENT HANDLERS
                # =========================================================
                
                # Create user
                def handle_create_user(username, password, password_confirm, email, is_admin):
                    if not username or not password:
                        return "❌ Benutzername und Passwort sind erforderlich"
                    
                    if password != password_confirm:
                        return "❌ Passwörter stimmen nicht überein"
                    
                    if len(password) < 8:
                        return "⚠️ Warnung: Passwort sollte mindestens 8 Zeichen haben"
                    
                    success, message = create_user(username, password, email, is_admin)
                    
                    if success:
                        # Clear form
                        return message
                    return message
                
                create_user_btn.click(
                    handle_create_user,
                    inputs=[
                        new_user_username,
                        new_user_password,
                        new_user_password_confirm,
                        new_user_email,
                        new_user_is_admin
                    ],
                    outputs=create_user_status
                )
                
                # Refresh users list
                def load_users_list(user_state=None):
                    if not user_state or not user_state.get("is_admin"):
                        return []
                    return get_all_users()
                        
                refresh_users_btn.click(
                    load_users_list, 
                    inputs=[session_state],
                    outputs=users_table
                )
                
                # Rename user
                def handle_rename_user(user_id, new_username):
                    if not new_username:
                        return "❌ Neuer Benutzername erforderlich", get_all_users()
                    
                    success, message = rename_user(int(user_id), new_username)
                    return message, get_all_users()
                
                rename_user_btn.click(
                    handle_rename_user,
                    inputs=[selected_user_id, rename_new_username],
                    outputs=[rename_status, users_table]
                )
                
                # Reset password
                def handle_reset_password(user_id, new_password, confirm_password):
                    if not new_password:
                        return "❌ Neues Passwort erforderlich"
                    
                    if new_password != confirm_password:
                        return "❌ Passwörter stimmen nicht überein"
                    
                    if len(new_password) < 8:
                        return "⚠️ Warnung: Passwort sollte mindestens 8 Zeichen haben"
                    
                    success, message = reset_user_password(int(user_id), new_password)
                    return message
                
                reset_password_btn.click(
                    handle_reset_password,
                    inputs=[selected_user_id, reset_new_password, reset_confirm_password],
                    outputs=reset_password_status
                )
                
                # Update email
                def handle_update_email(user_id, new_email):
                    success, message = update_user_email(int(user_id), new_email)
                    return message, get_all_users()
                
                update_email_btn.click(
                    handle_update_email,
                    inputs=[selected_user_id, update_email_input],
                    outputs=[update_email_status, users_table]
                )
                
                # Toggle admin status
                def handle_toggle_admin(user_id, user_state=None):
                    if not user_state or not user_state.get("id"):
                        return "❌ Nicht angemeldet", get_all_users()
                    
                    # Check if requesting user is actually admin
                    if not user_state.get("is_admin"):
                        return "⛔ Keine Berechtigung", get_all_users()

                    success, message = toggle_admin_status(int(user_id), user_state["id"])
                    return message, get_all_users()
                
                toggle_admin_btn.click(
                    handle_toggle_admin,
                    inputs=[selected_user_id, session_state],
                    outputs=[toggle_admin_status, users_table]
                )
                
                # Delete user
                def handle_delete_user(user_id, confirmation, user_state=None):
                    if confirmation != "LÖSCHEN":
                        return "❌ Bestätigung erforderlich: Tippe 'LÖSCHEN'", get_all_users()
                    
                    if not user_state or not user_state.get("id"):
                        return "❌ Nicht angemeldet", get_all_users()

                    if not user_state.get("is_admin"):
                        return "⛔ Keine Berechtigung", get_all_users()
                    
                    success, message = delete_user(int(user_id), user_state["id"])
                    return message, get_all_users()
                
                delete_user_btn.click(
                    handle_delete_user,
                    inputs=[selected_user_id, delete_confirm, session_state],
                    outputs=[delete_user_status, users_table]
                )

                # --- AUTO-LOAD ON TAB SWITCH ---
                trans_tab.select(fn=load_trans_data, outputs=[trans_history, trans_state])
                images_tab.select(fn=load_img_data, outputs=[images_history, img_state])
                prompts_tab.select(fn=load_prompts_data, outputs=[saved_prompts, prompt_state])
                
    def attempt_auto_login(json_str):
        """
        Called by JS on page load. 
        Input is a JSON string: '["user", "pass"]'
        """
        if not json_str or json_str == "":
            # Normal state (first visit or logged out)
            return (
                "", gr.update(visible=False), gr.update(visible=True), 
                "👤 Nicht angemeldet", gr.update(visible=False), 
                gr.update(visible=False), {"id": None, "username": None, "is_admin": False},
                gr.update(), gr.update(), gr.update()
            )

        # print(f"🔄 Auto-login RAW data: {json_str}") # Debug

        try:
            creds = json.loads(json_str)
            
            if not isinstance(creds, list) or len(creds) < 2:
                print(f"❌ Auto-login failed: Invalid format {type(creds)}")
                raise ValueError("Invalid format")
                
            username, password = creds[0], creds[1]
            print(f"🔑 Authenticating stored user: {username}")
            
            # Call the existing login handler
            return handle_login(username, password)
            
        except Exception as e:
            print(f"❌ Auto-login error: {e}")
            return (
                "", gr.update(visible=False), gr.update(visible=True), 
                "👤 Fehler beim Auto-Login", gr.update(visible=False), 
                gr.update(visible=False), {"id": None, "username": None, "is_admin": False},
                gr.update(), gr.update(), gr.update()
            )

    def handle_login(username, password):
        """Enhanced login with storage path initialization"""
        success, message, show_app, show_login, state_data = login_user(username, password)
        
        # status_text = f"👤 Angemeldet als: **{state_data['username']}**" if success else "👤 Nicht angemeldet"
        # SHORTER STATUS TEXT for mobile
        status_text = f"👤 {state_data['username']}" if success else "👤"
        
        show_admin_tab = state_data.get("is_admin", False)
        
        # Storage root for file browsers (user-specific)
        storage_root = None
        
        if success:
            try:
                # Create user storage directory
                if ensure_user_storage_dirs(username):
                    storage_root = get_user_storage_path(username)
                    logger.info(f"🗂️ Storage root for {username}: {storage_root}")
                else:
                    logger.warning(f"⚠️ Could not initialize storage for {username}")
            except Exception as e:
                logger.error(f"❌ Storage initialization error for {username}: {e}")
        
        # Return storage_root as an additional output to update FileExplorer components
        return (
            message,                        # login_message
            show_app,                       # main_app visibility
            show_login,                     # login_screen visibility
            status_text,                    # login_status
            gr.update(visible=True),        # logout_btn
            gr.update(visible=show_admin_tab), # admin_tab
            state_data,                     # session_state
            gr.update(root_dir=storage_root) if storage_root else gr.update(), # t_storage_browser
            gr.update(root_dir=storage_root) if storage_root else gr.update(), # v_storage_browser
            gr.update(root_dir=storage_root) if storage_root else gr.update()  # attach_sb_browser
        )
    
    def handle_logout():
        message, show_app, show_login, empty_state = logout_user()
        return message, show_app, show_login, "👤 Nicht angemeldet", gr.update(visible=False), gr.update(visible=False), empty_state
    
    # 1. Perform Python Login
    login_btn.click(
        fn=handle_login,
        inputs=[login_username, login_password],
        outputs=[
            login_message, main_app, login_screen, login_status, 
            logout_btn, admin_tab, session_state,
            t_storage_browser, v_storage_browser, attach_sb_browser
        ],
        # INLINE JS: Save credentials directly to LocalStorage
        # js="(u, p) => { localStorage.setItem('ak_user', u); localStorage.setItem('ak_pass', p); }"
        js="(u, p) => { localStorage.setItem('ak_user', u); localStorage.setItem('ak_pass', p); return [u, p]; }"
    
    )

    # 2. Logout Button Click
    # Python logout + JS clears LocalStorage
    logout_btn.click(
        fn=handle_logout,
        inputs=None,
        outputs=[login_message, main_app, login_screen, login_status, logout_btn, admin_tab, session_state],
        # INLINE JS: Clear LocalStorage
        # js="() => { localStorage.removeItem('ak_user'); localStorage.removeItem('ak_pass'); }"
        js="() => { localStorage.removeItem('ak_user'); localStorage.removeItem('ak_pass'); }"
    
    )

    # 3. AUTO-LOGIN ON LOAD (PWA Persistence)
    # Hidden component to act as bridge between JS and Python
    hidden_creds_bridge = gr.Textbox(visible=False, label="Creds Bridge")

    # On App Load: JS reads LocalStorage -> Feeds into hidden_creds -> Triggers Python Login
    demo.load(
        fn=None,
        inputs=None,
        outputs=[hidden_creds_bridge],
        #js="() => { const u = localStorage.getItem('ak_user'); const p = localStorage.getItem('ak_pass'); return (u && p) ? JSON.stringify([u, p]) : ''; }"
        js="() => { const u = localStorage.getItem('ak_user'); const p = localStorage.getItem('ak_pass'); return (u && p) ? JSON.stringify([u, p]) : ''; }"
    
        #js="async () => { await new Promise(r => setTimeout(r, 800)); return getCredsV2(); }"
    )

    # Python Trigger
    hidden_creds_bridge.change(
        fn=attempt_auto_login,
        inputs=[hidden_creds_bridge],
        outputs=[
            login_message, main_app, login_screen, login_status, 
            logout_btn, admin_tab, session_state,
            t_storage_browser, v_storage_browser, attach_sb_browser
        ]
    )

    # Assign ID for CSS Mobile Fix
    c_bot.elem_id = "chat_window"

# ==========================================
# 🚀 STARTUP SEQUENCE
# ==========================================

def initialize_application():
    """Complete application initialization with security checks"""
    
    print("=" * 60)
    print("🔐 KI SUITE - SECURE STARTUP")
    print("=" * 60)
    
    # 1. Database Schema
    print("\n1️⃣ Checking database schema...")
    ensure_database_schema()
    Base.metadata.create_all(bind=engine)  # Create new encrypted tables
    
    # 2. Encryption Check
    print("\n2️⃣ Checking encryption status...")
    ensure_encryption()  # Migrates old data if needed
    
    # 3. Crypto Health Check
    print("\n3️⃣ Cryptography status:")
    print(f"   Master Key: {'✅ Loaded' if crypto.master_key else '❌ Missing'}")
    print(f"   Post-Quantum: {'✅ Enabled (Kyber-512)' if HAS_PQ else '⚠️ Disabled (pqcrypto not installed)'}")
    
    # 4. Provider Mode
    print("\n4️⃣ Provider access mode:")
    print(f"   EU-Only Mode: {'✅ ENABLED (US providers blocked)' if EU_ONLY_MODE else '⚠️ DISABLED'}")
    
    # 5. Default Users
    print("\n5️⃣ Checking default users...")
    create_default_users()
    
    # 6. Security Recommendations
    print("\n6️⃣ Security recommendations:")
    if not EU_ONLY_MODE:
        print("   ⚠️ WARNING: US providers are enabled! Set EU_ONLY_MODE=True")
    if not HAS_PQ:
        print("   ⚠️ INFO: Install pqcrypto for post-quantum encryption")
    
    print("\n" + "=" * 60)
    print("✅ Secure initialization complete!")
    print("=" * 60 + "\n")

# Call this before launching
initialize_application()

# ==========================================
# 🚀 LAUNCH CONFIGURATION
# ==========================================
if __name__ == "__main__":
    from fastapi import FastAPI, Request, status
    from fastapi.responses import JSONResponse
    import uvicorn

    # 1. Configuration Constants
    APP_DIR = "/var/www/transkript_app"
    LOG_FILE = os.path.join(APP_DIR, "app.log")
    STATIC_DIR = os.path.join(APP_DIR, "static")
    IMAGES_DIR = os.path.join(APP_DIR, "generated_images")

    # 2. Ensure directories and permissions exist
    os.makedirs(IMAGES_DIR, exist_ok=True)
    
    if not os.path.exists(LOG_FILE):
        try:
            with open(LOG_FILE, 'w') as f: f.write("")
            os.chmod(LOG_FILE, 0o666) 
        except Exception as e:
            print(f"⚠️ Could not create log file: {e}")

    # 3. Define FastAPI Wrapper for Security
    app = FastAPI()

    @app.middleware("http")
    async def block_api_endpoints(request: Request, call_next):
        # Allow internal UI paths (/run, /queue), block external API access
        if request.url.path.startswith("/api"):
             return JSONResponse(
                status_code=status.HTTP_403_FORBIDDEN,
                content={"detail": "External API access is disabled."}
            )
        response = await call_next(request)
        return response

    print(f"🚀 Starting Server on Port 7860 (Fast Shutdown Enabled)...")
    print(f"📂 Serving files from: {APP_DIR}")

    # --- CRITICAL CSS FIX ---
    # We assign the CSS/Head directly to the properties because 
    # the constructor arguments are being ignored in this version.
    demo.css = CUSTOM_CSS
    demo.head = PWA_HEAD
    # ------------------------

    # 4. Mount Gradio
    print(f"🚀 Mounting Gradio app...")
    app = gr.mount_gradio_app(
        app, 
        demo, 
        path="/", 
        allowed_paths=[APP_DIR, STATIC_DIR, IMAGES_DIR, "/tmp/gradio"]
    )

    # 5. Run Server
    print(f"🚀 Starting Server on Port 7860...")
    config = uvicorn.Config(
        app, 
        host="0.0.0.0", 
        port=7860, 
        timeout_graceful_shutdown=1,
        log_level="info"
    )
    server = uvicorn.Server(config)
    server.run()