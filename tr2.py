#!/usr/bin/env python3
"""
Office Documents Translator with CTranslate2 and LLM Fallbacks

Supports: .pptx, .docx, .xlsx
Primary: CTranslate2 WMT21 model
Fallback: Multiple LLM providers
"""

import argparse
import logging
import sys
import os
import time
from pathlib import Path
from typing import List, Dict, Optional, Any
from dataclasses import dataclass
from collections import defaultdict
import requests

logger = logging.getLogger(__name__)

# === Quick Imports Check ===
print("🔍 Checking libraries...")

try:
    from pptx import Presentation
    HAS_PPTX = True
    print("  ✓ python-pptx")
except ImportError:
    HAS_PPTX = False
    print("  ✗ python-pptx")

try:
    from docx import Document
    HAS_DOCX = True
    print("  ✓ python-docx")
except ImportError:
    HAS_DOCX = False
    print("  ✗ python-docx")

try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
    print("  ✓ openpyxl")
except ImportError:
    HAS_OPENPYXL = False
    print("  ✗ openpyxl")

# Deferred imports for heavy libraries
HAS_CTRANSLATE2 = False
HAS_HF_HUB = False
HAS_TRANSFORMERS = False

def check_ctranslate2():
    global HAS_CTRANSLATE2
    try:
        import ctranslate2
        HAS_CTRANSLATE2 = True
        return True
    except ImportError:
        return False

def check_transformers():
    global HAS_TRANSFORMERS
    try:
        print("  ⏳ Loading transformers...", end='', flush=True)
        import transformers
        HAS_TRANSFORMERS = True
        print(" ✓")
        return True
    except ImportError:
        print(" ✗")
        return False

def check_hf_hub():
    global HAS_HF_HUB
    try:
        from huggingface_hub import snapshot_download
        HAS_HF_HUB = True
        return True
    except ImportError:
        return False

print("-" * 60)


@dataclass
class TextSegment:
    text: str
    metadata: Dict[str, Any]
    segment_type: str
    index: int


@dataclass
class TranslationResult:
    original: TextSegment
    translated_text: str
    provider: str


@dataclass
class ProviderConfig:
    name: str
    api_key: str
    models: List[str]
    rate_limit_rpm: int = 60
    base_url: Optional[str] = None


class CTranslate2Translator:
    """Primary translator using WMT21 model via CTranslate2"""
    
    def __init__(self, target_language: str = "de", model_path: Optional[str] = None):
        self.target_language = target_language
        self.model_repo = "cstr/wmt21ct2_int8"
        self.translator = None
        self.tokenizer = None
        self.available = False
        
        # Check dependencies
        logger.info("Checking CTranslate2 dependencies...")
        if not check_ctranslate2():
            logger.warning("ctranslate2 not installed")
            return
        print("  ✓ ctranslate2")
        
        if not check_hf_hub():
            logger.warning("huggingface_hub not installed")
            return
        print("  ✓ huggingface_hub")
        
        if not check_transformers():
            logger.warning("transformers not installed")
            return
        
        try:
            if model_path:
                self._load_from_path(model_path)
            else:
                self._load_or_download()
            self.available = True
        except Exception as e:
            logger.error(f"CTranslate2 initialization failed: {e}")
    
    def _check_model_cached(self) -> Optional[str]:
        """Check if model is in cache by following the refs/main pointer"""
        cache_base = Path.home() / '.cache' / 'huggingface' / 'hub'
        model_dir = cache_base / f"models--{self.model_repo.replace('/', '--')}"
        
        ref_path = model_dir / 'refs' / 'main'
        if ref_path.exists():
            with open(ref_path, 'r') as f:
                commit_hash = f.read().strip()
            
            snapshot_path = model_dir / 'snapshots' / commit_hash
            if (snapshot_path / 'model.bin').exists():
                return str(snapshot_path)
        return None
    
    def _load_from_path(self, model_path: str):
        """Load model from provided path"""
        import ctranslate2
        import transformers
        
        path = Path(model_path)
        if not path.exists():
            raise FileNotFoundError(f"Model path not found: {path}")
        
        logger.info(f"Loading from: {path}")
        print(f"  ⏳ Loading CTranslate2...", end='', flush=True)
        self.translator = ctranslate2.Translator(str(path), device="auto")
        print(" ✓")
        
        self._load_tokenizer()
    
    def _load_or_download(self):
        """Load cached model or download using fast hf_transfer"""
        import ctranslate2
        from huggingface_hub import snapshot_download
        
        cached_path = self._check_model_cached()
        if cached_path:
            model_path = cached_path
        else:
            import os
            os.environ["HF_HUB_ENABLE_HF_TRANSFER"] = "1" # Fast mode
            print(f"  ⏳ Downloading model (Max Speed Enabled)...")
            model_path = snapshot_download(
                repo_id=self.model_repo,
                resume_download=True,
                max_workers=16
            )
        
        self.translator = ctranslate2.Translator(model_path, device="auto")
        # Note: We call _load_tokenizer later in OfficeTranslator
    
    def _load_tokenizer(self, source_lang: str = "en"):
        """Load the tokenizer with dynamic source/target languages"""
        import transformers
        
        tokenizer_name = "facebook/wmt21-dense-24-wide-en-x"
        print(f"  ⏳ Loading tokenizer ({source_lang} -> {self.target_language})...", end='', flush=True)
        
        self.tokenizer = transformers.AutoTokenizer.from_pretrained(tokenizer_name)
        # Crucial: WMT21 requires these to be set for the internal prefix logic
        self.tokenizer.src_lang = source_lang
        self.tokenizer.tgt_lang = self.target_language
        print(" ✓")

    def translate_batch(self, texts: List[str]) -> List[Optional[str]]:
        if not self.available or not texts:
            return [None] * len(texts)
        
        try:
            source_batches = []
            for text in texts:
                tokens = self.tokenizer.tokenize(text)
                if self.tokenizer.src_lang not in tokens:
                    tokens = [self.tokenizer.src_lang] + tokens + [self.tokenizer.eos_token]
                source_batches.append(tokens)

            target_prefix = [self.tokenizer.lang_code_to_token[self.target_language]]
            
            results = self.translator.translate_batch(
                source_batches, 
                target_prefix=[target_prefix] * len(texts),
                beam_size=5,
                repetition_penalty=1.5
            )

            translated_results = []
            for res in results:
                target_tokens = res.hypotheses[0][len(target_prefix):]
                decoded = self.tokenizer.decode(
                    self.tokenizer.convert_tokens_to_ids(target_tokens),
                    skip_special_tokens=True
                )
                # CLEANUP: Remove leading punctuation artifacts
                clean = decoded.strip().lstrip('.,- ')
                translated_results.append(clean)
            
            return translated_results
        except Exception as e:
            logger.error(f"Batch translation failed: {e}")
            return [None] * len(texts)
    
    def translate_text(self, text: str) -> Optional[str]:
        """
        Translates text using CTranslate2 with the WMT21 model.
        Fixes the 'prefix leak' and 'no translation' issues by correctly 
        handling the target_prefix and using clean decoding.
        """
        if not self.available or not text.strip():
            return text if text.strip() else None
        
        try:
            # --- PHASE 1: SEGMENTATION (Preserve Newlines & Length Control) ---
            segments = []
            newline_sequences = []
            segment = ""
            
            i = 0
            while i < len(text):
                if text[i] == '\n':
                    newline_sequence = '\n'
                    while i + 1 < len(text) and text[i + 1] == '\n':
                        newline_sequence += '\n'
                        i += 1
                    if segment:
                        segments.append(segment)
                        segment = ""
                    newline_sequences.append(newline_sequence)
                else:
                    segment += text[i]
                    if len(segment) >= 500:
                        # Look for sentence boundaries within the limit
                        end_index = max(
                            segment.rfind('.', 0, 500),
                            segment.rfind('?', 0, 500),
                            segment.rfind('!', 0, 500)
                        )
                        if end_index != -1:
                            segments.append(segment[:end_index+1])
                            segment = segment[end_index+1:].lstrip()
                i += 1
            
            if segment:
                segments.append(segment)
            
            # --- PHASE 2: TRANSLATION (The Fix) ---
            translated_segments = []
            for seg in segments:
                if not seg.strip():
                    translated_segments.append(seg)
                    continue
                
                # Encode and prepare target prefix
                source = self.tokenizer.convert_ids_to_tokens(self.tokenizer.encode(seg))
                target_token = self.tokenizer.lang_code_to_token[self.target_language]
                target_prefix = [target_token]
                
                # Translate
                results = self.translator.translate_batch([source], target_prefix=[target_prefix])
                
                # FIX: Slice away the prefix token and decode cleanly
                target_tokens = results[0].hypotheses[0][len(target_prefix):]
                translated = self.tokenizer.decode(
                    self.tokenizer.convert_tokens_to_ids(target_tokens),
                    skip_special_tokens=True
                )
                translated_segments.append(translated)
            
            # --- PHASE 3: REASSEMBLY ---
            result = ""
            for i, seg in enumerate(translated_segments):
                result += seg
                if i < len(newline_sequences):
                    result += newline_sequences[i]
            
            logger.debug(f"CT2: '{text[:50]}...' -> '{result[:50]}...'")
            return result.strip()
            
        except Exception as e:
            logger.error(f"CTranslate2 translation failed: {e}")
            return None


class LLMTranslator:
    """Fallback translator using free LLM APIs"""
    
    def __init__(self, target_language: str = "German", source_language: str = "English"):
        self.target_language = target_language
        self.source_language = source_language
        self.providers = self._initialize_providers()
        self.current_provider_index = 0
        self.rate_trackers = defaultdict(lambda: {"requests": 0, "last_reset": time.time()})
        
        if not self.providers:
            logger.info("No LLM providers configured")
    
    def _initialize_providers(self) -> Dict[str, ProviderConfig]:
        """Initialize available LLM providers"""
        providers = {}
        
        # OpenRouter (free tier)
        if os.getenv("OPENROUTER_API_KEY"):
            providers["openrouter"] = ProviderConfig(
                name="OpenRouter",
                api_key=os.getenv("OPENROUTER_API_KEY"),
                models=["deepseek/deepseek-chat:free"],
                rate_limit_rpm=20,
                base_url="https://openrouter.ai/api/v1"
            )
            logger.info("✓ OpenRouter configured")
        
        # Groq (free tier)
        if os.getenv("GROQ_API_KEY"):
            providers["groq"] = ProviderConfig(
                name="Groq",
                api_key=os.getenv("GROQ_API_KEY"),
                models=["llama-3.1-8b-instant"],
                rate_limit_rpm=30,
                base_url="https://api.groq.com/openai/v1"
            )
            logger.info("✓ Groq configured")
        
        # Cerebras (free tier)
        if os.getenv("CEREBRAS_API_KEY"):
            providers["cerebras"] = ProviderConfig(
                name="Cerebras",
                api_key=os.getenv("CEREBRAS_API_KEY"),
                models=["llama3.1-8b"],
                rate_limit_rpm=30,
                base_url="https://api.cloud.cerebras.ai/v1"
            )
            logger.info("✓ Cerebras configured")
        
        # Ollama (local)
        if self._check_ollama():
            providers["ollama"] = ProviderConfig(
                name="Ollama",
                api_key="none",
                models=["llama3.2"],
                rate_limit_rpm=1000,
                base_url="http://localhost:11434"
            )
            logger.info("✓ Ollama detected")
        
        return providers
    
    def _check_ollama(self) -> bool:
        """Check if Ollama is running"""
        try:
            response = requests.get("http://localhost:11434/api/tags", timeout=2)
            if response.status_code == 200:
                data = response.json()
                models = data.get("models", [])
                return len(models) > 0
            return False
        except:
            return False
    
    def _can_make_request(self, provider_name: str) -> bool:
        """Check rate limits"""
        tracker = self.rate_trackers[provider_name]
        current_time = time.time()
        
        if current_time - tracker["last_reset"] >= 60:
            tracker["requests"] = 0
            tracker["last_reset"] = current_time
        
        config = self.providers[provider_name]
        return tracker["requests"] < config.rate_limit_rpm
    
    def _call_api(self, config: ProviderConfig, text: str) -> Optional[str]:
        if config.name == "Ollama":
            return self._call_ollama(config, text)
        
        headers = {"Content-Type": "application/json", "Authorization": f"Bearer {config.api_key}"}
        messages = [
            {"role": "system", "content": f"Translate from {self.source_language} to {self.target_language}. Return ONLY the translation."},
            {"role": "user", "content": text}
        ]
        
        payload = {
            "model": config.models[0],
            "messages": messages,
            "temperature": 0.3,
            "max_tokens": 1000
        }
        
        try:
            response = requests.post(
                f"{config.base_url}/chat/completions",
                headers=headers,
                json=payload,
                timeout=30
            )
            
            if response.status_code == 200:
                result = response.json()
                content = result["choices"][0]["message"]["content"].strip()
                logger.debug(f"{config.name}: '{text[:30]}' -> '{content[:30]}'")
                return content
            else:
                logger.debug(f"{config.name} error: {response.status_code}")
                return None
        except Exception as e:
            logger.debug(f"{config.name} failed: {e}")
            return None
    
    def _call_ollama(self, config: ProviderConfig, text: str) -> Optional[str]:
        """Call Ollama API"""
        try:
            # Get first available model
            response = requests.get(f"{config.base_url}/api/tags", timeout=2)
            if response.status_code != 200:
                return None
            
            models = response.json().get("models", [])
            if not models:
                return None
            
            model_name = models[0]["name"]
            
            payload = {
                "model": model_name,
                "prompt": f"Translate from {self.source_language} to {self.target_language}. Only return the translation:\n\n{text}",
                "stream": False
            }
            
            response = requests.post(
                f"{config.base_url}/api/generate",
                json=payload,
                timeout=60
            )
            
            if response.status_code == 200:
                result = response.json()
                content = result.get("response", "").strip()
                logger.debug(f"Ollama: '{text[:30]}' -> '{content[:30]}'")
                return content
            
            return None
        except Exception as e:
            logger.debug(f"Ollama failed: {e}")
            return None
    
    def translate_text(self, text: str) -> Optional[str]:
        """Translate text using available providers"""
        if not text.strip() or not self.providers:
            return text if text.strip() else None
        
        # Try each provider
        provider_names = list(self.providers.keys())
        for _ in range(len(provider_names) * 2):
            provider_name = provider_names[self.current_provider_index]
            self.current_provider_index = (self.current_provider_index + 1) % len(provider_names)
            
            if not self._can_make_request(provider_name):
                continue
            
            config = self.providers[provider_name]
            result = self._call_api(config, text)
            
            if result:
                self.rate_trackers[provider_name]["requests"] += 1
                return result
            
            time.sleep(0.5)
        
        return None


class PowerPointExtractor:
    """Extract and rebuild PowerPoint"""
    
    def __init__(self):
        if not HAS_PPTX:
            raise ImportError("python-pptx required")
    
    def extract_segments(self, file_path: Path) -> List[TextSegment]:
        logger.info(f"Extracting from: {file_path.name}")
        prs = Presentation(str(file_path))
        segments = []
        index = 0
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    if not paragraph.text.strip():
                        continue
                    
                    segments.append(TextSegment(
                        text=paragraph.text,
                        metadata={
                            'slide_idx': slide_idx,
                            'shape_idx': shape_idx,
                            'para_idx': para_idx,
                        },
                        segment_type='paragraph',
                        index=index
                    ))
                    index += 1
            
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    segments.append(TextSegment(
                        text=notes,
                        metadata={'slide_idx': slide_idx, 'is_notes': True},
                        segment_type='notes',
                        index=index
                    ))
                    index += 1
        
        logger.info(f"Found {len(segments)} segments")
        return segments
    
    def rebuild_document(self, original_path: Path, translations: List[TranslationResult], output_path: Path):
        logger.info("Rebuilding PowerPoint...")
        prs = Presentation(str(original_path))
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for trans in translations:
                        meta = trans.original.metadata
                        if (trans.original.segment_type == 'paragraph' and
                            meta.get('slide_idx') == slide_idx and
                            meta.get('shape_idx') == shape_idx and
                            meta.get('para_idx') == para_idx):
                            
                            if paragraph.runs:
                                paragraph.runs[0].text = trans.translated_text
                                for run in paragraph.runs[1:]:
                                    run.text = ""
                            else:
                                paragraph.text = trans.translated_text
                            break
            
            if slide.has_notes_slide:
                for trans in translations:
                    meta = trans.original.metadata
                    if (trans.original.segment_type == 'notes' and
                        meta.get('slide_idx') == slide_idx):
                        slide.notes_slide.notes_text_frame.text = trans.translated_text
                        break
        
        prs.save(str(output_path))
        logger.info(f"Saved: {output_path}")


class WordExtractor:
    """Extract and rebuild Word"""
    
    def __init__(self):
        if not HAS_DOCX:
            raise ImportError("python-docx required")
    
    def extract_segments(self, file_path: Path) -> List[TextSegment]:
        logger.info(f"Extracting from: {file_path.name}")
        doc = Document(str(file_path))
        segments = []
        index = 0
        
        for para_idx, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                segments.append(TextSegment(
                    text=paragraph.text,
                    metadata={'para_idx': para_idx},
                    segment_type='paragraph',
                    index=index
                ))
                index += 1
        
        for table_idx, table in enumerate(doc.tables):
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    if cell.text.strip():
                        segments.append(TextSegment(
                            text=cell.text,
                            metadata={
                                'table_idx': table_idx,
                                'row_idx': row_idx,
                                'cell_idx': cell_idx,
                            },
                            segment_type='table_cell',
                            index=index
                        ))
                        index += 1
        
        logger.info(f"Found {len(segments)} segments")
        return segments
    
    def rebuild_document(self, original_path: Path, translations: List[TranslationResult], output_path: Path):
        logger.info("Rebuilding Word...")
        doc = Document(str(original_path))
        
        for trans in translations:
            if trans.original.segment_type == 'paragraph':
                para_idx = trans.original.metadata['para_idx']
                if para_idx < len(doc.paragraphs):
                    para = doc.paragraphs[para_idx]
                    if para.runs:
                        para.runs[0].text = trans.translated_text
                        for run in para.runs[1:]:
                            run.text = ""
                    else:
                        para.text = trans.translated_text
        
        for trans in translations:
            if trans.original.segment_type == 'table_cell':
                meta = trans.original.metadata
                if meta['table_idx'] < len(doc.tables):
                    table = doc.tables[meta['table_idx']]
                    if meta['row_idx'] < len(table.rows):
                        row = table.rows[meta['row_idx']]
                        if meta['cell_idx'] < len(row.cells):
                            cell = row.cells[meta['cell_idx']]
                            if cell.paragraphs:
                                cell.paragraphs[0].text = trans.translated_text
        
        doc.save(str(output_path))
        logger.info(f"Saved: {output_path}")


class ExcelExtractor:
    """Extract and rebuild Excel"""
    
    def __init__(self):
        if not HAS_OPENPYXL:
            raise ImportError("openpyxl required")
    
    def extract_segments(self, file_path: Path) -> List[TextSegment]:
        logger.info(f"Extracting from: {file_path.name}")
        wb = load_workbook(str(file_path))
        segments = []
        index = 0
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row_idx, row in enumerate(sheet.iter_rows()):
                for col_idx, cell in enumerate(row):
                    if cell.value and isinstance(cell.value, str) and cell.value.strip():
                        segments.append(TextSegment(
                            text=cell.value,
                            metadata={'sheet': sheet_name, 'row': row_idx, 'col': col_idx},
                            segment_type='cell',
                            index=index
                        ))
                        index += 1
        
        wb.close()
        logger.info(f"Found {len(segments)} segments")
        return segments
    
    def rebuild_document(self, original_path: Path, translations: List[TranslationResult], output_path: Path):
        logger.info("Rebuilding Excel...")
        wb = load_workbook(str(original_path))
        
        for trans in translations:
            if trans.original.segment_type == 'cell':
                meta = trans.original.metadata
                if meta['sheet'] in wb.sheetnames:
                    sheet = wb[meta['sheet']]
                    cell = sheet.cell(row=meta['row']+1, column=meta['col']+1)
                    cell.value = trans.translated_text
        
        wb.save(str(output_path))
        wb.close()
        logger.info(f"Saved: {output_path}")


class OfficeTranslator:
    """Main translator orchestrator"""
    
    def __init__(self, target_language: str = "de", source_language: str = "en", use_ctranslate: bool = True, use_llm: bool = True, model_path: Optional[str] = None):
        self.target_language = target_language
        self.source_language = source_language
        
        # Map language codes
        lang_map = {
            'de': 'German', 'es': 'Spanish', 'fr': 'French',
            'it': 'Italian', 'ja': 'Japanese', 'zh': 'Chinese',
            'ru': 'Russian', 'pt': 'Portuguese', 'nl': 'Dutch',
        }
        lang_name = lang_map.get(target_language, target_language.title())
        src_name = lang_map.get(source_language, source_language.title())
        
        # Initialize translators
        self.ct2_translator = None
        self.llm_translator = None
        
        if use_ctranslate:
            try:
                logger.info("Initializing CTranslate2...")
                self.ct2_translator = CTranslate2Translator(target_language, model_path)
                # Pass the source language into the tokenizer loader
                self.ct2_translator._load_tokenizer(source_lang=source_language)
                if self.ct2_translator.available:
                    logger.info("✓ CTranslate2 ready")
            except KeyboardInterrupt:
                logger.info("CTranslate2 download cancelled, will use LLM fallback")
            except Exception as e:
                logger.warning(f"CTranslate2 not available: {e}")
        
        if use_llm:
            try:
                logger.info("Initializing LLM translator...")
                # Pass BOTH names to the LLM
                self.llm_translator = LLMTranslator(target_language=lang_name, source_language=src_name)
                if self.llm_translator.providers:
                    logger.info(f"✓ LLM ready ({len(self.llm_translator.providers)} providers)")
            except Exception as e:
                logger.warning(f"LLM not available: {e}")
        
        self.extractors = {
            '.pptx': PowerPointExtractor() if HAS_PPTX else None,
            '.docx': WordExtractor() if HAS_DOCX else None,
            '.xlsx': ExcelExtractor() if HAS_OPENPYXL else None,
        }
    
    def translate_document(self, input_file: Path, output_file: Path):
        print(f"\n{'='*60}\nFile: {input_file.name}\nTarget: {self.target_language}\n{'='*60}\n")
        
        suffix = input_file.suffix.lower()
        extractor = self.extractors.get(suffix)
        if not extractor: raise ValueError(f"Unsupported: {suffix}")
        
        segments = extractor.extract_segments(input_file)
        if not segments: return

        unique_texts = list(set(seg.text for seg in segments if seg.text.strip()))
        translation_cache = {}
        
        # SMART ROUTING: Separate Sentences from Fragments/Names
        heavy_tasks = []  # Long sentences for CT2
        light_tasks = []  # Short fragments/Names for LLM
        
        for text in unique_texts:
            # Logic: If it's 3 words or less, it's likely a name, date, or fragment
            if len(text.split()) <= 3:
                light_tasks.append(text)
            else:
                heavy_tasks.append(text)

        print(f"Routing: {len(heavy_tasks)} long segments to CT2, {len(light_tasks)} short segments to LLM")

        # Step 1: Process Sentences with CTranslate2
        if heavy_tasks and self.ct2_translator:
            print(f"  ⏳ Running CTranslate2 Batch...")
            ct2_results = self.ct2_translator.translate_batch(heavy_tasks)
            for original, translated in zip(heavy_tasks, ct2_results):
                # Validation: If CT2 just echoed the original or returned garbage, don't cache it
                if translated and translated.strip().lower() != original.strip().lower() and len(translated) > 2:
                    translation_cache[original] = (translated, "ctranslate2")

        # Step 2: Everything else to LLM (Light tasks + CT2 failures)
        remaining = [t for t in unique_texts if t not in translation_cache]
        if remaining and self.llm_translator:
            print(f"  ⏳ Running LLM Fallback/Short-string handling for {len(remaining)} segments...")
            for i, text in enumerate(remaining):
                translated = self.llm_translator.translate_text(text)
                if translated:
                    translation_cache[text] = (translated, "llm")
                    if logger.level <= logging.DEBUG:
                        print(f"    DEBUG: LLM Translated '{text[:30]}' -> '{translated[:30]}'")

        # Step 3: Final Mapping
        final_results = []
        stats = defaultdict(int)
        for seg in segments:
            translated_text, provider = translation_cache.get(seg.text, (seg.text, "original"))
            stats[provider] += 1
            final_results.append(TranslationResult(seg, translated_text, provider))
            
            if logger.level <= logging.DEBUG:
                print(f"  [{provider}] {seg.text[:40]} -> {translated_text[:40]}")

        # Step 4: Rebuild
        print("\nRebuilding...")
        extractor.rebuild_document(input_file, final_results, output_file)
        
        print(f"\n{'='*60}\n✓ Complete! Stats: {dict(stats)}\nOutput: {output_file}\n{'='*60}\n")


def setup_logging(verbose: bool):
    level = logging.DEBUG if verbose else logging.INFO
    handler = logging.StreamHandler()
    handler.setFormatter(logging.Formatter('%(levelname)s: %(message)s'))
    logger.setLevel(level)
    logger.addHandler(handler)


def main():
    parser = argparse.ArgumentParser(
        description='Office Documents Translator with CTranslate2 and LLM Fallbacks',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
═══════════════════════════════════════════════════════════════
OPTION 1: AUTO-DOWNLOAD (Default - Easiest)
═══════════════════════════════════════════════════════════════
Let the script download the WMT21 model automatically:

python %(prog)s input.pptx output.pptx

First run will download ~4.7GB (takes 10-30 min), then cached.

═══════════════════════════════════════════════════════════════
OPTION 2: MANUAL DOWNLOAD (Faster with aria2c)
═══════════════════════════════════════════════════════════════
Download manually with aria2c (10-20x faster):

# Install aria2c
brew install aria2  # macOS
sudo apt install aria2  # Linux

# Create directory
mkdir -p ~/.cache/huggingface/hub/models--cstr--wmt21ct2_int8/snapshots/main
cd ~/.cache/huggingface/hub/models--cstr--wmt21ct2_int8/snapshots/main

# Download all files (CORRECT NAMES!)
aria2c -x 16 -s 16 https://huggingface.co/cstr/wmt21ct2_int8/resolve/main/model.bin
aria2c -x 16 -s 16 https://huggingface.co/cstr/wmt21ct2_int8/resolve/main/config.json
aria2c -x 16 -s 16 https://huggingface.co/cstr/wmt21ct2_int8/resolve/main/sentencepiece.bpe.model
aria2c -x 16 -s 16 https://huggingface.co/cstr/wmt21ct2_int8/resolve/main/shared_vocabulary.json
aria2c -x 16 -s 16 https://huggingface.co/cstr/wmt21ct2_int8/resolve/main/special_tokens_map.json
aria2c -x 16 -s 16 https://huggingface.co/cstr/wmt21ct2_int8/resolve/main/tokenizer_config.json
aria2c -x 16 -s 16 https://huggingface.co/cstr/wmt21ct2_int8/resolve/main/vocab.json

# Then use normally
python %(prog)s input.pptx output.pptx

═══════════════════════════════════════════════════════════════
OPTION 3: LLM-ONLY (No Download)
═══════════════════════════════════════════════════════════════
Skip CTranslate2 and use free LLM APIs:

# Get free key from https://openrouter.ai
export OPENROUTER_API_KEY="sk-or-v1-..."

python %(prog)s input.pptx output.pptx --no-ctranslate

═══════════════════════════════════════════════════════════════
EXAMPLES:
═══════════════════════════════════════════════════════════════
# Test setup
python %(prog)s --test

# Translate to German
python %(prog)s presentation.pptx output_de.pptx -l de

# Translate to Spanish (verbose)
python %(prog)s doc.docx doc_es.docx -l es -v

# Use LLM only
python %(prog)s file.xlsx file_fr.xlsx -l fr --no-ctranslate

═══════════════════════════════════════════════════════════════
SUPPORTED:
═══════════════════════════════════════════════════════════════
Languages: de, es, fr, it, ja, zh, ru, pt, nl
Files: .pptx, .docx, .xlsx
Providers: CTranslate2 (primary), OpenRouter, Groq, Cerebras, Ollama
        """
    )
    
    parser.add_argument('input_file', nargs='?', help='Input file')
    parser.add_argument('output_file', nargs='?', help='Output file')
    parser.add_argument('-s', '--source', default='en', help='Source language (default: en)') 
    parser.add_argument('-l', '--language', default='de', help='Target language (default: de)')
    parser.add_argument('-v', '--verbose', action='store_true', help='Verbose output')
    parser.add_argument('--model-path', type=str, help='Path to CTranslate2 model (optional)')
    parser.add_argument('--no-ctranslate', action='store_true', help='Skip CTranslate2, use LLM only')
    parser.add_argument('--no-llm', action='store_true', help='Skip LLM fallback')
    parser.add_argument('--test', action='store_true', help='Test setup')
    
    args = parser.parse_args()
    setup_logging(args.verbose)
    
    if args.test:
        print("\n" + "="*60)
        print("TESTING SETUP")
        print("="*60 + "\n")
        
        # Test CTranslate2
        if not args.no_ctranslate:
            print("1. Testing CTranslate2...")
            try:
                ct2 = CTranslate2Translator('de', args.model_path)
                if ct2.available:
                    result = ct2.translate_text("Hello world!")
                    print(f"   ✓ CTranslate2 working!")
                    print(f"     Test: 'Hello world!' -> '{result}'")
                else:
                    print("   ✗ CTranslate2 not available")
            except KeyboardInterrupt:
                print("   ⊘ Download cancelled")
            except Exception as e:
                print(f"   ✗ Error: {e}")
        else:
            print("1. Skipping CTranslate2 (--no-ctranslate)")
        
        # Test LLM
        print("\n2. Testing LLM Providers...")
        translator = LLMTranslator("German")
        
        if translator.providers:
            print(f"   ✓ Found {len(translator.providers)} providers:")
            for name in translator.providers.keys():
                print(f"     - {name}")
            
            print("\n   Testing translation...")
            result = translator.translate_text("Hello world!")
            
            if result:
                print(f"   ✓ Translation working!")
                print(f"     'Hello world!' -> '{result}'")
            else:
                print("   ✗ Translation failed")
        else:
            print("   ✗ No providers configured")
            print("\n   Setup:")
            print("     export OPENROUTER_API_KEY='your-key'")
            print("     Get key at: https://openrouter.ai")
        
        # Test document libraries
        print("\n3. Document Support:")
        print(f"   PowerPoint: {'✓' if HAS_PPTX else '✗ (pip install python-pptx)'}")
        print(f"   Word: {'✓' if HAS_DOCX else '✗ (pip install python-docx)'}")
        print(f"   Excel: {'✓' if HAS_OPENPYXL else '✗ (pip install openpyxl)'}")
        
        print("\n" + "="*60)
        return
    
    if not args.input_file or not args.output_file:
        parser.print_help()
        return
    
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"Error: File not found: {input_path}")
        sys.exit(1)
    
    output_path = Path(args.output_file)
    
    try:
        translator = OfficeTranslator(
            target_language=args.language,
            source_language=args.source,
            use_ctranslate=not args.no_ctranslate,
            use_llm=not args.no_llm,
            model_path=args.model_path
        )
        
        if not translator.ct2_translator and not translator.llm_translator:
            print("Error: No translators available!")
            print("\nOptions:")
            print("  1. Let CTranslate2 auto-download (run without --no-ctranslate)")
            print("  2. Use LLM: export OPENROUTER_API_KEY='your-key'")
            print("  3. Manual download: see --help for aria2c instructions")
            sys.exit(1)
        
        translator.translate_document(input_path, output_path)
        
    except KeyboardInterrupt:
        print("\n\nInterrupted!")
        sys.exit(1)
    except Exception as e:
        logger.error(f"Failed: {e}", exc_info=args.verbose)
        sys.exit(1)


if __name__ == "__main__":
    main()