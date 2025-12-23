#!/usr/bin/env python3
"""
Office Documents Translator with Robust Paragraph-Level Translation
+ CLI limits for testing (--max-chars, --max-words, --max-paras)
"""

import argparse
import logging
import sys
import os
import time
from pathlib import Path
from typing import List, Dict, Optional, Any, Tuple
from dataclasses import dataclass, field
from collections import defaultdict
from enum import Enum
import requests

logger = logging.getLogger(__name__)

# === Library Checks ===
print("🔍 Checking libraries...")

try:
    from pptx import Presentation
    from pptx.util import Pt
    HAS_PPTX = True
    print("  ✓ python-pptx")
except ImportError:
    HAS_PPTX = False
    print("  ✗ python-pptx")

try:
    from docx import Document
    from docx.shared import Pt as DocxPt, RGBColor
    HAS_DOCX = True
    print("  ✓ python-docx")
except ImportError:
    HAS_DOCX = False
    print("  ✗ python-docx")

HAS_CTRANSLATE2 = False
HAS_TRANSFORMERS = False
HAS_HF_HUB = False

def check_ctranslate2():
    global HAS_CTRANSLATE2
    try:
        import ctranslate2
        HAS_CTRANSLATE2 = True
        return True
    except ImportError:
        return False

def check_transformers():
    global HAS_TRANSFORMERS
    try:
        print("  ⏳ Loading transformers...", end='', flush=True)
        import transformers
        HAS_TRANSFORMERS = True
        print(" ✓")
        return True
    except ImportError:
        print(" ✗")
        return False

def check_hf_hub():
    global HAS_HF_HUB
    try:
        from huggingface_hub import snapshot_download
        HAS_HF_HUB = True
        return True
    except ImportError:
        return False

print("-" * 60)


# ============================================================================
# SIMPLIFIED DATA STRUCTURES - Paragraph-level translation
# ============================================================================

@dataclass
class FormatRun:
    """A run of text with consistent formatting"""
    text: str
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    font_name: Optional[str] = None
    font_size: Optional[float] = None
    font_color: Optional[str] = None


@dataclass
class TranslatableParagraph:
    """A paragraph with formatting runs"""
    runs: List[FormatRun]
    position: Dict[str, Any]
    metadata: Dict[str, Any] = field(default_factory=dict)  # NEW: separate metadata
    
    def get_text(self) -> str:
        """Get plain text from all runs"""
        return ''.join(run.text for run in self.runs)
    
    def get_dominant_formatting(self) -> Dict[str, Any]:
        """Get the most common formatting in the paragraph"""
        non_empty_runs = [r for r in self.runs if r.text and r.text.strip()]
        if not non_empty_runs:
            return {}
        
        from collections import Counter
        format_counts = Counter()
        
        for run in non_empty_runs:
            key = (run.bold, run.italic, run.underline, run.font_name)
            format_counts[key] += len(run.text)
        
        most_common = format_counts.most_common(1)[0][0]
        
        return {
            'bold': most_common[0],
            'italic': most_common[1],
            'underline': most_common[2],
            'font_name': most_common[3]
        }
    
    def has_uniform_formatting(self) -> bool:
        """Check if all runs have same formatting"""
        non_empty_runs = [r for r in self.runs if r.text and r.text.strip()]
        if not non_empty_runs:
            return True
        
        first = non_empty_runs[0]
        return all(
            r.bold == first.bold and 
            r.italic == first.italic and 
            r.underline == first.underline and
            r.font_name == first.font_name
            for r in non_empty_runs
        )
    
    def set_text(self, translated_text: str):
        """
        Set translated text - PARAGRAPH-LEVEL formatting only.
        Inline formatting (specific words italicized) cannot be reliably preserved.
        """
        if not translated_text:
            for run in self.runs:
                run.text = ""
            return
        
        # Get dominant formatting
        non_empty = [r for r in self.runs if r.text and r.text.strip()]
        
        if not non_empty:
            self.runs = [FormatRun(text=translated_text)]
            return
        
        # Check if all text has same formatting (paragraph-level)
        first = non_empty[0]
        is_uniform = all(
            r.bold == first.bold and 
            r.italic == first.italic and 
            r.underline == first.underline
            for r in non_empty
        )
        
        if is_uniform:
            # All text has same formatting - apply to translation
            self.runs = [FormatRun(
                text=translated_text,
                bold=first.bold,
                italic=first.italic,
                underline=first.underline,
                font_name=first.font_name,
                font_size=first.font_size,
                font_color=first.font_color
            )]
            logger.debug(f"Applied uniform formatting: bold={first.bold}, italic={first.italic}")
            return
        
        # Mixed formatting - cannot reliably preserve across languages
        # Calculate what percentage is formatted
        total_chars = sum(len(r.text) for r in non_empty)
        italic_chars = sum(len(r.text) for r in non_empty if r.italic)
        bold_chars = sum(len(r.text) for r in non_empty if r.bold)
        
        italic_pct = italic_chars / total_chars if total_chars > 0 else 0
        bold_pct = bold_chars / total_chars if total_chars > 0 else 0
        
        logger.warning(
            f"Mixed inline formatting detected ({italic_pct:.0%} italic, {bold_pct:.0%} bold). "
            f"Inline formatting cannot be reliably preserved across translation."
        )
        
        # If >80% is formatted, apply to whole paragraph
        if italic_pct > 0.8 or bold_pct > 0.8:
            logger.info(f"Applying majority formatting to entire paragraph")
            self.runs = [FormatRun(
                text=translated_text,
                bold=bold_pct > 0.8,
                italic=italic_pct > 0.8,
                font_name=first.font_name
            )]
        else:
            # Mixed formatting on minority of text - use plain text
            logger.info("Dropping inline formatting (cannot preserve across languages)")
            self.runs = [FormatRun(
                text=translated_text,
                font_name=first.font_name
            )]


@dataclass
class DocumentStructure:
    """Document as list of paragraphs"""
    paragraphs: List[TranslatableParagraph] = field(default_factory=list)
    
    def add(self, paragraph: TranslatableParagraph):
        self.paragraphs.append(paragraph)
    
    def __len__(self):
        return len(self.paragraphs)
    
    def __iter__(self):
        return iter(self.paragraphs)
    
    def get_stats(self):
        """Get document statistics"""
        total_chars = sum(len(p.get_text()) for p in self.paragraphs)
        total_words = sum(len(p.get_text().split()) for p in self.paragraphs)
        return {
            'paragraphs': len(self.paragraphs),
            'chars': total_chars,
            'words': total_words
        }


# ============================================================================
# EXTRACTION - Paragraph-level
# ============================================================================

class ParagraphExtractor:
    """Extract paragraphs with formatting runs"""
    
    @staticmethod
    def extract_from_runs(runs, position: Dict) -> TranslatableParagraph:
        """Extract runs with formatting - CAPTURE ALL FORMATTING"""
        format_runs = []
        
        for run in runs:
            # Include empty runs to preserve structure
            format_runs.append(FormatRun(
                text=run.text,
                bold=run.font.bold if hasattr(run.font, 'bold') else None,
                italic=run.font.italic if hasattr(run.font, 'italic') else None,
                underline=run.font.underline if hasattr(run.font, 'underline') else None,
                font_name=run.font.name if hasattr(run.font, 'name') and run.font.name else None,
                font_size=run.font.size.pt if hasattr(run.font, 'size') and run.font.size else None,
            ))
        
        return TranslatableParagraph(runs=format_runs, position=position)


# ============================================================================
# TRANSLATION BACKENDS
# ============================================================================

class CTranslate2Translator:
    """CTranslate2 with better error handling"""
    
    MODELS = {
        'en_to_x': 'cstr/wmt21ct2_int8',
        'x_to_en': 'cstr/wmt21-x-en-ct2-int8',
    }
    
    SUPPORTED_LANGUAGES = ['de', 'es', 'fr', 'it', 'ja', 'zh', 'ru', 'pt', 'nl']
    
    def __init__(self, target_language: str, source_language: str, model_path: Optional[str] = None):
        self.target_language = target_language
        self.source_language = source_language
        self.translator = None
        self.tokenizer = None
        self.available = False
        self.direction = None
        
        if source_language == 'en' and target_language in self.SUPPORTED_LANGUAGES:
            self.direction = 'en_to_x'
            self.model_repo = self.MODELS['en_to_x']
            self.tokenizer_name = "facebook/wmt21-dense-24-wide-en-x"
        elif target_language == 'en' and source_language in self.SUPPORTED_LANGUAGES:
            self.direction = 'x_to_en'
            self.model_repo = self.MODELS['x_to_en']
            self.tokenizer_name = "facebook/wmt21-dense-24-wide-x-en"
        else:
            logger.info(f"⊘ CTranslate2: {source_language.upper()}→{target_language.upper()} not supported")
            return
        
        if not all([check_ctranslate2(), check_hf_hub(), check_transformers()]):
            return
        
        try:
            if model_path:
                self._load_from_path(model_path)
            else:
                self._load_or_download()
            self.available = True
            logger.info(f"✓ CTranslate2 ready ({source_language.upper()}→{target_language.upper()})")
        except Exception as e:
            logger.error(f"CTranslate2 init failed: {e}")
    
    def _check_model_cached(self) -> Optional[str]:
        cache_base = Path.home() / '.cache' / 'huggingface' / 'hub'
        model_dir = cache_base / f"models--{self.model_repo.replace('/', '--')}"
        ref_path = model_dir / 'refs' / 'main'
        if ref_path.exists():
            with open(ref_path, 'r') as f:
                commit_hash = f.read().strip()
            snapshot_path = model_dir / 'snapshots' / commit_hash
            if (snapshot_path / 'model.bin').exists():
                return str(snapshot_path)
        return None
    
    def _load_from_path(self, model_path: str):
        import ctranslate2
        path = Path(model_path)
        if not path.exists():
            raise FileNotFoundError(f"Model not found: {path}")
        self.translator = ctranslate2.Translator(str(path), device="auto")
    
    def _load_or_download(self):
        import ctranslate2
        from huggingface_hub import snapshot_download
        
        cached_path = self._check_model_cached()
        if cached_path:
            model_path = cached_path
        else:
            import os
            os.environ["HF_HUB_ENABLE_HF_TRANSFER"] = "1"
            print(f"  ⏳ Downloading {self.direction} model...")
            model_path = snapshot_download(
                repo_id=self.model_repo,
                resume_download=True,
                max_workers=16
            )
        self.translator = ctranslate2.Translator(model_path, device="auto")
    
    def load_tokenizer(self):
        import transformers
        print(f"  ⏳ Loading tokenizer ({self.source_language}→{self.target_language})...", end='', flush=True)
        self.tokenizer = transformers.AutoTokenizer.from_pretrained(self.tokenizer_name)
        self.tokenizer.src_lang = self.source_language
        self.tokenizer.tgt_lang = self.target_language
        print(" ✓")
    
    def translate_batch(self, texts: List[str]) -> List[str]:
        """
        Translate batch of texts.
        CRITICAL: Returns same count as input!
        """
        if not self.available or not texts:
            return texts
        
        try:
            logger.debug(f"CT2 batch: {len(texts)} texts")
            
            source_batches = []
            for text in texts:
                if not text.strip():
                    source_batches.append([self.tokenizer.src_lang])
                    continue
                
                # Tokenize
                tokens = self.tokenizer.tokenize(text)
                if self.tokenizer.src_lang not in tokens:
                    tokens = [self.tokenizer.src_lang] + tokens + [self.tokenizer.eos_token]
                source_batches.append(tokens)

            target_prefix = [self.tokenizer.lang_code_to_token[self.target_language]]
            
            # Translate
            results = self.translator.translate_batch(
                source_batches, 
                target_prefix=[target_prefix] * len(texts),
                beam_size=5,
                repetition_penalty=1.5,
                max_batch_size=32  # Smaller batches for stability
            )

            # Decode
            translated = []
            for i, res in enumerate(results):
                if not texts[i].strip():
                    translated.append(texts[i])
                    continue
                
                target_tokens = res.hypotheses[0][len(target_prefix):]
                decoded = self.tokenizer.decode(
                    self.tokenizer.convert_tokens_to_ids(target_tokens),
                    skip_special_tokens=True
                )
                translated.append(decoded.strip())
            
            logger.debug(f"CT2 batch complete: {len(translated)} translations")
            return translated
            
        except Exception as e:
            logger.error(f"CT2 batch failed: {e}")
            return texts  # Return originals on error


class LindatTranslator:
    """Lindat API (Charles University)"""
    
    SUPPORTED_PAIRS = {('en', 'cs'), ('cs', 'en'), ('en', 'uk'), ('uk', 'en'), ('cs', 'uk'), ('uk', 'cs')}
    
    def __init__(self, source_lang: str, target_lang: str):
        self.source_lang = source_lang
        self.target_lang = target_lang
        self.available = False
        
        if (source_lang, target_lang) not in self.SUPPORTED_PAIRS:
            logger.info(f"⊘ Lindat: {source_lang.upper()}→{target_lang.upper()} not supported")
            return
        
        try:
            response = requests.get("https://lindat.mff.cuni.cz/services/translation/api/v2/languages", timeout=5)
            if response.status_code == 200:
                self.available = True
                logger.info(f"✓ Lindat available ({source_lang.upper()}→{target_lang.upper()})")
        except:
            logger.warning("Lindat API unreachable")
    
    def translate_text(self, text: str) -> Optional[str]:
        if not self.available or not text.strip():
            return None
        try:
            model = f"{self.source_lang}-{self.target_lang}"
            url = f"https://lindat.mff.cuni.cz/services/translation/api/v2/models/{model}"
            response = requests.post(url, files={'input_text': ('input.txt', text, 'text/plain')}, timeout=60)
            if response.status_code == 200:
                result = response.json()
                return ' '.join(result) if isinstance(result, list) else None
            return None
        except Exception as e:
            logger.debug(f"Lindat failed: {e}")
            return None


class LLMTranslator:
    """LLM fallback"""
    
    def __init__(self, target_language: str, source_language: str):
        self.target_language = target_language
        self.source_language = source_language
        self.providers = self._initialize_providers()
        self.current_provider_index = 0
        
        if self.providers:
            logger.info(f"✓ LLM available ({len(self.providers)} providers)")
        else:
            logger.info("⊘ No LLM providers")
    
    def _initialize_providers(self) -> Dict[str, Any]:
        providers = {}
        
        if os.getenv("OPENAI_API_KEY"):
            providers["openai"] = {
                "name": "OpenAI",
                "api_key": os.getenv("OPENAI_API_KEY"),
                "model": "gpt-4o-mini",
                "base_url": "https://api.openai.com/v1"
            }
        
        if self._check_ollama():
            providers["ollama"] = {
                "name": "Ollama",
                "api_key": "none",
                "model": self._get_ollama_model(),
                "base_url": "http://localhost:11434"
            }
        
        return providers
    
    def _check_ollama(self) -> bool:
        try:
            response = requests.get("http://localhost:11434/api/tags", timeout=2)
            return response.status_code == 200 and len(response.json().get("models", [])) > 0
        except:
            return False
    
    def _get_ollama_model(self) -> str:
        try:
            response = requests.get("http://localhost:11434/api/tags", timeout=2)
            models = response.json().get("models", [])
            return models[0]["name"] if models else "llama3.2"
        except:
            return "llama3.2"
    
    def translate_text(self, text: str) -> Optional[str]:
        if not text.strip() or not self.providers:
            return None
        
        system_prompt = f"Translate from {self.source_language} to {self.target_language}. Return ONLY the translation."
        
        for _ in range(len(self.providers)):
            provider_name = list(self.providers.keys())[self.current_provider_index]
            self.current_provider_index = (self.current_provider_index + 1) % len(self.providers)
            
            config = self.providers[provider_name]
            
            if provider_name == "ollama":
                result = self._call_ollama(config, text, system_prompt)
            else:
                result = self._call_api(config, text, system_prompt)
            
            if result:
                return result
            time.sleep(0.5)
        
        return None
    
    def _call_api(self, config: Dict, text: str, system_prompt: str) -> Optional[str]:
        try:
            response = requests.post(
                f"{config['base_url']}/chat/completions",
                headers={"Content-Type": "application/json", "Authorization": f"Bearer {config['api_key']}"},
                json={
                    "model": config["model"],
                    "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": text}
                    ],
                    "temperature": 0.3,
                    "max_tokens": 4000
                },
                timeout=60
            )
            if response.status_code == 200:
                return response.json()["choices"][0]["message"]["content"].strip()
            return None
        except Exception as e:
            logger.debug(f"{config['name']} failed: {e}")
            return None
    
    def _call_ollama(self, config: Dict, text: str, system_prompt: str) -> Optional[str]:
        try:
            response = requests.post(
                f"{config['base_url']}/api/generate",
                json={"model": config["model"], "prompt": f"{system_prompt}\n\n{text}", "stream": False},
                timeout=120
            )
            if response.status_code == 200:
                return response.json().get("response", "").strip()
            return None
        except Exception as e:
            logger.debug(f"Ollama failed: {e}")
            return None


# ============================================================================
# TRANSLATOR WITH CASCADING FALLBACK
# ============================================================================

class DocumentTranslator:
    """Paragraph-level translator with proper chunking"""
    
    def __init__(self, ct2=None, lindat=None, llm=None):
        self.ct2 = ct2
        self.lindat = lindat
        self.llm = llm
    
    def translate_document(self, doc_structure: DocumentStructure, max_paras: Optional[int] = None) -> DocumentStructure:
        """Translate paragraphs maintaining 1:1 mapping"""
        paragraphs_to_translate = doc_structure.paragraphs[:max_paras] if max_paras else doc_structure.paragraphs
        
        logger.info(f"Translating {len(paragraphs_to_translate)} paragraphs...")
        
        translated_count = 0
        skipped_count = 0
        
        # Translate paragraph by paragraph to maintain mapping
        for i, para in enumerate(paragraphs_to_translate):
            original_text = para.get_text()
            
            if not original_text.strip():
                logger.debug(f"Para {i+1}/{len(paragraphs_to_translate)}: EMPTY (preserved)")
                skipped_count += 1
                continue
            
            logger.debug(f"Para {i+1}/{len(paragraphs_to_translate)}: Translating '{original_text[:80]}'...")
            logger.debug(f"  Runs: {len(para.runs)}, formatting: {[(r.bold, r.italic, r.font_name) for r in para.runs if r.text]}")
            
            # Translate this paragraph
            translated_text = self._translate_single_paragraph(original_text)
            
            logger.debug(f"  → '{translated_text[:80]}'")
            
            # Apply translation
            para.set_text(translated_text)
            translated_count += 1
        
        logger.info(f"Translated {translated_count} paragraphs, skipped {skipped_count} empty")
        return doc_structure
    
    def _translate_single_paragraph(self, text: str) -> str:
        """
        Translate a single paragraph with chunking if needed.
        CRITICAL: Always returns exactly ONE translation per input.
        """
        if not text.strip():
            return text
        
        # Try CTranslate2 first
        if self.ct2 and self.ct2.available:
            result = self._translate_with_ct2(text)
            if result:
                # Optional sanity check - WARNING ONLY, don't reject
                if not self._is_translated(text, result):
                    logger.warning(f"⚠ Translation may be unchanged: '{text[:50]}' → '{result[:50]}'")
                else:
                    logger.debug(f"✓ CT2: '{text[:50]}' → '{result[:50]}'")
                return result  # USE IT ANYWAY
        
        # Try Lindat
        if self.lindat and self.lindat.available:
            result = self.lindat.translate_text(text)
            if result:
                logger.debug(f"✓ Lindat: '{text[:50]}' → '{result[:50]}'")
                return result
        
        # Try LLM
        if self.llm and self.llm.providers:
            result = self.llm.translate_text(text)
            if result:
                logger.debug(f"✓ LLM: '{text[:50]}' → '{result[:50]}'")
                return result
        
        logger.error(f"❌ All methods failed for: '{text[:50]}'")
        return text  # Return original if all fail
    
    def _translate_with_ct2(self, text: str) -> Optional[str]:
        """
        Translate with CT2, handling long texts by sentence splitting.
        CRITICAL: Always returns exactly ONE string.
        """
        # Check length
        if len(text) <= 500:
            # Short enough - translate directly
            results = self.ct2.translate_batch([text])
            return results[0] if results else None
        
        # Long text - split by sentences
        logger.debug(f"Long paragraph ({len(text)} chars), splitting...")
        
        # Split into sentences
        import re
        sentences = re.split(r'([.!?]+\s+)', text)
        
        # Reconstruct sentence + punctuation pairs
        chunks = []
        i = 0
        while i < len(sentences):
            if i + 1 < len(sentences) and sentences[i+1].strip():
                # Sentence + punctuation
                chunk = sentences[i] + sentences[i+1]
                i += 2
            else:
                # Just sentence
                chunk = sentences[i]
                i += 1
            
            if chunk.strip():
                chunks.append(chunk)
        
        # Further split if chunks still too long
        final_chunks = []
        for chunk in chunks:
            if len(chunk) <= 500:
                final_chunks.append(chunk)
            else:
                # Split by commas or spaces
                words = chunk.split()
                current = []
                for word in words:
                    if sum(len(w) for w in current) + len(word) > 500:
                        final_chunks.append(' '.join(current))
                        current = [word]
                    else:
                        current.append(word)
                if current:
                    final_chunks.append(' '.join(current))
        
        logger.debug(f"Split into {len(final_chunks)} chunks")
        
        # Translate chunks
        translated_chunks = self.ct2.translate_batch(final_chunks)
        
        # Reconstruct single paragraph
        result = ' '.join(translated_chunks)
        return result
    
    def _is_translated(self, original: str, translated: str) -> bool:
        """
        Check if translation actually happened.
        IMPROVED: Handle short texts and proper nouns better.
        """
        if not original or not translated:
            return False
        
        # If texts are very similar in length and content, might not be translated
        orig_words = original.lower().split()
        trans_words = translated.lower().split()
        
        if not orig_words or not trans_words:
            return True
        
        # For very short texts (< 5 words), be more lenient
        if len(orig_words) < 5:
            # Just check if they're not identical
            return original.strip().lower() != translated.strip().lower()
        
        # For longer texts, check word overlap
        orig_set = set(orig_words)
        trans_set = set(trans_words)
        
        overlap = len(orig_set & trans_set) / len(orig_set)
        
        # More lenient threshold - proper nouns should stay the same
        return overlap < 0.8  # Was 0.5, now 0.8

# ============================================================================
# DOCUMENT HANDLERS
# ============================================================================

class WordHandler:
    """Word document handler"""
    
    def __init__(self):
        if not HAS_DOCX:
            raise ImportError("python-docx required")
    
    def extract_structure(self, file_path: Path, max_chars: Optional[int] = None, 
                        max_words: Optional[int] = None, max_paras: Optional[int] = None) -> DocumentStructure:
        """Extract with limits - including footnotes and EMPTY paragraphs"""
        logger.info(f"Extracting from: {file_path.name}")
        doc = Document(str(file_path))
        structure = DocumentStructure()
        
        total_chars = 0
        total_words = 0
        
        for para_idx, paragraph in enumerate(doc.paragraphs):
            if max_paras and len(structure) >= max_paras:
                logger.info(f"Reached paragraph limit: {max_paras}")
                break
            
            para_text = paragraph.text
            para_chars = len(para_text)
            para_words = len(para_text.split())
            
            if para_text.strip():
                if max_chars and total_chars + para_chars > max_chars:
                    logger.info(f"Reached character limit: {max_chars}")
                    break
                
                if max_words and total_words + para_words > max_words:
                    logger.info(f"Reached word limit: {max_words}")
                    break
                
                total_chars += para_chars
                total_words += para_words
            
            position = {'type': 'paragraph', 'para_idx': para_idx}
            trans_para = ParagraphExtractor.extract_from_runs(paragraph.runs, position)
            
            # Store metadata separately (NOT in position!)
            trans_para.metadata['style'] = paragraph.style.name if paragraph.style else None
            trans_para.metadata['alignment'] = paragraph.alignment
            
            structure.add(trans_para)
            
            logger.debug(f"Para {para_idx}: '{para_text[:60]}...' ({len(trans_para.runs)} runs)")
        
        # Tables
        for table_idx, table in enumerate(doc.tables):
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    if max_paras and len(structure) >= max_paras:
                        break
                    if not cell.paragraphs:
                        continue
                    
                    position = {'type': 'table_cell', 'table_idx': table_idx, 
                            'row_idx': row_idx, 'cell_idx': cell_idx}
                    trans_para = ParagraphExtractor.extract_from_runs(cell.paragraphs[0].runs, position)
                    structure.add(trans_para)
        
        # Footnotes
        try:
            if hasattr(doc, 'part') and hasattr(doc.part, 'footnotes_part'):
                footnotes = doc.part.footnotes_part.footnotes
                for footnote_idx, footnote in enumerate(footnotes):
                    if max_paras and len(structure) >= max_paras:
                        break
                    
                    for para_idx, paragraph in enumerate(footnote.paragraphs):
                        position = {'type': 'footnote', 'footnote_idx': footnote_idx, 
                                'para_idx': para_idx}
                        trans_para = ParagraphExtractor.extract_from_runs(paragraph.runs, position)
                        structure.add(trans_para)
                        
                        logger.debug(f"Footnote {footnote_idx}.{para_idx}: {paragraph.text[:50]}")
        except Exception as e:
            logger.warning(f"Could not extract footnotes: {e}")
        
        stats = structure.get_stats()
        logger.info(f"Extracted {stats['paragraphs']} paragraphs ({stats['words']} words, {stats['chars']} chars)")
        return structure


    def rebuild_document(self, original_path: Path, structure: DocumentStructure, output_path: Path):
        """Rebuild with translations - PRESERVE ALL STRUCTURE"""
        logger.info("Rebuilding document...")
        doc = Document(str(original_path))
        
        # Group by position
        by_position = {}
        for para in structure:
            key = tuple(sorted(para.position.items()))
            by_position[key] = para
        
        # Translate paragraphs
        for para_idx, paragraph in enumerate(doc.paragraphs):
            key = tuple(sorted({'type': 'paragraph', 'para_idx': para_idx}.items()))
            
            if key in by_position:
                trans_para = by_position[key]
                self._rebuild_paragraph(paragraph, trans_para)
        
        logger.info(f"Processed {len(doc.paragraphs)} paragraphs")
        
        # Handle tables
        for table_idx, table in enumerate(doc.tables):
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    key = tuple(sorted({'type': 'table_cell', 'table_idx': table_idx,
                                    'row_idx': row_idx, 'cell_idx': cell_idx}.items()))
                    if key in by_position and cell.paragraphs:
                        trans_para = by_position[key]
                        self._rebuild_paragraph(cell.paragraphs[0], trans_para)
        
        # Handle footnotes - DON'T REMOVE, TRANSLATE IN PLACE
        try:
            if hasattr(doc, 'part') and hasattr(doc.part, 'footnotes_part'):
                footnotes_part = doc.part.footnotes_part
                if footnotes_part and hasattr(footnotes_part, 'footnotes'):
                    footnotes = footnotes_part.footnotes
                    
                    for footnote_idx, footnote in enumerate(footnotes):
                        for para_idx, paragraph in enumerate(footnote.paragraphs):
                            key = tuple(sorted({'type': 'footnote', 'footnote_idx': footnote_idx,
                                            'para_idx': para_idx}.items()))
                            
                            if key in by_position:
                                trans_para = by_position[key]
                                self._rebuild_paragraph(paragraph, trans_para)
                                logger.debug(f"✓ Translated footnote {footnote_idx}.{para_idx}")
                    
                    logger.info(f"Processed {len(footnotes)} footnotes")
        except Exception as e:
            logger.warning(f"Could not process footnotes: {e}", exc_info=True)
        
        doc.save(str(output_path))
        logger.info(f"Saved: {output_path}")
    
    def _rebuild_paragraph(self, paragraph, trans_para: TranslatableParagraph):
        """Rebuild paragraph from translated runs"""
        
        # Apply paragraph-level style from METADATA
        if trans_para.metadata.get('style'):
            try:
                paragraph.style = trans_para.metadata['style']
            except:
                pass
        
        if trans_para.metadata.get('alignment'):
            paragraph.alignment = trans_para.metadata['alignment']
        
        # Clear existing runs
        for _ in range(len(paragraph.runs)):
            paragraph._element.remove(paragraph.runs[0]._element)
        
        has_any_text = any(run.text for run in trans_para.runs)
        if not has_any_text:
            logger.debug("Empty paragraph preserved")
            return
        
        logger.debug(f"Adding {len(trans_para.runs)} runs")
        for i, fmt_run in enumerate(trans_para.runs):
            if not fmt_run.text:
                continue
                
            run = paragraph.add_run(fmt_run.text)
            
            logger.debug(f"  Run {i}: '{fmt_run.text[:40]}' bold={fmt_run.bold} italic={fmt_run.italic}")
            
            if fmt_run.bold is not None:
                run.bold = fmt_run.bold
            if fmt_run.italic is not None:
                run.italic = fmt_run.italic
            if fmt_run.underline is not None:
                run.underline = fmt_run.underline
            if fmt_run.font_name:
                run.font.name = fmt_run.font_name
            if fmt_run.font_size:
                run.font.size = DocxPt(fmt_run.font_size)


class PowerPointHandler:
    """PowerPoint handler"""
    
    def __init__(self):
        if not HAS_PPTX:
            raise ImportError("python-pptx required")
    
    def extract_structure(self, file_path: Path, max_chars=None, max_words=None, max_paras=None):
        logger.info(f"Extracting from: {file_path.name}")
        prs = Presentation(str(file_path))
        structure = DocumentStructure()
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    if max_paras and len(structure) >= max_paras:
                        break
                    if not paragraph.text.strip():
                        continue
                    position = {'type': 'shape', 'slide_idx': slide_idx, 
                               'shape_idx': shape_idx, 'para_idx': para_idx}
                    trans_para = ParagraphExtractor.extract_from_runs(paragraph.runs, position)
                    structure.add(trans_para)
        
        stats = structure.get_stats()
        logger.info(f"Extracted {stats['paragraphs']} paragraphs")
        return structure
    
    def rebuild_document(self, original_path: Path, structure: DocumentStructure, output_path: Path):
        """Rebuild presentation with ONLY translated slides"""
        logger.info("Rebuilding presentation...")
        
        # Create new presentation
        prs = Presentation()
        prs.slide_width = Presentation(str(original_path)).slide_width
        prs.slide_height = Presentation(str(original_path)).slide_height
        
        # Group paragraphs by slide
        slides_content = defaultdict(list)
        for para in structure.paragraphs:
            if para.position['type'] == 'shape':
                slide_idx = para.position['slide_idx']
                slides_content[slide_idx].append(para)
        
        # Create slides with translated content
        for slide_idx in sorted(slides_content.keys()):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
            
            # Add text box with all paragraphs for this slide
            left = Pt(50)
            top = Pt(50)
            width = prs.slide_width - Pt(100)
            height = prs.slide_height - Pt(100)
            
            text_frame = slide.shapes.add_textbox(left, top, width, height).text_frame
            
            for para in slides_content[slide_idx]:
                p = text_frame.add_paragraph()
                self._rebuild_paragraph(p, para)
        
        prs.save(str(output_path))
        logger.info(f"Saved: {output_path}")
    
    def _rebuild_paragraph(self, paragraph, trans_para):
        for _ in range(len(paragraph.runs)):
            paragraph._element.remove(paragraph.runs[0]._element)
        for fmt_run in trans_para.runs:
            run = paragraph.add_run()
            run.text = fmt_run.text
            if fmt_run.bold is not None:
                run.font.bold = fmt_run.bold
            if fmt_run.italic is not None:
                run.font.italic = fmt_run.italic
            if fmt_run.underline is not None:
                run.font.underline = fmt_run.underline
            if fmt_run.font_name:
                run.font.name = fmt_run.font_name
            if fmt_run.font_size:
                run.font.size = Pt(fmt_run.font_size)


# ============================================================================
# MAIN ORCHESTRATOR
# ============================================================================

class OfficeTranslator:
    """Main orchestrator"""
    
    def __init__(self, target_language: str, source_language: str,
                 use_ctranslate: bool = True, use_lindat: bool = True,
                 use_llm: bool = True, model_path: Optional[str] = None):
        
        self.target_language = target_language
        self.source_language = source_language
        
        lang_map = {
            'de': 'German', 'es': 'Spanish', 'fr': 'French',
            'it': 'Italian', 'ja': 'Japanese', 'zh': 'Chinese',
            'ru': 'Russian', 'pt': 'Portuguese', 'nl': 'Dutch',
            'cs': 'Czech', 'uk': 'Ukrainian',
        }
        target_name = lang_map.get(target_language, target_language.title())
        source_name = lang_map.get(source_language, source_language.title())
        
        # Initialize backends
        ct2 = None
        lindat = None
        llm = None
        
        if use_ctranslate:
            logger.info("Initializing CTranslate2...")
            ct2 = CTranslate2Translator(target_language, source_language, model_path)
            if ct2.available:
                ct2.load_tokenizer()
        
        if use_lindat:
            logger.info("Initializing Lindat API...")
            lindat = LindatTranslator(source_language, target_language)
        
        if use_llm:
            logger.info("Initializing LLM...")
            llm = LLMTranslator(target_name, source_name)
        
        # Check availability
        has_method = any([
            ct2 and ct2.available,
            lindat and lindat.available,
            llm and llm.providers
        ])
        
        if not has_method:
            logger.error("❌ No translation methods available!")
            sys.exit(1)
        
        self.translator = DocumentTranslator(ct2, lindat, llm)
        
        self.handlers = {
            '.docx': WordHandler() if HAS_DOCX else None,
            '.pptx': PowerPointHandler() if HAS_PPTX else None,
        }
    
    def translate_document(self, input_file: Path, output_file: Path,
                          max_chars: Optional[int] = None,
                          max_words: Optional[int] = None,
                          max_paras: Optional[int] = None):
        
        print(f"\n{'='*60}")
        print(f"File: {input_file.name}")
        print(f"Direction: {self.source_language.upper()} → {self.target_language.upper()}")
        if max_chars:
            print(f"Limit: {max_chars} chars")
        if max_words:
            print(f"Limit: {max_words} words")
        if max_paras:
            print(f"Limit: {max_paras} paragraphs")
        print(f"{'='*60}\n")
        
        suffix = input_file.suffix.lower()
        handler = self.handlers.get(suffix)
        if not handler:
            raise ValueError(f"Unsupported: {suffix}")
        
        # Extract
        structure = handler.extract_structure(input_file, max_chars, max_words, max_paras)
        
        # Translate
        self.translator.translate_document(structure, max_paras)
        
        # Rebuild
        handler.rebuild_document(input_file, structure, output_file)
        
        print(f"\n{'='*60}\n✓ Complete!\nOutput: {output_file}\n{'='*60}\n")


# ============================================================================
# CLI
# ============================================================================

def setup_logging(verbose: bool):
    level = logging.DEBUG if verbose else logging.INFO
    handler = logging.StreamHandler()
    handler.setFormatter(logging.Formatter('%(levelname)s: %(message)s'))
    logger.setLevel(level)
    logger.addHandler(handler)


def main():
    parser = argparse.ArgumentParser(
        description='Office Documents Translator with Limits'
    )
    
    parser.add_argument('input_file', nargs='?', help='Input file')
    parser.add_argument('output_file', nargs='?', help='Output file')
    parser.add_argument('-s', '--source', default='en', help='Source language')
    parser.add_argument('-l', '--language', default='de', help='Target language')
    parser.add_argument('-v', '--verbose', action='store_true', help='Verbose')
    parser.add_argument('--model-path', type=str, help='CTranslate2 model path')
    parser.add_argument('--no-ctranslate', action='store_true', help='Skip CTranslate2')
    parser.add_argument('--no-lindat', action='store_true', help='Skip Lindat')
    parser.add_argument('--no-llm', action='store_true', help='Skip LLM')
    
    # LIMITS for testing
    parser.add_argument('--max-chars', type=int, help='Translate only first N characters')
    parser.add_argument('--max-words', type=int, help='Translate only first N words')
    parser.add_argument('--max-paras', type=int, help='Translate only first N paragraphs')
    
    args = parser.parse_args()
    setup_logging(args.verbose)
    
    if not args.input_file or not args.output_file:
        parser.print_help()
        return
    
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"Error: File not found: {input_path}")
        sys.exit(1)
    
    try:
        translator = OfficeTranslator(
            target_language=args.language,
            source_language=args.source,
            use_ctranslate=not args.no_ctranslate,
            use_lindat=not args.no_lindat,
            use_llm=not args.no_llm,
            model_path=args.model_path
        )
        
        translator.translate_document(
            input_path,
            Path(args.output_file),
            max_chars=args.max_chars,
            max_words=args.max_words,
            max_paras=args.max_paras
        )
        
    except KeyboardInterrupt:
        print("\n\nInterrupted!")
        sys.exit(1)
    except Exception as e:
        logger.error(f"Failed: {e}", exc_info=args.verbose)
        sys.exit(1)


if __name__ == "__main__":
    main()