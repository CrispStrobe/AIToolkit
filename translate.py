#!/usr/bin/env python3
"""
Office Documents Translator with CTranslate2 and LLM Fallbacks

Supports: .pptx, .docx, .xlsx
Primary: CTranslate2 WMT21 model (or smaller NLLB)
Fallback: Multiple LLM providers with rate limiting
"""

import argparse
import logging
import sys
import os
import time
import re
import json
import zipfile
import shutil
from pathlib import Path
from typing import List, Dict, Tuple, Optional, Any
from dataclasses import dataclass
from xml.etree import ElementTree as ET
from collections import defaultdict
import requests

# Configure logging
logger = logging.getLogger(__name__)

# === Quick Imports Check (non-blocking) ===
print("🔍 Checking for required libraries...")

# Core office libraries
try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    HAS_PPTX = True
    print("  ✓ python-pptx available")
except ImportError:
    HAS_PPTX = False
    print("  ✗ python-pptx NOT available")

try:
    from docx import Document
    from docx.shared import Pt as DocxPt
    HAS_DOCX = True
    print("  ✓ python-docx available")
except ImportError:
    HAS_DOCX = False
    print("  ✗ python-docx NOT available")

try:
    from openpyxl import load_workbook, Workbook
    HAS_OPENPYXL = True
    print("  ✓ openpyxl available")
except ImportError:
    HAS_OPENPYXL = False
    print("  ✗ openpyxl NOT available")

# Translation libraries - deferred import due to slow loading
HAS_CTRANSLATE2 = False
HAS_HF_HUB = False
HAS_TRANSFORMERS = False

def check_ctranslate2():
    """Lazy check for ctranslate2"""
    global HAS_CTRANSLATE2
    try:
        import ctranslate2
        HAS_CTRANSLATE2 = True
        return True
    except ImportError:
        return False

def check_transformers():
    """Lazy check for transformers"""
    global HAS_TRANSFORMERS
    try:
        print("  ⏳ Loading transformers (this may take a moment)...", end='', flush=True)
        import transformers
        HAS_TRANSFORMERS = True
        print(" ✓")
        return True
    except ImportError:
        print(" ✗")
        return False

def check_hf_hub():
    """Lazy check for huggingface_hub"""
    global HAS_HF_HUB
    try:
        from huggingface_hub import snapshot_download
        HAS_HF_HUB = True
        return True
    except ImportError:
        return False

# LLM providers (optional fallbacks) - also deferred
HAS_GROQ = False

def check_groq():
    global HAS_GROQ
    try:
        from groq import Groq
        HAS_GROQ = True
        return True
    except ImportError:
        return False

print("-" * 60)

# === XML Namespaces for Office Documents ===
NS_PPTX = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

NS_DOCX = {
    'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
}


# === Data Classes ===
@dataclass
class TextSegment:
    """Represents a translatable text segment with metadata"""
    text: str
    metadata: Dict[str, Any]
    segment_type: str
    index: int


@dataclass
class TranslationResult:
    """Stores translation with original segment info"""
    original: TextSegment
    translated_text: str
    provider: str


@dataclass
class ProviderConfig:
    """Configuration for LLM provider"""
    name: str
    api_key: str
    models: List[str]
    max_tokens: int
    rate_limit_rpm: int = 60
    rate_limit_tpm: int = 100000
    base_url: Optional[str] = None
    enabled: bool = True


# === CTranslate2 Translator ===
class CTranslate2Translator:
    """Primary translator using CTranslate2 models"""
    
    def __init__(self, target_language: str = "de", model_path: Optional[str] = None, use_small_model: bool = False):
        self.target_language = target_language
        self.translator = None
        self.tokenizer = None
        self.available = False
        self.model_path_provided = model_path
        
        # Select model and tokenizer based on size preference
        if use_small_model:
            self.model_repo = "michaelfeil/ct2fast-nllb-200-distilled-600M"
            self.tokenizer_name = "facebook/nllb-200-distilled-600M"
            self.model_size = "~600MB"
            # Map language codes for NLLB
            self.lang_code = self._map_to_nllb_code(target_language)
        else:
            self.model_repo = "cstr/wmt21ct2_int8"
            self.tokenizer_name = "facebook/wmt21-dense-24-wide-en-x"
            self.model_size = "~4.7GB"
            self.lang_code = target_language
        
        # Check for required libraries
        logger.info("Checking CTranslate2 dependencies...")
        if not check_ctranslate2():
            logger.warning("ctranslate2 not available")
            return
        print("  ✓ ctranslate2 available")
        
        if not model_path and not check_hf_hub():
            logger.warning("huggingface_hub not available")
            return
        if not model_path:
            print("  ✓ huggingface_hub available")
        
        if not check_transformers():
            logger.warning("transformers not available")
            return
        
        try:
            self._initialize()
            self.available = True
        except Exception as e:
            logger.error(f"Failed to initialize CTranslate2: {e}")
    
    def _map_to_nllb_code(self, lang_code: str) -> str:
        """Map simple language codes to NLLB codes"""
        nllb_map = {
            'de': 'deu_Latn',
            'es': 'spa_Latn',
            'fr': 'fra_Latn',
            'it': 'ita_Latn',
            'ja': 'jpn_Jpan',
            'zh': 'zho_Hans',
            'ru': 'rus_Cyrl',
        }
        return nllb_map.get(lang_code, f'{lang_code}_Latn')
    
    def _check_model_cached(self) -> Optional[str]:
        """Check if model is already cached"""
        cache_dir = Path.home() / '.cache' / 'huggingface' / 'hub'
        model_dir_pattern = f"models--{self.model_repo.replace('/', '--')}"
        
        if cache_dir.exists():
            for item in cache_dir.iterdir():
                if item.name.startswith(model_dir_pattern):
                    snapshots_dir = item / 'snapshots'
                    if snapshots_dir.exists():
                        for snapshot in snapshots_dir.iterdir():
                            model_bin = snapshot / 'model.bin'
                            if model_bin.exists():
                                logger.info(f"Found cached model at: {snapshot}")
                                return str(snapshot)
        return None
    
    def _initialize(self):
        """Download and initialize the CTranslate2 model"""
        import ctranslate2
        import transformers
        
        # Use provided path or download
        if self.model_path_provided:
            model_path = self.model_path_provided
            logger.info(f"Using provided model path: {model_path}")
            print(f"  ✓ Using model from: {model_path}")
        else:
            from huggingface_hub import snapshot_download
            
            # Check if model is already cached
            cached_path = self._check_model_cached()
            
            if cached_path:
                logger.info(f"Using cached model from: {cached_path}")
                print(f"  ✓ Using cached model (no download needed)")
                model_path = cached_path
            else:
                logger.warning(f"Model not cached - downloading {self.model_repo}...")
                print(f"\n{'='*60}")
                print(f"⚠️  FIRST-TIME MODEL DOWNLOAD")
                print(f"{'='*60}")
                print(f"Model: {self.model_repo}")
                print(f"Size: {self.model_size}")
                print(f"This will be cached for future use.")
                print(f"\n💡 TIP: For faster download, use aria2c:")
                print(f"   See --help for manual download instructions")
                print(f"\n💡 TIP: Or use --no-ctranslate for LLM-only mode (no download)")
                print(f"{'='*60}\n")
                
                response = input("Continue with download? [y/N]: ")
                if response.lower() != 'y':
                    logger.info("Download cancelled by user")
                    raise KeyboardInterrupt("Download cancelled by user")
                
                logger.info(f"Downloading model from {self.model_repo}...")
                print(f"  ⏳ Downloading {self.model_size} (this will take a while)...")
                print(f"     Press Ctrl+C to cancel - download will resume next time")
                
                try:
                    model_path = snapshot_download(
                        repo_id=self.model_repo,
                        resume_download=True
                    )
                    print(f"  ✓ Download complete")
                    logger.info(f"Model downloaded to: {model_path}")
                except KeyboardInterrupt:
                    print("\n\n⚠️  Download interrupted!")
                    print("   Next time you run, the download will resume from where it left off.")
                    raise
        
        print(f"  ⏳ Loading CTranslate2 model...", end='', flush=True)
        self.translator = ctranslate2.Translator(model_path, device="auto")
        print(" ✓")
        logger.info("CTranslate2 model loaded")
        
        print(f"  ⏳ Loading tokenizer...", end='', flush=True)
        self.tokenizer = transformers.AutoTokenizer.from_pretrained(self.tokenizer_name)
        
        # Set language codes based on model type
        if "nllb" in self.tokenizer_name.lower():
            self.tokenizer.src_lang = "eng_Latn"
            self.tokenizer.tgt_lang = self.lang_code
        else:
            self.tokenizer.src_lang = "en"
            self.tokenizer.tgt_lang = self.target_language
        
        print(" ✓")
        logger.info(f"Tokenizer loaded for en -> {self.target_language}")
    
    def translate_text(self, text: str) -> Optional[str]:
        if not self.available or not text.strip():
            return text if text.strip() else None
        
        try:
            # --- PHASE 1: SEGMENTATION ---
            segments = []
            newline_sequences = []
            segment = ""
            i = 0
            while i < len(text):
                if text[i] == '\n':
                    newline_sequence = '\n'
                    while i + 1 < len(text) and text[i + 1] == '\n':
                        newline_sequence += '\n'
                        i += 1
                    if segment:
                        segments.append(segment)
                        segment = ""
                    newline_sequences.append(newline_sequence)
                else:
                    segment += text[i]
                    if len(segment) >= 500:
                        end_index = max(segment.rfind('.', 0, 500), segment.rfind('?', 0, 500), segment.rfind('!', 0, 500))
                        if end_index != -1:
                            segments.append(segment[:end_index+1])
                            segment = segment[end_index+1:].lstrip()
                i += 1
            if segment:
                segments.append(segment)

            # --- PHASE 2: TRANSLATION ---
            translated_segments = []
            for seg in segments:
                if not seg.strip():
                    translated_segments.append(seg)
                    continue
                
                # FORCE language tokens for WMT21
                # Source must look like: ['__de__', 'token', ..., '</s>']
                source_tokens = self.tokenizer.tokenize(seg)
                if self.tokenizer.src_lang not in source_tokens:
                    source_tokens = [self.tokenizer.src_lang] + source_tokens + [self.tokenizer.eos_token]
                
                target_prefix = [self.tokenizer.lang_code_to_token[self.target_language]]
                
                results = self.translator.translate_batch([source_tokens], target_prefix=[target_prefix])
                
                # Extract and decode
                target_tokens = results[0].hypotheses[0][len(target_prefix):]
                translated = self.tokenizer.decode(
                    self.tokenizer.convert_tokens_to_ids(target_tokens),
                    skip_special_tokens=True
                )
                translated_segments.append(translated)
            
            # --- PHASE 3: REASSEMBLY ---
            result = ""
            for i, seg in enumerate(translated_segments):
                result += seg
                if i < len(newline_sequences):
                    result += newline_sequences[i]
            
            # This ensures logs show up when -v is used
            logger.debug(f"CT2: '{text[:40]}...' -> '{result[:40]}...'")
            return result.strip()
            
        except Exception as e:
            logger.error(f"CTranslate2 translation failed: {e}")
            return None


# === LLM Fallback Translator ===
class LLMTranslator:
    """Fallback translator using multiple LLM providers"""
    
    def __init__(self, target_language: str = "German", source_language: str = "English"):
        self.target_language = target_language
        self.source_language = source_language
        self.providers = self._initialize_providers()
        self.current_provider_index = 0
        self.rate_trackers = defaultdict(lambda: {"requests": 0, "last_reset": time.time()})
    
    def _initialize_providers(self) -> Dict[str, ProviderConfig]:
        """Initialize available LLM providers"""
        providers = {}
        
        # OpenRouter (free tier)
        if os.getenv("OPENROUTER_API_KEY"):
            providers["openrouter"] = ProviderConfig(
                name="OpenRouter",
                api_key=os.getenv("OPENROUTER_API_KEY"),
                models=["deepseek/deepseek-chat:free"],
                max_tokens=4000,
                rate_limit_rpm=20,
                base_url="https://openrouter.ai/api/v1"
            )
        
        # Groq
        if os.getenv("GROQ_API_KEY"):
            check_groq()
            if HAS_GROQ:
                providers["groq"] = ProviderConfig(
                    name="Groq",
                    api_key=os.getenv("GROQ_API_KEY"),
                    models=["llama-3.1-8b-instant"],
                    max_tokens=4000,
                    rate_limit_rpm=30
                )
        
        # Cerebras
        if os.getenv("CEREBRAS_API_KEY"):
            providers["cerebras"] = ProviderConfig(
                name="Cerebras",
                api_key=os.getenv("CEREBRAS_API_KEY"),
                models=["llama3.1-8b"],
                max_tokens=4000,
                rate_limit_rpm=30,
                base_url="https://api.cloud.cerebras.ai/v1"
            )
        
        # Ollama (local)
        if self._check_ollama():
            providers["ollama"] = ProviderConfig(
                name="Ollama",
                api_key="none",
                models=["llama3.2"],
                max_tokens=4000,
                rate_limit_rpm=1000,
                base_url="http://localhost:11434"
            )
        
        if providers:
            logger.info(f"Initialized {len(providers)} LLM providers: {list(providers.keys())}")
        
        return providers
    
    def _check_ollama(self) -> bool:
        """Check if Ollama is running"""
        try:
            response = requests.get("http://localhost:11434/api/tags", timeout=2)
            return response.status_code == 200
        except:
            return False
    
    def _can_make_request(self, provider_name: str) -> bool:
        """Check if we can make a request to this provider"""
        tracker = self.rate_trackers[provider_name]
        current_time = time.time()
        
        if current_time - tracker["last_reset"] >= 60:
            tracker["requests"] = 0
            tracker["last_reset"] = current_time
        
        config = self.providers[provider_name]
        return tracker["requests"] < config.rate_limit_rpm
    
    def _record_request(self, provider_name: str):
        """Record a request to provider"""
        self.rate_trackers[provider_name]["requests"] += 1
    
    def _call_openai_compatible(self, config: ProviderConfig, text: str, model: str) -> Optional[str]:
        """Call OpenAI-compatible API"""
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {config.api_key}"
        }
        
        if config.name == "OpenRouter":
            headers["HTTP-Referer"] = "https://github.com/office-translator"
        
        messages = [
            {
                "role": "system",
                "content": f"You are a professional translator. Translate the following text from English to {self.target_language}. Preserve formatting, line breaks, and meaning. Return ONLY the translation, no explanations."
            },
            {
                "role": "user",
                "content": text
            }
        ]
        
        payload = {
            "model": model,
            "messages": messages,
            "temperature": 0.3,
            "max_tokens": 1000
        }
        
        try:
            response = requests.post(
                f"{config.base_url}/chat/completions",
                headers=headers,
                json=payload,
                timeout=60
            )
            
            if response.status_code == 200:
                result = response.json()
                return result["choices"][0]["message"]["content"].strip()
            else:
                logger.error(f"{config.name} API error: {response.status_code}")
                return None
        except Exception as e:
            logger.error(f"{config.name} request failed: {e}")
            return None
    
    def _call_groq(self, config: ProviderConfig, text: str, model: str) -> Optional[str]:
        """Call Groq API"""
        try:
            from groq import Groq
            client = Groq(api_key=config.api_key)
            response = client.chat.completions.create(
                model=model,
                messages=[
                    {
                        "role": "system",
                        "content": f"Translate to {self.target_language}. Return only translation."
                    },
                    {"role": "user", "content": text}
                ],
                temperature=0.3,
                max_tokens=1000
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            logger.error(f"Groq failed: {e}")
            return None
    
    def _call_ollama(self, config: ProviderConfig, text: str, model: str) -> Optional[str]:
        """Call Ollama API"""
        try:
            payload = {
                "model": model,
                "messages": [
                    {
                        "role": "system",
                        "content": f"Translate to {self.target_language}. Return only translation."
                    },
                    {"role": "user", "content": text}
                ],
                "stream": False
            }
            
            response = requests.post(
                f"{config.base_url}/api/chat",
                json=payload,
                timeout=120
            )
            
            if response.status_code == 200:
                result = response.json()
                return result.get("message", {}).get("content", "").strip()
            return None
        except Exception as e:
            logger.error(f"Ollama failed: {e}")
            return None
    
    def translate_text(self, text: str) -> Optional[str]:
        """Translate text using available LLM provider"""
        if not text.strip():
            return text
        
        for _ in range(len(self.providers) * 2):
            provider_name = self.provider_names[self.current_provider_index]
            self.current_provider_index = (self.current_provider_index + 1) % len(self.provider_names)
            
            if not self._can_make_request(provider_name):
                logger.debug(f"{provider_name} rate limited, trying next")
                continue
            
            config = self.providers[provider_name]
            model = config.models[0]
            
            logger.debug(f"Trying {provider_name} with {model}")
            
            result = None
            if provider_name == "groq":
                result = self._call_groq(config, text, model)
            elif provider_name == "ollama":
                result = self._call_ollama(config, text, model)
            else:
                result = self._call_openai_compatible(config, text, model)
            
            if result:
                self._record_request(provider_name)
                logger.debug(f"LLM ({provider_name}) translated: '{text[:50]}...' -> '{result[:50]}...'")
                return result
            
            time.sleep(1)
        
        logger.warning(f"All LLM providers failed for text: {text[:50]}...")
        return None


# === Office Document Extractors ===
class PowerPointExtractor:
    """Extract and rebuild PowerPoint presentations"""
    
    def __init__(self):
        if not HAS_PPTX:
            raise ImportError("python-pptx is required for PowerPoint support")
    
    def extract_segments(self, file_path: Path) -> List[TextSegment]:
        """Extract all translatable text from PPTX"""
        logger.info(f"Extracting text from PowerPoint: {file_path}")
        prs = Presentation(str(file_path))
        segments = []
        index = 0
        
        for slide_idx, slide in enumerate(prs.slides):
            logger.debug(f"Processing slide {slide_idx + 1}/{len(prs.slides)}")
            
            # Extract text from shapes
            for shape_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    if not paragraph.text.strip():
                        continue
                    
                    segments.append(TextSegment(
                        text=paragraph.text,
                        metadata={
                            'slide_idx': slide_idx,
                            'shape_idx': shape_idx,
                            'para_idx': para_idx,
                            'font_size': paragraph.runs[0].font.size if paragraph.runs else None,
                            'bold': paragraph.runs[0].font.bold if paragraph.runs else None,
                            'italic': paragraph.runs[0].font.italic if paragraph.runs else None,
                        },
                        segment_type='paragraph',
                        index=index
                    ))
                    index += 1
            
            # Extract notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    segments.append(TextSegment(
                        text=notes_text,
                        metadata={'slide_idx': slide_idx, 'is_notes': True},
                        segment_type='notes',
                        index=index
                    ))
                    index += 1
        
        logger.info(f"Extracted {len(segments)} text segments from PowerPoint")
        return segments
    
    def rebuild_document(self, original_path: Path, translations: List[TranslationResult], output_path: Path):
        """Rebuild PowerPoint with translations"""
        logger.info(f"Rebuilding PowerPoint with translations")
        
        # Create translation lookup
        trans_map = {t.original.index: t.translated_text for t in translations}
        
        # Load original presentation
        prs = Presentation(str(original_path))
        
        for slide_idx, slide in enumerate(prs.slides):
            # Update text in shapes
            for shape_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    # Find matching translation
                    for trans in translations:
                        meta = trans.original.metadata
                        if (trans.original.segment_type == 'paragraph' and
                            meta.get('slide_idx') == slide_idx and
                            meta.get('shape_idx') == shape_idx and
                            meta.get('para_idx') == para_idx):
                            
                            # Update text while preserving formatting
                            if paragraph.runs:
                                first_run = paragraph.runs[0]
                                # Clear existing runs
                                for run in paragraph.runs[1:]:
                                    run.text = ""
                                first_run.text = trans.translated_text
                            else:
                                paragraph.text = trans.translated_text
                            break
            
            # Update notes
            if slide.has_notes_slide:
                for trans in translations:
                    meta = trans.original.metadata
                    if (trans.original.segment_type == 'notes' and
                        meta.get('slide_idx') == slide_idx and
                        meta.get('is_notes')):
                        slide.notes_slide.notes_text_frame.text = trans.translated_text
                        break
        
        prs.save(str(output_path))
        logger.info(f"Saved translated PowerPoint to: {output_path}")


class WordExtractor:
    """Extract and rebuild Word documents"""
    
    def __init__(self):
        if not HAS_DOCX:
            raise ImportError("python-docx is required for Word support")
    
    def extract_segments(self, file_path: Path) -> List[TextSegment]:
        """Extract all translatable text from DOCX"""
        logger.info(f"Extracting text from Word document: {file_path}")
        doc = Document(str(file_path))
        segments = []
        index = 0
        
        # Extract paragraphs
        for para_idx, paragraph in enumerate(doc.paragraphs):
            if not paragraph.text.strip():
                continue
            
            segments.append(TextSegment(
                text=paragraph.text,
                metadata={
                    'para_idx': para_idx,
                    'style': paragraph.style.name if paragraph.style else None,
                },
                segment_type='paragraph',
                index=index
            ))
            index += 1
        
        # Extract tables
        for table_idx, table in enumerate(doc.tables):
            for row_idx, row in enumerate(table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    if not cell.text.strip():
                        continue
                    
                    segments.append(TextSegment(
                        text=cell.text,
                        metadata={
                            'table_idx': table_idx,
                            'row_idx': row_idx,
                            'cell_idx': cell_idx,
                        },
                        segment_type='table_cell',
                        index=index
                    ))
                    index += 1
        
        logger.info(f"Extracted {len(segments)} text segments from Word document")
        return segments
    
    def rebuild_document(self, original_path: Path, translations: List[TranslationResult], output_path: Path):
        """Rebuild Word document with translations"""
        logger.info(f"Rebuilding Word document with translations")
        
        doc = Document(str(original_path))
        
        # Update paragraphs
        for trans in translations:
            if trans.original.segment_type == 'paragraph':
                para_idx = trans.original.metadata['para_idx']
                if para_idx < len(doc.paragraphs):
                    para = doc.paragraphs[para_idx]
                    # Preserve formatting by updating runs
                    if para.runs:
                        para.runs[0].text = trans.translated_text
                        for run in para.runs[1:]:
                            run.text = ""
                    else:
                        para.text = trans.translated_text
        
        # Update tables
        for trans in translations:
            if trans.original.segment_type == 'table_cell':
                meta = trans.original.metadata
                table_idx = meta['table_idx']
                row_idx = meta['row_idx']
                cell_idx = meta['cell_idx']
                
                if table_idx < len(doc.tables):
                    table = doc.tables[table_idx]
                    if row_idx < len(table.rows):
                        row = table.rows[row_idx]
                        if cell_idx < len(row.cells):
                            cell = row.cells[cell_idx]
                            # Update cell text
                            if cell.paragraphs:
                                cell.paragraphs[0].text = trans.translated_text
        
        doc.save(str(output_path))
        logger.info(f"Saved translated Word document to: {output_path}")


class ExcelExtractor:
    """Extract and rebuild Excel spreadsheets"""
    
    def __init__(self):
        if not HAS_OPENPYXL:
            raise ImportError("openpyxl is required for Excel support")
    
    def extract_segments(self, file_path: Path) -> List[TextSegment]:
        """Extract all translatable text from XLSX"""
        logger.info(f"Extracting text from Excel: {file_path}")
        wb = load_workbook(str(file_path))
        segments = []
        index = 0
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            logger.debug(f"Processing sheet: {sheet_name}")
            
            for row_idx, row in enumerate(sheet.iter_rows()):
                for col_idx, cell in enumerate(row):
                    if cell.value and isinstance(cell.value, str) and cell.value.strip():
                        segments.append(TextSegment(
                            text=cell.value,
                            metadata={
                                'sheet': sheet_name,
                                'row': row_idx,
                                'col': col_idx,
                                'coordinate': cell.coordinate,
                            },
                            segment_type='cell',
                            index=index
                        ))
                        index += 1
        
        wb.close()
        logger.info(f"Extracted {len(segments)} text segments from Excel")
        return segments
    
    def rebuild_document(self, original_path: Path, translations: List[TranslationResult], output_path: Path):
        """Rebuild Excel with translations"""
        logger.info(f"Rebuilding Excel with translations")
        
        wb = load_workbook(str(original_path))
        
        for trans in translations:
            if trans.original.segment_type == 'cell':
                meta = trans.original.metadata
                sheet_name = meta['sheet']
                row = meta['row']
                col = meta['col']
                
                if sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    cell = sheet.cell(row=row+1, column=col+1)  # openpyxl is 1-indexed
                    cell.value = trans.translated_text
        
        wb.save(str(output_path))
        wb.close()
        logger.info(f"Saved translated Excel to: {output_path}")


# === Main Translator Orchestrator ===

# === Main Translator Orchestrator ===
class OfficeTranslator:
    """Main orchestrator for translating office documents"""
    
    def __init__(self, target_language: str = "de", use_ctranslate: bool = True, use_llm: bool = True, model_path: Optional[str] = None, use_small_model: bool = False):
        self.target_language = target_language
        
        # Initialize translators
        self.ct2_translator = None
        self.llm_translator = None
        
        if use_ctranslate:
            try:
                logger.info("Initializing CTranslate2 translator...")
                self.ct2_translator = CTranslate2Translator(target_language, model_path, use_small_model)
                if self.ct2_translator.available:
                    logger.info("✓ CTranslate2 translator initialized")
            except Exception as e:
                logger.warning(f"CTranslate2 initialization failed: {e}")
        
        if use_llm:
            try:
                lang_map = {
                    'de': 'German',
                    'es': 'Spanish',
                    'fr': 'French',
                    'it': 'Italian',
                    'ja': 'Japanese',
                    'zh': 'Chinese',
                    'ru': 'Russian',
                }
                lang_name = lang_map.get(target_language, target_language)
                logger.info("Initializing LLM translator...")
                self.llm_translator = LLMTranslator(lang_name)
                if self.llm_translator.providers:
                    logger.info(f"✓ LLM translator initialized with {len(self.llm_translator.providers)} providers")
            except Exception as e:
                logger.warning(f"LLM initialization failed: {e}")
        
        # Initialize extractors
        self.extractors = {
            '.pptx': PowerPointExtractor() if HAS_PPTX else None,
            '.docx': WordExtractor() if HAS_DOCX else None,
            '.xlsx': ExcelExtractor() if HAS_OPENPYXL else None,
        }

    
    def translate_segment(self, segment: TextSegment) -> TranslationResult:
        """Translate a single segment using available translators"""
        logger.debug(f"Translating segment {segment.index}: '{segment.text[:50]}...'")
        
        translated = None
        provider = "none"
        
        # Try CTranslate2 first
        if self.ct2_translator and self.ct2_translator.available:
            translated = self.ct2_translator.translate_text(segment.text)
            if translated:
                provider = "ctranslate2"
                logger.debug(f"  ✓ CTranslate2 success")
        
        # Fallback to LLM
        if not translated and self.llm_translator:
            translated = self.llm_translator.translate_text(segment.text)
            if translated:
                provider = f"llm_{self.llm_translator.provider_names[self.llm_translator.current_provider_index]}"
                logger.debug(f"  ✓ LLM fallback success ({provider})")
        
        # Ultimate fallback: keep original
        if not translated:
            translated = segment.text
            provider = "original"
            logger.warning(f"  ✗ All translators failed, keeping original")
        
        return TranslationResult(
            original=segment,
            translated_text=translated,
            provider=provider
        )
    
    def translate_document(self, input_file: Path, output_file: Path):
        """Translate an entire office document"""
        logger.info(f"=" * 60)
        logger.info(f"Translating: {input_file}")
        logger.info(f"Output: {output_file}")
        logger.info(f"Target language: {self.target_language}")
        logger.info(f"=" * 60)
        
        # Determine file type
        suffix = input_file.suffix.lower()
        if suffix not in self.extractors:
            raise ValueError(f"Unsupported file type: {suffix}")
        
        extractor = self.extractors[suffix]
        if not extractor:
            raise ValueError(f"Extractor for {suffix} not available (missing library)")
        
        # Step 1: Extract segments
        logger.info("Step 1: Extracting text segments...")
        segments = extractor.extract_segments(input_file)
        
        if not segments:
            logger.warning("No text segments found in document!")
            return
        
        # Step 2: Translate segments
        logger.info(f"Step 2: Translating {len(segments)} segments...")
        translations = []
        
        for i, segment in enumerate(segments):
            logger.info(f"  Progress: {i+1}/{len(segments)} ({(i+1)/len(segments)*100:.1f}%)")
            result = self.translate_segment(segment)
            translations.append(result)
            
            # Show provider stats every 10 segments
            if (i + 1) % 10 == 0:
                provider_counts = defaultdict(int)
                for t in translations:
                    provider_counts[t.provider] += 1
                logger.info(f"  Provider usage: {dict(provider_counts)}")
        
        # Step 3: Rebuild document
        logger.info("Step 3: Rebuilding document with translations...")
        extractor.rebuild_document(input_file, translations, output_file)
        
        # Final statistics
        logger.info("=" * 60)
        logger.info("Translation complete!")
        provider_counts = defaultdict(int)
        for t in translations:
            provider_counts[t.provider] += 1
        logger.info(f"Provider usage: {dict(provider_counts)}")
        logger.info(f"Output saved to: {output_file}")
        logger.info("=" * 60)


# === CLI Interface ===
def setup_logging(verbose: bool):
    """Configure logging based on verbosity"""
    level = logging.DEBUG if verbose else logging.INFO
    
    formatter = logging.Formatter(
        '%(asctime)s - %(levelname)s - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )
    
    handler = logging.StreamHandler()
    handler.setFormatter(formatter)
    
    logger.setLevel(level)
    logger.addHandler(handler)


def main():
    parser = argparse.ArgumentParser(
        description='Office Documents Translator with CTranslate2 and LLM Fallbacks',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
QUICK START (No Downloads):
  1. Get free API key: https://openrouter.ai
  2. export OPENROUTER_API_KEY="sk-or-v1-..."
  3. python %(prog)s input.pptx output.pptx --no-ctranslate

MANUAL MODEL DOWNLOAD (Much Faster):
  # Install aria2c for fast parallel downloads
  brew install aria2  # macOS
  sudo apt install aria2  # Linux
  
  # Download model (replace with small or large model)
  mkdir -p ~/.cache/huggingface/hub/models--cstr--wmt21ct2_int8/snapshots/main
  cd ~/.cache/huggingface/hub/models--cstr--wmt21ct2_int8/snapshots/main
  
  aria2c -x 16 -s 16 https://huggingface.co/cstr/wmt21ct2_int8/resolve/main/model.bin
  aria2c -x 16 -s 16 https://huggingface.co/cstr/wmt21ct2_int8/resolve/main/config.json
  # ... download other files ...
  
  Then run: python %(prog)s input.pptx output.pptx

EXAMPLES:
  # LLM-only (recommended for quick start)
  %(prog)s presentation.pptx output.pptx -l de --no-ctranslate
  
  # Small CTranslate2 model (~600MB)
  %(prog)s document.docx translated.docx -l es --small-model
  
  # Large CTranslate2 model (~4.7GB, better quality)
  %(prog)s document.docx translated.docx -l es
  
  # Use local model path
  %(prog)s file.xlsx out.xlsx --model-path /path/to/model

Supported Languages:
  de - German       es - Spanish      fr - French
  it - Italian      ja - Japanese     zh - Chinese
  ru - Russian
        """
    )
    
    parser.add_argument('input_file', nargs='?', help='Input office document')
    parser.add_argument('output_file', nargs='?', help='Output translated document')
    parser.add_argument('-l', '--language', default='de', 
                       help='Target language code (default: de)')
    parser.add_argument('-v', '--verbose', action='store_true',
                       help='Enable verbose debug output')
    parser.add_argument('--no-ctranslate', action='store_true',
                       help='Disable CTranslate2, use only LLM')
    parser.add_argument('--no-llm', action='store_true',
                       help='Disable LLM fallback, use only CTranslate2')
    parser.add_argument('--small-model', action='store_true',
                       help='Use small NLLB model (~600MB) instead of WMT21 (~4.7GB)')
    parser.add_argument('--model-path', type=str,
                       help='Path to pre-downloaded CTranslate2 model directory')
    parser.add_argument('--test', action='store_true',
                       help='Test translation setup')
    
    args = parser.parse_args()
    
    # Setup logging
    setup_logging(args.verbose)
    
    # Test mode
    if args.test:
        print("\n" + "=" * 60)
        print("OFFICE TRANSLATOR - SETUP TEST")
        print("=" * 60)
        
        if not args.no_ctranslate:
            print("\n1. Testing CTranslate2...")
            if args.small_model:
                print("   Using SMALL model (~600MB)")
            else:
                print("   Using LARGE model (~4.7GB)")
            try:
                ct2 = CTranslate2Translator('de', args.model_path, args.small_model)
                if ct2.available:
                    test_text = "Hello world! This is a test."
                    result = ct2.translate_text(test_text)
                    print(f"  ✓ CTranslate2 working")
                    print(f"  Test: '{test_text}' -> '{result}'")
                else:
                    print(f"  ✗ CTranslate2 initialization failed")
            except Exception as e:
                print(f"  ✗ Error: {e}")
        else:
            print("\n1. Skipping CTranslate2 test (--no-ctranslate)")
        
        print("\n2. Testing LLM Providers...")
        llm = LLMTranslator('German')
        if llm.providers:
            print(f"  ✓ Found {len(llm.providers)} providers: {list(llm.providers.keys())}")
            test_text = "Hello world!"
            result = llm.translate_text(test_text)
            if result:
                print(f"  ✓ LLM working: '{test_text}' -> '{result}'")
            else:
                print(f"  ✗ LLM translation failed")
        else:
            print(f"  ✗ No LLM providers configured")
            print(f"\n  Get free API keys:")
            print(f"    OpenRouter: https://openrouter.ai (Recommended)")
            print(f"    Groq: https://console.groq.com")
        
        print("\n3. Checking Document Libraries...")
        print(f"  PowerPoint (.pptx): {'✓' if HAS_PPTX else '✗'}")
        print(f"  Word (.docx): {'✓' if HAS_DOCX else '✗'}")
        print(f"  Excel (.xlsx): {'✓' if HAS_OPENPYXL else '✗'}")
        
        print("\n" + "=" * 60)
        print("RECOMMENDATIONS:")
        print("  Quick start: Use --no-ctranslate with OpenRouter API")
        print("  Best quality: Download large model with aria2c")
        print("  Balanced: Use --small-model (~600MB)")
        print("=" * 60)
        return
    
    # Require input/output files
    if not args.input_file or not args.output_file:
        parser.print_help()
        return
    
    # Validate input file
    input_path = Path(args.input_file)
    if not input_path.exists():
        logger.error(f"Input file not found: {input_path}")
        sys.exit(1)
    
    output_path = Path(args.output_file)
    
    try:
        translator = OfficeTranslator(
            target_language=args.language,
            use_ctranslate=not args.no_ctranslate,
            use_llm=not args.no_llm,
            model_path=args.model_path,
            use_small_model=args.small_model
        )
        
        if not translator.ct2_translator and not translator.llm_translator:
            logger.error("No translators available!")
            print("\nSet up at least one translator:")
            print("  - LLM: export OPENROUTER_API_KEY='your-key'")
            print("  - CTranslate2: Use --small-model or download manually")
            sys.exit(1)
        
        translator.translate_document(input_path, output_path)
        
        print(f"\n✓ Success! Translated document saved to: {output_path}")
        
    except KeyboardInterrupt:
        print("\n\nTranslation interrupted by user")
        sys.exit(1)
    except Exception as e:
        logger.error(f"Translation failed: {e}", exc_info=args.verbose)
        sys.exit(1)


if __name__ == "__main__":
    main()