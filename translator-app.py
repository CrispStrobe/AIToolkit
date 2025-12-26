#!/usr/bin/env python3
"""
Gradio interface for Document Translator
Designed for Hugging Face Spaces deployment
"""

import gradio as gr
import asyncio
import os
import sys
import logging
import subprocess
import shutil
from pathlib import Path
from typing import Optional, Tuple
import tempfile

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Import from translator.py
from translator import (
    UltimateDocumentTranslator,
    TranslationMode,
    TranslationBackend,
    AlignerBackend
)

# ============================================================================
# ENVIRONMENT SETUP
# ============================================================================

def check_and_setup_environment():
    """
    Verify and setup required tools for Hugging Face Spaces.
    Returns status messages.
    """
    status_messages = []
    
    # 1. Check python-docx
    try:
        from docx import Document
        status_messages.append("✓ python-docx installed")
    except ImportError:
        status_messages.append("⚠ python-docx not found - installing...")
        try:
            subprocess.run([sys.executable, "-m", "pip", "install", "python-docx"], check=True)
            status_messages.append("✓ python-docx installed successfully")
        except Exception as e:
            status_messages.append(f"✗ python-docx installation failed: {e}")
    
    # 2. Check python-pptx 
    try:
        from pptx import Presentation
        status_messages.append("✓ python-pptx installed")
    except ImportError:
        status_messages.append("⚠ python-pptx not found - installing...")
        try:
            subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx"], check=True)
            status_messages.append("✓ python-pptx installed successfully")
        except Exception as e:
            status_messages.append(f"✗ python-pptx installation failed: {e}")
    
    # 3. Check CTranslate2
    try:
        import ctranslate2
        status_messages.append("✓ CTranslate2 installed")
    except ImportError:
        status_messages.append("⚠ CTranslate2 not found - installing...")
        try:
            subprocess.run([sys.executable, "-m", "pip", "install", "ctranslate2"], check=True)
            status_messages.append("✓ CTranslate2 installed successfully")
        except Exception as e:
            status_messages.append(f"✗ CTranslate2 installation failed: {e}")
    
    # 4. Check fast_align (optional, complex to build on HF Spaces)
    fast_align_path = shutil.which("fast_align")
    if fast_align_path:
        status_messages.append(f"✓ fast_align found at {fast_align_path}")
    else:
        status_messages.append("ℹ fast_align not available (optional - will use other aligners)")
    
    # 5. Check for API keys (optional)
    if os.getenv("OPENAI_API_KEY"):
        status_messages.append("✓ OpenAI API key detected")
    if os.getenv("ANTHROPIC_API_KEY"):
        status_messages.append("✓ Anthropic API key detected")
    if not os.getenv("OPENAI_API_KEY") and not os.getenv("ANTHROPIC_API_KEY"):
        status_messages.append("ℹ No LLM API keys found (LLM modes will be unavailable)")
    
    return "\n".join(status_messages)

# Run setup on startup
SETUP_STATUS = check_and_setup_environment()
logger.info(f"Setup complete:\n{SETUP_STATUS}")

# ============================================================================
# FILE TYPE DETECTION & VALIDATION 
# ============================================================================

def detect_and_validate_file(file_path: Path) -> Tuple[bool, str, str]:
    """
    Detect file type and validate it.
    
    Returns:
        Tuple of (is_valid, file_type, error_message)
        file_type: 'docx', 'pptx', or 'unknown'
    """
    if not file_path.exists():
        return False, 'unknown', f"File not found: {file_path}"
    
    suffix = file_path.suffix.lower()
    
    # Quick check by extension
    if suffix == '.docx':
        try:
            from docx import Document
            doc = Document(str(file_path))
            return True, 'docx', ""
        except Exception as e:
            return False, 'docx', f"Invalid Word document: {str(e)}"
    
    elif suffix == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            return True, 'pptx', ""
        except Exception as e:
            return False, 'pptx', f"Invalid PowerPoint presentation: {str(e)}"
    
    else:
        return False, 'unknown', "Unsupported format. Only .docx and .pptx files are supported."


def get_file_info(file_path: Path, file_type: str) -> str:
    """
    Get basic info about the uploaded file.
    
    Returns formatted string with file statistics.
    """
    try:
        if file_type == 'docx':
            from docx import Document
            doc = Document(str(file_path))
            para_count = len([p for p in doc.paragraphs if p.text.strip()])
            table_count = len(doc.tables)
            return f"📄 Word Document | {para_count} paragraphs | {table_count} tables"
        
        elif file_type == 'pptx':
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slide_count = len(prs.slides)
            shape_count = sum(len(slide.shapes) for slide in prs.slides)
            return f"📊 PowerPoint Presentation | {slide_count} slides | ~{shape_count} shapes"
    
    except Exception as e:
        return f"📎 {file_type.upper()} file"
    
    return "📎 Document"

# ============================================================================
# TRANSLATION FUNCTION
# ============================================================================

async def translate_document_async(
    input_file,
    source_lang: str,
    target_lang: str,
    mode: str,
    nmt_backend: str,
    nllb_size: str,
    aligner: str,
    llm_provider: Optional[str],
    progress=gr.Progress()
) -> Tuple[Optional[str], str]:
    """
    Asynchronous document translation with progress tracking. Handles both .docx and .pptx files
    
    Returns:
        Tuple of (output_file_path, log_messages)
    """
    
    if input_file is None:
        return None, "❌ Error: No file uploaded"
    
    # Create temp directory for processing
    temp_dir = Path(tempfile.mkdtemp())
    
    try:
        # Setup paths
        input_path = Path(input_file.name)
        
        # Detect file type and validate
        is_valid, file_type, error_msg = detect_and_validate_file(input_path)
        
        if not is_valid:
            return None, f"❌ Error: {error_msg}"
        
        # Dynamic output filename based on detected type
        output_extension = input_path.suffix  # Preserve original extension
        output_filename = f"{input_path.stem}_translated_{source_lang}_{target_lang}{output_extension}"
        output_path = temp_dir / output_filename
        
        # Get file info for display
        file_info = get_file_info(input_path, file_type)
        
        # Map UI selections to enums
        mode_map = {
            'NMT Only': TranslationMode.NMT_ONLY,
            'LLM with Alignment': TranslationMode.LLM_WITH_ALIGN,
            'LLM without Alignment': TranslationMode.LLM_WITHOUT_ALIGN,
            'Hybrid (Recommended)': TranslationMode.HYBRID
        }
        
        # Setup logging capture
        log_messages = []
        
        class LogCapture(logging.Handler):
            def emit(self, record):
                log_messages.append(self.format(record))
        
        log_handler = LogCapture()
        log_handler.setFormatter(logging.Formatter('%(levelname)s - %(message)s'))
        logging.getLogger().addHandler(log_handler)
        
        progress(0.1, desc=f"Initializing translator for {file_type.upper()}...")
        
        # Initialize translator
        translator = UltimateDocumentTranslator(
            src_lang=source_lang,
            tgt_lang=target_lang,
            mode=mode_map[mode],
            nmt_backend=nmt_backend.lower() if nmt_backend != "Auto" else "auto",
            llm_provider=llm_provider.lower() if llm_provider and llm_provider != "None" else None,
            aligner=aligner.lower() if aligner != "Auto" else "auto",
            nllb_model_size=nllb_size
        )
        
        # Progress message based on file type
        if file_type == 'docx':
            progress(0.2, desc="Processing Word document...")
        else:
            progress(0.2, desc="Processing PowerPoint slides...")
        
        # Use unified translate_file() method instead of translate_document()
        await translator.translate_file(input_path, output_path)
        
        progress(1.0, desc="Translation complete!")
        
        # Cleanup log handler
        logging.getLogger().removeHandler(log_handler)
        
        # Format-aware success message
        format_name = "Word Document" if file_type == 'docx' else "PowerPoint Presentation"
        
        # Format log output
        log_output = "\n".join(log_messages[-50:])  # Last 50 messages
        success_msg = f"""
✅ Translation Complete!

{file_info}
📄 Input: {input_path.name}
📄 Output: {output_filename}
🌍 Direction: {source_lang.upper()} → {target_lang.upper()}
⚙️ Mode: {mode}
🔧 Backend: {nmt_backend}

Recent Logs:
{log_output}
"""
        
        return str(output_path), success_msg
        
    except ImportError as e:
        error_msg = f"❌ Missing Library Error:\n{str(e)}\n\n"
        if 'docx' in str(e):
            error_msg += "Install with: pip install python-docx"
        elif 'pptx' in str(e):
            error_msg += "Install with: pip install python-pptx"
        logger.error(f"Import error: {e}", exc_info=True)
        return None, error_msg
        
    except Exception as e:
        error_msg = f"❌ Translation Error:\n{str(e)}\n\nPlease check your settings and try again."
        logger.error(f"Translation failed: {e}", exc_info=True)
        return None, error_msg

def translate_document_sync(*args, **kwargs):
    """Synchronous wrapper with explicit event loop management"""
    import asyncio
    import concurrent.futures
    from datetime import datetime
    
    print(f"\n{'='*60}")
    print(f"[{datetime.now().strftime('%H:%M:%S')}] 🔄 translate_document_sync called")
    
    # Define a helper to run the async function in a new thread's loop
    def run_in_new_loop(func, *args, **kwargs):
        """Create a fresh event loop in the new thread"""
        new_loop = asyncio.new_event_loop()
        asyncio.set_event_loop(new_loop)
        try:
            return new_loop.run_until_complete(func(*args, **kwargs))
        finally:
            new_loop.close()
    
    try:
        try:
            # Check if a loop is already running (Gradio/HF Spaces context)
            asyncio.get_running_loop()
            print(f"[DEBUG] ⚠️  Event loop running - Offloading to ThreadPool")
            
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
                # Pass the FUNCTION and ARGS separately
                future = executor.submit(
                    run_in_new_loop, 
                    translate_document_async, 
                    *args, 
                    **kwargs
                )
                result = future.result(timeout=600)
                print(f"[DEBUG] ✓ Thread execution completed")
                return result
                
        except RuntimeError:
            # No loop running (Standalone context)
            print(f"[DEBUG] ℹ️  No running loop - Using standard asyncio.run")
            result = asyncio.run(translate_document_async(*args, **kwargs))
            print(f"[DEBUG] ✓ asyncio.run() completed")
            return result
            
    except concurrent.futures.TimeoutError:
        error_msg = "❌ Translation timeout (>10 minutes)"
        print(f"[ERROR] {error_msg}")
        return None, error_msg
        
    except Exception as e:
        print(f"[ERROR] Critical failure: {e}")
        import traceback
        traceback.print_exc()
        return None, f"❌ Error: {str(e)}"
        
    finally:
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 🏁 Finished")
        print(f"{'='*60}\n")

# ============================================================================
# GRADIO INTERFACE
# ============================================================================

def create_interface():
    """
    Supports both .docx and .pptx files
    """
    
    # Language options (unchanged)
    languages = {
        "English": "en",
        "German": "de",
        "French": "fr",
        "Spanish": "es",
        "Italian": "it",
        "Portuguese": "pt",
        "Russian": "ru",
        "Chinese": "zh",
        "Japanese": "ja",
        "Korean": "ko",
        "Arabic": "ar",
        "Hindi": "hi",
        "Dutch": "nl",
        "Polish": "pl",
        "Turkish": "tr",
        "Czech": "cs",
        "Ukrainian": "uk",
        "Vietnamese": "vi"
    }
    
    with gr.Blocks(title="Document Translator") as demo:
        
        gr.Markdown("""
        # 🌍 Document Translator
        
        Translate Word documents and PowerPoint presentations while preserving formatting, 
        footnotes, styling, and layout.
        
        **Supported formats:** `.docx` (Word) and `.pptx` (PowerPoint)
        """)
        
        with gr.Row():
            with gr.Column(scale=1):
                gr.Markdown("### 📤 Input")
                
                input_file = gr.File(
                    label="Upload Document (.docx or .pptx)",
                    file_types=[".docx", ".pptx"],
                    type="filepath"
                )
                
                with gr.Row():
                    source_lang = gr.Dropdown(
                        choices=list(languages.keys()),
                        value="English",
                        label="Source Language"
                    )
                    target_lang = gr.Dropdown(
                        choices=list(languages.keys()),
                        value="German",
                        label="Target Language"
                    )
                
                gr.Markdown("### ⚙️ Settings")
                
                mode = gr.Dropdown(
                    choices=[
                        "Hybrid (Recommended)",
                        "NMT Only",
                        "LLM with Alignment",
                        "LLM without Alignment"
                    ],
                    value="Hybrid (Recommended)",
                    label="Translation Mode",
                    info="Hybrid uses NMT with optional LLM enhancement"
                )
                
                nmt_backend = gr.Dropdown(
                    choices=["NLLB", "Madlad", "Opus", "CT2", "Auto"],
                    value="NLLB",
                    label="NMT Backend",
                    info="NLLB: Fast & balanced | Madlad: Academic | Opus: Specialized pairs"
                )
                
                nllb_size = gr.Dropdown(
                    choices=["600M", "1.3B", "3.3B"],
                    value="600M",
                    label="NLLB Model Size",
                    info="600M recommended for Hugging Face Spaces (limited RAM)"
                )
                
                aligner = gr.Dropdown(
                    choices=["Auto", "Awesome", "SimAlign", "Lindat", "Heuristic"],
                    value="Auto",
                    label="Word Aligner",
                    info="Auto will select best available aligner"
                )
                
                llm_provider = gr.Dropdown(
                    choices=["None", "OpenAI", "Anthropic", "Ollama"],
                    value="None",
                    label="LLM Provider (Optional)",
                    info="Requires API key in environment variables"
                )
                
                translate_btn = gr.Button("🚀 Translate Document", variant="primary", size="lg")
            
            with gr.Column(scale=1):
                gr.Markdown("### 📥 Output")
                
                output_file = gr.File(
                    label="Translated Document",
                    interactive=False
                )
                
                log_output = gr.Textbox(
                    label="Translation Log",
                    lines=20,
                    max_lines=30,
                    interactive=False
                )

            
                file_preview = gr.Markdown(label="File Preview")
                input_file.change(
                     fn=preview_uploaded_file,
                     inputs=[input_file],
                     outputs=[file_preview]
                )
        
        gr.Markdown(f"### System Status\n```\n{SETUP_STATUS}\n```")
        
        gr.Markdown("""
        ### ✨ Features
        
        **Word Documents (.docx):**
        - Preserves formatting, footnotes, headers/footers
        - Maintains tables, styles, and paragraph formatting
        - Word-level alignment for accurate format transfer
        
        **PowerPoint Presentations (.pptx):**
        - Translates all text boxes and placeholders
        - Preserves slide layouts and positioning
        - Maintains tables, speaker notes, and formatting
        - Handles grouped shapes and complex layouts
        
        **Translation Engines:**
        - Multiple neural backends (NLLB, Madlad, Opus-MT, WMT21)
        - Optional LLM enhancement (OpenAI/Anthropic)
        - Word-level alignment for format preservation
        
        **Recommended Settings:**
        - Mode: Hybrid (best quality)
        - Backend: NLLB (fastest, good quality)
        - Size: 600M (good balance for Spaces)
                    
        ### 📖 Usage Tips
        
        **For Word Documents:**
        - Best for: Academic papers, reports, articles
        - Use "Hybrid" mode for complex formatting
        - Madlad backend excels with technical content
        
        **For PowerPoint:**
        - Best for: Business presentations, slides
        - Use "NMT Only" for speed
        - Check speaker notes in output file
        
        **General:**
        - Larger NLLB models (1.3B, 3.3B) improve quality but use more RAM
        - LLM modes require API keys and are slower
        - First translation may be slower (model download)
        
        ### ⚠️ Limitations
        
        - Only modern Office formats (.docx, .pptx) - not legacy .doc/.ppt
        - Large files may take several minutes to process
        - Complex formatting may require manual review after translation
        - LLM modes are slower and require API access
        - Embedded images and charts are not translated (text only)
        - On Hugging Face Spaces: Limited to 600M model due to RAM constraints
        
        ### 🔧 Advanced Settings Guide
        
        **Translation Modes:**
        - **NMT Only**: Pure neural translation, fastest, most reliable
        - **Hybrid**: Combines NMT with selective LLM enhancement (recommended)
        - **LLM with Alignment**: Uses LLM + word alignment for formatting
        - **LLM without Alignment**: LLM-only, best for natural output
        
        **Backends:**
        - **NLLB**: 200+ languages, fast, balanced quality
        - **Madlad**: Google's 3B model, excellent for academic/formal text
        - **Opus**: Specialized bilingual pairs, very fast for supported pairs
        - **CT2**: Dense models, best German/European quality
        
        **Aligners:**
        - **Awesome**: BERT-based, high precision (recommended for Mac/M1)
        - **SimAlign**: Heavy PyTorch BERT, good quality but slower
        - **Lindat**: Cloud API, no local resources needed
        - **Heuristic**: Simple fallback, fast but basic
        """)
        
        def handle_translate(input_f, src_lang_name, tgt_lang_name, mode, nmt, nllb_sz, algn, llm):
            src_code = languages.get(src_lang_name, "en")
            tgt_code = languages.get(tgt_lang_name, "de")
            return translate_document_sync(input_f, src_code, tgt_code, mode, nmt, nllb_sz, algn, llm)
        
        translate_btn.click(
            fn=handle_translate,
            inputs=[
                input_file,
                source_lang,
                target_lang,
                mode,
                nmt_backend,
                nllb_size,
                aligner,
                llm_provider
            ],
            outputs=[output_file, log_output]
        )
    
    return demo

# ============================================================================
# FILE PREVIEW
# ============================================================================

def preview_uploaded_file(file_path) -> str:
    """
    Generate a preview of the uploaded file content.
    Simplified version with better error handling.
    """
    if file_path is None:
        return "No file uploaded"
    
    try:
        input_path = Path(file_path.name)
        is_valid, file_type, error = detect_and_validate_file(input_path)
        
        if not is_valid:
            return f"❌ {error}"
        
        if file_type == 'docx':
            from docx import Document
            doc = Document(str(input_path))
            
            para_count = len([p for p in doc.paragraphs if p.text.strip()])
            table_count = len(doc.tables)
            
            preview_lines = [
                "📄 **Word Document Preview**\n",
                f"📊 Statistics:",
                f"  • {para_count} paragraphs with text",
                f"  • {table_count} tables",
                "\n**First paragraphs:**\n"
            ]
            
            count = 0
            for para in doc.paragraphs:
                if para.text.strip() and count < 3:
                    text = para.text.strip().replace('\n', ' ')
                    preview_lines.append(f"{count+1}. {text[:100]}{'...' if len(text) > 100 else ''}")
                    count += 1
            
            if count == 0:
                preview_lines.append("(No text content found)")
            
            return "\n".join(preview_lines)
        
        elif file_type == 'pptx':
            from pptx import Presentation
            
            try:
                prs = Presentation(str(input_path))
                slide_count = len(prs.slides)
                
                preview_lines = [
                    "📊 **PowerPoint Preview**\n",
                    f"📈 Statistics:",
                    f"  • {slide_count} slide{'s' if slide_count != 1 else ''}",
                ]
                
                # Count shapes across all slides
                total_shapes = 0
                text_shapes = 0
                
                for slide in prs.slides:
                    for shape in slide.shapes:
                        total_shapes += 1
                        try:
                            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                                if shape.text_frame.text.strip():
                                    text_shapes += 1
                        except:
                            pass
                
                preview_lines.append(f"  • {total_shapes} total shapes")
                preview_lines.append(f"  • {text_shapes} text boxes/placeholders")
                preview_lines.append("\n**Sample content from first slides:**\n")
                
                # Preview first 3 slides
                previewed = 0
                for idx in range(min(3, slide_count)):
                    try:
                        slide = prs.slides[idx]
                        preview_lines.append(f"**Slide {idx+1}:**")
                        
                        texts_found = 0
                        for shape in slide.shapes:
                            try:
                                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                                    text = shape.text_frame.text.strip()
                                    if text:
                                        text = text.replace('\n', ' ')
                                        preview_lines.append(f"  • {text[:80]}{'...' if len(text) > 80 else ''}")
                                        texts_found += 1
                                        if texts_found >= 2:  # Max 2 texts per slide
                                            break
                            except Exception as shape_err:
                                # Skip problematic shapes silently
                                continue
                        
                        if texts_found == 0:
                            preview_lines.append("  (No text content)")
                        
                        preview_lines.append("")
                        previewed += 1
                        
                    except Exception as slide_err:
                        logger.debug(f"Could not preview slide {idx+1}: {slide_err}")
                        preview_lines.append(f"**Slide {idx+1}:** (Could not access)\n")
                        continue
                
                if previewed == 0:
                    preview_lines.append("(Could not access slide content)")
                
                return "\n".join(preview_lines)
                
            except Exception as prs_error:
                return f"📊 PowerPoint file detected\n\n❌ Preview unavailable: {str(prs_error)}\n\nFile appears valid and should translate successfully."
    
    except Exception as e:
        logger.error(f"Preview generation failed: {e}", exc_info=True)
        return f"❌ Preview error: {str(e)}"
    
    return "Could not generate preview"

# ============================================================================
# MAIN
# ============================================================================

if __name__ == "__main__":
    demo = create_interface()
    demo.launch(
        server_name="0.0.0.0",
        server_port=7860,
        share=False,
        theme=gr.themes.Soft()  
    )