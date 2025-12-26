#!/usr/bin/env python3
import argparse
import asyncio
import logging
import os
import re
import sys
import gc
from collections import defaultdict
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional, Set, Tuple, Any
from enum import Enum
import requests
from tqdm import tqdm

# --- Core Library Diagnostics ---
print("🔍 System Check...")

def check_library(name, import_stmt):
    try:
        exec(import_stmt, globals())
        print(f"  ✓ {name}")
        return True
    except ImportError:
        print(f"  ⊘ {name} (optional)")
        return False
    except Exception as e:
        print(f"  ✗ {name} error: {e}")
        return False

HAS_DOCX = check_library("python-docx", "from docx import Document; from docx.shared import Pt, RGBColor; from docx.text.paragraph import Paragraph; from docx.oxml.shared import OxmlElement; from docx.oxml.ns import qn")
HAS_PPTX = check_library("python-pptx", "from pptx import Presentation; from pptx.util import Pt, Inches; from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR; from pptx.dml.color import RGBColor")
HAS_TORCH = check_library("torch", "import torch")
HAS_CT2 = check_library("CTranslate2", "import ctranslate2; from huggingface_hub import snapshot_download")
HAS_TRANSFORMERS = check_library("Transformers", "from transformers import AutoTokenizer, AutoModelForSeq2SeqLM")
HAS_OPENAI = check_library("OpenAI", "from openai import AsyncOpenAI")
HAS_ANTHROPIC = check_library("Anthropic", "from anthropic import AsyncAnthropic")
HAS_SIMALIGN = check_library("simalign", "from simalign import SentenceAligner")

# --- Device & Backend Configuration ---
def get_torch_device():
    if not HAS_TORCH: return "cpu"
    import torch
    if torch.cuda.is_available(): return torch.device("cuda")
    if hasattr(torch.backends, "mps") and torch.backends.mps.is_available():
        return torch.device("mps")
    return torch.device("cpu")

def get_ct2_settings():
    """Optimized for M1 Mac (Accelerate CPU) vs NVIDIA (CUDA)."""
    if HAS_TORCH:
        import torch
        if torch.cuda.is_available(): return "cuda", "float16"
    # Mac M1/M2/M3: MUST use 'cpu' for CT2 and 'int8' for peak ARM64 optimization
    return "cpu", "int8"

def check_fast_align():
    script_dir = Path(__file__).parent
    binary_locations = ["../fast_align/build/fast_align", "./fast_align/build/fast_align", "fast_align"]
    for loc in binary_locations:
        path = script_dir / loc if not loc.startswith('/') else Path(loc)
        if os.path.isfile(path) and os.access(path, os.X_OK): return True
    return False

HAS_FAST_ALIGN = check_fast_align()
print(f"  {'✓' if HAS_FAST_ALIGN else '⊘'} fast_align")
print("-" * 60)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# ============================================================================
# ENUMS FOR CONFIGURATION
# ============================================================================

class TranslationBackend(Enum):
    """Available translation backends"""
    CT2 = "ct2"
    NLLB = "nllb"
    OPENAI = "openai"
    ANTHROPIC = "anthropic"
    OLLAMA = "ollama"
    AUTO = "auto"


class AlignerBackend(Enum):
    """Available alignment backends"""
    LINDAT = "lindat"
    FAST_ALIGN = "fast_align"
    SIMALIGN = "simalign"
    HEURISTIC = "heuristic"
    AUTO = "auto"


class TranslationMode(Enum):
    """Translation modes"""
    NMT_ONLY = "nmt"
    LLM_WITH_ALIGN = "llm-align"
    LLM_WITHOUT_ALIGN = "llm-plain"
    HYBRID = "hybrid"


# ============================================================================
# DATA STRUCTURES
# ============================================================================

@dataclass
class FormatRun:
    """A run of text with formatting"""
    text: str
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    font_name: Optional[str] = None
    font_size: Optional[float] = None
    font_color: Optional[Tuple[int, int, int]] = None


@dataclass
class TranslatableParagraph:
    """Paragraph with formatting and metadata"""
    runs: List[FormatRun] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    def get_text(self) -> str:
        return ''.join(run.text for run in self.runs)
    
    def get_words(self) -> List[str]:
        """Clean tokenization for the ALIGNER only: strips punctuation."""
        import re
        text = self.get_text()
        # Extracts only alphanumeric sequences (e.g., "Theology" from "Theology:")
        return re.findall(r"\w+", text)
    
    def get_formatted_word_indices(self) -> Dict[str, Set[int]]:
        """Maps formatting to clean alphanumeric word indices."""
        formatted = {'italic': set(), 'bold': set(), 'italic_bold': set()}
        text = self.get_text()
        words = self.get_words() # Uses the same list the aligner sees
        
        if not words:
            return formatted
        
        char_to_word = {}
        last_found = 0
        for word_idx, word in enumerate(words):
            # Find the word in the text, starting search after the previous word
            start = text.find(word, last_found)
            if start != -1:
                for i in range(start, start + len(word)):
                    char_to_word[i] = word_idx
                last_found = start + len(word)
        
        char_pos = 0
        for run in self.runs:
            if not run.text:
                continue
            for char in run.text:
                if char_pos in char_to_word:
                    word_idx = char_to_word[char_pos]
                    if not char.isspace():
                        if run.bold and run.italic:
                            formatted['italic_bold'].add(word_idx)
                        elif run.italic:
                            formatted['italic'].add(word_idx)
                        elif run.bold:
                            formatted['bold'].add(word_idx)
                char_pos += 1
        return formatted


# ============================================================================
# TRANSLATION BACKENDS (keeping previous implementations)
# ============================================================================

class CTranslate2Translator:
    """WMT21 Backend: High-quality dense model."""
    
    MODELS = {
        'en_to_x': 'cstr/wmt21ct2_int8',
        'x_to_en': 'cstr/wmt21-ct2-x-en-int8',
    }
    
    SUPPORTED_LANGS = ['de', 'es', 'fr', 'it', 'ja', 'zh', 'ru', 'pt', 'nl', 'cs', 'uk']
    
    def __init__(self, src_lang: str, tgt_lang: str):
        self.src_lang, self.tgt_lang = src_lang, tgt_lang
        self.available = False
        self.ct2_dev, self.ct2_compute = get_ct2_settings()
        
        if not HAS_CT2: return
        
        if src_lang == 'en' and tgt_lang in self.SUPPORTED_LANGS:
            self.direction, self.tokenizer_name = 'en_to_x', "facebook/wmt21-dense-24-wide-en-x"
        elif tgt_lang == 'en' and src_lang in self.SUPPORTED_LANGS:
            self.direction, self.tokenizer_name = 'x_to_en', "facebook/wmt21-dense-24-wide-x-en"
        else: return
        
        try:
            logger.info(f"Loading WMT21-CT2 ({self.direction})...")
            model_path = self._get_or_download_model()
            self.translator = ctranslate2.Translator(
                model_path, device=self.ct2_dev, compute_type=self.ct2_compute
            )
            self.tokenizer = AutoTokenizer.from_pretrained(self.tokenizer_name)
            self.available = True
            logger.info(f"✓ WMT21 CT2 ready on {self.ct2_dev}")
        except Exception as e:
            logger.error(f"WMT21 initialization failure: {e}")
    
    def _get_or_download_model(self) -> str:
        cache_base = Path.home() / '.cache' / 'huggingface' / 'hub'
        model_repo = self.MODELS[self.direction]
        model_dir = cache_base / f"models--{model_repo.replace('/', '--')}"
        ref_path = model_dir / 'refs' / 'main'
        
        if ref_path.exists():
            with open(ref_path) as f:
                commit_hash = f.read().strip()
            snapshot_path = model_dir / 'snapshots' / commit_hash
            if (snapshot_path / 'model.bin').exists():
                return str(snapshot_path)
        
        logger.info("Downloading CT2 model...")
        return snapshot_download(repo_id=model_repo)
    
    def translate_batch(self, texts: List[str]) -> List[str]:
        if not self.available or not texts:
            return texts
        
        try:
            # 1. Ensure language tokens are correctly formatted
            # WMT21 models expect the target lang as a prefix
            target_prefix = [[self.tokenizer.lang_code_to_token[self.tgt_lang]]] * len(texts)
            
            # 2. Tokenize with specific padding/truncation
            source_tokens = [self.tokenizer.convert_ids_to_tokens(self.tokenizer.encode(t)) for t in texts]
            
            # 3. Translate with proper beam search and repetition penalty
            results = self.translator.translate_batch(
                source_tokens,
                target_prefix=target_prefix,
                beam_size=5,
                max_batch_size=16,
                repetition_penalty=1.2,
                # Prevent the model from just 'copying' the source if it gets confused
                disable_unk=True 
            )
            
            translated = []
            for i, res in enumerate(results):
                # Strip the language token from the start of the result
                tokens = res.hypotheses[0]
                if self.tokenizer.lang_code_to_token[self.tgt_lang] in tokens:
                    tokens = tokens[tokens.index(self.tokenizer.lang_code_to_token[self.tgt_lang]) + 1:]
                
                decoded = self.tokenizer.decode(self.tokenizer.convert_tokens_to_ids(tokens), skip_special_tokens=True)
                translated.append(decoded.strip())
                
            return translated
        except Exception as e:
            logger.error(f"CT2 Critical Error: {e}")
            return texts
class OpusMTTranslator:
    """Opus-MT Backend: Tiny, specialized bilingual models. Standalone logic."""
    def __init__(self, src_lang: str, tgt_lang: str):
        self.src_lang, self.tgt_lang = src_lang, tgt_lang
        self.available = False
        self.ct2_dev, self.ct2_compute = get_ct2_settings()
        
        # The standard name format for Opus-MT
        original_repo = f"Helsinki-NLP/opus-mt-{src_lang}-{tgt_lang}"
        # Your custom optimized repo
        self.custom_repo = f"cstr/opus-mt-{src_lang}-{tgt_lang}-ct2-int8"

        try:
            logger.info(f"NMT | Loading weights from {self.custom_repo}...")
            model_path = snapshot_download(repo_id=self.custom_repo)
            
            self.translator = ctranslate2.Translator(model_path, device=self.ct2_dev, compute_type=self.ct2_compute)
            
            # FIXED: Load tokenizer from the ORIGINAL repo name to get the correct Marian model_type
            # Transformers library will use the tiny cached JSON files from the original.
            self.tokenizer = AutoTokenizer.from_pretrained(original_repo)
            
            self.available = True
            logger.info(f"✓ NMT | Opus-MT ready (Weights: cstr / Tokenizer: original)")
        except Exception as e:
            logger.warning(f"NMT | Opus-MT primary load failed: {e}. Trying michaelfeil fallback...")
            try:
                fallback = f"michaelfeil/ct2fast-opus-mt-{src_lang}-{tgt_lang}"
                model_path = snapshot_download(repo_id=fallback)
                self.translator = ctranslate2.Translator(model_path, device=self.ct2_dev, compute_type=self.ct2_compute)
                self.tokenizer = AutoTokenizer.from_pretrained(original_repo)
                self.available = True
                logger.info(f"✓ NMT | Opus-MT ready using {fallback}")
            except:
                logger.error("NMT | All Opus-MT paths failed.")

    def translate_batch(self, texts: List[str]) -> List[str]:
        if not self.available or not texts: return texts
        source_tokens = [self.tokenizer.convert_ids_to_tokens(self.tokenizer.encode(t)) for t in texts]
        results = self.translator.translate_batch(source_tokens, beam_size=5)
        return [self.tokenizer.decode(self.tokenizer.convert_tokens_to_ids(r.hypotheses[0]), skip_special_tokens=True) for r in results]
    
class Madlad400Translator:
    """Madlad-400 Backend: Google's 3B powerhouse. Optimized for your cstr/ repo."""
    def __init__(self, src_lang: str, tgt_lang: str, model_size: str = "3b"):
        self.src_lang, self.tgt_lang = src_lang, tgt_lang
        self.available = False
        self.ct2_dev, self.ct2_compute = get_ct2_settings()
        
        original_repo = f"google/madlad400-{model_size}-mt"
        self.custom_repo = f"cstr/madlad400-{model_size}-ct2-int8"
        self.tgt_prefix = f"<2{tgt_lang}>"

        try:
            logger.info(f"NMT | Loading Madlad-400 from {self.custom_repo}...")
            model_path = snapshot_download(repo_id=self.custom_repo)
            
            self.translator = ctranslate2.Translator(model_path, device=self.ct2_dev, compute_type=self.ct2_compute)
            
            # FIXED: Point tokenizer to Google's repo to resolve the T5 architecture correctly
            self.tokenizer = AutoTokenizer.from_pretrained(original_repo)
            
            self.available = True
            logger.info(f"✓ NMT | Madlad-400 ready.")
        except Exception as e:
            logger.error(f"NMT | Madlad-400 load failed: {e}")

    def translate_batch(self, texts: List[str]) -> List[str]:
        if not self.available or not texts: return texts
        # Prepends <2de> (or similar) to every sentence in the batch
        source_tokens = [self.tokenizer.convert_ids_to_tokens(self.tokenizer.encode(f"{self.tgt_prefix} {t}")) for t in texts]
        results = self.translator.translate_batch(source_tokens, beam_size=1, repetition_penalty=2.0)
        return [self.tokenizer.decode(self.tokenizer.convert_tokens_to_ids(r.hypotheses[0]), skip_special_tokens=True) for r in results]
class NLLBTranslator:
    """NLLB-200 translator using CTranslate2 for 4x speedup and low memory"""
    
    # Map sizes to specific CTranslate2 optimized repositories
    REPOS = {
        "600M": "JustFrederik/nllb-200-distilled-600M-ct2-int8",
        "1.3B": "OpenNMT/nllb-200-distilled-1.3B-ct2-int8",
        "3.3B": "OpenNMT/nllb-200-3.3B-ct2-int8"
    }

    LANG_CODES = {
        'en': 'eng_Latn', 'de': 'deu_Latn', 'fr': 'fra_Latn', 
        'es': 'spa_Latn', 'it': 'ita_Latn', 'pt': 'por_Latn',
        'ru': 'rus_Cyrl', 'zh': 'zho_Hans', 'ja': 'jpn_Jpan',
        'ko': 'kor_Hang', 'ar': 'arb_Arab', 'hi': 'hin_Deva',
        'nl': 'nld_Latn', 'pl': 'pol_Latn', 'tr': 'tur_Latn',
        'cs': 'ces_Latn', 'uk': 'ukr_Cyrl', 'vi': 'vie_Latn',
    }
    
    def __init__(self, src_lang: str, tgt_lang: str, model_size: str = "600M"):
        self.src_lang, self.tgt_lang = src_lang, tgt_lang
        self.available = False
        self.mode = None
        self.device = get_torch_device()
        self.ct2_dev, self.ct2_compute = get_ct2_settings()
        
        self.src_code = self.LANG_CODES.get(src_lang)
        self.tgt_code = self.LANG_CODES.get(tgt_lang)
        
        if not self.src_code or not self.tgt_code: return

        ct2_repo = self.REPOS.get(model_size, self.REPOS["600M"])
        standard_repo = f"facebook/nllb-200-distilled-{model_size}"

        if HAS_CT2:
            try:
                logger.info(f"Loading NLLB-CT2 from {ct2_repo}...")
                model_path = snapshot_download(repo_id=ct2_repo)
                try:
                    self.translator = ctranslate2.Translator(
                        model_path, device=self.ct2_dev, compute_type=self.ct2_compute
                    )
                except Exception:
                    # Specific fallback for Mac CPU optimization mismatch
                    self.translator = ctranslate2.Translator(
                        model_path, device="cpu", compute_type="int8"
                    )

                self.tokenizer = AutoTokenizer.from_pretrained(standard_repo, src_lang=self.src_code)
                self.mode, self.available = "ct2", True
                logger.info(f"✓ NLLB-{model_size} CT2 ready.")
                return
            except Exception as e:
                logger.warning(f"NLLB-CT2 failed: {e}")

        if HAS_TRANSFORMERS and HAS_TORCH:
            try:
                logger.info(f"Fallback: Loading standard PyTorch NLLB...")
                self.model = AutoModelForSeq2SeqLM.from_pretrained(standard_repo)
                self.tokenizer = HFTokenizer.from_pretrained(standard_repo, src_lang=self.src_code)
                if self.device.type == "mps": self.model = self.model.half()
                self.model.to(self.device).eval()
                self.mode, self.available = "torch", True
                logger.info(f"✓ NLLB PyTorch ready on {self.device}")
            except Exception as e:
                logger.error(f"Critical: All NLLB paths failed: {e}")
    
    def translate_batch(self, texts: List[str], batch_size: int = 16) -> List[str]:
        """Translate batch using CT2 efficient beam search"""
        if not self.available or not texts:
            return texts
        
        try:
            results = []
            # NLLB requires the target language code as the first token (target prefix)
            target_prefix = [self.tgt_code]
            
            for i in range(0, len(texts), batch_size):
                batch = texts[i:i + batch_size]
                
                # Tokenize
                source_tokens = [
                    self.tokenizer.convert_ids_to_tokens(self.tokenizer.encode(t)) 
                    for t in batch
                ]
                
                # Translate
                step_results = self.translator.translate_batch(
                    source_tokens,
                    target_prefix=[target_prefix] * len(batch),
                    beam_size=4,
                    max_batch_size=batch_size,
                    repetition_penalty=1.1
                )
                
                # Decode (skipping the lang code prefix in the output)
                for res in step_results:
                    tokens = res.hypotheses[0]
                    # The first token is usually the target lang code
                    if tokens[0] == self.tgt_code:
                        tokens = tokens[1:]
                    
                    decoded = self.tokenizer.decode(
                        self.tokenizer.convert_tokens_to_ids(tokens), 
                        skip_special_tokens=True
                    )
                    results.append(decoded.strip())
            
            return results
        except Exception as e:
            logger.error(f"NLLB-CT2 translation failed: {e}")
            return texts
        
class LLMTranslator:
    """LLM translator (OpenAI/Anthropic/Ollama)"""
    
    def __init__(self, src_lang: str, tgt_lang: str, preferred_provider: Optional[str] = None):
        self.src_lang = src_lang
        self.tgt_lang = tgt_lang
        self.providers = self._init_providers(preferred_provider)
        
        if self.providers:
            logger.info(f"✓ LLM available ({list(self.providers.keys())})")
    
    def _init_providers(self, preferred: Optional[str] = None) -> Dict[str, Any]:
        providers = {}
        
        if HAS_OPENAI and os.getenv("OPENAI_API_KEY"):
            providers["openai"] = {
                "client": AsyncOpenAI(api_key=os.getenv("OPENAI_API_KEY")),
                "model": os.getenv("OPENAI_MODEL", "gpt-4o-mini")
            }
        
        if HAS_ANTHROPIC and os.getenv("ANTHROPIC_API_KEY"):
            providers["anthropic"] = {
                "client": AsyncAnthropic(api_key=os.getenv("ANTHROPIC_API_KEY")),
                "model": os.getenv("ANTHROPIC_MODEL", "claude-3-5-sonnet-20241022")
            }
        
        if self._check_ollama():
            providers["ollama"] = {
                "url": "http://localhost:11434/api/generate",
                "model": self._get_ollama_model()
            }
        
        if preferred and preferred in providers:
            return {preferred: providers[preferred]}
        
        return providers
    
    def _check_ollama(self) -> bool:
        try:
            r = requests.get("http://localhost:11434/api/tags", timeout=2)
            return r.status_code == 200 and len(r.json().get("models", [])) > 0
        except:
            return False
    
    def _get_ollama_model(self) -> str:
        try:
            r = requests.get("http://localhost:11434/api/tags", timeout=2)
            models = r.json().get("models", [])
            return models[0]["name"] if models else "llama3.2"
        except:
            return "llama3.2"
    
    async def translate_text(self, text: str, use_alignment: bool = True) -> Optional[str]:
        """Translate single text"""
        if not text.strip() or not self.providers:
            return None
        
        if use_alignment:
            prompt = (
                f"Translate the following text from {self.src_lang} to {self.tgt_lang}. "
                f"Preserve the word order as much as possible for alignment purposes. "
                f"Return ONLY the translation:\n\n{text}"
            )
        else:
            prompt = (
                f"Translate the following text from {self.src_lang} to {self.tgt_lang}. "
                f"Provide a natural, fluent translation. "
                f"Return ONLY the translation:\n\n{text}"
            )
        
        for provider_name, provider in self.providers.items():
            try:
                if provider_name == "openai":
                    response = await provider["client"].chat.completions.create(
                        model=provider["model"],
                        messages=[{"role": "user", "content": prompt}],
                        temperature=0.3,
                        max_tokens=4000
                    )
                    return response.choices[0].message.content.strip()
                
                elif provider_name == "anthropic":
                    response = await provider["client"].messages.create(
                        model=provider["model"],
                        messages=[{"role": "user", "content": prompt}],
                        temperature=0.3,
                        max_tokens=4000
                    )
                    return response.content[0].text.strip()
                
                elif provider_name == "ollama":
                    r = requests.post(
                        provider["url"],
                        json={"model": provider["model"], "prompt": prompt, "stream": False},
                        timeout=120
                    )
                    if r.status_code == 200:
                        return r.json().get("response", "").strip()
            
            except Exception as e:
                logger.debug(f"{provider_name} failed: {e}")
                continue
        
        return None
    
    async def translate_batch(self, texts: List[str], use_alignment: bool = True) -> List[str]:
        """Translate batch"""
        tasks = [self.translate_text(text, use_alignment) for text in texts]
        results = await asyncio.gather(*tasks)
        return [res if res else text for res, text in zip(results, texts)]


# ============================================================================
# WORD ALIGNERS
# ============================================================================

class LindatAligner:
    """Lindat word alignment API: Zero RAM usage fallback."""
    
    def __init__(self, src_lang: str, tgt_lang: str):
        self.src_lang, self.tgt_lang = src_lang, tgt_lang
        self.available = self._check_available()
        if self.available:
            logger.info(f"✓ ALIGN | Lindat API available ({src_lang}-{tgt_lang})")
    
    def _check_available(self) -> bool:
        try:
            r = requests.get(f"https://lindat.cz/services/text-aligner/align/{self.src_lang}-{self.tgt_lang}", timeout=3)
            return r.status_code in [200, 405]
        except: return False
    
    def align(self, src_words: List[str], tgt_words: List[str]) -> List[Tuple[int, int]]:
        if not self.available or not src_words or not tgt_words: return []
        try:
            r = requests.post(
                f"https://lindat.cz/services/text-aligner/align/{self.src_lang}-{self.tgt_lang}",
                headers={'Content-Type': 'application/json'},
                json={'src_tokens': [src_words], 'trg_tokens': [tgt_words]},
                timeout=15
            )
            if r.status_code == 200:
                alignment = r.json()["alignment"][0]
                return [(int(a[0]), int(a[1])) for a in alignment]
        except Exception as e:
            logger.debug(f"ALIGN | Lindat request failed: {e}")
        return []

class AwesomeAlignAligner:
    """BERT Aligner: Uses custom CT2-INT8 model from HuggingFace."""
    
    def __init__(self):
        self.available = False
        self.mode = None
        self.device = get_torch_device()
        self.ct2_dev, self.ct2_compute = get_ct2_settings()
        self.ct2_repo = "cstr/bert-base-multilingual-cased-ct2-int8"
        self.standard_repo = "bert-base-multilingual-cased"

        if HAS_CT2:
            try:
                logger.info(f"Loading CT2 Aligner from {self.ct2_repo}...")
                model_path = snapshot_download(repo_id=self.ct2_repo)
                self.encoder = ctranslate2.Encoder(
                    model_path, device=self.ct2_dev, compute_type=self.ct2_compute,
                    intra_threads=1 
                )
                # Load tokenizer from string to avoid local JSON collision
                self.tokenizer = AutoTokenizer.from_pretrained(self.standard_repo)
                self.mode, self.available = "ct2", True
                logger.info("✓ Awesome-Align CT2 ready.")
                return
            except Exception as e:
                logger.warning(f"CT2 Aligner load failed: {e}")

        if HAS_TRANSFORMERS and HAS_TORCH:
            try:
                from transformers import BertModel, BertTokenizer
                self.model = BertModel.from_pretrained(self.standard_repo)
                self.tokenizer = BertTokenizer.from_pretrained(self.standard_repo)
                if self.device.type == "mps": self.model = self.model.half()
                self.model.to(self.device).eval()
                self.mode, self.available = "torch", True
                logger.info(f"✓ Awesome-Align PyTorch ready on {self.device}")
            except Exception as e:
                logger.error(f"Critical: Aligner fallback failed: {e}")
    
    def align(self, src_words: List[str], tgt_words: List[str]) -> List[Tuple[int, int]]:
        """
        Extracts high-precision word alignments using BERT embeddings.
        Uses Mutual Argmax (Intersection) logic for 1-to-1 precision.
        Compatible with Mac CTranslate2 (no return_all_layers).
        """
        if not self.available or not src_words or not tgt_words:
            return []
        
        import numpy as np
        src_out, tgt_out = None, None
        
        try:
            # 1. PRE-PROCESSING: Subword tokenization
            def get_tokens_and_map(words):
                subtokens, word_map = [], []
                for i, w in enumerate(words):
                    tokens = self.tokenizer.tokenize(w) or [self.tokenizer.unk_token]
                    subtokens.extend(tokens)
                    word_map.extend([i] * len(tokens))
                return subtokens, word_map

            src_subtokens, src_word_map = get_tokens_and_map(src_words)
            tgt_subtokens, tgt_word_map = get_tokens_and_map(tgt_words)

            # 2. EMBEDDING EXTRACTION
            if self.mode == "ct2":
                src_input = [["[CLS]"] + src_subtokens + ["[SEP]"]]
                tgt_input = [["[CLS]"] + tgt_subtokens + ["[SEP]"]]
                
                # Use standard batch forward for cross-version compatibility
                res_src = self.encoder.forward_batch(src_input)
                res_tgt = self.encoder.forward_batch(tgt_input)
                
                # Extract Layer 12 (last_hidden_state)
                src_out = np.array(res_src.last_hidden_state)[0, 1:-1]
                tgt_out = np.array(res_tgt.last_hidden_state)[0, 1:-1]
                
            else: # PyTorch Mode
                import torch
                def to_ids(tokens):
                    ids = self.tokenizer.convert_tokens_to_ids(tokens)
                    return torch.tensor([self.tokenizer.cls_token_id] + ids + [self.tokenizer.sep_token_id]).to(self.device)
                
                with torch.no_grad():
                    # PyTorch uses Layer 8 (sweet spot)
                    out_s = self.model(to_ids(src_subtokens).unsqueeze(0), output_hidden_states=True)[2][8][0, 1:-1]
                    out_t = self.model(to_ids(tgt_subtokens).unsqueeze(0), output_hidden_states=True)[2][8][0, 1:-1]
                    src_out = out_s.detach().cpu().float().numpy()
                    tgt_out = out_t.detach().cpu().float().numpy()

            # 3. STABLE ALIGNMENT LOGIC: Mutual Argmax
            # Normalize vectors for cosine similarity
            src_norm = src_out / np.linalg.norm(src_out, axis=-1, keepdims=True)
            tgt_norm = tgt_out / np.linalg.norm(tgt_out, axis=-1, keepdims=True)
            similarity = np.dot(src_norm, tgt_norm.T)

            # Find best matches in both directions
            best_tgt_for_src = np.argmax(similarity, axis=1) # Shape: (src_len,)
            best_src_for_tgt = np.argmax(similarity, axis=0) # Shape: (tgt_len,)

            threshold = 1e-3
            align_words = set()

            # Mutual Agreement (The standard working method)
            for i, j in enumerate(best_tgt_for_src):
                # If source i picked target j, AND target j picked source i...
                if best_src_for_tgt[j] == i and similarity[i, j] > threshold:
                    # Map subword indices back to word indices
                    align_words.add((src_word_map[i], tgt_word_map[j]))
            
            final_alignments = sorted(list(align_words))
            
            # VERBOSE CLI LOG
            logger.debug(f"TRACE | Awesome-Align | Mode: {self.mode.upper()} | Links: {len(final_alignments)}")

            return final_alignments
            
        except Exception as e:
            logger.error(f"ALIGN | Logic failure: {e}")
            return []
    
    def align_old(self, src_words: List[str], tgt_words: List[str]) -> List[Tuple[int, int]]:
        if not self.available or not src_words or not tgt_words:
            return []
        
        try:
            import numpy as np
            
            # 1. Tokenize
            def get_tokens_and_map(words):
                tokens = []
                word_map = []
                for i, w in enumerate(words):
                    subwords = self.tokenizer.tokenize(w)
                    tokens.extend(subwords)
                    word_map.extend([i] * len(subwords))
                return tokens, word_map

            src_tokens, src_map = get_tokens_and_map(src_words)
            tgt_tokens, tgt_map = get_tokens_and_map(tgt_words)

            # 2. Extract Embeddings using CTranslate2
            # We add [CLS] and [SEP] just like BERT expects
            src_input = [["[CLS]"] + src_tokens + ["[SEP]"]]
            tgt_input = [["[CLS]"] + tgt_tokens + ["[SEP]"]]

            # forward_batch returns a StorageView; we convert to numpy for easy math
            src_out = np.array(self.encoder.forward_batch(src_input).last_hidden_state)[0, 1:-1]
            tgt_out = np.array(self.encoder.forward_batch(tgt_input).last_hidden_state)[0, 1:-1]

            # 3. Compute Similarity (Dot Product)
            # Normalizing vectors first ensures we are doing Cosine Similarity
            src_out /= np.linalg.norm(src_out, axis=-1, keepdims=True)
            tgt_out /= np.linalg.norm(tgt_out, axis=-1, keepdims=True)
            
            similarity = np.dot(src_out, tgt_out.T)

            # 4. Extract Alignments (Mutual Argmax / Threshold)
            threshold = 1e-3
            best_src = np.argmax(similarity, axis=1)
            best_tgt = np.argmax(similarity, axis=0)

            align_words = set()
            for i, j in enumerate(best_src):
                if best_tgt[j] == i and similarity[i, j] > threshold:
                    align_words.add((src_map[i], tgt_map[j]))
            
            return sorted(list(align_words))
            
        except Exception as e:
            logger.debug(f"CT2 Alignment failed: {e}")
            return []

class FastAlignAligner:
    """fast_align local aligner - uses temp file for binary"""
    
    def __init__(self):
        self.available = False
        self.mode = None
        self.binary_path = None
        self.atools_path = None
        
        # Check for Python package first
        try:
            import fast_align
            self.available = True
            self.mode = "python"
            return
        except ImportError:
            pass
        
        # Binary search logic
        script_dir = Path(__file__).parent
        search_dirs = [script_dir / "../fast_align/build", script_dir / "./fast_align/build"]
        
        for d in search_dirs:
            fa = d / "fast_align"
            at = d / "atools"
            if fa.exists() and os.access(fa, os.X_OK):
                self.binary_path = str(fa)
                self.atools_path = str(at) if at.exists() else None
                self.available = True
                self.mode = "binary"
                return
    
    def align(self, src_words: List[str], tgt_words: List[str]) -> List[Tuple[int, int]]:
        if not self.available or not src_words or not tgt_words:
            return []
        
        try:
            if self.mode == "python":
                from fast_align import align
                src_text = ' '.join(src_words)
                tgt_text = ' '.join(tgt_words)
                result = align([f"{src_text} ||| {tgt_text}"], forward=True)
                return [tuple(map(int, p.split('-'))) for p in result[0].split()]
            
            elif self.mode == "binary":
                import subprocess
                import tempfile
                
                input_str = f"{' '.join(src_words)} ||| {' '.join(tgt_words)}\n"
                
                with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
                    f.write(input_str)
                    temp_path = f.name
                
                try:
                    # -d: diagonal, -o: optimize tension, -v: variational bayes
                    # -I 1: single pass is enough for alignment of a known translation
                    cmd = [self.binary_path, '-i', temp_path, '-d', '-o', '-v', '-I', '1']
                    result = subprocess.run(cmd, capture_output=True, text=True, timeout=10)
                    
                    if result.returncode != 0:
                        return []
                    
                    # fast_align binary outputs to stdout in Pharaoh format (i-j)
                    alignments = []
                    output = result.stdout.strip()
                    if output:
                        for pair in output.split():
                            if '-' in pair:
                                i, j = map(int, pair.split('-'))
                                alignments.append((i, j))
                    return alignments
                finally:
                    if os.path.exists(temp_path): os.unlink(temp_path)
        except Exception as e:
            logger.debug(f"fast_align execution failed: {e}")
            return []

class SimAlignAligner:
    """SimAlign Aligner: Heavy PyTorch BERT (1.2GB RAM)."""
    def __init__(self):
        self.available = globals().get('HAS_SIMALIGN', False)
        if self.available:
            try:
                from simalign import SentenceAligner
                # FIXED: matching_methods must be a short string mapping code
                # "m" maps to "mwmf" internally in simalign.simalign.py
                self.aligner = SentenceAligner(model="bert", token_type="bpe", matching_methods="m", device="cpu")
                logger.info("✓ ALIGN | SimAlign (Standard BERT) ready.")
            except Exception as e:
                logger.error(f"ALIGN | SimAlign init failed: {e}")
                self.available = False
    
    def align(self, src_words: List[str], tgt_words: List[str]) -> List[Tuple[int, int]]:
        if not self.available or not src_words or not tgt_words: return []
        try:
            # returns dict of lists: {'mwmf': [(0,0), ...]}
            result = self.aligner.get_word_aligns(src_words, tgt_words)
            alignments = result.get('mwmf', [])
            return [(int(a[0]), int(a[1])) for a in alignments]
        except Exception as e:
            logger.error(f"ALIGN | SimAlign failed: {e}")
            return []

class HeuristicAligner:
    """Heuristic fallback - align words that appear in both"""
    
    def __init__(self):
        logger.info("✓ Heuristic aligner (fallback)")
    
    def align(self, src_words: List[str], tgt_words: List[str]) -> List[Tuple[int, int]]:
        """Simple heuristic alignment"""
        alignments = []
        
        src_lower = [w.lower().strip('.,!?;:') for w in src_words]
        tgt_lower = [w.lower().strip('.,!?;:') for w in tgt_words]
        
        for i, src_word in enumerate(src_lower):
            for j, tgt_word in enumerate(tgt_lower):
                if src_word == tgt_word and len(src_word) > 2:
                    alignments.append((i, j))
        
        return alignments


class MultiAligner:
    """Try multiple aligners with proper fallback"""
    
    def __init__(self, src_lang: str, tgt_lang: str, preferred: Optional[str] = None):
        self.aligners = []
        
        # If a specific aligner was requested
        if preferred and preferred != "auto":
            if preferred == "lindat":
                lindat = LindatAligner(src_lang, tgt_lang)
                if lindat.available:
                    self.aligners.append(("Lindat", lindat))
                else:
                    logger.warning(f"Requested aligner '{preferred}' not available, will try others")
            
            elif preferred == "awesome":
                awesome = AwesomeAlignAligner()
                if awesome.available:
                    self.aligners.append(("awesome-align", awesome))
            
            elif preferred == "fast_align":
                fast_align = FastAlignAligner()
                if HAS_FAST_ALIGN: 
                    self.aligners.append(("fast_align", fast_align))
                else:
                    logger.warning(f"fast_align requested but binary/package not found.")

            elif preferred == "simalign":
                simalign = SimAlignAligner()
                if simalign.available:
                    self.aligners.append(("SimAlign", simalign))
                else:
                    logger.warning(f"Requested aligner '{preferred}' not available, will try others")
            
            elif preferred == "heuristic":
                self.aligners.append(("Heuristic", HeuristicAligner()))
        
        # Auto mode or fallback: try all available aligners
        if not self.aligners or preferred == "auto":
            # Try Lindat first
            lindat = LindatAligner(src_lang, tgt_lang)
            if lindat.available:
                self.aligners.append(("Lindat", lindat))

            awesome = AwesomeAlignAligner()
            if awesome.available:
                self.aligners.append(("awesome-align", awesome))
            
            # Then SimAlign
            simalign = SimAlignAligner()
            if simalign.available:
                self.aligners.append(("SimAlign", simalign))
            
            # Then fast_align but will work right only with snapshot
            fast_align = FastAlignAligner()
            if fast_align.available:
                self.aligners.append(("fast_align", fast_align))
        
        # Always add heuristic as final fallback
        if not any(name == "Heuristic" for name, _ in self.aligners):
            self.aligners.append(("Heuristic", HeuristicAligner()))
        
        logger.info(f"Aligner chain: {[name for name, _ in self.aligners]}")
    
    def align(self, src_words: List[str], tgt_words: List[str]) -> List[Tuple[int, int]]:
        """Try aligners in order until one succeeds"""
        for name, aligner in self.aligners:
            result = aligner.align(src_words, tgt_words)
            if result:
                logger.debug(f"✓ {name} alignment: {len(result)} links")
                return result
        
        logger.debug("Using heuristic fallback (no quality alignments found)")
        return []

# ============================================================================
# PPTX HANDLING
# ============================================================================

# ============================================================================
# POWERPOINT-SPECIFIC STRUCTURES
# ============================================================================

@dataclass
class TranslatableTextFrame:
    """Text frame with positioning and shape metadata"""
    paragraphs: List[TranslatableParagraph] = field(default_factory=list)
    shape_metadata: Dict[str, Any] = field(default_factory=dict)
    
    def get_all_text(self) -> str:
        return '\n'.join(p.get_text() for p in self.paragraphs)


@dataclass
class SlideMetadata:
    """Slide-level properties"""
    layout: Any = None
    background: Any = None
    notes: Optional[str] = None




# ============================================================================
# DOCUMENT TRANSLATOR
# ============================================================================

class UltimateDocumentTranslator:
    """Document translator with configurable backends and M1 optimization."""
    
    def __init__(
        self,
        src_lang: str,
        tgt_lang: str,
        mode: TranslationMode = TranslationMode.HYBRID,
        nmt_backend: Optional[str] = "nllb",
        llm_provider: Optional[str] = None,
        aligner: Optional[str] = None,
        nllb_model_size: str = "600M"
    ):
        # Main initialization code
        self.src_lang, self.tgt_lang, self.mode = src_lang, tgt_lang, mode
        self.ct2 = None
        self.nllb = None
        self.llm = None
        self.aligner = None
        self.opus = None
        self.madlad = None
        
        # File type tracking
        self.current_file_type = None

        # Rest of initialization unchanged...
        logger.info(f"INIT | Starting Translator ({src_lang}→{tgt_lang})")
        logger.info(f"INIT | Mode: {mode.value} | NMT: {nmt_backend} | Aligner: {aligner or 'auto'}")
        self.log_memory("Initialization Start")

        # 1. NMT BACKEND SELECTION (Exclusive)
        if mode in [TranslationMode.NMT_ONLY, TranslationMode.HYBRID]:
            if nmt_backend == "nllb":
                self.nllb = NLLBTranslator(src_lang, tgt_lang, nllb_model_size)
                if not self.nllb.available:
                    logger.error("INIT | CRITICAL: NLLB specifically requested but failed to load.")
            
            elif nmt_backend == "opus":
                self.opus = OpusMTTranslator(src_lang, tgt_lang)
            
            elif nmt_backend == "madlad":
                self.madlad = Madlad400Translator(src_lang, tgt_lang, "3b")
            
            elif nmt_backend == "ct2":
                self.ct2 = CTranslate2Translator(src_lang, tgt_lang)
                if not self.ct2.available:
                    logger.error("INIT | CRITICAL: CTranslate2/WMT specifically requested but failed to load.")
            
            elif nmt_backend == "auto" or nmt_backend is None:
                # Priority: NLLB (Lightweight) -> CT2 (Dense)
                self.nllb = NLLBTranslator(src_lang, tgt_lang, nllb_model_size)
                if not (self.nllb and self.nllb.available):
                    logger.info("INIT | NLLB unavailable, trying CTranslate2/WMT...")
                    self.ct2 = CTranslate2Translator(src_lang, tgt_lang)

        # 2. LLM BACKEND (Hybrid or LLM modes)
        if mode in [TranslationMode.LLM_WITH_ALIGN, TranslationMode.LLM_WITHOUT_ALIGN, TranslationMode.HYBRID]:
            self.llm = LLMTranslator(src_lang, tgt_lang, llm_provider)
            if not self.llm.providers:
                logger.warning("INIT | Mode requires LLM but no providers (OpenAI/Anthropic/Ollama) available.")

        # 3. ALIGNER SELECTION (Priority: Awesome-Align)
        if mode in [TranslationMode.LLM_WITH_ALIGN, TranslationMode.HYBRID, TranslationMode.NMT_ONLY]:
            # We force 'awesome' as the preferred auto-aligner for M1 precision
            aligner_choice = aligner if aligner and aligner != "auto" else "awesome"
            self.aligner = MultiAligner(src_lang, tgt_lang, aligner_choice)
            
        self.log_memory("Initialization Complete")

        # FINAL VERIFICATION: Ensure at least one engine is ready
        engines_ready = any([
            self.nllb and self.nllb.available,
            self.ct2 and self.ct2.available,
            self.opus and self.opus.available,
            self.madlad and self.madlad.available,
            self.llm and self.llm.providers
        ])
        
        if not engines_ready:
            logger.error("INIT | CRITICAL FAILURE: No translation engines were able to load.")
            raise RuntimeError("No translation backends available. Check your model paths and API keys.")

    def log_memory(self, stage: str):
        """Log current RAM usage (requires psutil)."""
        try:
            import psutil
            process = psutil.Process(os.getpid())
            mem_mb = process.memory_info().rss / (1024 * 1024)
            logger.debug(f"MEM | Stage: {stage} | Usage: {mem_mb:.2f} MB")
        except ImportError:
            pass
        except Exception as e:
            logger.debug(f"Memory log failed: {e}")
    
    # with format specific error messages
    async def translate_file(self, input_path: Path, output_path: Path):
        """Main entry point with enhanced error handling"""
        try:
            is_valid, file_type, error = self.validate_file(input_path)
            if not is_valid:
                raise ValueError(error)
            
            self.current_file_type = file_type
            logger.info(f"Processing {file_type.upper()} file: {input_path.name}")
            
            if file_type == 'docx':
                await self.translate_document(input_path, output_path)
            elif file_type == 'pptx':
                await self.translate_presentation(input_path, output_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        
        except ImportError as e:
            if 'docx' in str(e) and self.current_file_type == 'docx':
                logger.error("python-docx not installed. Install with: pip install python-docx")
            elif 'pptx' in str(e) and self.current_file_type == 'pptx':
                logger.error("python-pptx not installed. Install with: pip install python-pptx")
            raise
        
        except Exception as e:
            format_name = "Word document" if self.current_file_type == 'docx' else "PowerPoint presentation"
            logger.error(f"Failed to translate {format_name}: {e}")
            raise
    
    async def translate_text(self, text: str) -> str:
        """Routes text through the active neural engine chain."""
        if not text.strip(): return text
        
        # 1. Try NLLB (Local CT2)
        if self.nllb and self.nllb.available:
            result = self.nllb.translate_batch([text])[0]
            if result and result.strip() != text.strip(): return result
            
        # 2. Try WMT21 (Local CT2 - 'ct2' backend)
        if self.ct2 and self.ct2.available:
            result = self.ct2.translate_batch([text])[0]
            if result and result.strip() != text.strip(): return result

        # 3. Try OPUS-MT (Local CT2)
        if self.opus and self.opus.available:
            result = self.opus.translate_batch([text])[0]
            if result and result.strip() != text.strip(): return result

        # 4. Try MADLAD-400 (Local CT2)
        if self.madlad and self.madlad.available:
            result = self.madlad.translate_batch([text])[0]
            if result and result.strip() != text.strip(): return result

        # 5. Try LLM (Hybrid/LLM modes)
        if self.llm and self.llm.providers:
            # use_alignment depends on mode
            use_align = (self.mode != TranslationMode.LLM_WITHOUT_ALIGN)
            result = await self.llm.translate_text(text, use_alignment=use_align)
            if result: return result
            
        return text
    
    def extract_paragraph(self, para: Paragraph) -> TranslatableParagraph:
        """Extracts paragraph with resolved font hierarchy to prevent Theme reversion."""
        # 1. Resolve the "Base Font" for this paragraph
        def get_resolved_font_name(p):
            # Check runs first
            for r in p.runs:
                if r.font.name: return r.font.name
            # Check style hierarchy
            curr_style = p.style
            while curr_style:
                if curr_style.font.name: return curr_style.font.name
                curr_style = curr_style.base_style
            return "Times New Roman" # Fallback if everything is 'None'

        resolved_base_font = get_resolved_font_name(para)
        
        runs = []
        for run in para.runs:
            f_color = None
            try:
                if run.font.color and run.font.color.rgb:
                    f_color = (run.font.color.rgb[0], run.font.color.rgb[1], run.font.color.rgb[2])
            except: pass

            runs.append(FormatRun(
                text=run.text,
                bold=run.bold,
                italic=run.italic,
                underline=run.underline,
                # If run has no font name, use the resolved base font we found
                font_name=run.font.name if run.font.name else resolved_base_font,
                font_size=run.font.size.pt if run.font.size else (para.style.font.size.pt if para.style.font.size else 12.0),
                font_color=f_color
            ))
        
        trans_para = TranslatableParagraph(runs=runs)
        trans_para.metadata['style'] = para.style
        trans_para.metadata['alignment'] = para.alignment
        
        pf = para.paragraph_format
        trans_para.metadata['layout'] = {
            'left_indent': pf.left_indent,
            'right_indent': pf.right_indent,
            'first_line_indent': pf.first_line_indent,
            'line_spacing': pf.line_spacing,
            'space_before': pf.space_before,
            'space_after': pf.space_after
        }
        return trans_para
    
    def log_memory(self, stage: str):
        """Log current RAM usage (requires psutil)."""
        try:
            import psutil
            process = psutil.Process(os.getpid())
            mem_mb = process.memory_info().rss / (1024 * 1024)
            logger.debug(f"MEM | Stage: {stage} | Usage: {mem_mb:.2f} MB")
        except ImportError:
            pass

    def copy_font_properties(self, target_run, source_run: FormatRun):
        """Forces Font Name, Size, and Color into XML. Does NOT touch bold/italic."""
        try:
            from docx.shared import Pt, RGBColor
            from docx.oxml.ns import qn
            if source_run.font_size:
                target_run.font.size = Pt(source_run.font_size)
            if source_run.font_color:
                target_run.font.color.rgb = RGBColor(*source_run.font_color)
            if source_run.underline is not None:
                target_run.font.underline = source_run.underline

            if source_run.font_name:
                # The XML 'rFonts' injection to bypass Theme defaults
                r = target_run._element
                rPr = r.get_or_add_rPr()
                rFonts = rPr.get_or_add_rFonts()
                rFonts.set(qn('w:ascii'), source_run.font_name)
                rFonts.set(qn('w:hAnsi'), source_run.font_name)
                rFonts.set(qn('w:eastAsia'), source_run.font_name)
                rFonts.set(qn('w:cs'), source_run.font_name)
                target_run.font.name = source_run.font_name
        except Exception as e:
            logger.debug(f"TRACE | Font Force Error: {e}")
    
    def apply_aligned_formatting(self, para: Paragraph, trans_para: TranslatableParagraph, translated_text: str, alignment: List[Tuple[int, int]]):
        """
        Reconstructs the Word paragraph by mapping source formatting onto translated text
        using neural word alignments. Preserves layout, font themes, and footnote anchors.
        """
        # 1. INITIAL TRACE & ANCHOR EXTRACTION
        self.log_para_trace(para, "INPUT")
        p_element = para._p
        
        # Extract all footnote reference elements (w:footnoteReference) or 
        # auto-numbered markers (w:footnoteRef) before we clear the paragraph.
        footnote_refs = p_element.xpath('.//w:r[w:footnoteReference or w:footnoteRef]')
        if footnote_refs:
            logger.debug(f"TRACE | Found {len(footnote_refs)} footnote anchors to re-attach.")

        # 2. RESTORE PARAGRAPH-LEVEL METADATA (Layout & Style)
        if trans_para.metadata.get('style'):
            para.style = trans_para.metadata['style']
        para.alignment = trans_para.metadata.get('alignment')
        
        layout = trans_para.metadata.get('layout', {})
        pf = para.paragraph_format
        try:
            # We explicitly set these to bypass Word's 'Normal' style defaults
            if layout.get('left_indent') is not None: pf.left_indent = layout['left_indent']
            if layout.get('right_indent') is not None: pf.right_indent = layout['right_indent']
            if layout.get('first_line_indent') is not None: pf.first_line_indent = layout['first_line_indent']
            if layout.get('line_spacing') is not None: pf.line_spacing = layout['line_spacing']
            if layout.get('space_before') is not None: pf.space_before = layout['space_before']
            if layout.get('space_after') is not None: pf.space_after = layout['space_after']
            logger.debug(f"TRACE | Layout metrics restored: Indent={pf.left_indent}, Spacing={pf.line_spacing}")
        except Exception as e:
            logger.error(f"TRACE | Layout Restore Error: {e}")

        # 3. CLEAR EXISTING CONTENT
        # We must remove runs one-by-one to keep the paragraph container intact
        for run in para.runs:
            p_element.remove(run._element)

        # 4. HANDLE FOOTNOTE-SPECIFIC STARTUP
        # For paragraphs that ARE footnote text, the marker must come first.
        is_footnote_content_para = "footnote" in str(para.style.name).lower()
        if is_footnote_content_para and footnote_refs:
            for ref in footnote_refs: 
                p_element.append(ref)
            para.add_run("\u00A0") # Add a non-breaking space after the number
            logger.debug("TRACE | Footnote number marker re-attached to start.")

        # 5. MAPPING LOGIC PREPARATION
        src_clean_words = trans_para.get_words()
        tgt_raw_units = translated_text.split() 
        formatted_indices = trans_para.get_formatted_word_indices()
        
        # Map Clean Aligner indices to Raw whitespace-split units
        clean_to_raw_tgt = {}
        clean_idx = 0
        for raw_idx, unit in enumerate(tgt_raw_units):
            # Only count as a 'word' if it contains alphanumeric characters
            if re.search(r'\w', unit):
                clean_to_raw_tgt[clean_idx] = raw_idx
                clean_idx += 1

        # Use the first source run as the baseline "Aesthetic" (Font Name/Size/Color)
        font_template = trans_para.runs[0] if trans_para.runs else None

        # 6. RECONSTRUCT RUNS WITH INLINE STYLES
        logger.debug(f"TRACE | Reconstructing {len(tgt_raw_units)} target units...")
        
        for i, unit in enumerate(tgt_raw_units):
            # Maintain original spacing
            run_text = unit + (" " if i < len(tgt_raw_units)-1 else "")
            run = para.add_run(run_text)
            
            # Determine Bold/Italic/Underline for this specific unit
            style_type = None
            matched_src_indices = [s_idx for s_idx, t_idx in alignment if clean_to_raw_tgt.get(t_idx) == i]
            
            if matched_src_indices:
                for s_idx in matched_src_indices:
                    # Check style priority: Bold+Italic > Bold > Italic
                    if s_idx in formatted_indices['italic_bold']:
                        style_type = 'italic_bold'; break
                    elif s_idx in formatted_indices['bold']:
                        style_type = 'bold'
                    elif s_idx in formatted_indices['italic'] and style_type != 'bold':
                        style_type = 'italic'
            
            # Apply Aligned Inline Styles
            if style_type == 'italic_bold': 
                run.bold = run.italic = True
            elif style_type == 'bold': 
                run.bold = True
            elif style_type == 'italic': 
                run.italic = True
            
            # Apply Baseline Aesthetics (The "Look" of the document)
            if font_template:
                self.copy_font_properties(run, font_template)

        # 7. RE-ANCHOR BODY FOOTNOTES
        # For main text, footnote citations usually go at the end of the translated block
        if not is_footnote_content_para and footnote_refs:
            logger.debug(f"TRACE | Body Footnote | Re-anchoring {len(footnote_refs)} refs to end of paragraph.")
            for ref in footnote_refs: 
                p_element.append(ref)

        self.log_para_trace(para, "OUTPUT")
    
    def is_paragraph_safe_to_translate(self, para: Paragraph) -> bool:
        """Check if paragraph can be safely translated"""
        if not para.text or not para.text.strip():
            return False
        
        text = para.text.strip()
        if not text:
            return False
        
        if len(para.text.strip()) <= 1 and not any(run.text.strip() for run in para.runs):
            return False
        
        try:
            for run in para.runs:
                if run._element.xpath('.//w:drawing | .//w:pict'):
                    logger.debug(f"Skipping paragraph with drawing/image")
                    return False
        except:
            pass
        
        try:
            field_chars = para._element.xpath('.//w:fldChar')
            if field_chars:
                is_footnote_field = para._element.xpath('.//w:footnoteReference')
                if not is_footnote_field:
                    logger.debug(f"Skipping paragraph with field")
                    return False
        except:
            pass
        
        return True
    
    async def translate_paragraph(self, para: Paragraph):
        """Paragraph translation with multi-backend routing and alignment."""
        if not para.text or not para.text.strip() or not self.is_paragraph_safe_to_translate(para):
            return
        
        try:
            self.log_memory("Pre-Paragraph")
            trans_para = self.extract_paragraph(para)
            original_text = trans_para.get_text()
            
            # 1. Split into sentences for higher NMT quality
            sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', original_text) if s.strip()]
            if not sentences: return

            # 2. FIXED: Use the central router for each sentence
            translated_sentences = []
            for s in sentences:
                t = await self.translate_text(s)
                translated_sentences.append(t)
            
            translated_text = " ".join(translated_sentences)
            
            # 3. Preparation for Aligner
            src_words = trans_para.get_words() 
            tgt_words_clean = re.findall(r"\w+", translated_text)
            
            # Log the translation for the CLI trace
            logger.info("-" * 30)
            logger.info(f"TRANS | Out: {translated_text[:60]}...")
            
            if not src_words or not tgt_words_clean:
                para.text = translated_text
                return
                
            # 4. Alignment Pass
            alignment = []
            if self.aligner:
                alignment = self.aligner.align(src_words, tgt_words_clean)
                logger.info(f"ALIGN | Matches: {len(alignment)}")
            
            # 5. Reconstruction
            self.apply_aligned_formatting(para, trans_para, translated_text, alignment)
            self.log_memory("Post-Paragraph")
                    
        except Exception as e:
            logger.error(f"PARA | Translation Failed: {e}", exc_info=True)
    
    def get_footnotes(self, doc: Document) -> List[Paragraph]:
        """Extract footnotes. We pass doc as parent to avoid Part attribute error"""
        try:
            document_part = doc.part
            footnote_part = None
            
            for rel in document_part.rels.values():
                if "relationships/footnotes" in rel.reltype:
                    footnote_part = rel.target_part
                    break
            
            if not footnote_part:
                return []
            
            from docx.oxml import parse_xml
            root = parse_xml(footnote_part.blob)
            ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
            
            paragraphs = []
            for footnote in root.xpath('//w:footnote', namespaces=ns):
                f_id = footnote.get(f'{{{ns["w"]}}}id')
                # Skip internal Word markers (id 0 and -1)
                if f_id and int(f_id) <= 0:
                    continue
                
                for p_elem in footnote.xpath('.//w:p', namespaces=ns):
                    # FIX: Pass 'doc' instead of 'None' so para.part is valid
                    para = Paragraph(p_elem, doc)
                    if para.text.strip():
                        paragraphs.append(para)
            
            self._footnote_root = root
            self._footnote_part = footnote_part
            return paragraphs
            
        except Exception as e:
            logger.warning(f"Could not extract footnotes: {e}")
            return []
    
    def get_all_paragraphs(self, doc: Document) -> List[Tuple[Paragraph, str]]:
        """Aggregate all paragraphs from document"""
        all_paras = []
        
        # Main body
        for para in doc.paragraphs:
            all_paras.append((para, "body"))
        
        # Tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        all_paras.append((para, "table"))
        
        # Footnotes - with error handling
        try:
            footnote_paras = self.get_footnotes(doc)
            for para in footnote_paras:
                all_paras.append((para, "footnote"))
        except Exception as e:
            logger.warning(f"Skipping footnotes due to error: {e}")
        
        # Headers/Footers
        for section in doc.sections:
            for para in section.header.paragraphs:
                all_paras.append((para, "header"))
            for para in section.footer.paragraphs:
                all_paras.append((para, "footer"))
        
        return all_paras
    
    def log_document_info(self, doc: Document, label: str):
        """Logs global document properties and style inventory."""
        logger.debug(f"{'='*20} DOCUMENT INFO [{label}] {'='*20}")
        for i, section in enumerate(doc.sections):
            logger.debug(f"Section {i} | Size: {section.page_width.pt:.1f}x{section.page_height.pt:.1f}pt")
            logger.debug(f"Section {i} | Margins: L:{section.left_margin.pt:.1f} R:{section.right_margin.pt:.1f} T:{section.top_margin.pt:.1f} B:{section.bottom_margin.pt:.1f}")
            logger.debug(f"Section {i} | Gutter: {section.gutter.pt:.1f}pt | Header-Dist: {section.header_distance.pt:.1f}pt")
        
        # Log all paragraph styles present in the document
        style_names = [s.name for s in doc.styles if s.type == 1]
        logger.debug(f"Style Inventory ({len(style_names)}): {', '.join(style_names)}")

    def log_para_trace(self, para: Paragraph, label: str):
        """Detailed debug log showing style, font, and layout metrics."""
        pf = para.paragraph_format
        style = para.style
        
        # Check if first run has inline formatting
        has_bold = any(r.bold for r in para.runs)
        has_ital = any(r.italic for r in para.runs)
        
        f_name = para.runs[0].font.name if para.runs and para.runs[0].font.name else "Inherited"
        f_size = f"{para.runs[0].font.size.pt if para.runs and para.runs[0].font.size else 'Default'}pt"
        
        logger.debug(f"PARA [{label}] | Style: '{style.name}'")
        logger.debug(f"  > Font: {f_name} @ {f_size} | Bold-Any: {has_bold} | Ital-Any: {has_ital}")
        logger.debug(f"  > Indent-L: {pf.left_indent.pt if pf.left_indent else 0:.1f}pt | Spacing-A: {pf.space_after.pt if pf.space_after else 0:.1f}pt")

    async def translate_document(self, input_path: Path, output_path: Path):
        """
        Full Word document translation lifecycle with robust XML commitment and verification logs.
        Called internally by translate_file().
        """
        self.log_memory("Initialization")
        doc = Document(str(input_path))
        
        if logger.isEnabledFor(logging.DEBUG):
            self.log_document_info(doc, "INPUT")

        # Gather footnotes properly to initialize _footnote_root
        _ = self.get_footnotes(doc)
        all_paras = self.get_all_paragraphs(doc)
        translatable = [(p, l) for p, l in all_paras if p.text.strip() and self.is_paragraph_safe_to_translate(p)]
        
        logger.info(f"Processing {len(translatable)} paragraphs across {len(all_paras)} total entities.")
        
        for para, location in tqdm(translatable, desc="Translating"):
            try:
                await self.translate_paragraph(para)
            except Exception as e:
                logger.error(f"Error in {location} translation: {e}")

        # FOOTNOTE XML COMMITMENT
        if hasattr(self, '_footnote_part') and hasattr(self, '_footnote_root'):
            try:
                from lxml import etree
                # Standard etree tostring ensures namespace 'w' is preserved correctly
                updated_xml = etree.tostring(self._footnote_root, encoding='utf-8', xml_declaration=True)
                self._footnote_part._blob = updated_xml
                logger.info("✓ Success | Footnote XML blob successfully serialized.")
            except Exception as e:
                logger.error(f"Error | Footnote commitment failed: {e}")

        logger.info(f"Saving file to {output_path}")
        doc.save(str(output_path))
        
        if logger.isEnabledFor(logging.DEBUG):
            self.log_document_info(Document(str(output_path)), "OUTPUT")
        logger.info("✓ Document Translation Complete.")

    # ============================================================================
    # PPTX EXTRACTION METHODS
    # ============================================================================

    def extract_text_frame(self, shape) -> TranslatableTextFrame:
        """
        Extract text frame with shape positioning and paragraph hierarchy.
        PPT equivalent of extract_paragraph().
        """
        from pptx.util import Pt, Inches
        
        trans_frame = TranslatableTextFrame()
        
        # Capture shape-level metadata (positioning, size, rotation)
        trans_frame.shape_metadata = {
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
            'rotation': shape.rotation,
            'shape_type': shape.shape_type,
            'name': shape.name
        }
        
        # Capture text frame properties
        if hasattr(shape, 'text_frame'):
            tf = shape.text_frame
            trans_frame.shape_metadata['text_frame'] = {
                'margin_left': tf.margin_left,
                'margin_right': tf.margin_right,
                'margin_top': tf.margin_top,
                'margin_bottom': tf.margin_bottom,
                'vertical_anchor': tf.vertical_anchor,
                'word_wrap': tf.word_wrap,
                'auto_size': tf.auto_size
            }
            
            # Extract each paragraph with runs
            for para in tf.paragraphs:
                trans_para = self.extract_ppt_paragraph(para)
                trans_frame.paragraphs.append(trans_para)
        
        return trans_frame


    def extract_ppt_paragraph(self, para) -> TranslatableParagraph:
        """
        Extract PowerPoint paragraph with run-level formatting.
        Similar to Word's extract_paragraph but for PPT-specific properties.
        """
        from pptx.util import Pt
        
        # Resolve base font from theme/master
        def get_resolved_ppt_font(p):
            # Check runs first
            for r in p.runs:
                if r.font.name:
                    return r.font.name
            # Check theme defaults
            try:
                if hasattr(p, '_element') and hasattr(p._element, 'pPr'):
                    # Theme font resolution logic here
                    pass
            except:
                pass
            return "Calibri"  # PPT default
        
        resolved_font = get_resolved_ppt_font(para)
        
        runs = []
        for run in para.runs:
            f_color = None
            try:
                if run.font.color and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    f_color = (rgb[0], rgb[1], rgb[2])
            except:
                pass
            
            runs.append(FormatRun(
                text=run.text,
                bold=run.font.bold,
                italic=run.font.italic,
                underline=run.font.underline,
                font_name=run.font.name if run.font.name else resolved_font,
                font_size=run.font.size.pt if run.font.size else 18.0,
                font_color=f_color
            ))
        
        trans_para = TranslatableParagraph(runs=runs)
        
        # Capture paragraph-level properties
        trans_para.metadata['alignment'] = para.alignment
        trans_para.metadata['level'] = para.level  # Indentation level
        trans_para.metadata['line_spacing'] = para.line_spacing
        trans_para.metadata['space_before'] = para.space_before
        trans_para.metadata['space_after'] = para.space_after
        
        return trans_para


    def extract_table_from_slide(self, table) -> List[List[TranslatableTextFrame]]:
        """
        Extract table structure with cell-level text frames.
        PPT tables are similar to Word but stored differently.
        """
        table_data = []
        
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_frame = self.extract_text_frame(cell)
                row_data.append(cell_frame)
            table_data.append(row_data)
        
        return table_data


    def get_speaker_notes(self, slide) -> Optional[str]:
        """
        Extract speaker notes from slide.
        PPT equivalent of footnotes in Word.
        """
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                return text_frame.text if text_frame.text.strip() else None
        except:
            return None
        
    # ============================================================================
    # PPTX SAFETY CHECKS 
    # ============================================================================

    def is_shape_safe_to_translate(self, shape) -> bool:
        """
        Check if PowerPoint shape can be safely translated.
        Equivalent to is_paragraph_safe_to_translate() for Word.
        """
        # Skip shapes without text frames
        if not hasattr(shape, 'text_frame'):
            return False
        
        try:
            text_frame = shape.text_frame
            
            # Skip empty text frames
            if not text_frame.text or not text_frame.text.strip():
                return False
            
            # Skip placeholder shapes with no actual content
            if len(text_frame.text.strip()) <= 1:
                return False
            
            # Skip shapes that are likely logos or decorative
            if shape.name and any(keyword in shape.name.lower() 
                                for keyword in ['logo', 'watermark', 'decoration', 'icon']):
                logger.debug(f"Skipping decorative shape: {shape.name}")
                return False
            
            # Skip very small shapes (likely decorative)
            if shape.width < 100000 or shape.height < 100000:  # Less than ~0.14 inches
                logger.debug(f"Skipping tiny shape: {shape.width}x{shape.height}")
                return False
            
            return True
            
        except Exception as e:
            logger.debug(f"Shape safety check failed: {e}")
            return False


    def is_slide_master_or_layout(self, slide) -> bool:
        """
        Detect if this is a master slide or layout (should not be translated).
        """
        try:
            # Master slides don't have a slide_id in the normal sense
            if not hasattr(slide, 'slide_id'):
                return True
            # Additional checks could go here
            return False
        except:
            return True


    # ============================================================================
    # PPTX RECONSTRUCTION METHODS
    # ============================================================================

    def apply_text_frame_formatting(
        self, 
        shape,
        trans_frame: TranslatableTextFrame,
        translated_paragraphs: List[Tuple[str, List[Tuple[int, int]]]]
    ):
        """
        Reconstruct text frame with aligned formatting.
        PPT equivalent of apply_aligned_formatting().
        
        Args:
            shape: PowerPoint shape object
            trans_frame: Original extracted text frame
            translated_paragraphs: List of (translated_text, alignment) tuples
        """
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
        
        tf = shape.text_frame
        
        # Restore text frame properties
        tf_meta = trans_frame.shape_metadata.get('text_frame', {})
        if tf_meta:
            try:
                tf.margin_left = tf_meta.get('margin_left', tf.margin_left)
                tf.margin_right = tf_meta.get('margin_right', tf.margin_right)
                tf.margin_top = tf_meta.get('margin_top', tf.margin_top)
                tf.margin_bottom = tf_meta.get('margin_bottom', tf.margin_bottom)
                tf.vertical_anchor = tf_meta.get('vertical_anchor', tf.vertical_anchor)
                tf.word_wrap = tf_meta.get('word_wrap', tf.word_wrap)
                tf.auto_size = tf_meta.get('auto_size', tf.auto_size)
            except Exception as e:
                logger.debug(f"Text frame property restoration failed: {e}")
        
        # Clear existing paragraphs
        for _ in range(len(tf.paragraphs)):
            tf._element.remove(tf.paragraphs[0]._element)
        
        # Reconstruct paragraphs
        for i, (trans_para, (translated_text, alignment)) in enumerate(
            zip(trans_frame.paragraphs, translated_paragraphs)
        ):
            para = tf.add_paragraph()
            
            # Restore paragraph properties
            para.alignment = trans_para.metadata.get('alignment')
            para.level = trans_para.metadata.get('level', 0)
            
            if trans_para.metadata.get('line_spacing'):
                para.line_spacing = trans_para.metadata['line_spacing']
            if trans_para.metadata.get('space_before'):
                para.space_before = trans_para.metadata['space_before']
            if trans_para.metadata.get('space_after'):
                para.space_after = trans_para.metadata['space_after']
            
            # Apply aligned formatting to runs
            self.apply_ppt_paragraph_formatting(
                para, trans_para, translated_text, alignment
            )


    def apply_ppt_paragraph_formatting(
        self,
        para,
        trans_para: TranslatableParagraph,
        translated_text: str,
        alignment: List[Tuple[int, int]]
    ):
        """
        Apply aligned formatting to PowerPoint paragraph runs.
        Core formatting transfer logic - same as Word but for PPT runs.
        """
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        
        src_clean_words = trans_para.get_words()
        tgt_raw_units = translated_text.split()
        formatted_indices = trans_para.get_formatted_word_indices()
        
        # Map clean indices to raw units
        clean_to_raw_tgt = {}
        clean_idx = 0
        for raw_idx, unit in enumerate(tgt_raw_units):
            if re.search(r'\w', unit):
                clean_to_raw_tgt[clean_idx] = raw_idx
                clean_idx += 1
        
        # Get font template
        font_template = trans_para.runs[0] if trans_para.runs else None
        
        # Reconstruct runs
        for i, unit in enumerate(tgt_raw_units):
            run_text = unit + (" " if i < len(tgt_raw_units)-1 else "")
            run = para.add_run()
            run.text = run_text
            
            # Determine style from alignment
            style_type = None
            matched_src = [s for s, t in alignment if clean_to_raw_tgt.get(t) == i]
            
            if matched_src:
                for s_idx in matched_src:
                    if s_idx in formatted_indices['italic_bold']:
                        style_type = 'italic_bold'
                        break
                    elif s_idx in formatted_indices['bold']:
                        style_type = 'bold'
                    elif s_idx in formatted_indices['italic'] and style_type != 'bold':
                        style_type = 'italic'
            
            # Apply inline styles
            if style_type == 'italic_bold':
                run.font.bold = run.font.italic = True
            elif style_type == 'bold':
                run.font.bold = True
            elif style_type == 'italic':
                run.font.italic = True
            
            # Apply baseline aesthetics
            if font_template:
                self.copy_ppt_font_properties(run, font_template)


    def copy_ppt_font_properties(self, target_run, source_run: FormatRun):
        """
        Force font properties in PowerPoint run.
        PPT equivalent of copy_font_properties().
        """
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        
        try:
            if source_run.font_name:
                target_run.font.name = source_run.font_name
            
            if source_run.font_size:
                target_run.font.size = Pt(source_run.font_size)
            
            if source_run.font_color:
                target_run.font.color.rgb = RGBColor(*source_run.font_color)
            
            if source_run.underline is not None:
                target_run.font.underline = source_run.underline
                
        except Exception as e:
            logger.debug(f"PPT font property copy failed: {e}")


    def restore_table_to_slide(
        self,
        table,
        table_data: List[List[TranslatableTextFrame]],
        translated_cells: List[List[Tuple[str, List[Tuple[int, int]]]]]
    ):
        """
        Restore translated content to PowerPoint table.
        """
        for i, row in enumerate(table.rows):
            for j, cell in enumerate(row.cells):
                if i < len(table_data) and j < len(table_data[i]):
                    trans_frame = table_data[i][j]
                    trans_paragraphs = translated_cells[i][j]
                    
                    # Treat cell as a shape with text frame
                    self.apply_text_frame_formatting(
                        cell, trans_frame, trans_paragraphs
                    )


    def set_speaker_notes(self, slide, translated_notes: str):
        """
        Set translated speaker notes.
        """
        try:
            if not slide.has_notes_slide:
                notes_slide = slide.notes_slide  # Creates if doesn't exist
            else:
                notes_slide = slide.notes_slide
            
            text_frame = notes_slide.notes_text_frame
            text_frame.clear()
            text_frame.text = translated_notes
        except Exception as e:
            logger.warning(f"Could not set speaker notes: {e}")


    # ============================================================================
    # PPTX SLIDE PROCESSING
    # ============================================================================

    async def translate_shape(self, shape):
        """
        Translate a single shape (text box, placeholder, etc.)
        """
        if not shape.has_text_frame:
            return
        
        try:
            # Extract
            trans_frame = self.extract_text_frame(shape)
            
            if not trans_frame.paragraphs:
                return
            
            # Translate each paragraph
            translated_paragraphs = []
            for trans_para in trans_frame.paragraphs:
                original_text = trans_para.get_text()
                if not original_text.strip():
                    translated_paragraphs.append(("", []))
                    continue
                
                # Translate
                translated_text = await self.translate_text(original_text)
                
                # Align
                src_words = trans_para.get_words()
                tgt_words = re.findall(r"\w+", translated_text)
                alignment = []
                if self.aligner and src_words and tgt_words:
                    alignment = self.aligner.align(src_words, tgt_words)
                
                translated_paragraphs.append((translated_text, alignment))
            
            # Reconstruct
            self.apply_text_frame_formatting(shape, trans_frame, translated_paragraphs)
            
        except Exception as e:
            logger.error(f"Shape translation failed: {e}", exc_info=True)


    async def translate_slide(self, slide):
        """
        Translate all content in a slide.
        """
        # Process shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                await self.translate_shape(shape)
            
            # Handle tables
            if shape.has_table:
                await self.translate_table_in_slide(shape.table)
            
            # Handle groups (recursive)
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                for sub_shape in shape.shapes:
                    if sub_shape.has_text_frame:
                        await self.translate_shape(sub_shape)
        
        # Process speaker notes
        notes_text = self.get_speaker_notes(slide)
        if notes_text:
            translated_notes = await self.translate_text(notes_text)
            self.set_speaker_notes(slide, translated_notes)


    async def translate_table_in_slide(self, table):
        """
        Translate table content in slide.
        """
        table_data = self.extract_table_from_slide(table)
        translated_cells = []
        
        for row_data in table_data:
            translated_row = []
            for cell_frame in row_data:
                cell_paragraphs = []
                for trans_para in cell_frame.paragraphs:
                    text = trans_para.get_text()
                    if text.strip():
                        translated = await self.translate_text(text)
                        src_words = trans_para.get_words()
                        tgt_words = re.findall(r"\w+", translated)
                        alignment = self.aligner.align(src_words, tgt_words) if self.aligner else []
                        cell_paragraphs.append((translated, alignment))
                    else:
                        cell_paragraphs.append(("", []))
                translated_row.append(cell_paragraphs)
            translated_cells.append(translated_row)
        
        self.restore_table_to_slide(table, table_data, translated_cells)


    async def translate_presentation(self, input_path: Path, output_path: Path):
        """
        Main presentation translation lifecycle.
        PPT equivalent of translate_document().
        """
        from pptx import Presentation
        
        prs = Presentation(str(input_path))
        
        logger.info(f"Processing {len(prs.slides)} slides")
        
        for slide_num, slide in enumerate(tqdm(prs.slides, desc="Translating slides"), 1):
            logger.info(f"Processing slide {slide_num}")
            await self.translate_slide(slide)
        
        logger.info(f"Saving presentation to {output_path}")
        prs.save(str(output_path))
        logger.info("✓ Presentation Translation Complete.")

    # ============================================================================
    # FILE TYPE DETECTION
    # ============================================================================

    def detect_file_type(self, file_path: Path) -> str:
        """
        Detect if file is Word (.docx) or PowerPoint (.pptx).
        Returns: 'docx', 'pptx', or 'unknown'
        """
        suffix = file_path.suffix.lower()
        
        if suffix == '.docx':
            return 'docx'
        elif suffix == '.pptx':
            return 'pptx'
        else:
            # Try to detect by magic bytes
            try:
                with open(file_path, 'rb') as f:
                    header = f.read(4)
                    # Both formats are ZIP files (PK header)
                    if header[:2] == b'PK':
                        # Try to open as docx first
                        try:
                            from docx import Document
                            Document(str(file_path))
                            return 'docx'
                        except:
                            pass
                        # Try as pptx
                        try:
                            from pptx import Presentation
                            Presentation(str(file_path))
                            return 'pptx'
                        except:
                            pass
            except:
                pass
        
        return 'unknown'


    def validate_file(self, file_path: Path) -> Tuple[bool, str, str]:
        """
        Validate input file and return (is_valid, file_type, error_message).
        
        Returns:
            Tuple of (success, file_type, error_msg)
        """
        if not file_path.exists():
            return False, 'unknown', f"File not found: {file_path}"
        
        file_type = self.detect_file_type(file_path)
        
        if file_type == 'unknown':
            return False, 'unknown', f"Unsupported file format. Only .docx and .pptx are supported."
        
        # Verify we can actually open it
        try:
            if file_type == 'docx':
                from docx import Document
                doc = Document(str(file_path))
                # Basic sanity check
                if not hasattr(doc, 'paragraphs'):
                    return False, file_type, "Invalid .docx file structure"
            
            elif file_type == 'pptx':
                from pptx import Presentation
                prs = Presentation(str(file_path))
                # Basic sanity check
                if not hasattr(prs, 'slides'):
                    return False, file_type, "Invalid .pptx file structure"
            
            return True, file_type, ""
        
        except Exception as e:
            return False, file_type, f"Cannot open file: {str(e)}"
        




# ============================================================================
# CLI
# ============================================================================

async def main():
    parser = argparse.ArgumentParser(
        description='🚀 Ultimate Document Translator - Multi-Backend Production Version',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Backends Comparison:
  nllb       - Distilled 600M (Default). Best speed/RAM ratio. Great general support.
  madlad     - Google's 3B model. Superior academic/formal quality (~3GB RAM).
  opus       - Specialized bilingual models. Tiny (~200MB), extremely fast, literal.
  ct2 (wmt)  - Dense Facebook models. Peak German/European quality (~6GB RAM).

Supported Formats:
  .docx      - Microsoft Word documents (paragraphs, tables, footnotes, headers/footers)
  .pptx      - Microsoft PowerPoint presentations (slides, text boxes, tables, notes)

Examples:
  # Translate Word document (NLLB-600M)
  %(prog)s input.docx output.docx -s en -t de
  
  # Translate PowerPoint presentation
  %(prog)s presentation.pptx translated.pptx -s en -t es
  
  # High-quality academic translation (Madlad-400)
  %(prog)s thesis.docx thesis_de.docx -s en -t de --nmt madlad
  
  # Translate slides with LLM (Claude)
  %(prog)s slides.pptx slides_fr.pptx -s en -t fr --mode llm-align --llm anthropic

Environment Variables:
  OPENAI_API_KEY, ANTHROPIC_API_KEY - Required for LLM backends.
        """
    )
    
    # 1. POSITIONAL ARGUMENTS, HELP TEXT
    parser.add_argument('input', help='Input file (.docx or .pptx)')
    parser.add_argument('output', help='Output file (.docx or .pptx)')
    
    parser.add_argument('-s', '--source', default='en', help='Source language code (default: en)')
    parser.add_argument('-t', '--target', default='de', help='Target language code (default: de)')
    
    parser.add_argument(
        '--mode',
        choices=['nmt', 'llm-align', 'llm-plain', 'hybrid'],
        default='hybrid',
        help='Translation strategy (default: hybrid)'
    )
    
    parser.add_argument(
        '--nmt',
        choices=['nllb', 'madlad', 'opus', 'ct2', 'auto'],
        default='nllb',
        help='Local NMT Engine: NLLB (general), Madlad (academic), Opus (specialized), CT2 (dense)'
    )
    
    parser.add_argument(
        '--nllb-size',
        choices=['600M', '1.3B', '3.3B'],
        default='600M',
        help='NLLB variant only: 600M (fastest), 1.3B (balanced), 3.3B (heavy)'
    )
    
    parser.add_argument(
        '--llm',
        choices=['openai', 'anthropic', 'ollama'],
        help='LLM provider for hybrid/llm modes'
    )
    
    parser.add_argument(
        '--aligner',
        choices=['awesome', 'simalign', 'lindat', 'fast_align', 'heuristic', 'auto'],
        default='auto',
        help='Word Aligner: awesome (M1 optimized), simalign (heavy), lindat (API)'
    )
    
    parser.add_argument('-v', '--verbose', action='store_true', help='Enable detailed TRACE logging')
    
    args = parser.parse_args()
    
    # Set logging level
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    # Enhanced validation
    input_path = Path(args.input)
    output_path = Path(args.output)
    
    if not input_path.exists():
        logger.error(f"File not found: {input_path}")
        sys.exit(1)
    
    # Validate file type compatibility
    input_type = input_path.suffix.lower()
    output_type = output_path.suffix.lower()
    
    if input_type not in ['.docx', '.pptx']:
        logger.error(f"Unsupported input format: {input_type}. Only .docx and .pptx are supported.")
        sys.exit(1)
    
    # Warn if output extension doesn't match input
    if output_type != input_type:
        logger.warning(f"Output extension ({output_type}) doesn't match input ({input_type}). Using {input_type}.")
        output_path = output_path.with_suffix(input_type)
    
    # Mode mapping unchanged
    mode_map = {
        'nmt': TranslationMode.NMT_ONLY,
        'llm-align': TranslationMode.LLM_WITH_ALIGN,
        'llm-plain': TranslationMode.LLM_WITHOUT_ALIGN,
        'hybrid': TranslationMode.HYBRID
    }
    
    # Updated status header
    file_type_name = "Word Document" if input_type == '.docx' else "PowerPoint Presentation"
    
    print(f"\n{'='*60}")
    print(f"🌍 DOCUMENT TRANSLATOR - PRODUCTION v12")
    print(f"{'='*60}")
    print(f"Format:     {file_type_name}")
    print(f"Input:      {input_path.name}")
    print(f"Output:     {output_path.name}")
    print(f"Direction:  {args.source.upper()} → {args.target.upper()}")
    print(f"Mode:       {args.mode.upper()}")
    print(f"NMT Engine: {args.nmt.upper()} {'('+args.nllb_size+')' if args.nmt=='nllb' else ''}")
    print(f"Aligner:    {args.aligner.upper()}")
    if args.llm:
        print(f"LLM:        {args.llm.upper()}")
    print(f"{'='*60}\n")
    
    # Initialize translator
    translator = UltimateDocumentTranslator(
        src_lang=args.source,
        tgt_lang=args.target,
        mode=mode_map[args.mode],
        nmt_backend=args.nmt,
        llm_provider=args.llm,
        aligner=args.aligner,
        nllb_model_size=args.nllb_size
    )
    
    # Use unified translate_file method
    try:
        await translator.translate_file(input_path, output_path)
        
        print(f"\n{'='*60}")
        print(f"✅ Success! {file_type_name} processed in {args.mode} mode.")
        print(f"💾 File saved to: {output_path}")
        print(f"{'='*60}\n")
    except Exception as e:
        logger.error(f"FAILED | Translation aborted: {e}", exc_info=args.verbose)
        sys.exit(1)


if __name__ == "__main__":
    asyncio.run(main())